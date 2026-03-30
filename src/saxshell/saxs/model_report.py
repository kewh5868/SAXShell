from __future__ import annotations

import json
import textwrap
from collections.abc import Callable, Sequence
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from tempfile import TemporaryDirectory

import numpy as np
from matplotlib import colormaps, rc_context
from matplotlib.colors import to_hex
from matplotlib.figure import Figure

from saxshell.saxs.dream import (
    DreamModelPlotData,
    DreamParameterEntry,
    DreamRunSettings,
    DreamSummary,
    DreamViolinPlotData,
)
from saxshell.saxs.prefit import PrefitEvaluation, PrefitParameterEntry
from saxshell.saxs.prefit.cluster_geometry import ClusterGeometryMetadataRow
from saxshell.saxs.project_manager import PowerPointExportSettings
from saxshell.saxs.project_manager.prior_plot import plot_md_prior_histogram
from saxshell.saxs.solution_scattering_estimator import (
    SolutionScatteringEstimate,
)

_SLIDE_WIDTH_INCHES = 13.333
_SLIDE_HEIGHT_INCHES = 7.5
_SLIDE_LEFT_INCHES = 0.45
_SLIDE_TOP_INCHES = 1.1
_SLIDE_CONTENT_WIDTH_INCHES = 12.43
_SLIDE_CONTENT_HEIGHT_INCHES = 5.85
_TABLE_TOP_INCHES = 1.42
_TABLE_WIDTH_INCHES = 12.43
_TABLE_HEIGHT_INCHES = 5.32
_TABLE_NOTE_TOP_INCHES = 6.88
_THICK_RULE_HEIGHT_INCHES = 0.06
_TEXT_WRAP_FULL = 98
_TEXT_WRAP_HALF = 43
_TEXT_WRAP_SIDE = 37
_TEXT_LINES_FULL = 24
_TEXT_LINES_HALF = 23
_TEXT_LINES_SIDE = 20
_PREFIT_WINDOW_EXPERIMENTAL_COLOR = to_hex("black", keep_alpha=False)
_PREFIT_WINDOW_MODEL_COLOR = to_hex("tab:red", keep_alpha=False)
_PREFIT_WINDOW_SOLVENT_COLOR = to_hex("green", keep_alpha=False)
_PREFIT_WINDOW_RESIDUAL_COLOR = to_hex("tab:blue", keep_alpha=False)
_DREAM_OUTPUT_EXPERIMENTAL_COLOR = to_hex("black", keep_alpha=False)
_DREAM_OUTPUT_MODEL_COLOR = to_hex("tab:red", keep_alpha=False)
_DREAM_OUTPUT_SOLVENT_COLOR = to_hex("green", keep_alpha=False)
_DREAM_OUTPUT_STRUCTURE_FACTOR_COLOR = to_hex(
    "tab:purple",
    keep_alpha=False,
)

ReportProgressCallback = Callable[[int, int, str], None]


@dataclass(frozen=True, slots=True)
class ReportComponentSeries:
    label: str
    q_values: np.ndarray
    intensities: np.ndarray
    color: str


@dataclass(frozen=True, slots=True)
class ReportComponentPlotData:
    title: str
    selected_q_min: float | None
    selected_q_max: float | None
    use_experimental_grid: bool
    log_x: bool
    log_y: bool
    experimental_q_values: np.ndarray | None
    experimental_intensities: np.ndarray | None
    solvent_q_values: np.ndarray | None
    solvent_intensities: np.ndarray | None
    component_series: tuple[ReportComponentSeries, ...]


@dataclass(frozen=True, slots=True)
class PriorHistogramRequest:
    title: str
    json_path: Path
    mode: str
    cmap: str
    secondary_element: str | None = None


@dataclass(frozen=True, slots=True)
class DreamFilterReportView:
    title: str
    description: str
    filter_mode: str
    is_active: bool
    summary: DreamSummary
    model_plot: DreamModelPlotData
    violin_plot: DreamViolinPlotData
    violin_payload: dict[str, object]
    weights_violin_payload: dict[str, object]
    effective_radii_violin_payload: dict[str, object]


@dataclass(frozen=True, slots=True)
class DreamModelReportContext:
    output_path: Path
    asset_dir: Path
    project_name: str
    project_dir: Path
    generated_at: datetime
    powerpoint_settings: PowerPointExportSettings
    user_q_range_text: str
    supported_q_range_text: str | None
    q_sampling_text: str
    template_name: str
    template_display_name: str
    template_module_path: Path | None
    model_equation_text: str | None
    model_context_lines: tuple[str, ...]
    model_definition_lines: tuple[str, ...]
    model_reference_lines: tuple[str, ...]
    prior_histograms: tuple[PriorHistogramRequest, ...]
    component_plot_without_solvent: ReportComponentPlotData | None
    component_plot_with_solvent: ReportComponentPlotData | None
    prefit_evaluation: PrefitEvaluation | None
    prefit_parameter_entries: tuple[PrefitParameterEntry, ...]
    prefit_statistics: dict[str, object]
    cluster_geometry_rows: tuple[ClusterGeometryMetadataRow, ...]
    solution_scattering_estimate: SolutionScatteringEstimate | None
    dream_settings: DreamRunSettings
    dream_summary: DreamSummary
    dream_model_plot: DreamModelPlotData
    dream_violin_plot: DreamViolinPlotData
    dream_violin_payload: dict[str, object]
    dream_parameter_map_entries: tuple[DreamParameterEntry, ...]
    dream_filter_assessments: tuple[dict[str, object], ...]
    dream_filter_views: tuple[DreamFilterReportView, ...]
    output_summary_lines: tuple[str, ...]
    directory_lines: tuple[str, ...]


@dataclass(frozen=True, slots=True)
class ModelReportExportResult:
    report_path: Path
    manifest_path: Path | None
    figure_paths: tuple[Path, ...]


class _ReportProgressTracker:
    def __init__(
        self,
        total_steps: int,
        callback: ReportProgressCallback | None,
    ) -> None:
        self.total_steps = max(int(total_steps), 1)
        self._callback = callback
        self._processed = 0
        self.emit("Generating DREAM model report PowerPoint. Please wait...")

    def emit(self, message: str) -> None:
        if self._callback is None:
            return
        self._callback(self._processed, self.total_steps, str(message))

    def advance(self, message: str) -> None:
        self._processed = min(self._processed + 1, self.total_steps)
        self.emit(message)


def export_dream_model_report_pptx(
    context: DreamModelReportContext,
    *,
    progress_callback: ReportProgressCallback | None = None,
) -> ModelReportExportResult:
    Presentation, Inches, Pt = _load_pptx_api()
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN

    export_settings = PowerPointExportSettings.from_dict(
        context.powerpoint_settings.to_dict()
    )
    context.output_path.parent.mkdir(parents=True, exist_ok=True)
    temporary_figure_dir: TemporaryDirectory[str] | None = None
    if (
        export_settings.generate_manifest
        or export_settings.export_figure_assets
    ):
        context.asset_dir.mkdir(parents=True, exist_ok=True)
    if export_settings.export_figure_assets:
        figure_dir = context.asset_dir / "figures"
        figure_dir.mkdir(parents=True, exist_ok=True)
    else:
        temporary_figure_dir = TemporaryDirectory()
        figure_dir = Path(temporary_figure_dir.name)

    prefit_parameter_rows = _prefit_parameter_rows(
        context.prefit_parameter_entries
    )
    geometry_parameter_rows = _cluster_geometry_rows(
        context.cluster_geometry_rows
    )
    dream_prior_rows = _dream_prior_rows(context.dream_parameter_map_entries)
    prior_histogram_pages = _chunked(context.prior_histograms, 4)
    prefit_table_pages = _table_row_chunks(prefit_parameter_rows, 12)
    geometry_table_pages = _table_row_chunks(geometry_parameter_rows, 11)
    dream_prior_table_pages = _table_row_chunks(dream_prior_rows, 11)
    prefit_summary_pages = (
        _paginate_text_lines(
            _prefit_summary_lines(
                context.prefit_evaluation,
                context.prefit_statistics,
            ),
            max_lines=_TEXT_LINES_SIDE,
            wrap_at=_TEXT_WRAP_SIDE,
        )
        if context.prefit_evaluation is not None
        else []
    )
    estimator_pages = (
        _paginate_text_lines(
            _solution_estimate_lines(context.solution_scattering_estimate),
            max_lines=_TEXT_LINES_FULL,
            wrap_at=_TEXT_WRAP_FULL,
        )
        if context.solution_scattering_estimate is not None
        else []
    )
    dream_settings_pages = _paginate_text_lines(
        _dream_settings_lines(
            context.dream_settings,
            context.dream_summary,
        ),
        max_lines=_TEXT_LINES_HALF,
        wrap_at=_TEXT_WRAP_HALF,
    )
    dream_assessment_pages = _paginate_text_lines(
        _dream_assessment_lines(
            context.dream_filter_assessments,
            context.dream_settings,
        ),
        max_lines=_TEXT_LINES_HALF,
        wrap_at=_TEXT_WRAP_HALF,
    )
    dream_output_pages = _paginate_text_lines(
        _dream_output_lines(
            context.dream_settings,
            context.dream_summary,
            context.dream_model_plot,
        ),
        max_lines=_TEXT_LINES_SIDE,
        wrap_at=_TEXT_WRAP_SIDE,
    )
    report_summary_pages = _paginate_text_lines(
        list(context.output_summary_lines),
        max_lines=_TEXT_LINES_HALF,
        wrap_at=_TEXT_WRAP_HALF,
    )
    directory_pages = _paginate_text_lines(
        list(context.directory_lines),
        max_lines=_TEXT_LINES_HALF,
        wrap_at=_TEXT_WRAP_HALF,
    )
    model_detail_lines: list[str] = []
    if context.model_equation_text:
        model_detail_lines.extend(
            [
                "Model equation:",
                context.model_equation_text,
            ]
        )
    if context.model_definition_lines:
        if model_detail_lines:
            model_detail_lines.append("")
        model_detail_lines.append("Term definitions:")
        model_detail_lines.extend(context.model_definition_lines)
    if context.model_reference_lines:
        if model_detail_lines:
            model_detail_lines.append("")
        model_detail_lines.append("References:")
        model_detail_lines.extend(context.model_reference_lines)
    model_context_pages = _paginate_text_lines(
        list(context.model_context_lines),
        max_lines=_TEXT_LINES_HALF,
        wrap_at=_TEXT_WRAP_HALF,
    )
    model_detail_pages = _paginate_text_lines(
        model_detail_lines,
        max_lines=_TEXT_LINES_HALF,
        wrap_at=_TEXT_WRAP_HALF,
    )
    prefit_has_solvent_trace = (
        context.prefit_evaluation is not None
        and _has_prefit_solvent_trace(context.prefit_evaluation)
    )

    total_slides = 1
    if context.model_context_lines or model_detail_lines:
        total_slides += max(len(model_context_pages), len(model_detail_pages))
    if export_settings.include_prior_histograms and prior_histogram_pages:
        total_slides += len(prior_histogram_pages)
    if export_settings.include_initial_traces and (
        context.component_plot_without_solvent is not None
        or context.component_plot_with_solvent is not None
    ):
        total_slides += 1
    if (
        export_settings.include_prefit_model
        and context.prefit_evaluation is not None
    ):
        total_slides += len(prefit_summary_pages)
        total_slides += int(prefit_has_solvent_trace)
    if export_settings.include_prefit_parameters:
        total_slides += len(prefit_table_pages)
    if (
        export_settings.include_geometry_table
        and context.cluster_geometry_rows
    ):
        total_slides += len(geometry_table_pages)
    if (
        export_settings.include_estimator_metrics
        and context.solution_scattering_estimate is not None
    ):
        total_slides += len(estimator_pages)
    if export_settings.include_dream_settings:
        total_slides += max(
            len(dream_settings_pages),
            len(dream_assessment_pages),
        )
    if export_settings.include_dream_prior_table:
        total_slides += len(dream_prior_table_pages)
    if export_settings.include_dream_output_model:
        total_slides += len(dream_output_pages)
    if (
        export_settings.include_posterior_comparisons
        and context.dream_filter_views
    ):
        total_slides += 4
    if (
        export_settings.include_output_summary
        and export_settings.include_directory_summary
    ):
        total_slides += max(len(report_summary_pages), len(directory_pages))
    elif export_settings.include_output_summary:
        total_slides += len(report_summary_pages)
    elif export_settings.include_directory_summary:
        total_slides += len(directory_pages)

    progress = _ReportProgressTracker(
        _count_report_figures(context)
        + total_slides
        + 1
        + int(export_settings.generate_manifest),
        progress_callback,
    )

    figure_paths: list[Path] = []
    rendered_figures: dict[str, Path] = {}

    if export_settings.include_prior_histograms:
        for index, request in enumerate(context.prior_histograms, start=1):
            figure_path = figure_dir / f"{_slugify(request.title)}.png"
            rendered_figures[request.title] = _render_prior_histogram(
                request,
                figure_path,
                settings=export_settings,
            )
            if export_settings.export_figure_assets:
                figure_paths.append(figure_path)
            progress.advance(
                "Rendered prior histogram "
                f"{index}/{len(context.prior_histograms)}."
            )

    if (
        export_settings.include_initial_traces
        and context.component_plot_without_solvent is not None
    ):
        figure_path = figure_dir / "initial_traces_no_solvent.png"
        rendered_figures["initial_traces_no_solvent"] = _render_component_plot(
            context.component_plot_without_solvent,
            figure_path,
            settings=export_settings,
        )
        if export_settings.export_figure_assets:
            figure_paths.append(figure_path)
        progress.advance("Rendered initial SAXS traces without solvent.")

    if (
        export_settings.include_initial_traces
        and context.component_plot_with_solvent is not None
    ):
        figure_path = figure_dir / "initial_traces_with_solvent.png"
        rendered_figures["initial_traces_with_solvent"] = (
            _render_component_plot(
                context.component_plot_with_solvent,
                figure_path,
                settings=export_settings,
            )
        )
        if export_settings.export_figure_assets:
            figure_paths.append(figure_path)
        progress.advance("Rendered initial SAXS traces with solvent.")

    if (
        export_settings.include_prefit_model
        and context.prefit_evaluation is not None
    ):
        figure_path = figure_dir / "prefit_model_without_solvent.png"
        rendered_figures["prefit_model_without_solvent"] = _render_prefit_plot(
            context.prefit_evaluation,
            context.prefit_statistics,
            figure_path,
            settings=export_settings,
            include_solvent=False,
        )
        if export_settings.export_figure_assets:
            figure_paths.append(figure_path)
        progress.advance("Rendered prefit model plot without solvent trace.")

        if prefit_has_solvent_trace:
            figure_path = figure_dir / "prefit_model_with_solvent.png"
            rendered_figures["prefit_model_with_solvent"] = (
                _render_prefit_plot(
                    context.prefit_evaluation,
                    context.prefit_statistics,
                    figure_path,
                    settings=export_settings,
                    include_solvent=True,
                )
            )
            if export_settings.export_figure_assets:
                figure_paths.append(figure_path)
            progress.advance("Rendered prefit model plot with solvent trace.")

    if export_settings.include_dream_output_model:
        figure_path = figure_dir / "dream_model.png"
        rendered_figures["dream_model"] = _render_dream_model_plot(
            context.dream_model_plot,
            figure_path,
            settings=export_settings,
        )
        if export_settings.export_figure_assets:
            figure_paths.append(figure_path)
        progress.advance("Rendered DREAM output model plot.")

    if (
        export_settings.include_posterior_comparisons
        and context.dream_filter_views
    ):
        figure_path = figure_dir / "dream_filter_violin_comparison.png"
        rendered_figures["dream_filter_violin_comparison"] = (
            _render_filter_violin_comparison(
                context.dream_filter_views,
                figure_path,
                settings=export_settings,
            )
        )
        if export_settings.export_figure_assets:
            figure_paths.append(figure_path)
        progress.advance("Rendered posterior violin comparison plot.")

        figure_path = figure_dir / "dream_filter_violin_comparison_weights.png"
        rendered_figures["dream_filter_violin_comparison_weights"] = (
            _render_filter_violin_comparison(
                context.dream_filter_views,
                figure_path,
                settings=export_settings,
                payload_variant="weights",
            )
        )
        if export_settings.export_figure_assets:
            figure_paths.append(figure_path)
        progress.advance(
            "Rendered posterior violin comparison plot for weight parameters."
        )

        figure_path = (
            figure_dir / "dream_filter_violin_comparison_effective_radii.png"
        )
        rendered_figures["dream_filter_violin_comparison_effective_radii"] = (
            _render_filter_violin_comparison(
                context.dream_filter_views,
                figure_path,
                settings=export_settings,
                payload_variant="effective_radii",
            )
        )
        if export_settings.export_figure_assets:
            figure_paths.append(figure_path)
        progress.advance(
            "Rendered posterior violin comparison plot for effective radii."
        )

        figure_path = figure_dir / "dream_filter_fit_comparison.png"
        rendered_figures["dream_filter_fit_comparison"] = (
            _render_filter_fit_comparison(
                context.dream_filter_views,
                figure_path,
                settings=export_settings,
            )
        )
        if export_settings.export_figure_assets:
            figure_paths.append(figure_path)
        progress.advance("Rendered posterior fit comparison plot.")

    manifest_path: Path | None = None
    if export_settings.generate_manifest:
        context.asset_dir.mkdir(parents=True, exist_ok=True)
        manifest_path = context.asset_dir / "report_manifest.json"
        manifest_path.write_text(
            json.dumps(
                _manifest_payload(
                    context,
                    figure_paths=figure_paths,
                ),
                indent=2,
            )
            + "\n",
            encoding="utf-8",
        )
        progress.advance("Wrote report manifest.")

    presentation = Presentation()
    presentation.slide_width = Inches(_SLIDE_WIDTH_INCHES)
    presentation.slide_height = Inches(_SLIDE_HEIGHT_INCHES)
    blank_layout = presentation.slide_layouts[6]
    slide_index = 0

    def first_run(paragraph):
        return paragraph.runs[0] if paragraph.runs else paragraph.add_run()

    def apply_run_style(
        run,
        *,
        font_size: float,
        color: str | None = None,
        bold: bool = False,
    ) -> None:
        run.font.name = export_settings.font_family
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = _rgb_color(color or export_settings.text_color)

    def add_title(slide, title: str, subtitle: str | None = None) -> None:
        title_box = slide.shapes.add_textbox(
            Inches(_SLIDE_LEFT_INCHES),
            Inches(0.25),
            Inches(_SLIDE_CONTENT_WIDTH_INCHES),
            Inches(0.52),
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = title
        title_paragraph.space_after = Pt(0)
        apply_run_style(first_run(title_paragraph), font_size=24, bold=True)
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(_SLIDE_LEFT_INCHES),
                Inches(0.72),
                Inches(_SLIDE_CONTENT_WIDTH_INCHES),
                Inches(0.3),
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()
            subtitle_frame.word_wrap = True
            subtitle_paragraph = subtitle_frame.paragraphs[0]
            subtitle_paragraph.text = subtitle
            subtitle_paragraph.space_after = Pt(0)
            apply_run_style(
                first_run(subtitle_paragraph),
                font_size=11,
                color=export_settings.text_color,
            )

    def add_text_block(
        slide,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        lines: list[str],
        font_size: float = 13,
        bold_first: bool = False,
        align=PP_ALIGN.LEFT,
    ) -> None:
        box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        if not lines:
            lines = [""]
        for index, line in enumerate(lines):
            paragraph = (
                frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            )
            paragraph.text = str(line)
            paragraph.alignment = align
            paragraph.space_after = Pt(0)
            paragraph.space_before = Pt(0)
            apply_run_style(
                first_run(paragraph),
                font_size=font_size,
                bold=bool(bold_first and index == 0),
            )

    def add_picture(
        slide,
        image_path: Path,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        fitted_left, fitted_top, fitted_width, fitted_height = (
            _fit_image_in_box(
                image_path,
                left=left,
                top=top,
                max_width=width,
                max_height=height,
            )
        )
        slide.shapes.add_picture(
            str(image_path),
            Inches(fitted_left),
            Inches(fitted_top),
            width=Inches(fitted_width),
            height=Inches(fitted_height),
        )

    def add_table_header_rule(
        slide,
        *,
        left: float,
        top: float,
        width: float,
    ) -> None:
        rule = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(_THICK_RULE_HEIGHT_INCHES),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = _rgb_color(export_settings.table_rule_color)
        rule.line.fill.background()

    def style_table_cell(
        cell,
        *,
        text: str,
        font_size: float,
        fill_color: str,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
    ) -> None:
        cell.text = str(text)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb_color(fill_color)
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        frame = cell.text_frame
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)
        apply_run_style(first_run(paragraph), font_size=font_size, bold=bold)

    def register_slide(message: str) -> None:
        nonlocal slide_index
        slide_index += 1
        progress.advance(
            f"Built slide {slide_index}/{total_slides}: {message}"
        )

    def add_table_slides(
        title: str,
        columns: list[str],
        rows: list[list[str]],
        *,
        rows_per_slide: int = 12,
        subtitle: str | None = None,
        note: str | None = None,
        column_width_weights: Sequence[float] | None = None,
    ) -> None:
        if not rows:
            slide = presentation.slides.add_slide(blank_layout)
            add_title(slide, title, subtitle)
            add_text_block(
                slide,
                left=_SLIDE_LEFT_INCHES,
                top=1.3,
                width=_SLIDE_CONTENT_WIDTH_INCHES,
                height=1.0,
                lines=[note or "No data are available for this section."],
                font_size=14,
            )
            register_slide(title)
            return
        chunks = _table_row_chunks(rows, rows_per_slide)
        for chunk_index, chunk in enumerate(chunks, start=1):
            slide = presentation.slides.add_slide(blank_layout)
            effective_title = _page_title(title, chunk_index, len(chunks))
            add_title(slide, effective_title, subtitle)
            table_height = (
                _TABLE_HEIGHT_INCHES - 0.28 if note else _TABLE_HEIGHT_INCHES
            )
            add_table_header_rule(
                slide,
                left=_SLIDE_LEFT_INCHES,
                top=_TABLE_TOP_INCHES,
                width=_TABLE_WIDTH_INCHES,
            )
            table_shape = slide.shapes.add_table(
                len(chunk) + 1,
                len(columns),
                Inches(_SLIDE_LEFT_INCHES),
                Inches(_TABLE_TOP_INCHES),
                Inches(_TABLE_WIDTH_INCHES),
                Inches(table_height),
            )
            table = table_shape.table
            column_widths = _resolve_column_widths(
                columns,
                chunk,
                total_width=_TABLE_WIDTH_INCHES,
                column_width_weights=column_width_weights,
            )
            for column_index, column_width in enumerate(column_widths):
                table.columns[column_index].width = Inches(column_width)
            row_height = table_height / max(len(chunk) + 1, 1)
            for row_index in range(len(chunk) + 1):
                table.rows[row_index].height = Inches(row_height)
            for column_index, column_name in enumerate(columns):
                style_table_cell(
                    table.cell(0, column_index),
                    text=column_name,
                    font_size=11,
                    fill_color=export_settings.table_header_fill,
                    bold=True,
                    align=PP_ALIGN.CENTER,
                )
            for row_index, row in enumerate(chunk, start=1):
                row_fill = (
                    export_settings.table_even_row_fill
                    if row_index % 2 == 1
                    else export_settings.table_odd_row_fill
                )
                for column_index, value in enumerate(row):
                    style_table_cell(
                        table.cell(row_index, column_index),
                        text=str(value),
                        font_size=10,
                        fill_color=row_fill,
                    )
            if note and chunk_index == len(chunks):
                add_text_block(
                    slide,
                    left=_SLIDE_LEFT_INCHES,
                    top=_TABLE_NOTE_TOP_INCHES,
                    width=_SLIDE_CONTENT_WIDTH_INCHES,
                    height=0.24,
                    lines=[note],
                    font_size=10,
                )
            register_slide(effective_title)

    def add_full_width_text_slides(
        title: str,
        subtitle: str | None,
        pages: list[list[str]],
        *,
        font_size: float = 13,
    ) -> None:
        total_pages = len(pages)
        for page_index, page_lines in enumerate(pages, start=1):
            slide = presentation.slides.add_slide(blank_layout)
            effective_title = _page_title(title, page_index, total_pages)
            add_title(slide, effective_title, subtitle)
            add_text_block(
                slide,
                left=_SLIDE_LEFT_INCHES,
                top=1.25,
                width=_SLIDE_CONTENT_WIDTH_INCHES,
                height=_SLIDE_CONTENT_HEIGHT_INCHES,
                lines=page_lines,
                font_size=font_size,
            )
            register_slide(effective_title)

    def add_two_column_text_slides(
        title: str,
        subtitle: str | None,
        left_pages: list[list[str]],
        right_pages: list[list[str]],
    ) -> None:
        total_pages = max(len(left_pages), len(right_pages))
        for page_index in range(total_pages):
            slide = presentation.slides.add_slide(blank_layout)
            effective_title = _page_title(title, page_index + 1, total_pages)
            add_title(slide, effective_title, subtitle)
            add_text_block(
                slide,
                left=0.55,
                top=1.12,
                width=5.85,
                height=5.82,
                lines=(
                    left_pages[page_index]
                    if page_index < len(left_pages)
                    else []
                ),
                font_size=11.5,
            )
            add_text_block(
                slide,
                left=6.7,
                top=1.12,
                width=5.78,
                height=5.82,
                lines=(
                    right_pages[page_index]
                    if page_index < len(right_pages)
                    else []
                ),
                font_size=11.5,
            )
            register_slide(effective_title)

    def add_picture_with_summary_slides(
        title: str,
        subtitle: str | None,
        image_path: Path,
        summary_pages: list[list[str]],
    ) -> None:
        total_pages = len(summary_pages)
        for page_index, page_lines in enumerate(summary_pages, start=1):
            slide = presentation.slides.add_slide(blank_layout)
            effective_title = _page_title(title, page_index, total_pages)
            add_title(slide, effective_title, subtitle)
            if page_index == 1:
                add_picture(
                    slide,
                    image_path,
                    left=0.45,
                    top=1.16,
                    width=7.55,
                    height=5.5,
                )
                add_text_block(
                    slide,
                    left=8.2,
                    top=1.2,
                    width=4.22,
                    height=5.26,
                    lines=page_lines,
                    font_size=11.5,
                )
            else:
                add_text_block(
                    slide,
                    left=_SLIDE_LEFT_INCHES,
                    top=1.25,
                    width=_SLIDE_CONTENT_WIDTH_INCHES,
                    height=_SLIDE_CONTENT_HEIGHT_INCHES,
                    lines=page_lines,
                    font_size=13,
                )
            register_slide(effective_title)

    def add_full_width_picture_slide(
        title: str,
        subtitle: str | None,
        image_path: Path,
    ) -> None:
        slide = presentation.slides.add_slide(blank_layout)
        add_title(slide, title, subtitle)
        add_picture(
            slide,
            image_path,
            left=0.45,
            top=1.12,
            width=12.43,
            height=5.58,
        )
        register_slide(title)

    slide = presentation.slides.add_slide(blank_layout)
    add_title(slide, "SAXS Model Report", "DREAM fit export")
    add_text_block(
        slide,
        left=0.6,
        top=1.25,
        width=_SLIDE_CONTENT_WIDTH_INCHES,
        height=_SLIDE_CONTENT_HEIGHT_INCHES,
        lines=_paginate_text_lines(
            [
                f"Project: {context.project_name}",
                f"Project directory: {context.project_dir}",
                (
                    "Generated: "
                    f"{context.generated_at.strftime('%Y-%m-%d %H:%M:%S')}"
                ),
                f"Template: {context.template_name}",
                f"User selected q-range: {context.user_q_range_text}",
                (
                    "Supported component q-range: "
                    f"{context.supported_q_range_text or 'Unavailable'}"
                ),
                f"q-grid: {context.q_sampling_text}",
                f"DREAM run directory: {context.dream_summary.run_dir}",
                (
                    "Posterior filter: "
                    f"{_describe_posterior_filter(context.dream_settings)}"
                ),
                (
                    "Posterior samples kept: "
                    f"{context.dream_summary.posterior_sample_count}"
                ),
            ],
            max_lines=_TEXT_LINES_FULL,
            wrap_at=_TEXT_WRAP_FULL,
        )[0],
        font_size=15,
    )
    register_slide("SAXS Model Report")

    if context.model_context_lines or model_detail_lines:
        add_two_column_text_slides(
            "Model Information",
            (f"{context.template_display_name} " f"({context.template_name})"),
            model_context_pages,
            model_detail_pages,
        )

    if export_settings.include_prior_histograms:
        for page_index, request_page in enumerate(
            prior_histogram_pages,
            start=1,
        ):
            slide = presentation.slides.add_slide(blank_layout)
            effective_title = _page_title(
                "Prior Histograms",
                page_index,
                len(prior_histogram_pages),
            )
            add_title(
                slide,
                effective_title,
                "Configured palettes for regular and solvent-sort views.",
            )
            grid_positions = [
                (0.45, 1.08, 5.86, 2.52),
                (6.97, 1.08, 5.86, 2.52),
                (0.45, 4.0, 5.86, 2.52),
                (6.97, 4.0, 5.86, 2.52),
            ]
            for request, (left, top, width, height) in zip(
                request_page,
                grid_positions,
                strict=False,
            ):
                figure_path = rendered_figures.get(request.title)
                if figure_path is None:
                    continue
                add_picture(
                    slide,
                    figure_path,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                )
            register_slide(effective_title)

    if export_settings.include_initial_traces and (
        context.component_plot_without_solvent is not None
        or context.component_plot_with_solvent is not None
    ):
        slide = presentation.slides.add_slide(blank_layout)
        add_title(
            slide,
            "Initial SAXS Traces",
            "Dual-axis rescaled views of the selected q-range.",
        )
        component_figure_keys = [
            key
            for key in (
                "initial_traces_no_solvent",
                "initial_traces_with_solvent",
            )
            if key in rendered_figures
        ]
        if len(component_figure_keys) == 1:
            add_picture(
                slide,
                rendered_figures[component_figure_keys[0]],
                left=0.65,
                top=1.2,
                width=12.0,
                height=5.35,
            )
        else:
            if "initial_traces_no_solvent" in rendered_figures:
                add_picture(
                    slide,
                    rendered_figures["initial_traces_no_solvent"],
                    left=0.45,
                    top=1.2,
                    width=5.88,
                    height=5.25,
                )
            if "initial_traces_with_solvent" in rendered_figures:
                add_picture(
                    slide,
                    rendered_figures["initial_traces_with_solvent"],
                    left=6.95,
                    top=1.2,
                    width=5.88,
                    height=5.25,
                )
        register_slide("Initial SAXS Traces")

    if (
        export_settings.include_prefit_model
        and context.prefit_evaluation is not None
        and "prefit_model_without_solvent" in rendered_figures
    ):
        add_picture_with_summary_slides(
            "Prefit Model",
            "Prefit window default view without the solvent trace.",
            rendered_figures["prefit_model_without_solvent"],
            prefit_summary_pages,
        )
        if "prefit_model_with_solvent" in rendered_figures:
            add_full_width_picture_slide(
                "Prefit Model With Solvent",
                "Same prefit fit with the solvent contribution trace enabled.",
                rendered_figures["prefit_model_with_solvent"],
            )

    if export_settings.include_prefit_parameters:
        add_table_slides(
            "Prefit Parameters",
            [
                "Parameter",
                "Category",
                "Value",
                "Vary",
                "Min",
                "Max",
                "Structure",
                "Motif",
            ],
            prefit_parameter_rows,
            rows_per_slide=12,
            column_width_weights=[
                1.65,
                1.2,
                0.8,
                0.72,
                0.82,
                0.82,
                1.05,
                1.05,
            ],
        )

    if (
        export_settings.include_geometry_table
        and context.cluster_geometry_rows
    ):
        add_table_slides(
            "Computed Geometry Parameters",
            [
                "Cluster",
                "Mapped Parameter",
                "Approx.",
                "Eff. Radius",
                "Rg",
                "Max Radius",
                "Anisotropy",
                "Axes (a/b/c)",
            ],
            geometry_parameter_rows,
            rows_per_slide=11,
            note=(
                "Geometry metrics come from the prefit estimator table "
                "saved with the active project."
            ),
            column_width_weights=[1.0, 1.6, 0.8, 0.9, 0.8, 0.95, 0.95, 1.7],
        )

    if (
        export_settings.include_estimator_metrics
        and context.solution_scattering_estimate is not None
    ):
        add_full_width_text_slides(
            "Estimator Metrics",
            "Volume fraction, attenuation, fluorescence, and number density.",
            estimator_pages,
            font_size=12.5,
        )

    if export_settings.include_dream_settings:
        add_two_column_text_slides(
            "DREAM Settings",
            "Sampler settings, posterior filtering, and active DREAM summary.",
            dream_settings_pages,
            dream_assessment_pages,
        )

    if export_settings.include_dream_prior_table:
        add_table_slides(
            "DREAM Prior Distributions",
            [
                "Parameter",
                "Type",
                "Structure",
                "Motif",
                "Value",
                "Vary",
                "Distribution",
                "Distribution Params",
            ],
            dream_prior_rows,
            rows_per_slide=11,
            column_width_weights=[1.2, 1.0, 0.9, 0.9, 0.78, 0.72, 1.08, 2.3],
        )

    if (
        export_settings.include_dream_output_model
        and "dream_model" in rendered_figures
    ):
        add_picture_with_summary_slides(
            "DREAM Output Model",
            "Best-fit model and posterior summary statistics.",
            rendered_figures["dream_model"],
            dream_output_pages,
        )

    if (
        export_settings.include_posterior_comparisons
        and "dream_filter_violin_comparison" in rendered_figures
    ):
        slide = presentation.slides.add_slide(blank_layout)
        add_title(
            slide,
            "Posterior Violin Comparison",
            "Default posterior filters with the active selection labeled.",
        )
        add_picture(
            slide,
            rendered_figures["dream_filter_violin_comparison"],
            left=0.45,
            top=1.12,
            width=12.43,
            height=5.58,
        )
        register_slide("Posterior Violin Comparison")

    if (
        export_settings.include_posterior_comparisons
        and "dream_filter_violin_comparison_weights" in rendered_figures
    ):
        slide = presentation.slides.add_slide(blank_layout)
        add_title(
            slide,
            "Posterior Violin Comparison - Weights",
            "Weight parameters (w##) shown on a dedicated y-axis scale.",
        )
        add_picture(
            slide,
            rendered_figures["dream_filter_violin_comparison_weights"],
            left=0.45,
            top=1.12,
            width=12.43,
            height=5.58,
        )
        register_slide("Posterior Violin Comparison - Weights")

    if (
        export_settings.include_posterior_comparisons
        and "dream_filter_violin_comparison_effective_radii"
        in rendered_figures
    ):
        slide = presentation.slides.add_slide(blank_layout)
        add_title(
            slide,
            "Posterior Violin Comparison - Effective Radii",
            "Effective-radius parameters shown on a dedicated y-axis scale.",
        )
        add_picture(
            slide,
            rendered_figures["dream_filter_violin_comparison_effective_radii"],
            left=0.45,
            top=1.12,
            width=12.43,
            height=5.58,
        )
        register_slide("Posterior Violin Comparison - Effective Radii")

    if (
        export_settings.include_posterior_comparisons
        and "dream_filter_fit_comparison" in rendered_figures
    ):
        slide = presentation.slides.add_slide(blank_layout)
        add_title(
            slide,
            "Filter Fit Comparison",
            "Corresponding fits and metrics for each posterior filter view.",
        )
        add_picture(
            slide,
            rendered_figures["dream_filter_fit_comparison"],
            left=0.45,
            top=1.12,
            width=12.43,
            height=5.58,
        )
        register_slide("Filter Fit Comparison")

    if (
        export_settings.include_output_summary
        and export_settings.include_directory_summary
    ):
        add_two_column_text_slides(
            "Output Summary",
            "Where to find the exported report and related figure data.",
            report_summary_pages,
            directory_pages,
        )
    elif export_settings.include_output_summary:
        add_full_width_text_slides(
            "Output Summary",
            "Summary information for the exported report.",
            report_summary_pages,
        )
    elif export_settings.include_directory_summary:
        add_full_width_text_slides(
            "Output Directories",
            "Project and report paths related to this export.",
            directory_pages,
        )

    presentation.save(str(context.output_path))
    progress.advance("Saved PowerPoint report.")
    result = ModelReportExportResult(
        report_path=context.output_path,
        manifest_path=manifest_path,
        figure_paths=tuple(figure_paths),
    )
    if temporary_figure_dir is not None:
        temporary_figure_dir.cleanup()
    return result


def _load_pptx_api():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError(
            "PowerPoint export requires the optional dependency "
            "`python-pptx`. Install it and retry."
        ) from exc
    return Presentation, Inches, Pt


def _count_report_figures(context: DreamModelReportContext) -> int:
    settings = PowerPointExportSettings.from_dict(
        context.powerpoint_settings.to_dict()
    )
    return (
        (
            len(context.prior_histograms)
            if settings.include_prior_histograms
            else 0
        )
        + (
            int(context.component_plot_without_solvent is not None)
            + int(context.component_plot_with_solvent is not None)
            if settings.include_initial_traces
            else 0
        )
        + (
            int(context.prefit_evaluation is not None)
            + int(
                context.prefit_evaluation is not None
                and _has_prefit_solvent_trace(context.prefit_evaluation)
            )
            if settings.include_prefit_model
            else 0
        )
        + int(settings.include_dream_output_model)
        + (
            4
            if settings.include_posterior_comparisons
            and context.dream_filter_views
            else 0
        )
    )


def _chunked(values: Sequence[object], size: int) -> list[list[object]]:
    if size <= 0:
        raise ValueError("chunk size must be positive")
    if not values:
        return []
    return [
        list(values[index : index + size])
        for index in range(0, len(values), size)
    ]


def _table_row_chunks(
    rows: Sequence[list[str]],
    rows_per_slide: int,
) -> list[list[list[str]]]:
    if rows_per_slide <= 0:
        raise ValueError("rows_per_slide must be positive")
    if not rows:
        return [[]]
    return [
        list(rows[index : index + rows_per_slide])
        for index in range(0, len(rows), rows_per_slide)
    ]


def _paginate_text_lines(
    lines: Sequence[str],
    *,
    max_lines: int,
    wrap_at: int,
) -> list[list[str]]:
    if max_lines <= 0:
        raise ValueError("max_lines must be positive")
    wrapped_lines: list[str] = []
    for raw_line in lines:
        line = str(raw_line).strip()
        if not line:
            wrapped_lines.append("")
            continue
        wrapped_lines.extend(
            textwrap.wrap(
                line,
                width=max(int(wrap_at), 1),
                break_long_words=True,
                break_on_hyphens=False,
            )
            or [line]
        )
    if not wrapped_lines:
        return [[]]
    return [
        wrapped_lines[index : index + max_lines]
        for index in range(0, len(wrapped_lines), max_lines)
    ]


def _page_title(title: str, page_index: int, total_pages: int) -> str:
    if total_pages <= 1:
        return title
    return f"{title} ({page_index}/{total_pages})"


def _resolve_column_widths(
    columns: Sequence[str],
    rows: Sequence[Sequence[str]],
    *,
    total_width: float,
    column_width_weights: Sequence[float] | None = None,
) -> list[float]:
    if not columns:
        return []
    if column_width_weights is not None:
        weights = [max(float(weight), 0.2) for weight in column_width_weights]
    else:
        weights = []
        for column_index, column_name in enumerate(columns):
            max_length = len(str(column_name))
            for row in rows:
                if column_index >= len(row):
                    continue
                max_length = max(max_length, len(str(row[column_index])))
            weights.append(max(0.8, min(float(max_length) ** 0.78, 3.2)))
    if len(weights) != len(columns):
        raise ValueError("column width weights must match the column count")
    scale = total_width / sum(weights)
    widths = [weight * scale for weight in weights]
    widths[-1] += total_width - sum(widths)
    return widths


def _fit_image_in_box(
    image_path: Path,
    *,
    left: float,
    top: float,
    max_width: float,
    max_height: float,
) -> tuple[float, float, float, float]:
    try:
        from PIL import Image
    except ImportError:
        return left, top, max_width, max_height

    with Image.open(image_path) as image:
        width_px, height_px = image.size
    if width_px <= 0 or height_px <= 0:
        return left, top, max_width, max_height
    image_ratio = width_px / height_px
    box_ratio = max_width / max_height if max_height > 0 else image_ratio
    if image_ratio >= box_ratio:
        fitted_width = max_width
        fitted_height = fitted_width / image_ratio
    else:
        fitted_height = max_height
        fitted_width = fitted_height * image_ratio
    fitted_left = left + (max_width - fitted_width) / 2.0
    fitted_top = top + (max_height - fitted_height) / 2.0
    return fitted_left, fitted_top, fitted_width, fitted_height


def _report_figure_context(
    settings: PowerPointExportSettings,
    *,
    compact: bool = False,
):
    return rc_context(
        {
            "font.family": settings.font_family,
            "font.sans-serif": [
                settings.font_family,
                "DejaVu Sans",
                "Liberation Sans",
                "sans-serif",
            ],
            "axes.titlesize": 10 if compact else 12,
            "axes.labelsize": 9 if compact else 10.5,
            "xtick.labelsize": 8 if compact else 9,
            "ytick.labelsize": 8 if compact else 9,
            "legend.fontsize": 7 if compact else 8,
            "figure.facecolor": "white",
            "axes.facecolor": "white",
        }
    )


def _save_figure(fig: Figure, output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(output_path, dpi=240, bbox_inches="tight")
    return output_path


def _render_prior_histogram(
    request: PriorHistogramRequest,
    output_path: Path,
    *,
    settings: PowerPointExportSettings,
) -> Path:
    with _report_figure_context(settings):
        figure = Figure(figsize=(6.8, 3.4))
        axis = figure.add_subplot(111)
        try:
            plot_md_prior_histogram(
                request.json_path,
                mode=request.mode,
                secondary_element=request.secondary_element,
                cmap=request.cmap,
                ax=axis,
            )
        except Exception as exc:
            axis.text(
                0.5,
                0.5,
                f"Unable to render {request.title}.\n{exc}",
                ha="center",
                va="center",
            )
            axis.set_axis_off()
        figure.tight_layout()
    return _save_figure(figure, output_path)


def _render_component_plot(
    plot_data: ReportComponentPlotData,
    output_path: Path,
    *,
    settings: PowerPointExportSettings,
) -> Path:
    with _report_figure_context(settings):
        figure = Figure(figsize=(7.2, 4.4))
        experimental_axis = figure.add_subplot(111)
        component_axis = (
            experimental_axis.twinx() if plot_data.component_series else None
        )
        legend_lines: list[object] = []
        legend_labels: list[str] = []

        if (
            plot_data.experimental_q_values is not None
            and plot_data.experimental_intensities is not None
        ):
            q_values = np.asarray(plot_data.experimental_q_values, dtype=float)
            intensities = np.asarray(
                plot_data.experimental_intensities,
                dtype=float,
            )
            (full_line,) = experimental_axis.plot(
                q_values,
                intensities,
                color=settings.experimental_trace_color,
                alpha=0.32,
                linewidth=1.2,
                label="Experimental data",
            )
            legend_lines.append(full_line)
            legend_labels.append("Experimental data")
            selected_mask = _selected_q_mask(
                q_values,
                plot_data.selected_q_min,
                plot_data.selected_q_max,
                use_experimental_grid=plot_data.use_experimental_grid,
            )
            if np.any(selected_mask) and not np.all(selected_mask):
                (selected_line,) = experimental_axis.plot(
                    q_values[selected_mask],
                    intensities[selected_mask],
                    color=settings.experimental_trace_color,
                    linewidth=1.8,
                    label="Selected q-range",
                )
                legend_lines.append(selected_line)
                legend_labels.append("Selected q-range")

        if (
            plot_data.solvent_q_values is not None
            and plot_data.solvent_intensities is not None
        ):
            q_values = np.asarray(plot_data.solvent_q_values, dtype=float)
            intensities = np.asarray(
                plot_data.solvent_intensities, dtype=float
            )
            (solvent_line,) = experimental_axis.plot(
                q_values,
                intensities,
                color=settings.solvent_trace_color,
                alpha=0.42,
                linewidth=1.2,
                label="Solvent data",
            )
            legend_lines.append(solvent_line)
            legend_labels.append("Solvent data")
            selected_mask = _selected_q_mask(
                q_values,
                plot_data.selected_q_min,
                plot_data.selected_q_max,
                use_experimental_grid=plot_data.use_experimental_grid,
            )
            if np.any(selected_mask) and not np.all(selected_mask):
                (selected_solvent_line,) = experimental_axis.plot(
                    q_values[selected_mask],
                    intensities[selected_mask],
                    color=settings.solvent_trace_color,
                    linewidth=1.8,
                    label="Selected solvent q-range",
                )
                legend_lines.append(selected_solvent_line)
                legend_labels.append("Selected solvent q-range")

        if plot_data.component_series:
            target_axis = component_axis or experimental_axis
            for series in plot_data.component_series:
                (line,) = target_axis.plot(
                    np.asarray(series.q_values, dtype=float),
                    np.asarray(series.intensities, dtype=float),
                    color=series.color,
                    linewidth=1.35,
                    label=series.label,
                )
                legend_lines.append(line)
                legend_labels.append(series.label)

        _apply_axis_scale(
            experimental_axis,
            log_x=plot_data.log_x,
            log_y=plot_data.log_y,
        )
        if component_axis is not None:
            _apply_axis_scale(
                component_axis,
                log_x=plot_data.log_x,
                log_y=plot_data.log_y,
            )
            component_axis.set_ylabel("Component intensity (arb. units)")
            _autoscale_component_plot_to_model_range(
                experimental_axis,
                component_axis,
                plot_data,
            )
        else:
            q_bounds = _component_q_bounds(plot_data)
            if q_bounds is not None:
                experimental_axis.set_xlim(*q_bounds)

        experimental_axis.set_xlabel("q (1/A)")
        experimental_axis.set_ylabel("Intensity (arb. units)")
        experimental_axis.set_title(plot_data.title, fontsize=12)
        experimental_axis.tick_params(labelsize=9)
        if component_axis is not None:
            component_axis.tick_params(labelsize=9)
        if legend_lines:
            columns = max(1, int(np.ceil(len(legend_lines) / 6.0)))
            experimental_axis.legend(
                legend_lines,
                legend_labels,
                fontsize=8,
                loc="upper right",
                ncols=columns,
                framealpha=0.92,
            )
        figure.tight_layout()
    return _save_figure(figure, output_path)


def _render_prefit_plot(
    evaluation: PrefitEvaluation,
    statistics: dict[str, object],
    output_path: Path,
    *,
    settings: PowerPointExportSettings,
    include_solvent: bool,
) -> Path:
    with _report_figure_context(settings):
        figure = Figure(figsize=(7.6, 4.95))
        has_experimental = evaluation.experimental_intensities is not None
        has_residuals = evaluation.residuals is not None
        if has_experimental and has_residuals:
            grid = figure.add_gridspec(2, 1, height_ratios=[3, 1])
            top_axis = figure.add_subplot(grid[0, 0])
            bottom_axis = figure.add_subplot(grid[1, 0], sharex=top_axis)
        else:
            top_axis = figure.add_subplot(111)
            bottom_axis = None
        _draw_prefit_plot_axes(
            top_axis,
            bottom_axis,
            evaluation,
            include_solvent=include_solvent,
            settings=settings,
        )
        figure.tight_layout()
    return _save_figure(figure, output_path)


def _render_dream_model_plot(
    model_plot: DreamModelPlotData,
    output_path: Path,
    *,
    settings: PowerPointExportSettings,
) -> Path:
    with _report_figure_context(settings):
        figure = Figure(figsize=(7.6, 4.95))
        grid = figure.add_gridspec(2, 1, height_ratios=[3, 1])
        top_axis = figure.add_subplot(grid[0, 0])
        bottom_axis = figure.add_subplot(grid[1, 0], sharex=top_axis)
        _draw_model_fit_axis(
            top_axis,
            q_values=np.asarray(model_plot.q_values, dtype=float),
            experimental=np.asarray(
                model_plot.experimental_intensities,
                dtype=float,
            ),
            model=np.asarray(model_plot.model_intensities, dtype=float),
            solvent=(
                None
                if model_plot.solvent_contribution is None
                else np.asarray(model_plot.solvent_contribution, dtype=float)
            ),
            structure_factor=(
                None
                if model_plot.structure_factor_trace is None
                else np.asarray(model_plot.structure_factor_trace, dtype=float)
            ),
            title=f"DREAM output ({model_plot.bestfit_method})",
            metrics_lines=[
                f"Template: {model_plot.template_name}",
                f"RMSE: {model_plot.rmse:.6g}",
                ("Mean |res|: " f"{model_plot.mean_abs_residual:.6g}"),
                f"R^2: {model_plot.r_squared:.6g}",
            ],
            show_legend=True,
            settings=settings,
        )
        residuals = np.asarray(
            model_plot.model_intensities - model_plot.experimental_intensities,
            dtype=float,
        )
        bottom_axis.axhline(0.0, color="#6b7280", linewidth=1.0)
        bottom_axis.plot(
            np.asarray(model_plot.q_values, dtype=float),
            residuals,
            color=settings.residual_trace_color,
            linewidth=1.2,
        )
        bottom_axis.set_xscale("log")
        bottom_axis.set_xlabel("q (1/A)")
        bottom_axis.set_ylabel("Residual")
        top_axis.tick_params(labelsize=9)
        bottom_axis.tick_params(labelsize=9)
        figure.tight_layout()
    return _save_figure(figure, output_path)


def _render_filter_violin_comparison(
    filter_views: tuple[DreamFilterReportView, ...],
    output_path: Path,
    *,
    settings: PowerPointExportSettings,
    payload_variant: str = "default",
) -> Path:
    with _report_figure_context(settings, compact=True):
        rows, columns = _comparison_grid_shape(len(filter_views))
        figure = Figure(figsize=(12.0, 6.8))
        axes = figure.subplots(rows, columns, squeeze=False)
        axes_flat = list(axes.ravel())
        payloads = _comparison_violin_payloads(
            filter_views,
            payload_variant=payload_variant,
        )
        for axis, view, payload in zip(
            axes_flat,
            filter_views,
            payloads,
            strict=False,
        ):
            _draw_violin_axis(
                axis,
                view,
                settings=settings,
                payload=payload,
            )
        for axis in axes_flat[len(filter_views) :]:
            axis.set_axis_off()
        figure.tight_layout()
    return _save_figure(figure, output_path)


def _render_filter_fit_comparison(
    filter_views: tuple[DreamFilterReportView, ...],
    output_path: Path,
    *,
    settings: PowerPointExportSettings,
) -> Path:
    with _report_figure_context(settings, compact=True):
        rows, columns = _comparison_grid_shape(len(filter_views))
        figure = Figure(figsize=(12.0, 6.8))
        axes = figure.subplots(rows, columns, squeeze=False)
        axes_flat = list(axes.ravel())
        for axis, view in zip(axes_flat, filter_views, strict=False):
            _draw_model_fit_axis(
                axis,
                q_values=np.asarray(view.model_plot.q_values, dtype=float),
                experimental=np.asarray(
                    view.model_plot.experimental_intensities,
                    dtype=float,
                ),
                model=np.asarray(
                    view.model_plot.model_intensities, dtype=float
                ),
                solvent=None,
                structure_factor=(
                    None
                    if view.model_plot.structure_factor_trace is None
                    else np.asarray(
                        view.model_plot.structure_factor_trace,
                        dtype=float,
                    )
                ),
                title=view.title,
                metrics_lines=[
                    f"RMSE: {view.model_plot.rmse:.5g}",
                    (
                        "Mean |res|: "
                        f"{view.model_plot.mean_abs_residual:.5g}"
                    ),
                    f"R^2: {view.model_plot.r_squared:.5g}",
                    f"Samples: {view.summary.posterior_sample_count}",
                ],
                show_legend=False,
                compact=True,
                dream_output_style=True,
                settings=settings,
            )
        for axis in axes_flat[len(filter_views) :]:
            axis.set_axis_off()
        figure.tight_layout()
    return _save_figure(figure, output_path)


def _draw_prefit_plot_axes(
    top_axis,
    bottom_axis,
    evaluation: PrefitEvaluation,
    *,
    include_solvent: bool,
    settings: PowerPointExportSettings,
) -> None:
    del settings
    q_values = np.asarray(evaluation.q_values, dtype=float)
    plotted_lines: list[object] = []
    if evaluation.experimental_intensities is not None:
        experimental = np.asarray(
            evaluation.experimental_intensities,
            dtype=float,
        )
        (experimental_line,) = top_axis.plot(
            q_values,
            experimental,
            color=_PREFIT_WINDOW_EXPERIMENTAL_COLOR,
            label="Experimental",
        )
        plotted_lines.append(experimental_line)
    if include_solvent and evaluation.solvent_contribution is not None:
        solvent_values = np.asarray(
            evaluation.solvent_contribution,
            dtype=float,
        )
        solvent_mask = np.isfinite(solvent_values) & (solvent_values > 0.0)
        if np.any(solvent_mask):
            (solvent_line,) = top_axis.plot(
                q_values[solvent_mask],
                solvent_values[solvent_mask],
                color=_PREFIT_WINDOW_SOLVENT_COLOR,
                linewidth=1.5,
                label="Solvent contribution",
            )
            plotted_lines.append(solvent_line)
    (model_line,) = top_axis.plot(
        q_values,
        np.asarray(evaluation.model_intensities, dtype=float),
        color=_PREFIT_WINDOW_MODEL_COLOR,
        label="Model",
    )
    plotted_lines.append(model_line)
    top_axis.set_xscale("log")
    top_axis.set_yscale("log")
    top_axis.set_ylabel("Intensity (arb. units)")
    top_axis.set_title(
        "Prefit model + solvent trace" if include_solvent else "Prefit model"
    )
    top_axis.text(
        0.02,
        0.02,
        "\n".join(_prefit_metric_lines(evaluation)),
        transform=top_axis.transAxes,
        ha="left",
        va="bottom",
        fontsize=9,
        bbox={
            "boxstyle": "round,pad=0.35",
            "facecolor": "white",
            "edgecolor": "0.6",
            "alpha": 0.85,
        },
    )
    if plotted_lines:
        top_axis.legend(
            plotted_lines,
            [str(line.get_label()) for line in plotted_lines],
            fontsize=8,
            loc="best",
            framealpha=0.92,
        )
    if bottom_axis is not None and evaluation.residuals is not None:
        bottom_axis.axhline(0.0, color="0.5", linewidth=1.0)
        bottom_axis.plot(
            np.asarray(evaluation.q_values, dtype=float),
            np.asarray(evaluation.residuals, dtype=float),
            color=_PREFIT_WINDOW_RESIDUAL_COLOR,
            linewidth=1.2,
        )
        bottom_axis.set_xscale("log")
        bottom_axis.set_xlabel("q (1/A)")
        bottom_axis.set_ylabel("Residual")
        bottom_axis.tick_params(labelsize=9)
        top_axis.set_xlabel("")
    else:
        top_axis.set_xlabel("q (1/A)")
    top_axis.tick_params(labelsize=9)
    _autoscale_visible_axis_limits(
        top_axis,
        log_x=True,
        log_y=True,
    )
    if bottom_axis is not None:
        _autoscale_visible_axis_limits(
            bottom_axis,
            log_x=True,
            log_y=False,
        )


def _draw_model_fit_axis(
    axis,
    *,
    q_values: np.ndarray,
    experimental: np.ndarray | None,
    model: np.ndarray,
    solvent: np.ndarray | None,
    structure_factor: np.ndarray | None,
    title: str,
    metrics_lines: list[str],
    show_legend: bool,
    compact: bool = False,
    dream_output_style: bool = False,
    settings: PowerPointExportSettings,
) -> None:
    q_values = np.asarray(q_values, dtype=float)
    model = np.asarray(model, dtype=float)
    legend_lines: list[object] = []
    legend_labels: list[str] = []
    experimental_color = (
        _DREAM_OUTPUT_EXPERIMENTAL_COLOR
        if dream_output_style
        else settings.experimental_trace_color
    )
    model_color = (
        _DREAM_OUTPUT_MODEL_COLOR
        if dream_output_style
        else settings.model_trace_color
    )
    solvent_color = (
        _DREAM_OUTPUT_SOLVENT_COLOR
        if dream_output_style
        else settings.solvent_trace_color
    )
    structure_factor_color = (
        _DREAM_OUTPUT_STRUCTURE_FACTOR_COLOR
        if dream_output_style
        else settings.structure_factor_color
    )
    if experimental is not None:
        experimental = np.asarray(experimental, dtype=float)
        artist = axis.scatter(
            q_values,
            experimental,
            color=experimental_color,
            s=8 if compact else 14,
            zorder=3,
            label="Experimental",
        )
        legend_lines.append(artist)
        legend_labels.append("Experimental")
    if solvent is not None:
        solvent_mask = np.isfinite(solvent) & (solvent > 0.0)
        if np.any(solvent_mask):
            (solvent_line,) = axis.plot(
                q_values[solvent_mask],
                solvent[solvent_mask],
                color=solvent_color,
                linewidth=1.0 if compact else 1.4,
                label="Solvent contribution",
            )
            legend_lines.append(solvent_line)
            legend_labels.append("Solvent")
    if structure_factor is not None:
        structure_mask = np.isfinite(structure_factor)
        if np.any(structure_mask):
            twin_axis = axis.twinx()
            twin_axis.set_xscale("log")
            (structure_line,) = twin_axis.plot(
                q_values[structure_mask],
                structure_factor[structure_mask],
                color=structure_factor_color,
                linestyle="--",
                linewidth=1.0 if compact else 1.2,
                label=(
                    "Structure factor S(q)"
                    if dream_output_style
                    else "Structure factor"
                ),
            )
            twin_axis.set_ylabel("S(q)", color=structure_factor_color)
            twin_axis.tick_params(
                axis="y",
                colors=structure_factor_color,
            )
            twin_axis.spines["right"].set_color(structure_factor_color)
            legend_lines.append(structure_line)
            legend_labels.append("Structure factor")
    (model_line,) = axis.plot(
        q_values,
        model,
        color=model_color,
        linewidth=1.5 if compact else 2.0,
        label="Model" if not dream_output_style else "Model",
    )
    legend_lines.append(model_line)
    legend_labels.append("Model")
    axis.set_xscale("log")
    axis.set_yscale("log")
    axis.set_ylabel("Intensity")
    axis.set_title(title, fontsize=10 if compact else None)
    axis.tick_params(labelsize=8 if compact else 9)
    metrics_font = 7 if compact else 9
    axis.text(
        0.02,
        0.02,
        "\n".join(metrics_lines),
        transform=axis.transAxes,
        ha="left",
        va="bottom",
        fontsize=metrics_font,
        bbox={
            "boxstyle": "round,pad=0.35",
            "facecolor": "white",
            "edgecolor": "#9ca3af",
            "alpha": 0.9,
        },
    )
    if show_legend and legend_lines:
        axis.legend(
            legend_lines,
            legend_labels,
            fontsize=7 if compact else 8,
            loc="best",
            framealpha=0.92,
        )
    axis.set_xlabel("q (1/A)")


def _draw_violin_axis(
    axis,
    view: DreamFilterReportView,
    *,
    settings: PowerPointExportSettings,
    payload: dict[str, object] | None = None,
) -> None:
    violin_payload = view.violin_payload if payload is None else payload
    samples = np.asarray(violin_payload.get("samples", []), dtype=float)
    if samples.ndim == 1:
        samples = samples.reshape(1, -1)
    display_names = [
        str(label) for label in violin_payload.get("display_names", [])
    ]
    selected_values = np.asarray(
        violin_payload.get("selected_values", []),
        dtype=float,
    )
    interval_low_values = np.asarray(
        violin_payload.get("interval_low_values", []),
        dtype=float,
    )
    interval_high_values = np.asarray(
        violin_payload.get("interval_high_values", []),
        dtype=float,
    )
    if samples.size == 0 or not display_names:
        axis.text(
            0.5,
            0.5,
            "No violin data available.",
            ha="center",
            va="center",
        )
        axis.set_axis_off()
        return
    positions = np.arange(1, len(display_names) + 1)
    violin_parts = axis.violinplot(
        samples,
        positions=positions,
        showmedians=True,
    )
    body_colors = _gradient_colors(
        settings.component_color_map,
        len(display_names),
    )
    for body, color in zip(violin_parts["bodies"], body_colors, strict=False):
        body.set_facecolor(color)
        body.set_edgecolor("#374151")
        body.set_alpha(0.62)
        body.set_linewidth(0.7)
    for key in ("cbars", "cmins", "cmaxes"):
        artist = violin_parts.get(key)
        if artist is not None:
            artist.set_color("#4b5563")
            artist.set_linewidth(1.0)
    median_artist = violin_parts.get("cmedians")
    if median_artist is not None:
        median_artist.set_color("#111827")
        median_artist.set_linewidth(1.2)
    axis.vlines(
        positions,
        interval_low_values,
        interval_high_values,
        color="#4b5563",
        linewidth=1.3,
    )
    axis.scatter(
        positions,
        selected_values,
        color=settings.model_trace_color,
        s=12,
        zorder=3,
    )
    axis.set_xticks(positions)
    axis.set_xticklabels(display_names, rotation=45, ha="right", fontsize=7)
    axis.set_ylabel(str(violin_payload.get("ylabel", "Parameter value")))
    axis.set_title(view.title, fontsize=10)
    axis.tick_params(labelsize=8)
    y_limits = violin_payload.get("y_limits")
    if y_limits is not None:
        axis.set_ylim(float(y_limits[0]), float(y_limits[1]))
    axis.grid(True, axis="y", alpha=0.15)


def _comparison_grid_shape(count: int) -> tuple[int, int]:
    if count <= 1:
        return (1, 1)
    if count <= 2:
        return (1, 2)
    return (2, 2)


def _comparison_violin_payloads(
    filter_views: tuple[DreamFilterReportView, ...],
    *,
    payload_variant: str,
) -> list[dict[str, object]]:
    payload_attribute = {
        "default": "violin_payload",
        "weights": "weights_violin_payload",
        "effective_radii": "effective_radii_violin_payload",
    }.get(payload_variant)
    if payload_attribute is None:
        raise ValueError(
            "Unknown violin payload variant: "
            f"{payload_variant}. Expected 'default', 'weights', "
            "or 'effective_radii'."
        )
    payloads = [
        _copy_violin_payload(getattr(view, payload_attribute))
        for view in filter_views
    ]
    if payload_variant in {"weights", "effective_radii"}:
        _apply_shared_violin_y_limits(payloads)
    return payloads


def _copy_violin_payload(payload: dict[str, object]) -> dict[str, object]:
    copied: dict[str, object] = {}
    for key, value in payload.items():
        if isinstance(value, np.ndarray):
            copied[key] = np.asarray(value, dtype=float).copy()
        elif isinstance(value, list):
            copied[key] = list(value)
        elif isinstance(value, tuple):
            copied[key] = tuple(value)
        else:
            copied[key] = value
    return copied


def _apply_shared_violin_y_limits(
    payloads: Sequence[dict[str, object]]
) -> None:
    if not payloads:
        return
    if any(payload.get("y_limits") is not None for payload in payloads):
        return
    y_limits = _shared_violin_y_limits(payloads)
    if y_limits is None:
        return
    for payload in payloads:
        payload["y_limits"] = y_limits


def _shared_violin_y_limits(
    payloads: Sequence[dict[str, object]],
) -> tuple[float, float] | None:
    lower_bounds: list[float] = []
    upper_bounds: list[float] = []
    for payload in payloads:
        for key in (
            "samples",
            "selected_values",
            "interval_low_values",
            "interval_high_values",
        ):
            values = np.asarray(payload.get(key, []), dtype=float)
            finite_values = values[np.isfinite(values)]
            if finite_values.size == 0:
                continue
            lower_bounds.append(float(np.min(finite_values)))
            upper_bounds.append(float(np.max(finite_values)))
    if not lower_bounds or not upper_bounds:
        return None
    lower = min(lower_bounds)
    upper = max(upper_bounds)
    if np.isclose(lower, upper):
        padding = max(abs(lower) * 0.08, 1.0)
    else:
        padding = (upper - lower) * 0.08
    shared_lower = lower - padding
    shared_upper = upper + padding
    if lower >= 0.0 and shared_lower < 0.0:
        shared_lower = 0.0
    if shared_upper <= shared_lower:
        shared_upper = shared_lower + 1.0
    return (shared_lower, shared_upper)


def _apply_axis_scale(axis, *, log_x: bool, log_y: bool) -> None:
    axis.set_xscale("log" if log_x else "linear")
    axis.set_yscale("log" if log_y else "linear")


def _has_prefit_solvent_trace(evaluation: PrefitEvaluation) -> bool:
    if evaluation.solvent_contribution is None:
        return False
    solvent = np.asarray(evaluation.solvent_contribution, dtype=float)
    return bool(np.any(np.isfinite(solvent) & (solvent > 0.0)))


def _autoscale_visible_axis_limits(
    axis,
    *,
    log_x: bool,
    log_y: bool,
) -> None:
    try:
        axis.relim(visible_only=True)
        axis.autoscale_view()
    except Exception:
        pass
    x_values: list[np.ndarray] = []
    y_values: list[np.ndarray] = []
    for line in axis.get_lines():
        if not line.get_visible():
            continue
        line_x = np.asarray(line.get_xdata(orig=False), dtype=float)
        line_y = np.asarray(line.get_ydata(orig=False), dtype=float)
        mask = np.isfinite(line_x) & np.isfinite(line_y)
        if log_x:
            mask &= line_x > 0.0
        if log_y:
            mask &= line_y > 0.0
        if not np.any(mask):
            continue
        x_values.append(line_x[mask])
        y_values.append(line_y[mask])
    if x_values:
        axis.set_xlim(
            *_autoscale_axis_limits(np.concatenate(x_values), log_scale=log_x)
        )
    if y_values:
        axis.set_ylim(
            *_autoscale_axis_limits(np.concatenate(y_values), log_scale=log_y)
        )


def _autoscale_axis_limits(
    values: np.ndarray,
    *,
    log_scale: bool,
) -> tuple[float, float]:
    data = np.asarray(values, dtype=float)
    finite = data[np.isfinite(data)]
    if finite.size == 0:
        return (0.0, 1.0)
    lower = float(np.nanmin(finite))
    upper = float(np.nanmax(finite))
    if np.isclose(lower, upper):
        if log_scale:
            lower = lower / 1.15
            upper = upper * 1.15
            return (lower, upper)
        padding = max(abs(lower) * 0.05, 1e-12)
        return (lower - padding, upper + padding)
    if log_scale:
        return (lower / 1.05, upper * 1.05)
    padding = 0.05 * (upper - lower)
    return (lower - padding, upper + padding)


def _component_q_bounds(
    plot_data: ReportComponentPlotData,
) -> tuple[float, float] | None:
    q_segments: list[np.ndarray] = []
    if plot_data.experimental_q_values is not None:
        q_segments.append(
            np.asarray(plot_data.experimental_q_values, dtype=float)
        )
    if plot_data.solvent_q_values is not None:
        q_segments.append(np.asarray(plot_data.solvent_q_values, dtype=float))
    for series in plot_data.component_series:
        q_segments.append(np.asarray(series.q_values, dtype=float))
    if not q_segments:
        return None
    merged = np.concatenate(q_segments)
    finite = merged[np.isfinite(merged)]
    if finite.size == 0:
        return None
    return (float(np.nanmin(finite)), float(np.nanmax(finite)))


def _autoscale_component_plot_to_model_range(
    experimental_axis,
    component_axis,
    plot_data: ReportComponentPlotData,
) -> None:
    model_q_bounds = _component_model_q_bounds(plot_data)
    if model_q_bounds is None:
        return
    q_min, q_max = model_q_bounds
    component_axis.set_xlim(q_min, q_max)
    if experimental_axis is not None:
        experimental_axis.set_xlim(q_min, q_max)
        _autoscale_axis_y_for_plot(
            experimental_axis,
            q_min,
            q_max,
            log_scale=plot_data.log_y,
        )
        _normalize_component_axis_to_experimental(
            experimental_axis,
            component_axis,
            plot_data,
        )
        return
    _autoscale_axis_y_for_plot(
        component_axis,
        q_min,
        q_max,
        log_scale=plot_data.log_y,
    )


def _component_model_q_bounds(
    plot_data: ReportComponentPlotData,
) -> tuple[float, float] | None:
    q_segments = [
        np.asarray(series.q_values, dtype=float)
        for series in plot_data.component_series
    ]
    if not q_segments:
        return None
    merged = np.concatenate(q_segments)
    finite = merged[np.isfinite(merged)]
    if finite.size == 0:
        return None
    return (float(np.nanmin(finite)), float(np.nanmax(finite)))


def _autoscale_axis_y_for_plot(
    axis,
    q_min: float,
    q_max: float,
    *,
    log_scale: bool,
) -> None:
    y_segments: list[np.ndarray] = []
    for line in axis.get_lines():
        if not line.get_visible():
            continue
        x_data = np.asarray(line.get_xdata(orig=False), dtype=float)
        y_data = np.asarray(line.get_ydata(orig=False), dtype=float)
        mask = (
            np.isfinite(x_data)
            & np.isfinite(y_data)
            & (x_data >= q_min)
            & (x_data <= q_max)
        )
        if log_scale:
            mask &= y_data > 0.0
        if np.any(mask):
            y_segments.append(y_data[mask])
    if not y_segments:
        return
    y_values = np.concatenate(y_segments)
    y_min = float(np.nanmin(y_values))
    y_max = float(np.nanmax(y_values))
    if np.isclose(y_min, y_max):
        padding = max(abs(y_min) * 0.05, 1e-12)
        axis.set_ylim(y_min - padding, y_max + padding)
        return
    if log_scale:
        axis.set_ylim(y_min / 1.15, y_max * 1.15)
    else:
        padding = 0.05 * (y_max - y_min)
        axis.set_ylim(y_min - padding, y_max + padding)


def _normalize_component_axis_to_experimental(
    experimental_axis,
    component_axis,
    plot_data: ReportComponentPlotData,
) -> None:
    if (
        plot_data.experimental_q_values is None
        or plot_data.experimental_intensities is None
    ):
        return
    component_values = [
        np.asarray(series.intensities, dtype=float)
        for series in plot_data.component_series
    ]
    if not component_values:
        return
    component_data = np.concatenate(component_values)
    component_data = component_data[np.isfinite(component_data)]
    if plot_data.log_y:
        component_data = component_data[component_data > 0.0]
    if component_data.size == 0:
        return
    experimental_q_values = np.asarray(
        plot_data.experimental_q_values,
        dtype=float,
    )
    experimental_intensities = np.asarray(
        plot_data.experimental_intensities,
        dtype=float,
    )
    experimental_mask = _selected_q_mask(
        experimental_q_values,
        plot_data.selected_q_min,
        plot_data.selected_q_max,
        use_experimental_grid=plot_data.use_experimental_grid,
    )
    if not np.any(experimental_mask):
        return
    filtered_q = experimental_q_values[experimental_mask]
    filtered_i = experimental_intensities[experimental_mask]
    model_q_bounds = _component_model_q_bounds(plot_data)
    if model_q_bounds is None:
        return
    overlap_mask = (filtered_q >= model_q_bounds[0]) & (
        filtered_q <= model_q_bounds[1]
    )
    if np.any(overlap_mask):
        filtered_i = filtered_i[overlap_mask]
    filtered_i = filtered_i[np.isfinite(filtered_i)]
    if plot_data.log_y:
        filtered_i = filtered_i[filtered_i > 0.0]
    if filtered_i.size == 0:
        return
    left_limits = experimental_axis.get_ylim()
    right_limits = _aligned_y_limits(
        left_limits,
        float(np.nanmin(filtered_i)),
        float(np.nanmax(filtered_i)),
        float(np.nanmin(component_data)),
        float(np.nanmax(component_data)),
        log_scale=plot_data.log_y,
    )
    component_axis.set_ylim(right_limits)


def _aligned_y_limits(
    left_limits: tuple[float, float],
    experimental_min: float,
    experimental_max: float,
    component_min: float,
    component_max: float,
    *,
    log_scale: bool,
) -> tuple[float, float]:
    if log_scale:
        if (
            min(
                left_limits[0],
                left_limits[1],
                experimental_min,
                experimental_max,
                component_min,
                component_max,
            )
            <= 0.0
        ):
            log_scale = False
    if not log_scale:
        left_low, left_high = left_limits
        exp_low, exp_high = sorted((experimental_min, experimental_max))
        comp_low, comp_high = sorted((component_min, component_max))
        if np.isclose(left_high, left_low) or np.isclose(exp_high, exp_low):
            padding = max(abs(comp_low) * 0.1, 1e-12)
            return comp_low - padding, comp_high + padding
        p0 = (exp_low - left_low) / (left_high - left_low)
        p1 = (exp_high - left_low) / (left_high - left_low)
        if np.isclose(p1, p0):
            padding = max(abs(comp_low) * 0.1, 1e-12)
            return comp_low - padding, comp_high + padding
        delta = (comp_high - comp_low) / (p1 - p0)
        right_low = comp_low - p0 * delta
        right_high = right_low + delta
        return right_low, right_high
    left_logs = np.log10(np.asarray(left_limits, dtype=float))
    exp_logs = np.log10(
        np.asarray(sorted((experimental_min, experimental_max)), dtype=float)
    )
    comp_logs = np.log10(
        np.asarray(sorted((component_min, component_max)), dtype=float)
    )
    if np.isclose(left_logs[1], left_logs[0]) or np.isclose(
        exp_logs[1], exp_logs[0]
    ):
        return component_min / 1.2, component_max * 1.2
    p0 = (exp_logs[0] - left_logs[0]) / (left_logs[1] - left_logs[0])
    p1 = (exp_logs[1] - left_logs[0]) / (left_logs[1] - left_logs[0])
    if np.isclose(p1, p0):
        return component_min / 1.2, component_max * 1.2
    delta = (comp_logs[1] - comp_logs[0]) / (p1 - p0)
    right_low_log = comp_logs[0] - p0 * delta
    right_high_log = right_low_log + delta
    return 10**right_low_log, 10**right_high_log


def _selected_q_mask(
    q_values: np.ndarray,
    lower: float | None,
    upper: float | None,
    *,
    use_experimental_grid: bool,
) -> np.ndarray:
    values = np.asarray(q_values, dtype=float)
    if values.size == 0:
        return np.zeros(0, dtype=bool)
    if lower is None and upper is None:
        return np.ones_like(values, dtype=bool)
    lower_bound = lower if lower is not None else float(np.nanmin(values))
    upper_bound = upper if upper is not None else float(np.nanmax(values))
    if lower_bound > upper_bound:
        return np.zeros_like(values, dtype=bool)
    if use_experimental_grid:
        start_index = int(np.argmin(np.abs(values - lower_bound)))
        end_index = int(np.argmin(np.abs(values - upper_bound)))
        lo_index, hi_index = sorted((start_index, end_index))
        mask = np.zeros_like(values, dtype=bool)
        mask[lo_index : hi_index + 1] = True
        return mask
    return (values >= lower_bound) & (values <= upper_bound)


def _gradient_colors(cmap_name: str, count: int) -> list[str]:
    try:
        cmap = colormaps[cmap_name]
    except Exception:
        cmap = colormaps["viridis"]
    if count <= 1:
        return [to_hex(cmap(0.72), keep_alpha=False)]
    positions = np.linspace(0.22, 0.9, count)
    return [
        to_hex(cmap(float(position)), keep_alpha=False)
        for position in positions
    ]


def _prefit_summary_lines(
    evaluation: PrefitEvaluation,
    statistics: dict[str, object],
) -> list[str]:
    lines = [
        f"Points: {len(np.asarray(evaluation.q_values, dtype=float))}",
        (
            "q-range: "
            f"{float(np.min(evaluation.q_values)):.6g} to "
            f"{float(np.max(evaluation.q_values)):.6g}"
        ),
    ]
    lines.extend(_prefit_metric_lines(evaluation))
    if statistics.get("method"):
        lines.append(f"Method: {statistics['method']}")
    if statistics.get("nfev") is not None:
        lines.append(f"Function evals: {statistics['nfev']}")
    if statistics.get("chi_square") is not None:
        lines.append(f"Chi^2: {float(statistics['chi_square']):.6g}")
    if statistics.get("reduced_chi_square") is not None:
        lines.append(
            f"Reduced chi^2: {float(statistics['reduced_chi_square']):.6g}"
        )
    if statistics.get("r_squared") is not None:
        lines.append(f"Saved-fit R^2: {float(statistics['r_squared']):.6g}")
    if statistics.get("saved_at"):
        lines.append(f"Saved state timestamp: {statistics['saved_at']}")
    return lines


def _prefit_metric_lines(evaluation: PrefitEvaluation) -> list[str]:
    if (
        evaluation.experimental_intensities is None
        or evaluation.residuals is None
    ):
        return [
            "Model Only Mode",
            "Experimental fit metrics unavailable",
        ]
    experimental = np.asarray(evaluation.experimental_intensities, dtype=float)
    model = np.asarray(evaluation.model_intensities, dtype=float)
    residuals = np.asarray(model - experimental, dtype=float)
    rmse = float(np.sqrt(np.mean(residuals**2)))
    mean_abs = float(np.mean(np.abs(residuals)))
    mean_experimental = float(np.mean(experimental))
    total_sum_squares = float(np.sum((experimental - mean_experimental) ** 2))
    residual_sum_squares = float(np.sum(residuals**2))
    r_squared = (
        float(1.0 - (residual_sum_squares / total_sum_squares))
        if total_sum_squares > 0.0
        else 1.0
    )
    return [
        f"RMSE: {rmse:.6g}",
        f"Mean |res|: {mean_abs:.6g}",
        f"R^2: {r_squared:.6g}",
    ]


def _prefit_parameter_rows(
    entries: tuple[PrefitParameterEntry, ...],
) -> list[list[str]]:
    rows: list[list[str]] = []
    for entry in entries:
        rows.append(
            [
                str(entry.name),
                str(entry.category),
                f"{float(entry.value):.6g}",
                "Yes" if entry.vary else "No",
                f"{float(entry.minimum):.6g}",
                f"{float(entry.maximum):.6g}",
                str(entry.structure or "-"),
                str(entry.motif or "-"),
            ]
        )
    return rows


def _cluster_geometry_rows(
    rows: tuple[ClusterGeometryMetadataRow, ...],
) -> list[list[str]]:
    table_rows: list[list[str]] = []
    for row in rows:
        axes_text = (
            f"{float(row.active_semiaxis_a):.4g} / "
            f"{float(row.active_semiaxis_b):.4g} / "
            f"{float(row.active_semiaxis_c):.4g}"
        )
        table_rows.append(
            [
                str(row.cluster_id),
                str(row.mapped_parameter or "-"),
                str(row.sf_approximation),
                f"{float(row.effective_radius):.5g}",
                f"{float(row.mean_radius_of_gyration):.5g}",
                f"{float(row.mean_max_radius):.5g}",
                f"{float(row.anisotropy_metric):.5g}",
                axes_text,
            ]
        )
    return table_rows


def _solution_estimate_lines(
    estimate: SolutionScatteringEstimate,
) -> list[str]:
    lines = [
        "Solution scattering estimator summary",
        f"Incident energy: {float(estimate.settings.beam.incident_energy_kev):.6g} keV",
    ]
    if estimate.number_density_estimate is not None:
        lines.extend(
            [
                "",
                "Number density:",
                (
                    "Atoms/A^3: "
                    f"{estimate.number_density_estimate.number_density_a3:.6g}"
                ),
                (
                    "Atoms/cm^3: "
                    f"{estimate.number_density_estimate.number_density_cm3:.6g}"
                ),
                (
                    "Total atoms: "
                    f"{estimate.number_density_estimate.total_atoms:.6g}"
                ),
            ]
        )
    if estimate.volume_fraction_estimate is not None:
        lines.extend(
            [
                "",
                "Physical volume fraction:",
                (
                    "Physical solute-associated volume fraction: "
                    f"{float(estimate.volume_fraction_estimate.solute_volume_fraction):.6f}"
                ),
                (
                    "Physical solvent-associated volume fraction: "
                    f"{float(estimate.volume_fraction_estimate.solvent_volume_fraction):.6f}"
                ),
            ]
        )
    if estimate.interaction_contrast_estimate is not None:
        lines.extend(
            [
                "",
                "SAXS-effective interaction ratio:",
                (
                    "Contrast weight factor: "
                    f"{float(estimate.interaction_contrast_estimate.contrast_weight_factor):.6g}"
                ),
                (
                    "SAXS-effective solute interaction ratio: "
                    f"{float(estimate.interaction_contrast_estimate.saxs_effective_solute_interaction_ratio):.6f}"
                ),
                (
                    "SAXS-effective solvent background ratio: "
                    f"{float(estimate.interaction_contrast_estimate.saxs_effective_solvent_background_ratio):.6f}"
                ),
            ]
        )
    if estimate.attenuation_estimate is not None:
        lines.extend(
            [
                "",
                "Attenuation:",
                (
                    "Sample transmission: "
                    f"{float(estimate.attenuation_estimate.sample_transmission):.6f}"
                ),
                (
                    "Neat-solvent transmission: "
                    f"{float(estimate.attenuation_estimate.neat_solvent_transmission):.6f}"
                ),
                (
                    "Solvent scale factor: "
                    f"{float(estimate.attenuation_estimate.solvent_scattering_scale_factor):.6f}"
                ),
            ]
        )
        if estimate.interaction_contrast_estimate is not None:
            lines.append(
                "Single-weight solvent multiplier: "
                f"{float(estimate.attenuation_estimate.solvent_scattering_scale_factor * estimate.interaction_contrast_estimate.saxs_effective_solvent_background_ratio):.6f}"
            )
        if (
            estimate.attenuation_estimate.neat_solvent_to_sample_ratio
            is not None
        ):
            lines.append(
                "Neat-solvent/sample-solvent ratio: "
                f"{float(estimate.attenuation_estimate.neat_solvent_to_sample_ratio):.6g}"
            )
    if estimate.fluorescence_estimate is not None:
        lines.extend(
            [
                "",
                "Fluorescence:",
                (
                    "Primary yield proxy: "
                    f"{float(estimate.fluorescence_estimate.total_primary_detected_yield):.6g}"
                ),
                (
                    "Secondary yield proxy: "
                    f"{float(estimate.fluorescence_estimate.total_secondary_detected_yield):.6g}"
                ),
            ]
        )
        for line_estimate in estimate.fluorescence_estimate.line_estimates[:5]:
            lines.append(
                f"{line_estimate.element} {line_estimate.family}: "
                f"{float(line_estimate.total_detected_yield):.6g}"
            )
    return lines


def _dream_settings_lines(
    settings: DreamRunSettings,
    summary: DreamSummary,
) -> list[str]:
    return [
        f"Best-fit method: {summary.bestfit_method}",
        ("Posterior filter: " f"{_describe_posterior_filter(settings)}"),
        f"Posterior samples kept: {summary.posterior_sample_count}",
        (
            "Credible interval: "
            f"{summary.credible_interval_low:g} - "
            f"{summary.credible_interval_high:g}"
        ),
        f"MAP location: chain {summary.map_chain + 1}, step {summary.map_step + 1}",
        f"nchains: {settings.nchains}",
        f"niterations: {settings.niterations}",
        f"burn-in (%): {settings.burnin_percent}",
        f"nseedchains: {settings.nseedchains}",
        f"crossover burn-in: {settings.crossover_burnin}",
        f"run label: {settings.run_label}",
        f"violin mode: {settings.violin_parameter_mode}",
        f"violin sample source: {settings.violin_sample_source}",
        f"violin weight order: {settings.violin_weight_order}",
        f"violin y-scale: {settings.violin_value_scale_mode}",
        (
            "Auto-select best filter after run: "
            f"{'on' if settings.auto_select_best_posterior_filter else 'off'}"
        ),
    ]


def _dream_assessment_lines(
    assessments: tuple[dict[str, object], ...],
    settings: DreamRunSettings,
) -> list[str]:
    lines = [
        "Posterior filtering assessment",
        ("Active selection: " f"{_describe_posterior_filter(settings)}"),
    ]
    if not assessments:
        lines.append("No saved assessment metrics are available.")
        return lines
    for assessment in assessments:
        lines.extend(
            [
                "",
                str(assessment.get("description", "Unnamed filter")),
                f"RMSE: {float(assessment.get('rmse', 0.0)):.6g}",
                (
                    "Mean |res|: "
                    f"{float(assessment.get('mean_abs_residual', 0.0)):.6g}"
                ),
                f"R^2: {float(assessment.get('r_squared', 0.0)):.6g}",
                (
                    "Posterior samples: "
                    f"{int(assessment.get('posterior_sample_count', 0))}"
                ),
            ]
        )
    return lines


def _dream_prior_rows(
    entries: tuple[DreamParameterEntry, ...],
) -> list[list[str]]:
    rows: list[list[str]] = []
    for entry in entries:
        dist_params = ", ".join(
            f"{key}={float(value):.6g}"
            for key, value in sorted(entry.dist_params.items())
        )
        rows.append(
            [
                str(entry.param),
                str(entry.param_type),
                str(entry.structure or "-"),
                str(entry.motif or "-"),
                f"{float(entry.value):.6g}",
                "Yes" if entry.vary else "No",
                str(entry.distribution),
                dist_params,
            ]
        )
    return rows


def _dream_output_lines(
    settings: DreamRunSettings,
    summary: DreamSummary,
    model_plot: DreamModelPlotData,
) -> list[str]:
    lines = [
        f"Template: {model_plot.template_name}",
        f"Best-fit method: {settings.bestfit_method}",
        ("Posterior filter: " f"{_describe_posterior_filter(settings)}"),
        f"Posterior samples kept: {summary.posterior_sample_count}",
        f"RMSE: {model_plot.rmse:.6g}",
        f"Mean |res|: {model_plot.mean_abs_residual:.6g}",
        f"R^2: {model_plot.r_squared:.6g}",
        (
            "Credible interval: "
            f"{summary.credible_interval_low:g} - "
            f"{summary.credible_interval_high:g}"
        ),
        (
            "Active parameters: "
            f"{', '.join(summary.active_parameter_names) or 'None'}"
        ),
    ]
    for index, name in enumerate(summary.full_parameter_names[:10]):
        lines.append(
            f"{name}: {float(summary.bestfit_params[index]):.6g} "
            f"(p{summary.credible_interval_low:g}="
            f"{float(summary.interval_low_values[index]):.6g}, "
            f"p{summary.credible_interval_high:g}="
            f"{float(summary.interval_high_values[index]):.6g})"
        )
    if len(summary.full_parameter_names) > 10:
        lines.append(
            f"... {len(summary.full_parameter_names) - 10} additional parameters"
        )
    return lines


def _describe_posterior_filter(settings: DreamRunSettings) -> str:
    if settings.posterior_filter_mode == "top_percent_logp":
        return (
            f"top_percent_logp "
            f"(top {settings.posterior_top_percent:g}% by log-posterior)"
        )
    if settings.posterior_filter_mode == "top_n_logp":
        return (
            f"top_n_logp "
            f"(top {settings.posterior_top_n} samples by log-posterior)"
        )
    return "all_post_burnin"


def _manifest_payload(
    context: DreamModelReportContext,
    *,
    figure_paths: list[Path],
) -> dict[str, object]:
    return {
        "report_type": "dream_model_report_pptx",
        "generated_at": context.generated_at.isoformat(),
        "project_name": context.project_name,
        "project_dir": str(context.project_dir),
        "report_path": str(context.output_path),
        "asset_dir": str(context.asset_dir),
        "dream_run_dir": str(context.dream_summary.run_dir),
        "user_q_range": context.user_q_range_text,
        "supported_q_range": context.supported_q_range_text,
        "q_sampling": context.q_sampling_text,
        "template_name": context.template_name,
        "template_display_name": context.template_display_name,
        "template_module_path": (
            None
            if context.template_module_path is None
            else str(context.template_module_path)
        ),
        "model_equation": context.model_equation_text,
        "model_reference_lines": list(context.model_reference_lines),
        "figure_paths": [str(path) for path in figure_paths],
        "prior_histograms": [
            {
                "title": request.title,
                "mode": request.mode,
                "secondary_element": request.secondary_element,
                "cmap": request.cmap,
                "source": str(request.json_path),
            }
            for request in context.prior_histograms
        ],
        "prefit_statistics": {
            str(key): value
            for key, value in context.prefit_statistics.items()
            if not isinstance(value, Path)
        },
        "dream_settings": context.dream_settings.to_dict(),
        "dream_summary": {
            "bestfit_method": context.dream_summary.bestfit_method,
            "posterior_filter_mode": context.dream_summary.posterior_filter_mode,
            "posterior_sample_count": context.dream_summary.posterior_sample_count,
            "credible_interval_low": context.dream_summary.credible_interval_low,
            "credible_interval_high": context.dream_summary.credible_interval_high,
        },
        "dream_filter_views": [
            {
                "title": view.title,
                "description": view.description,
                "filter_mode": view.filter_mode,
                "is_active": view.is_active,
                "posterior_sample_count": view.summary.posterior_sample_count,
                "rmse": view.model_plot.rmse,
                "mean_abs_residual": view.model_plot.mean_abs_residual,
                "r_squared": view.model_plot.r_squared,
            }
            for view in context.dream_filter_views
        ],
    }


def _rgb_color(value: str):
    from pptx.dml.color import RGBColor

    red, green, blue = tuple(
        int(value[index : index + 2], 16) for index in (1, 3, 5)
    )
    return RGBColor(red, green, blue)


def _slugify(value: str) -> str:
    safe = "".join(
        character.lower() if character.isalnum() else "_"
        for character in value
    )
    while "__" in safe:
        safe = safe.replace("__", "_")
    return safe.strip("_") or "figure"


__all__ = [
    "DreamFilterReportView",
    "DreamModelReportContext",
    "ModelReportExportResult",
    "PriorHistogramRequest",
    "ReportComponentPlotData",
    "ReportComponentSeries",
    "export_dream_model_report_pptx",
]
