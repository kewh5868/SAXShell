from __future__ import annotations

import textwrap
from collections.abc import Callable, Sequence
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import TYPE_CHECKING

from matplotlib.figure import Figure
from matplotlib.text import Text

from saxshell.clusterdynamics import ClusterDynamicsResult
from saxshell.saxs.project_manager import (
    PowerPointExportSettings,
    build_project_paths,
)

if TYPE_CHECKING:
    from saxshell.clusterdynamicsml.workflow import ClusterDynamicsMLResult

ClusterReportProgressCallback = Callable[[int, int, str], None]

_SLIDE_WIDTH_INCHES = 13.333
_SLIDE_HEIGHT_INCHES = 7.5
_SLIDE_LEFT_INCHES = 0.45
_SLIDE_CONTENT_WIDTH_INCHES = 12.43
_TITLE_TOP_INCHES = 0.28
_TITLE_HEIGHT_INCHES = 0.42
_SUBTITLE_TOP_INCHES = 0.72
_SUBTITLE_HEIGHT_INCHES = 0.22
_TEXT_TOP_INCHES = 1.10
_TEXT_HEIGHT_INCHES = 5.92
_FIGURE_TOP_INCHES = 1.08
_FIGURE_HEIGHT_INCHES = 5.96
_TABLE_TOP_INCHES = 1.08
_TABLE_HEIGHT_INCHES = 5.98
_THICK_RULE_HEIGHT_INCHES = 0.04


@dataclass(slots=True)
class ClusterPowerPointExportResult:
    report_path: Path
    appended_to_existing: bool
    added_slide_count: int


@dataclass(slots=True)
class _TextSection:
    title: str
    subtitle: str | None
    pages: list[list[str]]
    placeholder: str
    font_size: float = 13.0


@dataclass(slots=True)
class _FigureSection:
    title: str
    subtitle: str | None
    figure: Figure | None
    placeholder: str


@dataclass(slots=True)
class _TableSection:
    title: str
    subtitle: str | None
    columns: tuple[str, ...]
    rows: list[list[str]]
    rows_per_slide: int
    column_width_weights: tuple[float, ...] | None = None
    alignments: tuple[str, ...] | None = None
    note: str | None = None
    header_font_size: float = 11.0
    row_font_size: float = 10.0


class _ReportProgressTracker:
    def __init__(
        self,
        total_steps: int,
        callback: ClusterReportProgressCallback | None,
        *,
        opening_message: str,
    ) -> None:
        self.total_steps = max(int(total_steps), 1)
        self._callback = callback
        self._processed = 0
        self.emit(opening_message)

    def emit(self, message: str) -> None:
        if self._callback is None:
            return
        self._callback(self._processed, self.total_steps, str(message))

    def advance(self, message: str) -> None:
        self._processed = min(self._processed + 1, self.total_steps)
        self.emit(message)


def latest_project_powerpoint_report(project_dir: str | Path) -> Path | None:
    reports_dir = build_project_paths(project_dir).reports_dir
    reports_dir.mkdir(parents=True, exist_ok=True)
    candidates = sorted(
        reports_dir.glob("*.pptx"),
        key=lambda path: (path.stat().st_mtime, path.name.lower()),
    )
    if not candidates:
        return None
    return candidates[-1]


def default_powerpoint_report_path(
    *,
    project_dir: str | Path | None,
    fallback_dir: str | Path,
    fallback_stem: str,
) -> Path:
    if project_dir is not None:
        project_dir_path = Path(project_dir).expanduser().resolve()
        existing = latest_project_powerpoint_report(project_dir_path)
        if existing is not None:
            return existing
        reports_dir = build_project_paths(project_dir_path).reports_dir
        reports_dir.mkdir(parents=True, exist_ok=True)
        return reports_dir / f"{project_dir_path.name}_results.pptx"
    fallback_dir_path = Path(fallback_dir).expanduser().resolve()
    fallback_dir_path.mkdir(parents=True, exist_ok=True)
    return fallback_dir_path / f"{fallback_stem}.pptx"


def normalize_powerpoint_output_path(path: str | Path) -> Path:
    resolved = Path(path).expanduser()
    if resolved.suffix.lower() != ".pptx":
        resolved = resolved.with_suffix(".pptx")
    return resolved.resolve()


def export_cluster_dynamics_report_pptx(
    *,
    result: ClusterDynamicsResult,
    selection_summary: str,
    result_summary: str,
    figure: Figure | None,
    output_path: str | Path,
    settings: PowerPointExportSettings | None = None,
    project_dir: str | Path | None = None,
    frames_dir: str | Path | None = None,
    progress_callback: ClusterReportProgressCallback | None = None,
) -> ClusterPowerPointExportResult:
    export_settings = _normalized_settings(settings)
    lifetime_columns = (
        "Label",
        "Size",
        "Mean lifetime (fs)",
        "Std lifetime (fs)",
        "Completed",
        "Window-truncated",
        "Assoc. rate (1/ps)",
        "Dissoc. rate (1/ps)",
        "Occupancy (%)",
        "Mean count/frame",
    )
    lifetime_rows = [
        [
            entry.label,
            str(entry.cluster_size),
            _format_optional_float(entry.mean_lifetime_fs),
            _format_optional_float(entry.std_lifetime_fs),
            str(entry.completed_lifetime_count),
            str(entry.window_truncated_lifetime_count),
            f"{entry.association_rate_per_ps:.3f}",
            f"{entry.dissociation_rate_per_ps:.3f}",
            f"{entry.occupancy_fraction * 100.0:.1f}",
            f"{entry.mean_count_per_frame:.3f}",
        ]
        for entry in result.lifetime_by_label
    ]
    generated_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    cover_lines = [
        f"Generated: {generated_at}",
        f"Frames analyzed: {result.analyzed_frames}",
        f"Time bins: {result.bin_count}",
        "Project directory: "
        + (
            "not set"
            if project_dir is None
            else str(Path(project_dir).resolve())
        ),
        "Frames folder: "
        + (
            "not set"
            if frames_dir is None
            else str(Path(frames_dir).resolve())
        ),
    ]
    text_sections = [
        _TextSection(
            title="ClusterDynamics Settings",
            subtitle="Selection preview and analysis inputs",
            pages=_paginate_text_lines(
                selection_summary.splitlines(),
                max_lines=18,
                wrap_at=60,
            ),
            placeholder="Selection settings are not available for this result.",
        ),
        _TextSection(
            title="ClusterDynamics Summary",
            subtitle="Observed result summary",
            pages=_paginate_text_lines(
                result_summary.splitlines(),
                max_lines=18,
                wrap_at=60,
            ),
            placeholder="Summary information is not available for this result.",
        ),
    ]
    figure_sections = [
        _FigureSection(
            title="Cluster Distribution Heatmap",
            subtitle="Current clusterdynamics plot settings",
            figure=figure,
            placeholder="No cluster-distribution plot is available for this result.",
        )
    ]
    table_sections = [
        _TableSection(
            title="Observed Cluster Lifetimes",
            subtitle="Lifetime statistics by stoichiometry label",
            columns=lifetime_columns,
            rows=lifetime_rows,
            rows_per_slide=11,
            column_width_weights=(
                1.55,
                0.55,
                1.05,
                0.95,
                0.7,
                1.0,
                1.0,
                1.0,
                0.95,
                1.05,
            ),
            alignments=(
                "left",
                "center",
                "center",
                "center",
                "center",
                "center",
                "center",
                "center",
                "center",
                "center",
            ),
            row_font_size=9.5,
        )
    ]
    return _export_cluster_report_pptx(
        report_title="ClusterDynamics Report",
        cover_subtitle="Time-binned cluster-distribution analysis",
        cover_lines=cover_lines,
        text_sections=text_sections,
        figure_sections=figure_sections,
        table_sections=table_sections,
        output_path=output_path,
        settings=export_settings,
        progress_callback=progress_callback,
    )


def export_cluster_dynamicsai_report_pptx(
    *,
    result: "ClusterDynamicsMLResult",
    selection_summary: str,
    result_summary: str,
    dynamics_figure: Figure | None,
    predicted_structures_figure: Figure | None,
    output_path: str | Path,
    settings: PowerPointExportSettings | None = None,
    project_dir: str | Path | None = None,
    frames_dir: str | Path | None = None,
    progress_callback: ClusterReportProgressCallback | None = None,
) -> ClusterPowerPointExportResult:
    export_settings = _normalized_settings(settings)
    lifetime_columns = (
        "Label",
        "Size",
        "Mean lifetime (fs)",
        "Std lifetime (fs)",
        "Completed",
        "Window-truncated",
        "Assoc. rate (1/ps)",
        "Dissoc. rate (1/ps)",
        "Occupancy (%)",
        "Mean count/frame",
    )
    lifetime_rows = [
        [
            entry.label,
            str(entry.cluster_size),
            _format_optional_float(entry.mean_lifetime_fs),
            _format_optional_float(entry.std_lifetime_fs),
            str(entry.completed_lifetime_count),
            str(entry.window_truncated_lifetime_count),
            f"{entry.association_rate_per_ps:.3f}",
            f"{entry.dissociation_rate_per_ps:.3f}",
            f"{entry.occupancy_fraction * 100.0:.1f}",
            f"{entry.mean_count_per_frame:.3f}",
        ]
        for entry in result.dynamics_result.lifetime_by_label
    ]
    prediction_columns = (
        "Target nodes",
        "Rank",
        "Label",
        "Share (%)",
        "Mean count/frame",
        "Mean lifetime (fs)",
        "Assoc. rate",
        "Dissoc. rate",
        "Source",
        "Notes",
    )
    prediction_rows = [
        [
            str(entry.target_node_count),
            str(entry.rank),
            entry.label,
            f"{entry.predicted_population_share * 100.0:.2f}",
            f"{entry.predicted_mean_count_per_frame:.4f}",
            f"{entry.predicted_mean_lifetime_fs:.3f}",
            f"{entry.predicted_association_rate_per_ps:.3f}",
            f"{entry.predicted_dissociation_rate_per_ps:.3f}",
            "" if entry.source_label is None else entry.source_label,
            entry.notes,
        ]
        for entry in result.predictions
    ]
    generated_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    cover_lines = [
        f"Generated: {generated_at}",
        f"Frames analyzed: {result.dynamics_result.analyzed_frames}",
        f"Observed node counts: {result.preview.observed_node_counts or ('n/a',)}",
        f"Predicted candidates: {len(result.predictions)}",
        "Project directory: "
        + (
            "not set"
            if project_dir is None
            else str(Path(project_dir).resolve())
        ),
        "Frames folder: "
        + (
            "not set"
            if frames_dir is None
            else str(Path(frames_dir).resolve())
        ),
    ]
    text_sections = [
        _TextSection(
            title="ClusterDynamicsML Settings",
            subtitle="Selection preview and prediction inputs",
            pages=_paginate_text_lines(
                selection_summary.splitlines(),
                max_lines=18,
                wrap_at=60,
            ),
            placeholder="Selection settings are not available for this result.",
        ),
        _TextSection(
            title="ClusterDynamicsML Summary",
            subtitle="Observed and predicted result summary",
            pages=_paginate_text_lines(
                result_summary.splitlines(),
                max_lines=18,
                wrap_at=60,
            ),
            placeholder="Summary information is not available for this result.",
        ),
    ]
    figure_sections = [
        _FigureSection(
            title="Observed Cluster Distribution",
            subtitle="Current clusterdynamics heatmap settings",
            figure=dynamics_figure,
            placeholder="No observed cluster-distribution plot is available.",
        ),
        _FigureSection(
            title="Predicted Structures SAXS Comparison",
            subtitle=(
                "Observed-only and observed + Predicted Structures SAXS traces"
            ),
            figure=predicted_structures_figure,
            placeholder=(
                "No Predicted Structures SAXS plot is available for this result."
            ),
        ),
    ]
    table_sections = [
        _TableSection(
            title="Observed Cluster Lifetimes",
            subtitle=(
                "Lifetime statistics used by the Predicted Structures workflow"
            ),
            columns=lifetime_columns,
            rows=lifetime_rows,
            rows_per_slide=11,
            column_width_weights=(
                1.55,
                0.55,
                1.05,
                0.95,
                0.7,
                1.0,
                1.0,
                1.0,
                0.95,
                1.05,
            ),
            alignments=(
                "left",
                "center",
                "center",
                "center",
                "center",
                "center",
                "center",
                "center",
                "center",
                "center",
            ),
            row_font_size=9.5,
        ),
        _TableSection(
            title="Predicted Larger Clusters",
            subtitle=(
                "Ranked candidates for Predicted Structures above the current threshold"
            ),
            columns=prediction_columns,
            rows=prediction_rows,
            rows_per_slide=8,
            column_width_weights=(
                0.75,
                0.48,
                1.1,
                0.78,
                1.0,
                1.02,
                0.82,
                0.88,
                0.95,
                3.65,
            ),
            alignments=(
                "center",
                "center",
                "left",
                "center",
                "center",
                "center",
                "center",
                "center",
                "left",
                "left",
            ),
            row_font_size=9.0,
        ),
    ]
    return _export_cluster_report_pptx(
        report_title="ClusterDynamicsML Report",
        cover_subtitle="Cluster extrapolation and Predicted Structures SAXS analysis",
        cover_lines=cover_lines,
        text_sections=text_sections,
        figure_sections=figure_sections,
        table_sections=table_sections,
        output_path=output_path,
        settings=export_settings,
        progress_callback=progress_callback,
    )


def _export_cluster_report_pptx(
    *,
    report_title: str,
    cover_subtitle: str,
    cover_lines: list[str],
    text_sections: Sequence[_TextSection],
    figure_sections: Sequence[_FigureSection],
    table_sections: Sequence[_TableSection],
    output_path: str | Path,
    settings: PowerPointExportSettings,
    progress_callback: ClusterReportProgressCallback | None,
) -> ClusterPowerPointExportResult:
    (
        Presentation,
        Inches,
        Pt,
        RGBColor,
        PP_ALIGN,
        MSO_VERTICAL_ANCHOR,
        MSO_AUTO_SHAPE_TYPE,
    ) = _load_pptx_api()

    def rgb_color(value: str):
        red, green, blue = tuple(
            int(value[index : index + 2], 16) for index in (1, 3, 5)
        )
        return RGBColor(red, green, blue)

    def first_run(paragraph):
        if paragraph.runs:
            return paragraph.runs[0]
        return paragraph.add_run()

    def apply_run_style(
        run,
        *,
        font_size: float,
        bold: bool = False,
        color: str | None = None,
    ) -> None:
        run.font.name = settings.font_family
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rgb_color(
            settings.text_color if color is None else color
        )

    def set_slide_background(slide) -> None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb_color("#FFFFFF")

    def add_title(slide, title: str, subtitle: str | None = None) -> None:
        title_box = slide.shapes.add_textbox(
            Inches(_SLIDE_LEFT_INCHES),
            Inches(_TITLE_TOP_INCHES),
            Inches(_SLIDE_CONTENT_WIDTH_INCHES),
            Inches(_TITLE_HEIGHT_INCHES),
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = False
        title_frame.clear()
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = str(title)
        title_paragraph.space_after = Pt(0)
        title_paragraph.space_before = Pt(0)
        apply_run_style(first_run(title_paragraph), font_size=22, bold=True)
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(_SLIDE_LEFT_INCHES),
                Inches(_SUBTITLE_TOP_INCHES),
                Inches(_SLIDE_CONTENT_WIDTH_INCHES),
                Inches(_SUBTITLE_HEIGHT_INCHES),
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = False
            subtitle_frame.clear()
            subtitle_paragraph = subtitle_frame.paragraphs[0]
            subtitle_paragraph.text = str(subtitle)
            subtitle_paragraph.space_after = Pt(0)
            subtitle_paragraph.space_before = Pt(0)
            apply_run_style(
                first_run(subtitle_paragraph),
                font_size=10,
                color="#4B5563",
            )

    def add_text_block(
        slide,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        lines: Sequence[str],
        font_size: float,
    ) -> None:
        textbox = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        frame = textbox.text_frame
        frame.word_wrap = True
        frame.clear()
        frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        for index, line in enumerate(lines):
            paragraph = (
                frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            )
            paragraph.text = str(line)
            paragraph.space_after = Pt(0)
            paragraph.space_before = Pt(0)
            apply_run_style(first_run(paragraph), font_size=font_size)

    def add_cover_slide(
        title: str, subtitle: str, lines: Sequence[str]
    ) -> None:
        slide = presentation.slides.add_slide(blank_layout)
        set_slide_background(slide)
        add_title(slide, title, subtitle)
        add_text_block(
            slide,
            left=_SLIDE_LEFT_INCHES,
            top=1.45,
            width=_SLIDE_CONTENT_WIDTH_INCHES,
            height=4.9,
            lines=lines,
            font_size=15,
        )
        register_slide(title)

    def add_picture(
        slide,
        image_path: Path,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        fitted_left, fitted_top, fitted_width, fitted_height = (
            _fit_image_in_box(
                image_path,
                left=left,
                top=top,
                max_width=width,
                max_height=height,
            )
        )
        slide.shapes.add_picture(
            str(image_path),
            Inches(fitted_left),
            Inches(fitted_top),
            width=Inches(fitted_width),
            height=Inches(fitted_height),
        )

    def add_table_header_rule(
        slide,
        *,
        left: float,
        top: float,
        width: float,
    ) -> None:
        rule = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(_THICK_RULE_HEIGHT_INCHES),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = rgb_color(settings.table_rule_color)
        rule.line.fill.background()

    def style_table_cell(
        cell,
        *,
        text: str,
        font_size: float,
        fill_color: str,
        bold: bool = False,
        align=None,
    ) -> None:
        cell.text = str(text)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb_color(fill_color)
        cell.margin_left = Inches(0.035)
        cell.margin_right = Inches(0.035)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        frame = cell.text_frame
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT if align is None else align
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)
        apply_run_style(first_run(paragraph), font_size=font_size, bold=bold)

    def resolve_alignment(value: str | None):
        if value == "center":
            return PP_ALIGN.CENTER
        if value == "right":
            return PP_ALIGN.RIGHT
        return PP_ALIGN.LEFT

    def add_text_section(section: _TextSection) -> None:
        pages = section.pages or [[]]
        if not pages:
            pages = [[]]
        total_pages = len(pages)
        for page_index, page in enumerate(pages, start=1):
            slide = presentation.slides.add_slide(blank_layout)
            set_slide_background(slide)
            effective_title = _page_title(
                section.title, page_index, total_pages
            )
            add_title(slide, effective_title, section.subtitle)
            content = page or [section.placeholder]
            add_text_block(
                slide,
                left=_SLIDE_LEFT_INCHES,
                top=_TEXT_TOP_INCHES,
                width=_SLIDE_CONTENT_WIDTH_INCHES,
                height=_TEXT_HEIGHT_INCHES,
                lines=content,
                font_size=section.font_size,
            )
            register_slide(effective_title)

    def add_figure_section(
        section: _FigureSection,
        temporary_dir: Path,
    ) -> None:
        slide = presentation.slides.add_slide(blank_layout)
        set_slide_background(slide)
        add_title(slide, section.title, section.subtitle)
        if section.figure is None:
            add_text_block(
                slide,
                left=_SLIDE_LEFT_INCHES,
                top=2.4,
                width=_SLIDE_CONTENT_WIDTH_INCHES,
                height=1.0,
                lines=[section.placeholder],
                font_size=15,
            )
            register_slide(section.title)
            return
        image_path = temporary_dir / f"{_slugify(section.title)}.png"
        _save_figure_image(
            section.figure,
            image_path,
            font_family=settings.font_family,
        )
        add_picture(
            slide,
            image_path,
            left=_SLIDE_LEFT_INCHES,
            top=_FIGURE_TOP_INCHES,
            width=_SLIDE_CONTENT_WIDTH_INCHES,
            height=_FIGURE_HEIGHT_INCHES,
        )
        register_slide(section.title)

    def add_table_section(section: _TableSection) -> None:
        chunks = _table_row_chunks(section.rows, section.rows_per_slide)
        if not chunks:
            chunks = [[]]
        total_pages = len(chunks)
        for page_index, chunk in enumerate(chunks, start=1):
            slide = presentation.slides.add_slide(blank_layout)
            set_slide_background(slide)
            effective_title = _page_title(
                section.title, page_index, total_pages
            )
            add_title(slide, effective_title, section.subtitle)
            if not chunk:
                add_text_block(
                    slide,
                    left=_SLIDE_LEFT_INCHES,
                    top=2.4,
                    width=_SLIDE_CONTENT_WIDTH_INCHES,
                    height=1.0,
                    lines=["No data are available for this section."],
                    font_size=15,
                )
                register_slide(effective_title)
                continue
            add_table_header_rule(
                slide,
                left=_SLIDE_LEFT_INCHES,
                top=_TABLE_TOP_INCHES,
                width=_SLIDE_CONTENT_WIDTH_INCHES,
            )
            table_height = _TABLE_HEIGHT_INCHES - (
                0.25 if section.note else 0.0
            )
            table_shape = slide.shapes.add_table(
                len(chunk) + 1,
                len(section.columns),
                Inches(_SLIDE_LEFT_INCHES),
                Inches(_TABLE_TOP_INCHES),
                Inches(_SLIDE_CONTENT_WIDTH_INCHES),
                Inches(table_height),
            )
            table = table_shape.table
            column_widths = _resolve_column_widths(
                section.columns,
                chunk,
                total_width=_SLIDE_CONTENT_WIDTH_INCHES,
                column_width_weights=section.column_width_weights,
            )
            for column_index, column_width in enumerate(column_widths):
                table.columns[column_index].width = Inches(column_width)
            row_height = table_height / max(len(chunk) + 1, 1)
            for row_index in range(len(chunk) + 1):
                table.rows[row_index].height = Inches(row_height)
            for column_index, column_name in enumerate(section.columns):
                style_table_cell(
                    table.cell(0, column_index),
                    text=column_name,
                    font_size=section.header_font_size,
                    fill_color=settings.table_header_fill,
                    bold=True,
                    align=PP_ALIGN.CENTER,
                )
            for row_index, row in enumerate(chunk, start=1):
                row_fill = (
                    settings.table_even_row_fill
                    if row_index % 2 == 1
                    else settings.table_odd_row_fill
                )
                for column_index, value in enumerate(row):
                    alignment = None
                    if section.alignments is not None and column_index < len(
                        section.alignments
                    ):
                        alignment = resolve_alignment(
                            section.alignments[column_index]
                        )
                    style_table_cell(
                        table.cell(row_index, column_index),
                        text=str(value),
                        font_size=section.row_font_size,
                        fill_color=row_fill,
                        align=alignment,
                    )
            if section.note and page_index == total_pages:
                add_text_block(
                    slide,
                    left=_SLIDE_LEFT_INCHES,
                    top=_TABLE_TOP_INCHES + table_height + 0.08,
                    width=_SLIDE_CONTENT_WIDTH_INCHES,
                    height=0.18,
                    lines=[section.note],
                    font_size=9.5,
                )
            register_slide(effective_title)

    output_path = normalize_powerpoint_output_path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    appended_to_existing = output_path.is_file()
    if appended_to_existing:
        presentation = Presentation(str(output_path))
    else:
        presentation = Presentation()
        presentation.slide_width = Inches(_SLIDE_WIDTH_INCHES)
        presentation.slide_height = Inches(_SLIDE_HEIGHT_INCHES)
    blank_layout = _best_blank_layout(presentation)
    initial_slide_count = len(presentation.slides)
    total_slide_count = (
        1
        + sum(max(len(section.pages), 1) for section in text_sections)
        + len(figure_sections)
        + sum(
            max(
                len(_table_row_chunks(section.rows, section.rows_per_slide)), 1
            )
            for section in table_sections
        )
    )
    progress = _ReportProgressTracker(
        total_slide_count + 2,
        progress_callback,
        opening_message=f"Generating {report_title}. Please wait...",
    )
    progress.advance(
        "Opened existing PowerPoint report."
        if appended_to_existing
        else "Created new PowerPoint report."
    )
    slide_index = 0

    def register_slide(message: str) -> None:
        nonlocal slide_index
        slide_index += 1
        progress.advance(
            f"Built slide {slide_index}/{total_slide_count}: {message}"
        )

    with TemporaryDirectory() as temporary_directory:
        temporary_dir = Path(temporary_directory)
        add_cover_slide(report_title, cover_subtitle, cover_lines)
        for section in text_sections:
            add_text_section(section)
        for section in figure_sections:
            add_figure_section(section, temporary_dir)
        for section in table_sections:
            add_table_section(section)
        presentation.save(str(output_path))
    progress.advance("Saved PowerPoint report.")
    return ClusterPowerPointExportResult(
        report_path=output_path,
        appended_to_existing=appended_to_existing,
        added_slide_count=len(presentation.slides) - initial_slide_count,
    )


def _best_blank_layout(presentation):
    for layout in presentation.slide_layouts:
        try:
            if len(layout.placeholders) == 0:
                return layout
        except Exception:
            continue
    if len(presentation.slide_layouts) > 6:
        return presentation.slide_layouts[6]
    return presentation.slide_layouts[-1]


def _load_pptx_api():
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError(
            "PowerPoint export requires the optional dependency "
            "`python-pptx`. Install it and retry."
        ) from exc
    return (
        Presentation,
        Inches,
        Pt,
        RGBColor,
        PP_ALIGN,
        MSO_VERTICAL_ANCHOR,
        MSO_AUTO_SHAPE_TYPE,
    )


def _normalized_settings(
    settings: PowerPointExportSettings | None,
) -> PowerPointExportSettings:
    if settings is None:
        return PowerPointExportSettings()
    return PowerPointExportSettings.from_dict(settings.to_dict())


def _paginate_text_lines(
    lines: Sequence[str],
    *,
    max_lines: int,
    wrap_at: int,
) -> list[list[str]]:
    if max_lines <= 0:
        raise ValueError("max_lines must be positive")
    wrapped_lines: list[str] = []
    for raw_line in lines:
        line = str(raw_line).strip()
        if not line:
            wrapped_lines.append("")
            continue
        wrapped_lines.extend(
            textwrap.wrap(
                line,
                width=max(int(wrap_at), 1),
                break_long_words=True,
                break_on_hyphens=False,
            )
            or [line]
        )
    if not wrapped_lines:
        return [[]]
    return [
        wrapped_lines[index : index + max_lines]
        for index in range(0, len(wrapped_lines), max_lines)
    ]


def _table_row_chunks(
    rows: Sequence[Sequence[str]],
    rows_per_slide: int,
) -> list[list[list[str]]]:
    if rows_per_slide <= 0:
        raise ValueError("rows_per_slide must be positive")
    if not rows:
        return []
    return [
        [list(value) for value in rows[index : index + rows_per_slide]]
        for index in range(0, len(rows), rows_per_slide)
    ]


def _page_title(title: str, page_index: int, total_pages: int) -> str:
    if total_pages <= 1:
        return title
    return f"{title} ({page_index}/{total_pages})"


def _resolve_column_widths(
    columns: Sequence[str],
    rows: Sequence[Sequence[str]],
    *,
    total_width: float,
    column_width_weights: Sequence[float] | None = None,
) -> list[float]:
    if not columns:
        return []
    if column_width_weights is not None:
        weights = [max(float(weight), 0.2) for weight in column_width_weights]
    else:
        weights = []
        for column_index, column_name in enumerate(columns):
            max_length = len(str(column_name))
            for row in rows:
                if column_index >= len(row):
                    continue
                max_length = max(max_length, len(str(row[column_index])))
            weights.append(max(0.8, min(float(max_length) ** 0.78, 3.2)))
    if len(weights) != len(columns):
        raise ValueError("column width weights must match the column count")
    scale = total_width / sum(weights)
    widths = [weight * scale for weight in weights]
    widths[-1] += total_width - sum(widths)
    return widths


def _fit_image_in_box(
    image_path: Path,
    *,
    left: float,
    top: float,
    max_width: float,
    max_height: float,
) -> tuple[float, float, float, float]:
    try:
        from PIL import Image
    except ImportError:
        return left, top, max_width, max_height

    with Image.open(image_path) as image:
        width_px, height_px = image.size
    if width_px <= 0 or height_px <= 0:
        return left, top, max_width, max_height
    image_ratio = width_px / height_px
    box_ratio = max_width / max_height if max_height > 0 else image_ratio
    if image_ratio >= box_ratio:
        fitted_width = max_width
        fitted_height = fitted_width / image_ratio
    else:
        fitted_height = max_height
        fitted_width = fitted_height * image_ratio
    fitted_left = left + (max_width - fitted_width) / 2.0
    fitted_top = top + (max_height - fitted_height) / 2.0
    return fitted_left, fitted_top, fitted_width, fitted_height


def _save_figure_image(
    figure: Figure,
    output_path: Path,
    *,
    font_family: str,
) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    text_artists = list(figure.findobj(match=Text))
    original_font_families = [
        artist.get_fontfamily() for artist in text_artists
    ]
    try:
        for artist in text_artists:
            artist.set_fontfamily(font_family)
        if figure.canvas is not None:
            figure.canvas.draw()
        figure.savefig(
            output_path,
            format="png",
            dpi=220,
            bbox_inches="tight",
            facecolor="white",
            edgecolor="white",
        )
    finally:
        for artist, font_family_value in zip(
            text_artists,
            original_font_families,
            strict=False,
        ):
            artist.set_fontfamily(font_family_value)
        if figure.canvas is not None:
            figure.canvas.draw_idle()


def _format_optional_float(value: float | None) -> str:
    if value is None:
        return "n/a"
    return f"{value:.6g}"


def _slugify(value: str) -> str:
    safe = "".join(
        character.lower() if character.isalnum() else "_"
        for character in value
    )
    return safe.strip("_") or "figure"
