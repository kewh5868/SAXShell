from __future__ import annotations

import csv
import os
from pathlib import Path

import numpy as np
import pytest
from PySide6.QtWidgets import QApplication

import saxshell.clusterdynamics.cli as clusterdynamics_cli_module
from saxshell import saxshell as saxshell_module
from saxshell.clusterdynamics import (
    ClusterDynamicsWorkflow,
    load_cluster_dynamics_dataset,
    save_cluster_dynamics_dataset,
)
from saxshell.clusterdynamics.ui.main_window import ClusterDynamicsMainWindow
from saxshell.saxs.project_manager import (
    SAXSProjectManager,
    build_project_paths,
)

ATOM_TYPE_DEFINITIONS = {
    "node": [("Pb", None)],
    "linker": [("I", None)],
    "shell": [("O", None)],
}
PAIR_CUTOFFS = {
    ("Pb", "I"): {0: 1.7},
    ("Pb", "O"): {1: 1.3},
}


@pytest.fixture(scope="module")
def qapp():
    os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
    app = QApplication.instance()
    if app is None:
        app = QApplication([])
    yield app


def _connected_xyz_lines() -> str:
    return (
        "5\n"
        "frame_connected\n"
        "Pb 0.0 0.0 0.0\n"
        "I 1.0 0.0 0.0\n"
        "Pb 2.0 0.0 0.0\n"
        "O 0.2 1.0 0.0\n"
        "H 0.2 1.7 0.0\n"
    )


def _disconnected_xyz_lines() -> str:
    return (
        "5\n"
        "frame_disconnected\n"
        "Pb 0.0 0.0 0.0\n"
        "I 5.0 0.0 0.0\n"
        "Pb 10.0 0.0 0.0\n"
        "O 0.2 1.0 0.0\n"
        "H 0.2 1.7 0.0\n"
    )


def _build_frames_dir(tmp_path: Path) -> Path:
    frames_dir = tmp_path / "splitxyz0001"
    frames_dir.mkdir()
    sequence = (
        _disconnected_xyz_lines(),
        _connected_xyz_lines(),
        _connected_xyz_lines(),
        _disconnected_xyz_lines(),
        _connected_xyz_lines(),
        _disconnected_xyz_lines(),
    )
    for index, content in enumerate(sequence):
        (frames_dir / f"frame_{index:04d}.xyz").write_text(content)
    return frames_dir


def _build_offset_frames_dir(tmp_path: Path) -> Path:
    frames_dir = tmp_path / "splitxyz_f847fs"
    frames_dir.mkdir()
    sequence = (
        _connected_xyz_lines(),
        _disconnected_xyz_lines(),
    )
    for index, content in enumerate(sequence, start=1866):
        (frames_dir / f"frame_{index:04d}.xyz").write_text(content)
    return frames_dir


def _write_energy_file(tmp_path: Path) -> Path:
    energy_path = tmp_path / "traj.ener"
    lines = [
        f"{step} {step * 10.0:.1f} {1.0 + step:.3f} {300.0 + step:.3f} {-20.0 - step:.3f}\n"
        for step in range(6)
    ]
    energy_path.write_text("".join(lines))
    return energy_path


def _presentation_text(presentation) -> str:
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text)
    return "\n".join(texts)


def test_cluster_dynamics_workflow_bins_clusters_and_lifetimes(tmp_path):
    frames_dir = _build_frames_dir(tmp_path)
    energy_path = _write_energy_file(tmp_path)
    workflow = ClusterDynamicsWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        shell_levels=(1,),
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=3,
        energy_file=energy_path,
    )

    preview = workflow.preview_selection()
    result = workflow.analyze()

    assert preview.selected_frames == 6
    assert preview.frames_per_colormap_timestep == 3
    assert preview.colormap_timestep_fs == pytest.approx(30.0)
    assert preview.bin_count == 2
    assert preview.analysis_start_fs == 0.0
    assert preview.analysis_stop_fs == 60.0
    assert result.bin_count == 2
    assert result.analyzed_frames == 6
    assert result.cluster_labels == ("I", "Pb", "Pb2I")

    label_index = {
        label: index for index, label in enumerate(result.cluster_labels)
    }
    assert result.raw_count_matrix[label_index["Pb2I"], :].tolist() == [
        2.0,
        1.0,
    ]
    assert result.raw_count_matrix[label_index["Pb"], :].tolist() == [2.0, 4.0]
    assert result.raw_count_matrix[label_index["I"], :].tolist() == [1.0, 2.0]

    pb2i_summary = next(
        entry for entry in result.lifetime_by_label if entry.label == "Pb2I"
    )
    assert pb2i_summary.cluster_size == 3
    assert pb2i_summary.completed_lifetime_count == 2
    assert pb2i_summary.window_truncated_lifetime_count == 0
    assert pb2i_summary.mean_lifetime_fs == pytest.approx(15.0)
    assert pb2i_summary.std_lifetime_fs == pytest.approx(5.0)
    assert pb2i_summary.association_events == 2
    assert pb2i_summary.dissociation_events == 2

    size1_summary = next(
        entry for entry in result.lifetime_by_size if entry.cluster_size == 1
    )
    assert size1_summary.completed_lifetime_count == 3
    assert size1_summary.window_truncated_lifetime_count == 6
    assert size1_summary.mean_lifetime_fs == pytest.approx(10.0)
    assert size1_summary.std_lifetime_fs == pytest.approx(0.0)

    assert result.energy_data is not None
    x_values, y_values, label = result.energy_series("temperature")
    assert len(x_values) == 6
    assert y_values.tolist() == pytest.approx(
        [300.0, 301.0, 302.0, 303.0, 304.0, 305.0]
    )
    assert label == "Temperature (K)"


def test_cluster_dynamics_preview_honors_explicit_time_window(tmp_path):
    frames_dir = _build_frames_dir(tmp_path)
    workflow = ClusterDynamicsWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=2,
        analysis_start_fs=10.0,
        analysis_stop_fs=40.0,
    )

    preview = workflow.preview_selection()

    assert preview.selected_frames == 4
    assert preview.first_selected_frame == "frame_0001.xyz"
    assert preview.last_selected_frame == "frame_0004.xyz"
    assert preview.first_selected_time_fs == pytest.approx(10.0)
    assert preview.last_selected_time_fs == pytest.approx(40.0)
    assert preview.colormap_timestep_fs == pytest.approx(20.0)
    assert preview.bin_count == 2


def test_cluster_dynamics_uses_frame_filename_indices_for_time_axis(tmp_path):
    frames_dir = _build_offset_frames_dir(tmp_path)
    workflow = ClusterDynamicsWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        frame_timestep_fs=0.5,
        frames_per_colormap_timestep=1,
    )

    preview = workflow.preview_selection()

    assert preview.folder_start_time_fs == pytest.approx(847.0)
    assert preview.first_frame_time_fs == pytest.approx(933.0)
    assert preview.first_selected_source_frame_index == 1866
    assert preview.last_selected_source_frame_index == 1867
    assert preview.time_source_label == "Frame filenames x timestep"
    assert any(
        "847.000 fs" in message and "933.000 fs" in message
        for message in preview.time_warnings
    )


def test_cluster_dynamics_dataset_round_trip(tmp_path):
    frames_dir = _build_frames_dir(tmp_path)
    energy_path = _write_energy_file(tmp_path)
    result = ClusterDynamicsWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        shell_levels=(1,),
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=3,
        energy_file=energy_path,
    ).analyze()

    saved = save_cluster_dynamics_dataset(
        result,
        tmp_path / "cluster_dynamics_saved.json",
        analysis_settings={
            "frame_timestep_fs": 10.0,
            "frames_per_colormap_timestep": 3,
            "project_dir": str(tmp_path),
        },
    )
    loaded = load_cluster_dynamics_dataset(saved.dataset_file)

    assert saved.dataset_file.exists()
    assert any(
        path.name.endswith("_cluster_distribution.csv")
        for path in saved.written_files
    )
    assert loaded.result.cluster_labels == result.cluster_labels
    assert (
        loaded.result.raw_count_matrix.shape == result.raw_count_matrix.shape
    )
    assert np.array_equal(
        loaded.result.raw_count_matrix, result.raw_count_matrix
    )
    assert loaded.analysis_settings["frame_timestep_fs"] == 10.0
    assert loaded.analysis_settings["frames_per_colormap_timestep"] == 3
    assert loaded.result.energy_data is not None


def test_cluster_dynamics_main_window_updates_preview_for_xyz_frames(
    qapp,
    tmp_path,
):
    del qapp
    frames_dir = _build_frames_dir(tmp_path)

    window = ClusterDynamicsMainWindow(initial_frames_dir=frames_dir)
    window.time_panel.frame_timestep_spin.setValue(10.0)
    window.time_panel.frames_per_colormap_timestep_spin.setValue(3)

    preview_text = window.run_panel.selection_box.toPlainText()

    assert window.trajectory_panel.mode_label.text() == "Mode: XYZ frames"
    assert window.run_panel.title() == "Run Analysis"
    assert window.dataset_panel.title() == "Saved Results"
    assert (
        window.dataset_panel.save_dataset_button.text()
        == "Save Current Result"
    )
    assert (
        window.dataset_panel.load_dataset_button.text() == "Open Saved Result"
    )
    assert window.definitions_panel.title() == "Cluster Definitions (XYZ mode)"
    assert window.definitions_panel.parentWidget() is not None
    assert "Frames selected: 6" in preview_text
    assert "Time bins: 2" in preview_text
    assert "Frame timestep: 10.000 fs" in preview_text
    assert "Frames per colormap timestep: 3" in preview_text
    assert "Colormap timestep: 30.000 fs" in preview_text
    assert window.time_panel.colormap_timestep_value.text() == "30.000"
    window.close()


def test_cluster_dynamics_main_window_inherits_project_dir_and_start_time(
    qapp,
    tmp_path,
):
    del qapp
    frames_dir = _build_offset_frames_dir(tmp_path)

    window = ClusterDynamicsMainWindow(
        initial_frames_dir=frames_dir,
        initial_project_dir=tmp_path,
    )

    preview_text = window.run_panel.selection_box.toPlainText()

    assert window.dataset_panel.project_dir() == tmp_path
    assert window.time_panel.frame_timestep_fs() == pytest.approx(0.5)
    assert window.time_panel.frames_per_colormap_timestep() == 1
    assert window.time_panel.colormap_timestep_fs() == pytest.approx(0.5)
    assert window.time_panel.folder_start_time_fs() == pytest.approx(847.0)
    assert "Time source: Frame filenames x timestep" in preview_text
    assert "Colormap timestep: 0.500 fs" in preview_text
    assert "Source frame index range: 1866 to 1867" in preview_text
    window.close()


def test_cluster_dynamics_main_window_exports_colormap_and_lifetime_csv(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    frames_dir = _build_frames_dir(tmp_path)
    result = ClusterDynamicsWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        shell_levels=(1,),
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=2,
    ).analyze()

    window = ClusterDynamicsMainWindow(initial_frames_dir=frames_dir)
    window._last_result = result
    window.plot_panel.set_result(result)
    window.plot_panel.display_mode_combo.setCurrentIndex(
        window.plot_panel.display_mode_combo.findData("mean_count")
    )
    window.plot_panel.time_unit_combo.setCurrentIndex(
        window.plot_panel.time_unit_combo.findData("ps")
    )

    colormap_path = tmp_path / "exported_colormap.csv"
    lifetime_path = tmp_path / "exported_lifetime.csv"
    selected_paths = iter((str(colormap_path), str(lifetime_path)))
    monkeypatch.setattr(
        "saxshell.clusterdynamics.ui.main_window.QFileDialog.getSaveFileName",
        lambda *args, **kwargs: (next(selected_paths), "CSV Files (*.csv)"),
    )

    window.save_colormap_data()
    window.save_lifetime_table()

    with colormap_path.open(newline="", encoding="utf-8") as handle:
        colormap_rows = list(csv.DictReader(handle))
    with lifetime_path.open(newline="", encoding="utf-8") as handle:
        lifetime_rows = list(csv.DictReader(handle))

    pb2i_lifetime = next(
        row for row in lifetime_rows if row["label"] == "Pb2I"
    )

    assert (
        window.dataset_panel.save_colormap_button.text()
        == "Save Colormap Data"
    )
    assert (
        window.dataset_panel.save_lifetime_button.text()
        == "Save Lifetime Table"
    )
    assert len(colormap_rows) == len(result.cluster_labels) * result.bin_count
    assert colormap_rows[0]["display_mode"] == "mean_count"
    assert colormap_rows[0]["time_unit"] == "ps"
    assert "colormap_value" in colormap_rows[0]
    assert pb2i_lifetime["mean_lifetime_fs"] == "15"
    assert pb2i_lifetime["std_lifetime_fs"] == "5"
    window.close()


def test_cluster_dynamics_main_window_exports_powerpoint_report(
    qapp,
    tmp_path,
    monkeypatch,
):
    pytest.importorskip("pptx")
    from pptx import Presentation

    del qapp
    frames_dir = _build_frames_dir(tmp_path)
    project_dir = tmp_path / "saxs_project"
    SAXSProjectManager().create_project(project_dir)
    result = ClusterDynamicsWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        shell_levels=(1,),
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=2,
    ).analyze()

    window = ClusterDynamicsMainWindow(
        initial_frames_dir=frames_dir,
        initial_project_dir=project_dir,
    )
    window._last_result = result
    window.plot_panel.set_result(result)
    window.run_panel.set_selection_summary(
        window._format_preview_text(result.preview)
    )
    window._populate_summary_box(result)

    report_path = (
        build_project_paths(project_dir).reports_dir
        / f"{project_dir.name}_results.pptx"
    )
    captured_default_path: dict[str, str] = {}

    def fake_get_save_file_name(*args, **kwargs):
        captured_default_path["value"] = str(args[2])
        return (str(report_path), "PowerPoint Files (*.pptx)")

    monkeypatch.setattr(
        "saxshell.clusterdynamics.ui.main_window.QFileDialog.getSaveFileName",
        fake_get_save_file_name,
    )

    window.save_powerpoint_report()

    presentation = Presentation(str(report_path))

    assert captured_default_path["value"] == str(report_path)
    assert (
        window.dataset_panel.save_powerpoint_button.text()
        == "Save PowerPoint Report"
    )
    assert len(presentation.slides) >= 5
    assert "ClusterDynamics Report" in _presentation_text(presentation)
    assert "Observed Cluster Lifetimes" in _presentation_text(presentation)
    window.close()


def test_saxshell_cli_forwards_to_clusterdynamics_subcommand(monkeypatch):
    captured: dict[str, object] = {}

    def fake_clusterdynamics_main(argv=None):
        captured["argv"] = argv
        return 27

    monkeypatch.setattr(
        clusterdynamics_cli_module,
        "main",
        fake_clusterdynamics_main,
    )

    exit_code = saxshell_module.main(
        ["clusterdynamics", "--", "frames", "--energy-file", "traj.ener"]
    )

    assert exit_code == 27
    assert captured["argv"] == ["frames", "--energy-file", "traj.ener"]
