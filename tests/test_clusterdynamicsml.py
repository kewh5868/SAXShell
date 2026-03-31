from __future__ import annotations

import csv
import json
import os
import time
from pathlib import Path
from threading import Event
from types import SimpleNamespace

import numpy as np
import pytest
from matplotlib import rcParams
from PySide6.QtCore import QPoint, QPointF, Qt
from PySide6.QtGui import QWheelEvent
from PySide6.QtTest import QTest
from PySide6.QtWidgets import QApplication

import saxshell.clusterdynamicsml.cli as clusterdynamicsml_cli_module
import saxshell.clusterdynamicsml.workflow as clusterdynamicsml_workflow_module
from saxshell import saxshell as saxshell_module
from saxshell.clusterdynamicsml import (
    ClusterDynamicsMLWorkflow,
    load_cluster_dynamicsai_dataset,
    save_cluster_dynamicsai_dataset,
)
from saxshell.clusterdynamicsml.ui.main_window import (
    _UI_REFRESH_DELAY_MS,
    ClusterDynamicsMLMainWindow,
    ClusterDynamicsMLSettingsPanel,
    _combined_model_weight_rows,
)
from saxshell.clusterdynamicsml.ui.plot_panel import (
    _build_population_histogram_payload,
    _distribution_entries,
)
from saxshell.saxs.debye import (
    compute_debye_intensity,
    compute_debye_intensity_with_debye_waller,
)
from saxshell.saxs.project_manager import (
    SAXSProjectManager,
    build_project_paths,
)
from saxshell.saxs.project_manager.prior_plot import (
    build_prior_histogram_export_payload,
    list_secondary_filter_elements,
)

ATOM_TYPE_DEFINITIONS = {
    "node": [("Pb", None)],
    "linker": [("I", None)],
}
PAIR_CUTOFFS = {
    ("Pb", "I"): {0: 1.2},
}


@pytest.fixture(scope="module")
def qapp():
    os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
    app = QApplication.instance()
    if app is None:
        app = QApplication([])
    yield app


def _disconnected_xyz_lines() -> str:
    return (
        "5\n"
        "frame_disconnected\n"
        "Pb 0.0 0.0 0.0\n"
        "Pb 10.0 0.0 0.0\n"
        "Pb 20.0 0.0 0.0\n"
        "I 30.0 0.0 0.0\n"
        "I 40.0 0.0 0.0\n"
    )


def _pair_xyz_lines() -> str:
    return (
        "5\n"
        "frame_pair\n"
        "Pb 0.0 0.0 0.0\n"
        "Pb 2.0 0.0 0.0\n"
        "Pb 10.0 0.0 0.0\n"
        "I 1.0 0.0 0.0\n"
        "I 40.0 0.0 0.0\n"
    )


def _triple_xyz_lines() -> str:
    return (
        "5\n"
        "frame_triple\n"
        "Pb 0.0 0.0 0.0\n"
        "Pb 2.0 0.0 0.0\n"
        "Pb 4.0 0.0 0.0\n"
        "I 1.0 0.0 0.0\n"
        "I 3.0 0.0 0.0\n"
    )


def _build_frames_dir(tmp_path: Path) -> Path:
    frames_dir = tmp_path / "splitxyz_f0fs"
    frames_dir.mkdir()
    sequence = (
        _disconnected_xyz_lines(),
        _pair_xyz_lines(),
        _triple_xyz_lines(),
        _triple_xyz_lines(),
        _pair_xyz_lines(),
        _disconnected_xyz_lines(),
    )
    for index, content in enumerate(sequence):
        (frames_dir / f"frame_{index:04d}.xyz").write_text(content)
    return frames_dir


def _build_clusters_dir(tmp_path: Path) -> Path:
    clusters_dir = tmp_path / "clusters_training"
    (clusters_dir / "Pb").mkdir(parents=True)
    (clusters_dir / "Pb2I").mkdir(parents=True)
    (clusters_dir / "Pb3I2").mkdir(parents=True)
    single = "1\n" "pb_single\n" "Pb 0.0 0.0 0.0\n"
    pair_a = (
        "3\n"
        "pb2i_pair_a\n"
        "Pb 0.0 0.0 0.0\n"
        "Pb 2.0 0.0 0.0\n"
        "I 1.0 0.0 0.0\n"
    )
    pair_b = (
        "3\n"
        "pb2i_pair_b\n"
        "Pb 0.0 0.1 0.0\n"
        "Pb 2.1 -0.1 0.0\n"
        "I 1.0 0.2 0.1\n"
    )
    triple_a = (
        "5\n"
        "pb3i2_chain_a\n"
        "Pb 0.0 0.0 0.0\n"
        "Pb 2.0 0.0 0.0\n"
        "Pb 4.0 0.0 0.0\n"
        "I 1.0 0.2 0.0\n"
        "I 3.0 -0.2 0.0\n"
    )
    triple_b = (
        "5\n"
        "pb3i2_chain_b\n"
        "Pb 0.0 0.0 0.0\n"
        "Pb 2.1 0.1 0.0\n"
        "Pb 4.2 0.0 0.1\n"
        "I 1.1 0.3 0.1\n"
        "I 3.1 -0.2 -0.1\n"
    )
    (clusters_dir / "Pb" / "pb_0001.xyz").write_text(single)
    (clusters_dir / "Pb" / "pb_0002.xyz").write_text(single)
    (clusters_dir / "Pb2I" / "pb2i_0001.xyz").write_text(pair_a)
    (clusters_dir / "Pb2I" / "pb2i_0002.xyz").write_text(pair_b)
    (clusters_dir / "Pb3I2" / "pb3i2_0001.xyz").write_text(triple_a)
    (clusters_dir / "Pb3I2" / "pb3i2_0002.xyz").write_text(triple_b)
    return clusters_dir


def _wait_for(
    qapp,
    predicate,
    *,
    timeout: float = 5.0,
) -> None:
    deadline = time.monotonic() + timeout
    while time.monotonic() < deadline:
        qapp.processEvents()
        if predicate():
            return
        time.sleep(0.01)
    qapp.processEvents()
    assert predicate()


def _make_wheel_event(widget) -> QWheelEvent:
    center = widget.rect().center()
    local_pos = QPointF(center)
    global_pos = QPointF(widget.mapToGlobal(center))
    return QWheelEvent(
        local_pos,
        global_pos,
        QPoint(0, 0),
        QPoint(0, 120),
        Qt.MouseButton.NoButton,
        Qt.KeyboardModifier.NoModifier,
        Qt.ScrollPhase.ScrollUpdate,
        False,
    )


def _build_node_only_label_clusters_dir(tmp_path: Path) -> Path:
    clusters_dir = tmp_path / "clusters_training_node_only_labels"
    (clusters_dir / "Pb").mkdir(parents=True)
    (clusters_dir / "Pb2").mkdir(parents=True)
    (clusters_dir / "Pb3").mkdir(parents=True)
    single = "1\n" "pb_single\n" "Pb 0.0 0.0 0.0\n"
    pair = (
        "3\n"
        "pb2i_pair\n"
        "Pb 0.0 0.0 0.0\n"
        "Pb 2.0 0.0 0.0\n"
        "I 1.0 0.0 0.0\n"
    )
    triple = (
        "5\n"
        "pb3i2_chain\n"
        "Pb 0.0 0.0 0.0\n"
        "Pb 2.0 0.0 0.0\n"
        "Pb 4.0 0.0 0.0\n"
        "I 1.0 0.2 0.0\n"
        "I 3.0 -0.2 0.0\n"
    )
    (clusters_dir / "Pb" / "pb_0001.xyz").write_text(single)
    (clusters_dir / "Pb2" / "pb2_0001.xyz").write_text(pair)
    (clusters_dir / "Pb2" / "pb2_0002.xyz").write_text(pair)
    (clusters_dir / "Pb3" / "pb3_0001.xyz").write_text(triple)
    (clusters_dir / "Pb3" / "pb3_0002.xyz").write_text(triple)
    return clusters_dir


def _build_clusters_dir_with_secondary_atoms(tmp_path: Path) -> Path:
    clusters_dir = tmp_path / "clusters_training_secondary"
    (clusters_dir / "Pb").mkdir(parents=True)
    (clusters_dir / "Pb2I").mkdir(parents=True)
    (clusters_dir / "Pb3I2").mkdir(parents=True)
    single_a = "1\n" "pb_single_a\n" "Pb 0.0 0.0 0.0\n"
    single_b = "2\n" "pb_single_b_o\n" "Pb 0.0 0.0 0.0\n" "O 1.5 0.0 0.0\n"
    pair_a = (
        "3\n"
        "pb2i_pair_a\n"
        "Pb 0.0 0.0 0.0\n"
        "Pb 2.0 0.0 0.0\n"
        "I 1.0 0.0 0.0\n"
    )
    pair_b = (
        "5\n"
        "pb2i_pair_b_o2\n"
        "Pb 0.0 0.1 0.0\n"
        "Pb 2.1 -0.1 0.0\n"
        "I 1.0 0.2 0.1\n"
        "O 1.0 1.8 0.0\n"
        "O 1.0 -1.8 0.0\n"
    )
    triple_a = (
        "6\n"
        "pb3i2_chain_a_o\n"
        "Pb 0.0 0.0 0.0\n"
        "Pb 2.0 0.0 0.0\n"
        "Pb 4.0 0.0 0.0\n"
        "I 1.0 0.2 0.0\n"
        "I 3.0 -0.2 0.0\n"
        "O 2.0 2.0 0.0\n"
    )
    triple_b = (
        "5\n"
        "pb3i2_chain_b\n"
        "Pb 0.0 0.0 0.0\n"
        "Pb 2.1 0.1 0.0\n"
        "Pb 4.2 0.0 0.1\n"
        "I 1.1 0.3 0.1\n"
        "I 3.1 -0.2 -0.1\n"
    )
    (clusters_dir / "Pb" / "pb_0001.xyz").write_text(single_a)
    (clusters_dir / "Pb" / "pb_0002.xyz").write_text(single_b)
    (clusters_dir / "Pb2I" / "pb2i_0001.xyz").write_text(pair_a)
    (clusters_dir / "Pb2I" / "pb2i_0002.xyz").write_text(pair_b)
    (clusters_dir / "Pb3I2" / "pb3i2_0001.xyz").write_text(triple_a)
    (clusters_dir / "Pb3I2" / "pb3i2_0002.xyz").write_text(triple_b)
    return clusters_dir


def _write_experimental_data_file(tmp_path: Path) -> Path:
    q_values = np.linspace(0.05, 1.0, 60)
    intensities = np.exp(-2.0 * q_values) + 0.15
    output = tmp_path / "experimental.txt"
    with output.open("w", encoding="utf-8") as handle:
        for q_value, intensity in zip(q_values, intensities, strict=False):
            handle.write(f"{q_value:.6f} {intensity:.8f}\n")
    return output


def _write_energy_file(tmp_path: Path, name: str = "traj.ener") -> Path:
    energy_path = tmp_path / name
    energy_path.write_text(
        "# step time kinetic temperature potential\n"
        "1 0.0 1.0 300.0 -10.0\n"
        "2 5.0 1.1 301.0 -10.1\n",
        encoding="utf-8",
    )
    return energy_path


def _build_project_dir(
    tmp_path: Path,
    *,
    clusters_dir: Path,
    experimental_data_file: Path,
) -> Path:
    manager = SAXSProjectManager()
    project_dir = tmp_path / "saxs_project"
    settings = manager.create_project(project_dir)
    settings.clusters_dir = str(clusters_dir)
    settings.experimental_data_path = str(experimental_data_file)
    settings.copied_experimental_data_file = str(experimental_data_file)
    manager.save_project(settings)
    return project_dir


def _write_component_profile_file(
    output_path: Path,
    *,
    q_values: np.ndarray,
    intensity: np.ndarray,
) -> None:
    data = np.column_stack(
        [
            np.asarray(q_values, dtype=float),
            np.asarray(intensity, dtype=float),
            np.zeros_like(intensity, dtype=float),
            np.zeros_like(intensity, dtype=float),
        ]
    )
    np.savetxt(
        output_path,
        data,
        comments="",
        header="# Columns: q, S(q)_avg, S(q)_std, S(q)_se\n",
        fmt=["%.8f", "%.8f", "%.8f", "%.8f"],
    )


def _write_project_component_artifacts(
    project_dir: Path,
    *,
    q_values: np.ndarray,
) -> None:
    paths = build_project_paths(project_dir)
    component_specs = {
        "Pb": {"motif_0000": np.full_like(q_values, 1.25, dtype=float)},
        "Pb2I": {"motif_0000": np.full_like(q_values, 2.75, dtype=float)},
        "Pb3I2": {"motif_0000": np.full_like(q_values, 4.50, dtype=float)},
    }
    saxs_map: dict[str, dict[str, str]] = {}
    prior_structures: dict[str, dict[str, object]] = {}
    for structure, motifs in component_specs.items():
        saxs_map[structure] = {}
        prior_structures[structure] = {}
        motif_count = len(motifs)
        for motif, intensity in motifs.items():
            filename = f"{structure}_{motif}.txt"
            _write_component_profile_file(
                paths.scattering_components_dir / filename,
                q_values=q_values,
                intensity=intensity,
            )
            saxs_map[structure][motif] = filename
            prior_structures[structure][motif] = {
                "count": 1,
                "weight": 1.0 / motif_count,
                "profile_file": filename,
            }
    (project_dir / "md_saxs_map.json").write_text(
        json.dumps({"saxs_map": saxs_map}, indent=2) + "\n",
        encoding="utf-8",
    )
    (project_dir / "md_prior_weights.json").write_text(
        json.dumps(
            {
                "structures": prior_structures,
                "q_range": {
                    "qmin": float(q_values.min()),
                    "qmax": float(q_values.max()),
                    "points": int(q_values.size),
                },
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )


def _presentation_text(presentation) -> str:
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text)
    return "\n".join(texts)


def test_clusterdynamicsml_prediction_panel_tooltips(qapp):
    del qapp
    panel = ClusterDynamicsMLSettingsPanel()

    assert "stoichiometry label" in panel.clusters_dir_edit.toolTip()
    assert "Leave this blank" in panel.experimental_data_edit.toolTip()
    assert "Lowest node count to predict" in panel.target_start_spin.toolTip()
    assert "Every integer node count" in panel.target_stop_spin.toolTip()
    assert (
        "ranked candidate stoichiometries" in panel.candidates_spin.toolTip()
    )
    assert (
        "share among the predicted candidates"
        in panel.share_threshold_spin.toolTip()
    )
    assert "no experimental data file is loaded" in panel.q_min_spin.toolTip()
    assert "no experimental data file is loaded" in panel.q_max_spin.toolTip()
    assert "fallback SAXS grid" in panel.q_points_spin.toolTip()
    assert "larger node counts to predict" in panel.toolTip()

    panel.close()


def test_clusterdynamicsml_workflow_predicts_larger_clusters(tmp_path):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)
    project_dir = _build_project_dir(
        tmp_path,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
    )

    workflow = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        project_dir=project_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    )

    preview = workflow.preview_selection()
    result = workflow.analyze()

    assert preview.structure_label_count == 3
    assert preview.observed_node_counts == (1, 2, 3)
    assert preview.target_node_counts == (4, 5)

    assert result.max_observed_node_count == 3
    assert result.max_predicted_node_count in {4, 5}
    assert {entry.target_node_count for entry in result.predictions} == {4, 5}
    assert all(
        entry.label.startswith("Pb4") or entry.label.startswith("Pb5")
        for entry in result.predictions
    )
    assert {entry.label for entry in result.predictions}.isdisjoint(
        {"Pb4", "Pb5"}
    )
    assert sum(
        entry.predicted_population_share for entry in result.predictions
    ) == pytest.approx(1.0)
    assert all(
        entry.predicted_population_share > 0.0 for entry in result.predictions
    )
    assert all(
        len(entry.generated_elements) == sum(entry.element_counts.values())
        for entry in result.predictions
    )
    assert result.saxs_comparison is not None
    assert (
        result.saxs_comparison.experimental_data_path
        == experimental_data_file.resolve()
    )
    assert result.saxs_comparison.rmse is not None
    assert len(result.saxs_comparison.component_weights) >= 3


def test_clusterdynamicsml_estimates_debye_waller_pairs_and_uses_them_for_predicted_traces(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4,),
    ).analyze()

    observed_rows = [
        row
        for row in result.debye_waller_estimates
        if row.source == "observed"
    ]
    predicted_rows = [
        row
        for row in result.debye_waller_estimates
        if row.source == "predicted"
    ]
    assert result.saxs_comparison is not None
    target_prediction = max(
        result.predictions,
        key=lambda entry: entry.predicted_population_share,
    )
    prediction_sigma_lookup = {
        (row.element_a, row.element_b): row.sigma
        for row in predicted_rows
        if row.label == target_prediction.label
    }
    included_pair_indices = (
        clusterdynamicsml_workflow_module._first_shell_pair_indices(
            target_prediction.generated_coordinates,
            list(target_prediction.generated_elements),
            node_elements={"Pb"},
            pair_cutoff_definitions=PAIR_CUTOFFS,
        )
    )
    classic_trace = compute_debye_intensity(
        target_prediction.generated_coordinates,
        list(target_prediction.generated_elements),
        result.saxs_comparison.q_values,
    )
    debye_waller_trace = compute_debye_intensity_with_debye_waller(
        target_prediction.generated_coordinates,
        list(target_prediction.generated_elements),
        result.saxs_comparison.q_values,
        pair_sigma_by_element=prediction_sigma_lookup,
        included_pair_indices=included_pair_indices,
    )
    predicted_component = next(
        entry
        for entry in result.saxs_comparison.component_weights
        if entry.label == target_prediction.label
        and entry.source == "predicted"
    )
    assert predicted_component.profile_path is not None
    saved_profile = np.loadtxt(predicted_component.profile_path, comments="#")

    assert any(
        row.label == "Pb3I2"
        and row.method == "ensemble"
        and (row.element_a, row.element_b) == ("Pb", "Pb")
        and row.sigma > 0.0
        for row in observed_rows
    )
    assert predicted_rows
    assert prediction_sigma_lookup
    assert np.any(debye_waller_trace < classic_trace)
    np.testing.assert_allclose(saved_profile[:, 1], debye_waller_trace)


def test_clusterdynamicsml_estimates_debye_waller_from_first_shell_pairs_only(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = tmp_path / "clusters_local_debye_waller"
    structure_dir = clusters_dir / "PbI2"
    structure_dir.mkdir(parents=True, exist_ok=True)
    structure_a = structure_dir / "pbi2_0001.xyz"
    structure_b = structure_dir / "pbi2_0002.xyz"
    structure_a.write_text(
        "3\npbi2_a\nPb 0.0 0.0 0.0\nI 1.0 0.0 0.0\nI 3.0 0.0 0.0\n",
        encoding="utf-8",
    )
    structure_b.write_text(
        "3\npbi2_b\nPb 0.0 0.0 0.0\nI 1.1 0.0 0.0\nI 5.0 0.0 0.0\n",
        encoding="utf-8",
    )

    workflow = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(2,),
    )
    observation = (
        clusterdynamicsml_workflow_module.ClusterDynamicsMLTrainingObservation(
            label="PbI2",
            node_count=1,
            cluster_size=1,
            element_counts={"Pb": 1, "I": 2},
            file_count=2,
            representative_path=structure_a,
            structure_dir=structure_dir,
            motifs=("no_motif",),
            mean_atom_count=3.0,
            mean_radius_of_gyration=1.0,
            mean_max_radius=1.0,
            mean_semiaxis_a=1.0,
            mean_semiaxis_b=1.0,
            mean_semiaxis_c=1.0,
            total_observations=2,
            occupied_frames=2,
            mean_count_per_frame=1.0,
            occupancy_fraction=1.0,
            association_events=0,
            dissociation_events=0,
            association_rate_per_ps=0.0,
            dissociation_rate_per_ps=0.0,
            completed_lifetime_count=2,
            window_truncated_lifetime_count=0,
            mean_lifetime_fs=10.0,
            std_lifetime_fs=0.0,
        )
    )

    estimates = workflow._estimate_observed_debye_waller_pairs([observation])
    pair_row = next(
        row
        for row in estimates
        if row.label == "PbI2"
        and (row.element_a, row.element_b) == ("I", "Pb")
    )

    assert pair_row.aligned_pair_count == 1
    assert pair_row.sigma == pytest.approx(0.05, abs=1e-6)


def test_clusterdynamicsml_structure_observations_use_reference_atom_counts(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_node_only_label_clusters_dir(tmp_path)

    workflow = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4,),
    )

    observations = workflow._build_structure_observations(clusters_dir)
    counts_by_label = {
        entry.label: entry.element_counts for entry in observations
    }

    assert counts_by_label["Pb"] == {"Pb": 1}
    assert counts_by_label["Pb2"] == {"Pb": 2, "I": 1}
    assert counts_by_label["Pb3"] == {"Pb": 3, "I": 2}


def test_clusterdynamicsml_learns_reference_geometry_statistics(tmp_path):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir_with_secondary_atoms(tmp_path)

    workflow = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions={
            **ATOM_TYPE_DEFINITIONS,
            "solvent": [("O", None)],
        },
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4,),
    )
    result = workflow.analyze()
    statistics = workflow._collect_training_geometry_statistics(
        list(result.training_observations)
    )

    assert statistics.atom_type_by_element["Pb"] == "node"
    assert statistics.atom_type_by_element["I"] == "linker"
    assert statistics.atom_type_by_element["O"] == "shell"
    assert statistics.node_bond_length == pytest.approx(2.05, abs=0.20)
    assert statistics.bond_length_medians[("I", "Pb")] == pytest.approx(
        1.03,
        abs=0.20,
    )
    assert statistics.node_angle_medians[("node", "node")] > 170.0
    assert statistics.node_coordination_medians["node"] == pytest.approx(
        1.0,
        abs=0.1,
    )
    assert statistics.node_coordination_medians["linker"] == pytest.approx(
        1.0,
        abs=0.1,
    )
    assert statistics.contact_distance_medians[("I", "I")] == pytest.approx(
        2.05,
        abs=0.20,
    )
    assert statistics.contact_distance_medians[("I", "O")] == pytest.approx(
        2.03,
        abs=0.30,
    )
    assert statistics.non_node_node_coordination_medians["I"] == pytest.approx(
        2.0,
        abs=0.1,
    )
    assert statistics.non_node_node_coordination_medians["O"] == pytest.approx(
        1.0,
        abs=0.1,
    )
    assert statistics.atom_coordination_medians[
        ("shell", "linker")
    ] == pytest.approx(
        1.0,
        abs=0.1,
    )
    assert statistics.atom_coordination_medians[("linker", "linker")] > 0.0
    assert statistics.atom_coordination_medians[("linker", "shell")] > 0.0


def test_clusterdynamicsml_predicted_structure_follows_reference_geometry(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)

    workflow = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4,),
    )
    result = workflow.analyze()
    statistics = workflow._collect_training_geometry_statistics(
        list(result.training_observations)
    )
    prediction = max(
        (
            entry
            for entry in result.predictions
            if entry.target_node_count == 4
            and entry.element_counts.get("I", 0) > 0
        ),
        key=lambda entry: entry.predicted_population_share,
    )

    node_indices = [
        index
        for index, element in enumerate(prediction.generated_elements)
        if element == "Pb"
    ]
    linker_indices = [
        index
        for index, element in enumerate(prediction.generated_elements)
        if element == "I"
    ]
    node_coordinates = np.asarray(
        prediction.generated_coordinates[node_indices],
        dtype=float,
    )
    linker_coordinates = np.asarray(
        prediction.generated_coordinates[linker_indices],
        dtype=float,
    )
    scaffold_edges = clusterdynamicsml_workflow_module._node_scaffold_edges(
        node_coordinates,
        ["Pb"] * len(node_coordinates),
        pair_cutoff_definitions=PAIR_CUTOFFS,
    )
    edge_lengths = np.asarray(
        [
            np.linalg.norm(
                node_coordinates[index_a] - node_coordinates[index_b]
            )
            for index_a, index_b in scaffold_edges
        ],
        dtype=float,
    )
    scaffold_adjacency = (
        clusterdynamicsml_workflow_module._adjacency_from_edges(
            len(node_coordinates),
            scaffold_edges,
        )
    )
    node_angles = [
        clusterdynamicsml_workflow_module._angle_between_vectors(
            node_coordinates[neighbors[0]] - node_coordinates[node_index],
            node_coordinates[neighbors[1]] - node_coordinates[node_index],
        )
        for node_index, neighbors in (
            (index, sorted(entries))
            for index, entries in scaffold_adjacency.items()
            if len(entries) == 2
        )
    ]
    pb_i_distance = statistics.bond_length_medians[("I", "Pb")]
    linker_distances = np.asarray(
        [
            np.linalg.norm(node_coordinates - coordinate, axis=1)
            for coordinate in linker_coordinates
        ],
        dtype=float,
    )
    linker_bridge_counts = [
        int(np.sum(distances <= pb_i_distance * 1.25))
        for distances in linker_distances
    ]

    assert len(scaffold_edges) == len(node_coordinates) - 1
    assert np.median(edge_lengths) == pytest.approx(
        statistics.node_bond_length,
        abs=0.35,
    )
    assert node_angles
    assert min(angle for angle in node_angles if angle is not None) > 140.0
    assert np.median(np.min(linker_distances, axis=1)) == pytest.approx(
        pb_i_distance,
        abs=0.25,
    )
    assert max(linker_bridge_counts) >= 2


def test_clusterdynamicsml_predicted_structure_respects_non_node_contacts(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir_with_secondary_atoms(tmp_path)

    workflow = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions={
            **ATOM_TYPE_DEFINITIONS,
            "solvent": [("O", None)],
        },
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4,),
    )
    result = workflow.analyze()
    statistics = workflow._collect_training_geometry_statistics(
        list(result.training_observations)
    )
    source_observation = max(
        result.training_observations,
        key=lambda row: row.node_count,
    )
    generated_elements, generated_coordinates = (
        workflow._generate_predicted_structure(
            source_observation,
            target_counts={"Pb": 4, "I": 3, "O": 2},
            predicted_max_radius=max(
                source_observation.mean_max_radius * 1.3, 1.0
            ),
            geometry_statistics=statistics,
        )
    )

    coordinates = np.asarray(generated_coordinates, dtype=float)
    linker_indices = [
        index
        for index, element in enumerate(generated_elements)
        if element == "I"
    ]
    shell_indices = [
        index
        for index, element in enumerate(generated_elements)
        if element == "O"
    ]
    linker_coordinates = np.asarray(coordinates[linker_indices], dtype=float)
    shell_coordinates = np.asarray(coordinates[shell_indices], dtype=float)
    linker_linker_distances = [
        float(
            np.linalg.norm(
                linker_coordinates[index_a] - linker_coordinates[index_b]
            )
        )
        for index_a in range(len(linker_coordinates))
        for index_b in range(index_a + 1, len(linker_coordinates))
    ]
    shell_to_linker_distances = [
        float(np.min(np.linalg.norm(linker_coordinates - coordinate, axis=1)))
        for coordinate in shell_coordinates
    ]
    linker_shell_cutoff = (
        clusterdynamicsml_workflow_module._contact_distance_cutoff(
            "I",
            "O",
            geometry_statistics=statistics,
            default_distance=statistics.node_bond_length,
            pair_cutoff_definitions=PAIR_CUTOFFS,
        )
    )
    linker_linker_cutoff = (
        clusterdynamicsml_workflow_module._contact_distance_cutoff(
            "I",
            "I",
            geometry_statistics=statistics,
            default_distance=statistics.node_bond_length,
            pair_cutoff_definitions=PAIR_CUTOFFS,
        )
    )
    linker_linker_contact_counts = [
        int(
            np.sum(
                np.linalg.norm(linker_coordinates - coordinate, axis=1)
                <= linker_linker_cutoff
            )
            - 1
        )
        for coordinate in linker_coordinates
    ]

    assert linker_indices
    assert shell_indices
    assert min(linker_linker_distances) == pytest.approx(
        statistics.contact_distance_medians[("I", "I")],
        abs=0.35,
    )
    assert np.median(
        np.asarray(shell_to_linker_distances, dtype=float)
    ) == pytest.approx(
        statistics.contact_distance_medians[("I", "O")],
        abs=0.35,
    )
    assert all(
        distance <= linker_shell_cutoff
        for distance in shell_to_linker_distances
    )
    assert sum(count >= 1 for count in linker_linker_contact_counts) >= 2


def test_clusterdynamicsml_histogram_binning_keeps_secondary_oxygen_out_of_labels(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir_with_secondary_atoms(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions={
            **ATOM_TYPE_DEFINITIONS,
            "solvent": [("O", None)],
        },
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    oxygen_predictions = [
        entry
        for entry in result.predictions
        if entry.element_counts.get("O", 0) > 0
    ]
    assert oxygen_predictions

    payload = _build_population_histogram_payload(
        result,
        include_predictions=True,
    )

    assert payload is not None
    assert "O" in payload["available_elements"]
    assert all("O" not in label for label in payload["structures"])

    predicted_payloads = [
        motif_payload
        for motifs in payload["structures"].values()
        for motif_name, motif_payload in motifs.items()
        if motif_name.startswith("predicted_rank_")
    ]
    assert any(
        "O" in payload.get("secondary_atom_distributions", {})
        and any(
            int(segment) > 0
            for segment in payload["secondary_atom_distributions"]["O"]
        )
        for payload in predicted_payloads
    )


def test_clusterdynamicsml_solvent_sort_histograms_stay_normalized_with_predictions(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir_with_secondary_atoms(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions={
            **ATOM_TYPE_DEFINITIONS,
            "solvent": [("O", None)],
        },
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    payload = _build_population_histogram_payload(
        result,
        include_predictions=True,
    )

    assert payload is not None

    structure_fraction = build_prior_histogram_export_payload(
        payload,
        mode="solvent_sort_structure_fraction",
        value_mode="percent",
        secondary_element="O",
    )
    atom_fraction = build_prior_histogram_export_payload(
        payload,
        mode="solvent_sort_atom_fraction",
        value_mode="percent",
        secondary_element="O",
    )

    assert float(np.sum(structure_fraction["totals"])) == pytest.approx(100.0)
    assert float(np.sum(atom_fraction["totals"])) == pytest.approx(100.0)
    assert float(np.max(structure_fraction["totals"])) <= 100.0 + 1.0e-9
    assert float(np.max(atom_fraction["totals"])) <= 100.0 + 1.0e-9


def test_clusterdynamicsml_predictions_keep_required_linkers_and_drop_tiny_tail(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir_with_secondary_atoms(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions={
            **ATOM_TYPE_DEFINITIONS,
            "solvent": [("O", None)],
        },
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    minimum_iodides = {4: 2, 5: 3}

    assert result.predictions
    assert {entry.target_node_count for entry in result.predictions} == {4, 5}
    assert {entry.label for entry in result.predictions}.isdisjoint(
        {"Pb4O4", "Pb5O5"}
    )
    assert all(
        entry.element_counts.get("I", 0)
        >= minimum_iodides.get(entry.target_node_count, 0)
        for entry in result.predictions
    )
    assert all(
        entry.predicted_population_share
        >= result.prediction_population_share_threshold
        for entry in result.predictions
    )


def test_clusterdynamicsml_secondary_atom_predictions_preserve_reference_bonds(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir_with_secondary_atoms(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    workflow = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions={
            **ATOM_TYPE_DEFINITIONS,
            "solvent": [("O", None)],
        },
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
        prediction_population_share_threshold=0.0,
    )
    result = workflow.analyze()
    statistics = workflow._collect_training_geometry_statistics(
        list(result.training_observations)
    )

    assert any(
        entry.element_counts.get("O", 0) > 0 for entry in result.predictions
    )

    for prediction in result.predictions:
        node_indices = [
            index
            for index, element in enumerate(prediction.generated_elements)
            if element == "Pb"
        ]
        node_coordinates = np.asarray(
            prediction.generated_coordinates[node_indices],
            dtype=float,
        )
        scaffold_edges = (
            clusterdynamicsml_workflow_module._node_scaffold_edges(
                node_coordinates,
                ["Pb"] * len(node_coordinates),
                pair_cutoff_definitions=PAIR_CUTOFFS,
            )
        )
        edge_lengths = np.asarray(
            [
                np.linalg.norm(
                    node_coordinates[index_a] - node_coordinates[index_b]
                )
                for index_a, index_b in scaffold_edges
            ],
            dtype=float,
        )
        assert edge_lengths.size > 0
        assert np.median(edge_lengths) == pytest.approx(
            statistics.node_bond_length,
            abs=0.20,
        )

        for element in {"I", "O"}:
            atom_indices = [
                index
                for index, atom_element in enumerate(
                    prediction.generated_elements
                )
                if atom_element == element
            ]
            if not atom_indices:
                continue
            nearest_node_distances = np.asarray(
                [
                    np.min(
                        np.linalg.norm(
                            node_coordinates
                            - prediction.generated_coordinates[atom_index],
                            axis=1,
                        )
                    )
                    for atom_index in atom_indices
                ],
                dtype=float,
            )
            assert np.median(nearest_node_distances) == pytest.approx(
                statistics.bond_length_medians[(element, "Pb")],
                abs=0.20,
            )


def test_clusterdynamicsml_reuses_project_component_profiles_for_observed_saxs(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)
    project_dir = _build_project_dir(
        tmp_path,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
    )
    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.experimental_data_path = None
    settings.copied_experimental_data_file = None
    manager.save_project(settings)
    q_values = np.linspace(0.20, 0.60, 9)
    _write_project_component_artifacts(project_dir, q_values=q_values)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        project_dir=project_dir,
        experimental_data_file=None,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
        q_points=q_values.size,
    ).analyze()

    assert result.saxs_comparison is not None
    observed_entries = [
        entry
        for entry in result.saxs_comparison.component_weights
        if entry.source == "observed_project"
    ]
    assert len(observed_entries) == len(result.training_observations)
    assert result.saxs_comparison.q_values[0] == pytest.approx(
        float(q_values[0])
    )
    assert result.saxs_comparison.q_values[-1] == pytest.approx(
        float(q_values[-1])
    )

    expected_weights: list[float] = []
    expected_traces: list[np.ndarray] = []
    for row in result.training_observations:
        profile_entry = next(
            entry for entry in observed_entries if entry.label == row.label
        )
        assert profile_entry.profile_path is not None
        profile_data = np.loadtxt(profile_entry.profile_path, comments="#")
        expected_traces.append(np.asarray(profile_data[:, 1], dtype=float))
        expected_weights.append(
            max(float(row.mean_count_per_frame), 0.0)
            * max(float(row.occupancy_fraction), 0.05)
        )
    normalized = np.asarray(expected_weights, dtype=float)
    normalized = normalized / np.sum(normalized)
    expected_model = np.einsum(
        "i,ij->j", normalized, np.asarray(expected_traces)
    )
    assert np.allclose(
        result.saxs_comparison.observed_raw_model_intensity,
        expected_model,
    )


def test_clusterdynamicsml_maps_prediction_shares_into_combined_weights(
    tmp_path,
    monkeypatch,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    original_fit = (
        clusterdynamicsml_workflow_module.ClusterDynamicsMLWorkflow._fit_candidate_property_models
    )

    def _collapsed_count_models(
        self,
        training_observations,
        feature_matrix,
        weights,
        *,
        non_node_elements,
        node_elements,
    ):
        models = original_fit(
            self,
            training_observations,
            feature_matrix,
            weights,
            non_node_elements=non_node_elements,
            node_elements=node_elements,
        )
        models["mean_count_per_frame"] = (
            clusterdynamicsml_workflow_module._PropertyModel(
                coefficients=None,
                constant_value=0.0,
                transform="identity",
                default_value=0.0,
                lower_bound=0.0,
            )
        )
        models["occupancy_fraction"] = (
            clusterdynamicsml_workflow_module._PropertyModel(
                coefficients=None,
                constant_value=0.35,
                transform="identity",
                default_value=0.35,
                lower_bound=0.0,
                upper_bound=1.0,
            )
        )
        models["mean_lifetime_fs"] = (
            clusterdynamicsml_workflow_module._PropertyModel(
                coefficients=None,
                constant_value=25.0,
                transform="identity",
                default_value=25.0,
                lower_bound=self.frame_timestep_fs,
            )
        )
        return models

    monkeypatch.setattr(
        clusterdynamicsml_workflow_module.ClusterDynamicsMLWorkflow,
        "_fit_candidate_property_models",
        _collapsed_count_models,
    )

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    assert all(
        entry.predicted_mean_count_per_frame == pytest.approx(0.0)
        for entry in result.predictions
    )
    assert sum(
        entry.predicted_population_share for entry in result.predictions
    ) == pytest.approx(1.0)
    assert result.saxs_comparison is not None

    predicted_component_weights = {
        entry.label: float(entry.weight)
        for entry in result.saxs_comparison.component_weights
        if entry.source == "predicted"
    }
    assert set(predicted_component_weights) == {
        entry.label for entry in result.predictions
    }
    assert all(weight > 0.0 for weight in predicted_component_weights.values())

    normalized_predicted_model_weights = np.asarray(
        [
            predicted_component_weights[entry.label]
            for entry in result.predictions
        ],
        dtype=float,
    )
    normalized_predicted_model_weights = (
        normalized_predicted_model_weights
        / np.sum(normalized_predicted_model_weights)
    )
    expected_shares = np.asarray(
        [entry.predicted_population_share for entry in result.predictions],
        dtype=float,
    )
    expected_shares = expected_shares / np.sum(expected_shares)
    assert np.allclose(normalized_predicted_model_weights, expected_shares)

    combined_entries = _distribution_entries(result, include_predictions=True)
    predicted_entries = [
        entry
        for entry in combined_entries
        if int(entry["node_count"]) > result.max_observed_node_count
    ]
    assert len(predicted_entries) == len(result.predictions)
    assert all(
        float(entry["normalized_weight"]) > 0.0 for entry in predicted_entries
    )

    combined_payload = _build_population_histogram_payload(
        result,
        include_predictions=True,
    )
    assert combined_payload is not None
    assert {entry.label for entry in result.predictions}.issubset(
        set(combined_payload["structures"])
    )


def test_clusterdynamicsml_caps_predicted_weight_takeover_from_extreme_shares(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    assert result.predictions
    result.predictions[0].predicted_population_share = 0.98
    residual_share = 0.02 / max(len(result.predictions) - 1, 1)
    for entry in result.predictions[1:]:
        entry.predicted_population_share = residual_share
    result.predictions[0].predicted_mean_count_per_frame = 1.0e6
    result.predictions[0].predicted_occupancy_fraction = 1.0

    observed_weights, predicted_weights = (
        clusterdynamicsml_workflow_module._resolved_population_weights(
            result.training_observations,
            result.predictions,
            frame_timestep_fs=10.0,
        )
    )
    observed_size_totals: dict[int, float] = {}
    for observation, weight in zip(
        result.training_observations,
        observed_weights,
        strict=False,
    ):
        observed_size_totals[int(observation.node_count)] = (
            observed_size_totals.get(int(observation.node_count), 0.0)
            + float(weight)
        )

    assert float(np.sum(predicted_weights)) <= (
        float(observed_size_totals[max(observed_size_totals)]) + 1.0e-12
    )
    combined_total = float(
        np.sum(observed_weights) + np.sum(predicted_weights)
    )
    assert combined_total > 0.0
    assert float(np.max(predicted_weights) / combined_total) < 0.5


def test_clusterdynamicsml_writes_predicted_structure_xyz_and_component_profiles(
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    assert result.saxs_comparison is not None
    assert result.saxs_comparison.component_output_dir is not None
    assert result.saxs_comparison.predicted_structure_dir is not None
    assert result.saxs_comparison.component_output_dir.is_dir()
    assert result.saxs_comparison.predicted_structure_dir.is_dir()

    observed_entries = [
        entry
        for entry in result.saxs_comparison.component_weights
        if entry.source == "observed_direct"
    ]
    predicted_entries = [
        entry
        for entry in result.saxs_comparison.component_weights
        if entry.source == "predicted"
    ]
    assert observed_entries
    assert predicted_entries
    assert all(
        entry.profile_path is not None and entry.profile_path.is_file()
        for entry in observed_entries + predicted_entries
    )
    assert all(
        entry.structure_path is not None and entry.structure_path.is_file()
        for entry in predicted_entries
    )
    first_xyz = predicted_entries[0].structure_path.read_text(encoding="utf-8")
    assert predicted_entries[0].structure_path.suffix == ".xyz"
    assert first_xyz.splitlines()[0].strip().isdigit()


def test_clusterdynamicsml_writes_xyz_for_every_prediction(
    tmp_path,
    monkeypatch,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    original_resolve = (
        clusterdynamicsml_workflow_module._resolved_population_weights
    )

    def _drop_predicted_model_members(
        training_observations,
        predictions,
        *,
        frame_timestep_fs,
    ):
        observed_weights, predicted_weights = original_resolve(
            training_observations,
            predictions,
            frame_timestep_fs=frame_timestep_fs,
        )
        filtered = np.zeros_like(predicted_weights)
        for index, entry in enumerate(predictions):
            if entry.rank == 1:
                filtered[index] = predicted_weights[index]
        return observed_weights, filtered

    monkeypatch.setattr(
        clusterdynamicsml_workflow_module,
        "_resolved_population_weights",
        _drop_predicted_model_members,
    )

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    assert result.saxs_comparison is not None
    assert result.saxs_comparison.predicted_structure_dir is not None
    written_paths = sorted(
        result.saxs_comparison.predicted_structure_dir.glob("*.xyz")
    )

    assert len(result.predictions) > len(
        {entry.target_node_count for entry in result.predictions}
    )
    assert len(written_paths) == len(result.predictions)
    assert {path.name for path in written_paths} == {
        f"{entry.target_node_count:02d}_rank{entry.rank:02d}_{entry.label}.xyz"
        for entry in result.predictions
    }


def test_clusterdynamicsml_dataset_round_trip(tmp_path):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    saved = save_cluster_dynamicsai_dataset(
        result,
        tmp_path / "clusterdynamicsml_saved.json",
        analysis_settings={
            "frame_timestep_fs": 10.0,
            "target_node_counts": [4, 5],
        },
    )
    loaded = load_cluster_dynamicsai_dataset(saved.dataset_file)

    assert saved.dataset_file.exists()
    assert any(
        path.name.endswith("_cluster_distribution.csv")
        for path in saved.written_files
    )
    assert any(
        path.name.endswith("_lifetime.csv") for path in saved.written_files
    )
    assert any(
        path.name.endswith("_predictions.csv") for path in saved.written_files
    )
    assert any(
        path.name.endswith("_observed_histogram.csv")
        for path in saved.written_files
    )
    assert any(
        path.name.endswith("_observed_plus_predicted_structures_histogram.csv")
        for path in saved.written_files
    )
    assert any(path.suffix == ".xyz" for path in saved.written_files)
    assert any(
        path.parent.name.endswith("_saxs_components")
        for path in saved.written_files
    )
    assert (
        loaded.result.max_observed_node_count == result.max_observed_node_count
    )
    assert [entry.label for entry in loaded.result.predictions] == [
        entry.label for entry in result.predictions
    ]
    assert loaded.analysis_settings["target_node_counts"] == [4, 5]
    assert loaded.result.saxs_comparison is not None
    assert np.allclose(
        loaded.result.saxs_comparison.fitted_model_intensity,
        result.saxs_comparison.fitted_model_intensity,
    )
    assert np.allclose(
        loaded.result.saxs_comparison.observed_fitted_model_intensity,
        result.saxs_comparison.observed_fitted_model_intensity,
    )
    assert (
        loaded.result.saxs_comparison.component_output_dir
        == result.saxs_comparison.component_output_dir
    )
    assert [
        (
            entry.source,
            entry.method,
            entry.label,
            entry.element_a,
            entry.element_b,
        )
        for entry in loaded.result.debye_waller_estimates
    ] == [
        (
            entry.source,
            entry.method,
            entry.label,
            entry.element_a,
            entry.element_b,
        )
        for entry in result.debye_waller_estimates
    ]


def test_clusterdynamicsml_window_stores_runtime_training_history_and_updates_preview(
    qapp,
    tmp_path,
):
    del qapp
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)
    project_dir = _build_project_dir(
        tmp_path,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
    )

    window = ClusterDynamicsMLMainWindow(
        initial_frames_dir=frames_dir,
        initial_project_dir=project_dir,
        initial_clusters_dir=clusters_dir,
        initial_experimental_data_file=experimental_data_file,
    )
    config = window._build_job_config()
    preview = window._build_preview_workflow().preview_selection()

    window._store_runtime_training_example(
        config=config,
        preview=preview,
        runtime_seconds=18.0,
    )
    window._refresh_selection_preview()
    estimate = window._estimate_runtime_for_preview(preview, config=config)
    history_path = window._runtime_history_file(
        project_dir=project_dir,
        frames_dir=frames_dir,
    )

    assert history_path is not None
    assert history_path.is_file()
    assert estimate.seconds is not None
    assert estimate.sample_count == 1
    assert (
        "Estimated compute time:"
        in window.run_panel.selection_box.toPlainText()
    )
    assert (
        "learned from 1 previous run"
        in window.run_panel.selection_box.toPlainText()
    )
    payload = json.loads(history_path.read_text(encoding="utf-8"))
    assert len(payload["runs"]) == 1
    assert payload["runs"][0]["runtime_seconds"] == pytest.approx(18.0)
    window.close()


def test_clusterdynamicsml_window_autosaves_and_restores_project_result_bundle(
    qapp,
    tmp_path,
):
    del qapp
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)
    project_dir = _build_project_dir(
        tmp_path,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
    )

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        project_dir=project_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    window = ClusterDynamicsMLMainWindow(
        initial_frames_dir=frames_dir,
        initial_project_dir=project_dir,
    )
    window.time_panel.set_frame_timestep_fs(10.0)
    window.time_panel.set_frames_per_colormap_timestep(1)
    window.prediction_panel.set_target_node_counts((4, 5))
    window._on_run_finished(result)

    saved_results_dir = (
        build_project_paths(project_dir).exported_data_dir
        / "clusterdynamicsml"
        / "saved_results"
    )
    dataset_files = sorted(saved_results_dir.rglob("*_clusterdynamicsml.json"))
    assert dataset_files
    cached_dataset = dataset_files[-1]
    bundle_dir = cached_dataset.parent
    bundle_files = {path.name for path in bundle_dir.iterdir()}
    assert f"{cached_dataset.stem}_selection_preview.txt" in bundle_files
    assert f"{cached_dataset.stem}_summary.txt" in bundle_files
    assert f"{cached_dataset.stem}_saxs.csv" in bundle_files
    assert f"{cached_dataset.stem}_observed_histogram.csv" in bundle_files
    assert (
        f"{cached_dataset.stem}_observed_plus_predicted_structures_histogram.csv"
        in bundle_files
    )
    assert (bundle_dir / f"{cached_dataset.stem}_saxs_components").is_dir()

    window.close()

    reopened = ClusterDynamicsMLMainWindow(initial_project_dir=project_dir)

    assert reopened._last_result is not None
    assert reopened._last_dataset_file == cached_dataset
    assert reopened.trajectory_panel.get_frames_dir() == frames_dir
    assert reopened.lifetime_table.rowCount() == (
        len(result.training_observations) + len(result.predictions)
    )
    assert "SAXS components in mixture" in reopened.summary_box.toPlainText()
    observed_hist_patches = sum(
        len(axis.patches) for axis in reopened.histogram_panel.figure.axes
    )
    assert observed_hist_patches > 0
    reopened.close()


def test_clusterdynamicsml_window_compares_prediction_history_and_defaults_to_latest(
    qapp,
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)
    project_dir = _build_project_dir(
        tmp_path,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
    )

    result_one = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        project_dir=project_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4,),
        prediction_population_share_threshold=0.01,
    ).analyze()
    result_two = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        project_dir=project_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
        prediction_population_share_threshold=0.05,
    ).analyze()

    window = ClusterDynamicsMLMainWindow(
        initial_frames_dir=frames_dir,
        initial_project_dir=project_dir,
    )
    window.time_panel.set_frame_timestep_fs(10.0)
    window.time_panel.set_frames_per_colormap_timestep(1)
    window.prediction_panel.set_target_node_counts((4,))
    window.prediction_panel.set_prediction_population_share_threshold(0.01)
    window._on_run_finished(result_one)
    first_dataset = window._last_dataset_file

    window.prediction_panel.set_target_node_counts((4, 5))
    window.prediction_panel.set_prediction_population_share_threshold(0.05)
    window._on_run_finished(result_two)
    second_dataset = window._last_dataset_file

    assert first_dataset is not None
    assert second_dataset is not None
    assert first_dataset != second_dataset
    assert window.history_table.rowCount() == 2
    assert window._selected_history_dataset_file() == second_dataset
    assert window._last_result is not None
    assert window._last_result.preview.target_node_counts == (4, 5)

    older_row = next(
        row
        for row in range(window.history_table.rowCount())
        if window._history_dataset_file_for_row(row) == first_dataset
    )
    window.history_table.selectRow(older_row)
    window._load_selected_history_entry()
    _wait_for(qapp, lambda: window._dataset_load_thread is None)

    assert window._last_dataset_file == first_dataset
    assert window._selected_history_dataset_file() == first_dataset
    assert window._last_result is not None
    assert window._last_result.preview.target_node_counts == (4,)
    window.close()

    reopened = ClusterDynamicsMLMainWindow(initial_project_dir=project_dir)

    assert reopened._last_result is not None
    assert reopened._last_dataset_file == second_dataset
    assert reopened._selected_history_dataset_file() == second_dataset
    assert reopened._last_result.preview.target_node_counts == (4, 5)
    reopened.close()


def test_clusterdynamicsml_window_loads_history_entries_off_ui_thread(
    qapp,
    tmp_path,
    monkeypatch,
):
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)
    project_dir = _build_project_dir(
        tmp_path,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
    )

    result_one = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        project_dir=project_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4,),
        prediction_population_share_threshold=0.01,
    ).analyze()
    result_two = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        project_dir=project_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
        prediction_population_share_threshold=0.05,
    ).analyze()

    window = ClusterDynamicsMLMainWindow(
        initial_frames_dir=frames_dir,
        initial_project_dir=project_dir,
    )
    window.time_panel.set_frame_timestep_fs(10.0)
    window.time_panel.set_frames_per_colormap_timestep(1)
    window.prediction_panel.set_target_node_counts((4,))
    window.prediction_panel.set_prediction_population_share_threshold(0.01)
    window._on_run_finished(result_one)
    first_dataset = window._last_dataset_file

    window.prediction_panel.set_target_node_counts((4, 5))
    window.prediction_panel.set_prediction_population_share_threshold(0.05)
    window._on_run_finished(result_two)

    assert first_dataset is not None
    older_row = next(
        row
        for row in range(window.history_table.rowCount())
        if window._history_dataset_file_for_row(row) == first_dataset
    )
    gate = Event()
    original_loader = (
        "saxshell.clusterdynamicsml.ui.main_window."
        "load_cluster_dynamicsai_dataset"
    )

    def delayed_loader(dataset_file):
        gate.wait(timeout=2.0)
        return load_cluster_dynamicsai_dataset(dataset_file)

    monkeypatch.setattr(original_loader, delayed_loader)
    window.history_table.selectRow(older_row)
    window._load_selected_history_entry()
    qapp.processEvents()

    assert window._dataset_load_thread is not None
    assert not window.history_load_button.isEnabled()
    assert not window.history_table.isEnabled()
    assert "loading" in (window.history_status_label.text().lower())

    gate.set()
    _wait_for(qapp, lambda: window._dataset_load_thread is None)

    assert window._last_dataset_file == first_dataset
    assert window._last_result is not None
    assert window._last_result.preview.target_node_counts == (4,)
    assert window.history_load_button.isEnabled()
    assert window.history_table.isEnabled()
    window.close()


def test_clusterdynamicsml_window_blocks_accidental_field_scroll_and_escape(
    qapp,
):
    window = ClusterDynamicsMLMainWindow()
    window.show()
    qapp.processEvents()

    candidates_spin = window.prediction_panel.candidates_spin
    original_value = candidates_spin.value()
    candidates_spin.setFocus()
    qapp.processEvents()
    QApplication.sendEvent(candidates_spin, _make_wheel_event(candidates_spin))
    qapp.processEvents()

    assert candidates_spin.value() == original_value

    clusters_edit = window.prediction_panel.clusters_dir_edit
    clusters_edit.setFocus()
    qapp.processEvents()
    assert clusters_edit.hasFocus()
    QTest.keyClick(clusters_edit, Qt.Key.Key_Escape)
    qapp.processEvents()

    assert not clusters_edit.hasFocus()

    candidates_spin.lineEdit().setFocus()
    qapp.processEvents()
    assert candidates_spin.lineEdit().hasFocus()
    QTest.keyClick(candidates_spin.lineEdit(), Qt.Key.Key_Escape)
    qapp.processEvents()

    assert not candidates_spin.hasFocus()
    assert not candidates_spin.lineEdit().hasFocus()
    window.close()


def test_clusterdynamicsml_window_debounces_preview_refresh_while_typing(
    qapp,
    monkeypatch,
):
    window = ClusterDynamicsMLMainWindow()
    window.show()
    qapp.processEvents()
    QTest.qWait(_UI_REFRESH_DELAY_MS + 50)

    refresh_calls: list[str] = []

    def record_refresh():
        refresh_calls.append(window.prediction_panel.clusters_dir_edit.text())

    monkeypatch.setattr(window, "_refresh_selection_preview", record_refresh)

    window.prediction_panel.clusters_dir_edit.setText("/tmp/a")
    window.prediction_panel.clusters_dir_edit.setText("/tmp/ab")
    window.prediction_panel.clusters_dir_edit.setText("/tmp/abc")
    qapp.processEvents()

    assert refresh_calls == []

    QTest.qWait(_UI_REFRESH_DELAY_MS + 50)

    assert refresh_calls == ["/tmp/abc"]
    window.close()


def test_clusterdynamicsml_window_debounces_frames_dir_changes_while_typing(
    qapp,
    monkeypatch,
):
    window = ClusterDynamicsMLMainWindow()
    window.show()
    qapp.processEvents()
    QTest.qWait(_UI_REFRESH_DELAY_MS + 50)

    observed_frames_dirs: list[Path | None] = []

    def record_frames_dir_change(frames_dir: Path | None) -> None:
        observed_frames_dirs.append(frames_dir)

    monkeypatch.setattr(
        window, "_on_frames_dir_changed", record_frames_dir_change
    )

    window.trajectory_panel.frames_dir_edit.setText("/tmp/f")
    window.trajectory_panel.frames_dir_edit.setText("/tmp/fr")
    window.trajectory_panel.frames_dir_edit.setText("/tmp/frames")
    qapp.processEvents()

    assert observed_frames_dirs == []

    QTest.qWait(_UI_REFRESH_DELAY_MS + 50)

    assert observed_frames_dirs == [Path("/tmp/frames")]
    window.close()


def test_clusterdynamicsml_window_detailed_report_toggle_skips_preview_refresh(
    qapp,
    monkeypatch,
):
    window = ClusterDynamicsMLMainWindow()
    window.show()
    qapp.processEvents()
    QTest.qWait(_UI_REFRESH_DELAY_MS + 50)

    refresh_calls: list[bool] = []

    def record_refresh():
        refresh_calls.append(True)

    monkeypatch.setattr(window, "_refresh_selection_preview", record_refresh)

    window.run_panel.auto_report_checkbox.click()
    qapp.processEvents()
    QTest.qWait(_UI_REFRESH_DELAY_MS + 50)

    assert refresh_calls == []
    window.close()


def test_clusterdynamicsml_window_autofills_sibling_energy_file(
    qapp,
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    energy_path = _write_energy_file(tmp_path)

    window = ClusterDynamicsMLMainWindow()
    window.show()
    qapp.processEvents()

    window.trajectory_panel.frames_dir_edit.setText(str(frames_dir))
    _wait_for(
        qapp,
        lambda: window.run_panel.energy_file() == energy_path.resolve(),
        timeout=2.0,
    )

    assert window.run_panel.energy_file() == energy_path.resolve()
    window.close()


def test_clusterdynamicsml_window_keeps_manual_energy_file_when_frames_change(
    qapp,
    tmp_path,
):
    frames_dir = _build_frames_dir(tmp_path)
    _write_energy_file(tmp_path)
    manual_dir = tmp_path / "manual_energy"
    manual_dir.mkdir()
    manual_energy = _write_energy_file(manual_dir, name="manual.ener")

    window = ClusterDynamicsMLMainWindow()
    window.show()
    qapp.processEvents()

    window.run_panel.energy_path_edit.setText(str(manual_energy))
    window.trajectory_panel.frames_dir_edit.setText(str(frames_dir))
    QTest.qWait(_UI_REFRESH_DELAY_MS + 75)
    qapp.processEvents()

    assert window.run_panel.energy_file() == manual_energy.resolve()
    window.close()


def test_clusterdynamicsml_window_inherits_project_defaults(
    qapp,
    tmp_path,
):
    del qapp
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)
    project_dir = _build_project_dir(
        tmp_path,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
    )

    window = ClusterDynamicsMLMainWindow(
        initial_frames_dir=frames_dir,
        initial_project_dir=project_dir,
    )
    window.time_panel.set_frame_timestep_fs(10.0)
    window.time_panel.set_frames_per_colormap_timestep(1)
    window.prediction_panel.set_target_node_counts((4, 5))
    window._refresh_selection_preview()
    preview_text = window.run_panel.selection_box.toPlainText()

    assert window.dataset_panel.project_dir() == project_dir
    assert window.prediction_panel.clusters_dir() == clusters_dir
    assert (
        window.prediction_panel.experimental_data_file()
        == experimental_data_file
    )
    assert "Observed node counts: (1, 2, 3)" in preview_text
    assert "Target node counts: (4, 5)" in preview_text
    window.close()


def test_clusterdynamicsml_window_exports_colormap_and_lifetime_csv(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    window = ClusterDynamicsMLMainWindow(initial_frames_dir=frames_dir)
    window._last_result = result
    window.dynamics_plot_panel.set_result(result.dynamics_result)
    window.dynamics_plot_panel.display_mode_combo.setCurrentIndex(
        window.dynamics_plot_panel.display_mode_combo.findData("count")
    )
    window.dynamics_plot_panel.time_unit_combo.setCurrentIndex(
        window.dynamics_plot_panel.time_unit_combo.findData("ps")
    )

    colormap_path = tmp_path / "ai_colormap.csv"
    lifetime_path = tmp_path / "ai_lifetime.csv"
    selected_paths = iter((str(colormap_path), str(lifetime_path)))
    monkeypatch.setattr(
        "saxshell.clusterdynamicsml.ui.main_window.QFileDialog.getSaveFileName",
        lambda *args, **kwargs: (next(selected_paths), "CSV Files (*.csv)"),
    )

    window.save_colormap_data()
    window.save_lifetime_table()

    with colormap_path.open(newline="", encoding="utf-8") as handle:
        colormap_rows = list(csv.DictReader(handle))
    with lifetime_path.open(newline="", encoding="utf-8") as handle:
        lifetime_rows = list(csv.DictReader(handle))

    assert len(colormap_rows) == (
        len(result.dynamics_result.cluster_labels)
        * result.dynamics_result.bin_count
    )
    assert colormap_rows[0]["display_mode"] == "count"
    assert colormap_rows[0]["time_unit"] == "ps"
    assert any(row["Label"] == "Pb3I2" for row in lifetime_rows)
    assert any(row["Type"] == "Predicted" for row in lifetime_rows)
    window.close()


def test_clusterdynamicsml_window_shows_observed_lifetime_tab(
    qapp,
    tmp_path,
):
    del qapp
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    window = ClusterDynamicsMLMainWindow(initial_frames_dir=frames_dir)
    window._on_run_finished(result)

    tab_titles = [
        window.results_tabs.tabText(index)
        for index in range(window.results_tabs.count())
    ]
    lifetime_labels = [
        window.lifetime_table.item(row, 4).text()
        for row in range(window.lifetime_table.rowCount())
    ]
    lifetime_types = [
        window.lifetime_table.item(row, 0).text()
        for row in range(window.lifetime_table.rowCount())
    ]
    observed_only_weights = [
        window.lifetime_table.item(row, 5).text()
        for row in range(window.lifetime_table.rowCount())
        if window.lifetime_table.item(row, 0).text() == "Observed"
    ]
    debye_pairs = [
        window.debye_waller_table.item(row, 5).text()
        for row in range(window.debye_waller_table.rowCount())
    ]

    assert tab_titles == [
        "Summary",
        "Lifetimes",
        "Debye-Waller",
        "Histograms",
        "SAXS",
    ]
    assert window.lifetime_table.rowCount() == (
        len(result.training_observations) + len(result.predictions)
    )
    assert window.debye_waller_table.rowCount() == len(
        result.debye_waller_estimates
    )
    assert "Pb3I2" in lifetime_labels
    assert "Pb-Pb" in debye_pairs
    assert "Predicted" in lifetime_types
    assert all(weight != "n/a" for weight in observed_only_weights)
    assert window.lifetime_table.item(0, 8) is not None
    window.close()


def test_clusterdynamicsml_window_progress_messages_show_current_ml_step(
    qapp,
):
    del qapp
    window = ClusterDynamicsMLMainWindow()
    window._active_runtime_estimate = SimpleNamespace(
        seconds=90.0, sample_count=2
    )
    window._run_started_at_monotonic = time.monotonic() - 15.0

    window._on_worker_progress(
        "Step 4/7: Computing bond-length, bond-angle, and coordination distributions.\n"
        "Learning node, linker, shell, and non-node contact statistics."
    )

    assert "step 4/7" in window.run_panel.progress_label.text().lower()
    assert "remaining" in window.run_panel.progress_label.text().lower()
    assert window.run_panel.progress_bar.maximum() == 7
    assert window.run_panel.progress_bar.value() == 4
    assert (
        "Computing bond-length, bond-angle, and coordination distributions."
        in (window.run_panel.log_box.toPlainText())
    )
    window.close()


def test_clusterdynamicsml_window_shows_histogram_tabs_and_saxs_model_overlay(
    qapp,
    tmp_path,
):
    del qapp
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    window = ClusterDynamicsMLMainWindow(initial_frames_dir=frames_dir)
    window._on_run_finished(result)

    tab_titles = [
        window.results_tabs.tabText(index)
        for index in range(window.results_tabs.count())
    ]
    observed_hist_patches = sum(
        len(axis.patches) for axis in window.histogram_panel.figure.axes
    )
    window.histogram_panel.population_combo.setCurrentIndex(1)
    combined_hist_patches = sum(
        len(axis.patches) for axis in window.histogram_panel.figure.axes
    )
    observed_weight_sum = sum(
        float(entry["normalized_weight"])
        for entry in _distribution_entries(result, include_predictions=False)
    )
    combined_weight_sum = sum(
        float(entry["normalized_weight"])
        for entry in _distribution_entries(result, include_predictions=True)
    )
    line_labels = {
        line.get_label()
        for axis in window.saxs_panel.figure.axes
        for line in axis.lines
    }
    component_lines = [
        line
        for axis in window.saxs_panel.figure.axes
        for line in axis.lines
        if "component:" in str(line.get_gid() or "")
    ]
    initial_component_visibility = {
        line.get_label(): line.get_visible() for line in component_lines
    }
    component_colors = [str(line.get_color()) for line in component_lines]
    initial_axes = [
        {
            "xscale": axis.get_xscale(),
            "yscale": axis.get_yscale(),
            "ylabel": axis.get_ylabel(),
        }
        for axis in window.saxs_panel.figure.axes
    ]
    legend_line_colors = {
        key: str(handle.get_color())
        for key, handle in window.saxs_panel._legend_handle_lookup.items()
        if hasattr(handle, "get_color")
    }
    trace_line_colors = {
        key: str(line.get_color())
        for key, line in window.saxs_panel._trace_line_lookup.items()
    }
    prediction_size_ranks: dict[int, set[int]] = {}
    for row in range(window.lifetime_table.rowCount()):
        if window.lifetime_table.item(row, 0).text() != "Predicted":
            continue
        node_count = int(window.lifetime_table.item(row, 1).text())
        size_rank = int(window.lifetime_table.item(row, 2).text())
        prediction_size_ranks.setdefault(node_count, set()).add(size_rank)
    model_rows = _combined_model_weight_rows(result)
    combined_weight_percents = [
        float(window.lifetime_table.item(row, 6).text())
        for row in range(window.lifetime_table.rowCount())
    ]
    window.saxs_panel._toggle_all_component_traces()
    toggled_component_visibility = {
        line.get_label(): line.get_visible()
        for axis in window.saxs_panel.figure.axes
        for line in axis.lines
        if "component:" in str(line.get_gid() or "")
    }
    observed_trace_visibility = {
        line.get_label(): line.get_visible()
        for axis in window.saxs_panel.figure.axes
        for line in axis.lines
        if line.get_label() == "observed-only model"
        or line.get_label().startswith("observed component:")
    }
    predicted_trace_visibility = {
        line.get_label(): line.get_visible()
        for axis in window.saxs_panel.figure.axes
        for line in axis.lines
        if line.get_label() == "observed + predicted structures model"
        or line.get_label().startswith("predicted structure component:")
    }
    window.saxs_panel._toggle_observed_traces()
    hidden_observed_visibility = {
        line.get_label(): line.get_visible()
        for axis in window.saxs_panel.figure.axes
        for line in axis.lines
        if line.get_label() == "observed-only model"
        or line.get_label().startswith("observed component:")
    }
    window.saxs_panel._toggle_predicted_traces()
    hidden_predicted_visibility = {
        line.get_label(): line.get_visible()
        for axis in window.saxs_panel.figure.axes
        for line in axis.lines
        if line.get_label() == "observed + predicted structures model"
        or line.get_label().startswith("predicted structure component:")
    }
    window.saxs_panel._toggle_observed_traces()
    window.saxs_panel._toggle_predicted_traces()
    window.saxs_panel.model_range_button.setChecked(True)
    autoscaled_axes = window.saxs_panel.figure.axes
    model_lines = [
        line for line in autoscaled_axes[1].lines if line.get_visible()
    ]
    model_q_values = np.concatenate(
        [
            np.asarray(line.get_xdata(orig=False), dtype=float)
            for line in model_lines
        ]
    )

    assert tab_titles == [
        "Summary",
        "Lifetimes",
        "Debye-Waller",
        "Histograms",
        "SAXS",
    ]
    assert observed_hist_patches > 0
    assert combined_hist_patches > 0
    assert observed_weight_sum == pytest.approx(1.0)
    assert combined_weight_sum == pytest.approx(1.0)
    assert len(initial_axes) == 2
    assert initial_axes[0]["xscale"] == "log"
    assert initial_axes[0]["yscale"] == "log"
    assert initial_axes[1]["yscale"] == "log"
    assert initial_axes[0]["ylabel"] == "Intensity (arb. units)"
    assert initial_axes[1]["ylabel"] == "Model Intensity (arb. units)"
    assert window.right_splitter.orientation() == Qt.Orientation.Vertical
    assert window.right_splitter.count() == 3
    assert prediction_size_ranks == {4: {2}, 5: {1}}
    assert window.lifetime_table.rowCount() == (
        len(result.training_observations) + len(result.predictions)
    )
    assert sum(
        row["normalized_weight"] for row in model_rows
    ) == pytest.approx(1.0)
    assert sum(combined_weight_percents) == pytest.approx(100.0, abs=0.2)
    assert "observed-only model" in line_labels
    assert "observed + predicted structures model" in line_labels
    assert any(
        "predicted structure component:" in label for label in line_labels
    )
    assert trace_line_colors
    assert legend_line_colors
    assert legend_line_colors.items() <= trace_line_colors.items()
    assert initial_component_visibility
    assert not any(initial_component_visibility.values())
    default_cycle = list(
        rcParams["axes.prop_cycle"].by_key().get("color", ["#1f77b4"])
    )
    assert component_colors == [
        default_cycle[index % len(default_cycle)]
        for index in range(len(component_colors))
    ]
    assert toggled_component_visibility
    assert all(toggled_component_visibility.values())
    assert observed_trace_visibility
    assert predicted_trace_visibility
    assert any(observed_trace_visibility.values())
    assert any(predicted_trace_visibility.values())
    assert hidden_observed_visibility
    assert not any(hidden_observed_visibility.values())
    assert hidden_predicted_visibility
    assert not any(hidden_predicted_visibility.values())
    assert window.saxs_panel.observed_traces_button.text() == (
        "Hide Observed Traces"
    )
    assert window.saxs_panel.predicted_traces_button.text() == (
        "Hide Predicted Traces"
    )
    assert window.saxs_panel.component_traces_button.text() == (
        "Hide Component Traces"
    )
    assert window.saxs_panel.model_range_button.isChecked()
    assert autoscaled_axes[0].get_xlim() == pytest.approx(
        (
            float(np.nanmin(model_q_values)),
            float(np.nanmax(model_q_values)),
        )
    )
    window.close()


def test_clusterdynamicsml_histogram_tabs_match_project_setup_modes(
    qapp,
    tmp_path,
):
    del qapp
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir_with_secondary_atoms(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    observed_payload = _build_population_histogram_payload(
        result,
        include_predictions=False,
    )
    combined_payload = _build_population_histogram_payload(
        result,
        include_predictions=True,
    )

    assert observed_payload is not None
    assert combined_payload is not None
    assert list_secondary_filter_elements(observed_payload) == ["O"]
    assert set(combined_payload["structures"]) > set(
        observed_payload["structures"]
    )

    window = ClusterDynamicsMLMainWindow(initial_frames_dir=frames_dir)
    window._on_run_finished(result)

    observed_panel = window.histogram_panel
    mode_labels = [
        observed_panel.mode_combo.itemText(index)
        for index in range(observed_panel.mode_combo.count())
    ]
    assert mode_labels == [
        "Structure Fraction",
        "Atom Fraction",
        "Solvent Sort - Structure Fraction",
        "Solvent Sort - Atom Fraction",
    ]
    assert observed_panel.secondary_combo.count() == 1
    assert observed_panel.secondary_combo.itemText(0) == "O"

    observed_panel.mode_combo.setCurrentIndex(2)
    observed_axis = observed_panel.figure.axes[0]
    assert observed_axis.get_title() == (
        "Solvent-Sort Structure Fraction Prior Histogram (O)"
    )
    assert len(observed_axis.patches) > 0

    combined_panel = window.histogram_panel
    combined_panel.population_combo.setCurrentIndex(1)
    combined_panel.mode_combo.setCurrentIndex(3)
    combined_axis = combined_panel.figure.axes[0]
    combined_tick_labels = {
        tick.get_text() for tick in combined_axis.get_xticklabels()
    }

    assert (
        combined_axis.get_title()
        == "Solvent-Sort Atom Fraction Prior Histogram (O)"
    )
    assert len(combined_axis.patches) > 0
    assert any("Pb" in label for label in combined_tick_labels)
    window.close()


def test_clusterdynamicsml_window_appends_powerpoint_report_to_existing_project_report(
    qapp,
    tmp_path,
    monkeypatch,
):
    pytest.importorskip("pptx")
    from pptx import Presentation

    del qapp
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)
    experimental_data_file = _write_experimental_data_file(tmp_path)
    project_dir = _build_project_dir(
        tmp_path,
        clusters_dir=clusters_dir,
        experimental_data_file=experimental_data_file,
    )
    existing_report = (
        build_project_paths(project_dir).reports_dir
        / "existing_project_results.pptx"
    )
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[0])
    presentation.save(str(existing_report))
    initial_slide_count = len(presentation.slides)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        project_dir=project_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    window = ClusterDynamicsMLMainWindow(
        initial_frames_dir=frames_dir,
        initial_project_dir=project_dir,
    )
    window._last_result = result
    window.dynamics_plot_panel.set_result(result.dynamics_result)
    window.predicted_structures_plot_panel.set_result(result)
    window.run_panel.set_selection_summary(
        window._format_preview_text(result.preview)
    )
    window._populate_summary_box(result)

    captured_default_path: dict[str, str] = {}

    def fake_get_save_file_name(*args, **kwargs):
        captured_default_path["value"] = str(args[2])
        return (str(existing_report), "PowerPoint Files (*.pptx)")

    monkeypatch.setattr(
        "saxshell.clusterdynamicsml.ui.main_window.QFileDialog.getSaveFileName",
        fake_get_save_file_name,
    )

    window.save_powerpoint_report()

    updated_presentation = Presentation(str(existing_report))

    assert captured_default_path["value"] == str(existing_report)
    assert len(updated_presentation.slides) > initial_slide_count
    assert "ClusterDynamicsML Report" in _presentation_text(
        updated_presentation
    )
    assert "Predicted Larger Clusters" in _presentation_text(
        updated_presentation
    )
    window.close()


def test_clusterdynamicsml_window_auto_exports_detailed_report_after_run(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    frames_dir = _build_frames_dir(tmp_path)
    clusters_dir = _build_clusters_dir(tmp_path)

    result = ClusterDynamicsMLWorkflow(
        frames_dir,
        atom_type_definitions=ATOM_TYPE_DEFINITIONS,
        pair_cutoff_definitions=PAIR_CUTOFFS,
        clusters_dir=clusters_dir,
        frame_timestep_fs=10.0,
        frames_per_colormap_timestep=1,
        target_node_counts=(4, 5),
    ).analyze()

    window = ClusterDynamicsMLMainWindow(initial_frames_dir=frames_dir)
    output_path = tmp_path / "auto_cluster_dynamics_ml_report.pptx"
    captured: dict[str, object] = {}

    def fake_export_cluster_dynamicsai_report_pptx(**kwargs):
        captured["output_path"] = Path(kwargs["output_path"])
        captured["selection_summary"] = str(kwargs["selection_summary"])
        captured["result_summary"] = str(kwargs["result_summary"])
        captured["project_dir"] = kwargs["project_dir"]
        captured["frames_dir"] = kwargs["frames_dir"]
        return SimpleNamespace(
            appended_to_existing=False,
            report_path=Path(kwargs["output_path"]),
            added_slide_count=5,
        )

    monkeypatch.setattr(
        "saxshell.clusterdynamicsml.ui.main_window.export_cluster_dynamicsai_report_pptx",
        fake_export_cluster_dynamicsai_report_pptx,
    )
    monkeypatch.setattr(
        ClusterDynamicsMLMainWindow,
        "_default_powerpoint_report_file",
        lambda self: output_path,
    )

    assert not window.run_panel.auto_report_checkbox.isHidden()
    assert not window.run_panel.auto_report_enabled()

    window.run_panel.set_auto_report_enabled(True)
    window._on_run_finished(result)

    assert captured["output_path"] == output_path
    assert captured["frames_dir"] == frames_dir
    assert "Target node counts" in str(captured["selection_summary"])
    assert "Predicted candidates:" in str(captured["result_summary"])
    assert window.run_panel.progress_label.text() == (
        "Progress: detailed report saved"
    )
    assert str(output_path) in window.run_panel.log_box.toPlainText()
    window.close()


def test_saxshell_cli_forwards_to_clusterdynamicsml_subcommand(monkeypatch):
    captured: dict[str, object] = {}

    def fake_clusterdynamicsml_main(argv=None):
        captured["argv"] = argv
        return 31

    monkeypatch.setattr(
        clusterdynamicsml_cli_module,
        "main",
        fake_clusterdynamicsml_main,
    )

    exit_code = saxshell_module.main(
        [
            "clusterdynamicsml",
            "--",
            "frames",
            "--clusters-dir",
            "clusters",
            "--experimental-data",
            "exp.txt",
        ]
    )

    assert exit_code == 31
    assert captured["argv"] == [
        "frames",
        "--clusters-dir",
        "clusters",
        "--experimental-data",
        "exp.txt",
    ]
