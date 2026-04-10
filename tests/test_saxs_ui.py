from __future__ import annotations

import json
import multiprocessing as mp
import os
import pickle
import shutil
import threading
import time
import warnings
from pathlib import Path
from types import SimpleNamespace

import numpy as np
import pytest
from matplotlib import colormaps
from matplotlib.backends.backend_qtagg import NavigationToolbar2QT
from matplotlib.collections import LineCollection, PolyCollection
from matplotlib.colors import to_hex
from matplotlib.figure import Figure
from PySide6.QtCore import QObject, QRect, Qt, Signal
from PySide6.QtGui import QColor, QTextOption
from PySide6.QtTest import QTest
from PySide6.QtWidgets import (
    QAbstractScrollArea,
    QApplication,
    QCheckBox,
    QDialog,
    QFileDialog,
    QFormLayout,
    QHeaderView,
    QInputDialog,
    QLabel,
    QMessageBox,
    QPushButton,
    QScrollArea,
    QSizePolicy,
    QSplitter,
    QWidget,
)
from scipy import stats

import saxshell.saxs.project_manager.project as project_module
import saxshell.saxs.ui.prefit_tab as prefit_tab_module
from saxshell.clusterdynamicsml.workflow import (
    ClusterDynamicsMLTrainingObservation,
    PredictedClusterCandidate,
)
from saxshell.saxs._model_templates import (
    list_template_specs,
    load_template_module,
    load_template_spec,
)
from saxshell.saxs.contrast.debye import build_contrast_component_profiles
from saxshell.saxs.contrast.electron_density import (
    CONTRAST_SOLVENT_METHOD_DIRECT,
    CONTRAST_SOLVENT_METHOD_NEAT,
    ContrastGeometryDensitySettings,
    ContrastSolventDensitySettings,
    compute_contrast_geometry_and_electron_density,
)
from saxshell.saxs.contrast.representatives import (
    analyze_contrast_representatives,
)
from saxshell.saxs.contrast.settings import (
    COMPONENT_BUILD_MODE_BORN_APPROXIMATION,
    COMPONENT_BUILD_MODE_CONTRAST,
    COMPONENT_BUILD_MODE_NO_CONTRAST,
    ContrastRepresentativeSamplerSettings,
)
from saxshell.saxs.contrast.ui.main_window import ContrastModeMainWindow
from saxshell.saxs.debye.profiles import AveragedComponent, ClusterBin
from saxshell.saxs.dream import (
    DreamParameterEntry,
    DreamRunSettings,
    SAXSDreamResultsLoader,
    SAXSDreamWorkflow,
    load_dream_settings,
)
from saxshell.saxs.model_report import export_dream_model_report_pptx
from saxshell.saxs.prefit import (
    PrefitEvaluation,
    PrefitParameterEntry,
    SAXSPrefitWorkflow,
    compute_cluster_geometry_metadata,
)
from saxshell.saxs.prefit.workflow import PrefitFitResult
from saxshell.saxs.project_manager import (
    ClusterImportResult,
    ExperimentalDataSummary,
    PowerPointExportSettings,
    ProjectSettings,
    SAXSProjectManager,
    build_prior_histogram_export_payload,
    build_project_paths,
    load_experimental_data_file,
    plot_md_prior_histogram,
    project_artifact_paths,
)
from saxshell.saxs.solution_scattering_estimator import (
    SolutionScatteringEstimatorSettings,
    calculate_solution_scattering_estimate,
)
from saxshell.saxs.template_installation import install_template_candidate
from saxshell.saxs.ui.distribution_window import (
    INTERACTIVE_PREVIEW_THROTTLE_MS,
    DistributionSetupWindow,
)
from saxshell.saxs.ui.dream_tab import DreamTab
from saxshell.saxs.ui.experimental_data_loader import (
    ExperimentalDataHeaderDialog,
)
from saxshell.saxs.ui.main_window import (
    AUTO_SNAP_PANES_KEY,
    PROJECT_LOAD_TOTAL_STEPS,
    InstallModelDialog,
    RuntimeBundleOpener,
    SAXSMainWindow,
    TemplateInstallRequest,
    launch_saxs_ui,
)
from saxshell.saxs.ui.prefit_tab import PrefitTab, TableCellComboBox
from saxshell.saxs.ui.prior_histogram_window import PriorHistogramWindow
from saxshell.saxs.ui.project_setup_tab import ProjectSetupTab
from saxshell.saxs.ui.solute_volume_fraction_widget import (
    SOLUTE_VOLUME_FRACTION_CITATION_URL,
    SOLUTE_VOLUME_FRACTION_HELP_TEXT,
    SoluteVolumeFractionWidget,
)
from saxshell.version import __version__

POLY_LMA_HS_TEMPLATE = "template_pydream_poly_lma_hs"
POLY_LMA_HS_MIX_TEMPLATE = "template_pydream_poly_lma_hs_mix_approx"


def _table_column_index(table, label: str) -> int:
    for index in range(table.columnCount()):
        header_item = table.horizontalHeaderItem(index)
        if header_item is not None and header_item.text() == label:
            return index
    raise AssertionError(f"Column {label!r} was not found.")


def _plot_lines_by_gid(axis, gid: str):
    return [line for line in axis.get_lines() if line.get_gid() == gid]


def _pane_snap_resize_result(
    splitter: QSplitter,
    pane_snap_filter,
    *,
    target_index: int,
    desired_width: int | None = None,
    other_width: int = 220,
) -> tuple[list[int], list[int]]:
    pane_widget = splitter.widget(target_index)
    assert pane_widget is not None
    click_widget = pane_widget
    if isinstance(pane_widget, QScrollArea):
        click_widget = pane_widget.viewport()
    current_sizes = splitter.sizes()
    current_total = sum(size for size in current_sizes if size > 0)
    assert current_total > 0
    resolved_width = (
        max(720, int(current_total * 0.7))
        if desired_width is None
        else int(desired_width)
    )
    if target_index == 0:
        splitter.setSizes([250, max(current_total - 250, 1)])
    else:
        splitter.setSizes([max(current_total - 250, 1), 250])
    QApplication.processEvents()
    before_sizes = splitter.sizes()

    pane_snap_filter._preferred_width = lambda widget: (
        resolved_width
        if widget is splitter.widget(target_index)
        else int(other_width)
    )
    QTest.mouseClick(
        click_widget,
        Qt.MouseButton.LeftButton,
        Qt.KeyboardModifier.NoModifier,
        click_widget.rect().center(),
    )
    QApplication.processEvents()
    return before_sizes, splitter.sizes()


def _write_component_file(path, q_values, intensities):
    data = np.column_stack(
        [
            q_values,
            intensities,
            np.zeros_like(q_values),
            np.zeros_like(q_values),
        ]
    )
    np.savetxt(
        path,
        data,
        header="# Number of files: 1\n# Columns: q, S(q)_avg, S(q)_std, S(q)_se",
        comments="",
    )


def _build_saved_contrast_distribution_project(tmp_path):
    pytest.importorskip("xraydb")
    manager = SAXSProjectManager()
    project_dir = tmp_path / "saved_contrast_project"
    settings = manager.create_project(project_dir)
    paths = build_project_paths(project_dir)

    clusters_dir = paths.project_dir / "clusters"
    cluster_bin_dir = clusters_dir / "PbI2"
    cluster_bin_dir.mkdir(parents=True, exist_ok=True)
    (cluster_bin_dir / "frame_0001.xyz").write_text(
        "\n".join(
            [
                "5",
                "frame_0001",
                "Pb 0.0 0.0 0.0",
                "I 3.0 0.0 0.0",
                "I 0.7 1.8 0.0",
                "O 0.0 0.0 2.6",
                "O 0.0 0.0 5.1",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    (cluster_bin_dir / "frame_0002.xyz").write_text(
        "\n".join(
            [
                "5",
                "frame_0002",
                "Pb 0.0 0.0 0.0",
                "I 2.1 0.0 0.0",
                "I 0.0 2.1 0.0",
                "O 0.0 0.0 2.8",
                "O 0.0 0.0 5.2",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    settings.model_only_mode = True
    settings.clusters_dir = str(clusters_dir.resolve())
    settings.use_experimental_grid = False
    settings.q_min = 0.05
    settings.q_max = 0.30
    settings.q_points = 8
    settings.component_build_mode = COMPONENT_BUILD_MODE_CONTRAST
    settings.selected_model_template = (
        "template_pd_likelihood_monosq_decoupled"
    )
    manager.save_project(settings)
    build_result = manager.build_scattering_components(settings)
    artifact_paths = project_artifact_paths(
        settings,
        storage_mode="distribution",
        allow_legacy_fallback=False,
    )
    return project_dir, settings, artifact_paths, build_result


def _build_minimal_saxs_project(
    tmp_path,
    *,
    frames_dir: Path | None = None,
    pdb_frames_dir: Path | None = None,
):
    manager = SAXSProjectManager()
    project_dir = tmp_path / "saxs_project"
    settings = manager.create_project(project_dir)
    paths = build_project_paths(project_dir)

    q_values = np.linspace(0.05, 0.3, 8)
    component = np.linspace(10.0, 17.0, 8)
    template_name = "template_pd_likelihood_monosq_decoupled"
    template_module = load_template_module(template_name)
    experimental = template_module.lmfit_model_profile(
        q_values,
        np.zeros_like(q_values),
        [component],
        w0=0.6,
        solv_w=0.0,
        offset=0.05,
        eff_r=9.0,
        vol_frac=0.0,
        scale=5e-4,
    )

    experimental_path = paths.experimental_data_dir / "exp_demo.txt"
    np.savetxt(experimental_path, np.column_stack([q_values, experimental]))
    _write_component_file(
        paths.scattering_components_dir / "A_no_motif.txt",
        q_values,
        component,
    )
    (paths.project_dir / "md_prior_weights.json").write_text(
        json.dumps(
            {
                "origin": "clusters",
                "total_files": 1,
                "structures": {
                    "A": {
                        "no_motif": {
                            "count": 1,
                            "weight": 0.6,
                            "representative": "frame_0001.xyz",
                            "profile_file": "A_no_motif.txt",
                        }
                    }
                },
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    (paths.project_dir / "md_saxs_map.json").write_text(
        json.dumps(
            {"saxs_map": {"A": {"no_motif": "A_no_motif.txt"}}},
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )

    settings.experimental_data_path = str(experimental_path)
    settings.copied_experimental_data_file = str(experimental_path)
    settings.selected_model_template = template_name
    if frames_dir is not None:
        settings.frames_dir = str(frames_dir)
    if pdb_frames_dir is not None:
        settings.pdb_frames_dir = str(pdb_frames_dir)
    manager.save_project(settings)
    return project_dir, paths


def _write_predicted_structure_artifacts(
    paths,
    *,
    observed_weight: float = 0.75,
    predicted_weight: float = 0.25,
):
    q_values = np.linspace(0.05, 0.3, 8)
    observed = np.linspace(10.0, 17.0, 8)
    predicted = np.linspace(4.0, 11.0, 8)
    _write_component_file(
        paths.predicted_scattering_components_dir / "A_no_motif.txt",
        q_values,
        observed,
    )
    _write_component_file(
        paths.predicted_scattering_components_dir / "A2_predicted_rank01.txt",
        q_values,
        predicted,
    )
    dataset_dir = (
        paths.exported_data_dir
        / "clusterdynamicsml"
        / "saved_results"
        / "20260330_120000_demo"
    )
    dataset_dir.mkdir(parents=True, exist_ok=True)
    dataset_file = dataset_dir / "20260330_120000_demo_clusterdynamicsml.json"
    dataset_file.write_text(
        json.dumps({"predictions": [{"label": "A2", "rank": 1}]}, indent=2)
        + "\n",
        encoding="utf-8",
    )
    (
        paths.project_dir / "md_prior_weights_predicted_structures.json"
    ).write_text(
        json.dumps(
            {
                "origin": "clusters_predicted_structures",
                "total_files": 1,
                "available_elements": ["A"],
                "value_kind": "normalized_weight",
                "includes_predicted_structures": True,
                "prediction_dataset_file": str(dataset_file),
                "structures": {
                    "A": {
                        "no_motif": {
                            "count": 1,
                            "weight": observed_weight,
                            "normalized_weight": observed_weight,
                            "observed_only_weight": 1.0,
                            "representative": "frame_0001.xyz",
                            "profile_file": "A_no_motif.txt",
                            "source_kind": "cluster_dir",
                        }
                    },
                    "A2": {
                        "predicted_rank01": {
                            "count": 1,
                            "weight": predicted_weight,
                            "normalized_weight": predicted_weight,
                            "observed_only_weight": 0.0,
                            "representative": "02_rank01_A2.xyz",
                            "profile_file": "A2_predicted_rank01.txt",
                            "source_kind": "predicted_structure",
                        }
                    },
                },
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    (paths.project_dir / "md_saxs_map_predicted_structures.json").write_text(
        json.dumps(
            {
                "saxs_map": {
                    "A": {"no_motif": "A_no_motif.txt"},
                    "A2": {"predicted_rank01": "A2_predicted_rank01.txt"},
                }
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    return dataset_file


def _write_distribution_history_artifacts(
    artifact_paths,
    *,
    state_name: str,
    run_name: str,
):
    prefit_state_dir = artifact_paths.prefit_dir / state_name
    prefit_state_dir.mkdir(parents=True, exist_ok=True)
    (prefit_state_dir / "prefit_state.json").write_text(
        json.dumps({"saved_at": state_name, "parameter_entries": []}, indent=2)
        + "\n",
        encoding="utf-8",
    )
    run_dir = artifact_paths.dream_runtime_dir / run_name
    run_dir.mkdir(parents=True, exist_ok=True)
    np.save(
        run_dir / "dream_sampled_params.npy", np.zeros((2, 2), dtype=float)
    )
    np.save(run_dir / "dream_log_ps.npy", np.zeros((2,), dtype=float))
    (run_dir / "pd_settings.json").write_text(
        "{}\n",
        encoding="utf-8",
    )


def _seed_saved_distribution_from_root(
    project_dir: Path,
    *,
    include_cluster_geometry: bool = False,
    component_build_mode: str | None = None,
):
    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    if component_build_mode is not None:
        settings.component_build_mode = component_build_mode
    paths = build_project_paths(project_dir)
    artifact_paths = project_artifact_paths(
        settings,
        storage_mode="distribution",
        allow_legacy_fallback=False,
    )
    manager.ensure_project_dirs(paths)
    manager.ensure_artifact_dirs(artifact_paths)
    shutil.copytree(
        paths.scattering_components_dir,
        artifact_paths.component_dir,
        dirs_exist_ok=True,
    )
    shutil.copy2(
        paths.project_dir / "md_saxs_map.json",
        artifact_paths.component_map_file,
    )
    shutil.copy2(
        paths.project_dir / "md_prior_weights.json",
        artifact_paths.prior_weights_file,
    )
    manager._write_distribution_metadata(
        settings, artifact_paths=artifact_paths
    )
    if include_cluster_geometry:
        artifact_paths.cluster_geometry_metadata_file.parent.mkdir(
            parents=True,
            exist_ok=True,
        )
        artifact_paths.cluster_geometry_metadata_file.write_text(
            json.dumps({"rows": []}, indent=2) + "\n",
            encoding="utf-8",
        )
    return settings, artifact_paths


def _build_poly_lma_geometry_project(
    tmp_path,
    *,
    template_name: str = POLY_LMA_HS_MIX_TEMPLATE,
    single_atom: bool = False,
):
    manager = SAXSProjectManager()
    project_dir = tmp_path / "poly_lma_project"
    settings = manager.create_project(project_dir)
    paths = build_project_paths(project_dir)

    q_values = np.linspace(0.05, 0.3, 8)
    component = np.linspace(10.0, 17.0, 8)

    clusters_dir = paths.project_dir / "clusters"
    structure_dir = clusters_dir / "A"
    structure_dir.mkdir(parents=True, exist_ok=True)
    if single_atom:
        (structure_dir / "frame_0001.xyz").write_text(
            "1\nframe 1\nI 0.0 0.0 0.0\n",
            encoding="utf-8",
        )
        cluster_count = 1
    else:
        (structure_dir / "frame_0001.xyz").write_text(
            "4\nframe 1\nPb 0.0 0.0 0.0\nI 2.0 0.0 0.0\n"
            "I 0.0 2.0 0.0\nI 0.0 0.0 2.0\n",
            encoding="utf-8",
        )
        (structure_dir / "frame_0002.xyz").write_text(
            "4\nframe 2\nPb 0.0 0.0 0.0\nI 2.2 0.0 0.0\n"
            "I 0.0 1.8 0.0\nI 0.0 0.0 2.1\n",
            encoding="utf-8",
        )
        cluster_count = 2
    geometry_table = compute_cluster_geometry_metadata(
        clusters_dir,
        template_name=template_name,
    )
    effective_radius = geometry_table.rows[0].effective_radius

    template_module = load_template_module(template_name)
    experimental = template_module.lmfit_model_profile(
        q_values,
        np.zeros_like(q_values),
        [component],
        np.asarray([effective_radius], dtype=float),
        w0=1.0,
        phi_solute=0.02,
        phi_int=0.02,
        solvent_scale=1.0,
        scale=1.0,
        offset=0.0,
        log_sigma=-9.21,
    )

    experimental_path = paths.experimental_data_dir / "exp_demo.txt"
    np.savetxt(experimental_path, np.column_stack([q_values, experimental]))
    _write_component_file(
        paths.scattering_components_dir / "A_no_motif.txt",
        q_values,
        component,
    )
    (paths.project_dir / "md_prior_weights.json").write_text(
        json.dumps(
            {
                "origin": "clusters",
                "total_files": cluster_count,
                "structures": {
                    "A": {
                        "no_motif": {
                            "count": cluster_count,
                            "weight": 1.0,
                            "representative": "frame_0001.xyz",
                            "profile_file": "A_no_motif.txt",
                        }
                    }
                },
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    (paths.project_dir / "md_saxs_map.json").write_text(
        json.dumps(
            {"saxs_map": {"A": {"no_motif": "A_no_motif.txt"}}},
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )

    settings.experimental_data_path = str(experimental_path)
    settings.copied_experimental_data_file = str(experimental_path)
    settings.selected_model_template = template_name
    settings.clusters_dir = str(clusters_dir)
    manager.save_project(settings)
    return project_dir, paths


def _write_minimal_dream_results(
    project_dir,
    *,
    settings: DreamRunSettings | None = None,
    entries: list[DreamParameterEntry] | None = None,
):
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)
    parameter_entries = entries or workflow.create_default_parameter_map()
    bundle = workflow.create_runtime_bundle(
        settings=settings,
        entries=parameter_entries,
    )

    active_values = np.asarray(
        [entry.value for entry in parameter_entries if entry.vary],
        dtype=float,
    )
    sampled_params = []
    log_ps = []
    for chain_index in range(2):
        chain_samples = []
        chain_logps = []
        for step_index in range(4):
            adjustment = (chain_index + 1) * (step_index + 1) * 0.01
            chain_samples.append(active_values + adjustment)
            chain_logps.append(-5.0 + chain_index + step_index * 0.5)
        sampled_params.append(chain_samples)
        log_ps.append(chain_logps)
    np.save(
        bundle.run_dir / "dream_sampled_params.npy", np.asarray(sampled_params)
    )
    np.save(bundle.run_dir / "dream_log_ps.npy", np.asarray(log_ps))
    return bundle


def _write_weight_order_dream_results(tmp_path):
    run_dir = tmp_path / "dream_order_test"
    run_dir.mkdir(parents=True)
    metadata = {
        "settings": {"burnin_percent": 0},
        "template_name": "template_pd_likelihood_monosq_decoupled",
        "parameter_map": [
            {
                "structure": "PbI2O",
                "motif": "m3",
                "param_type": "Both",
                "param": "w2",
                "value": 0.25,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 0.25, "s": 0.1},
            },
            {
                "structure": "I2",
                "motif": "m1",
                "param_type": "Both",
                "param": "w0",
                "value": 0.35,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 0.35, "s": 0.1},
            },
            {
                "structure": "Pb2",
                "motif": "m2",
                "param_type": "Both",
                "param": "w1",
                "value": 0.4,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 0.4, "s": 0.1},
            },
            {
                "structure": "",
                "motif": "",
                "param_type": "SAXS",
                "param": "scale",
                "value": 1.0,
                "vary": False,
                "distribution": "norm",
                "dist_params": {"loc": 1.0, "scale": 0.1},
            },
        ],
        "active_parameter_entries": [
            {
                "structure": "PbI2O",
                "motif": "m3",
                "param_type": "Both",
                "param": "w2",
                "value": 0.25,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 0.25, "s": 0.1},
            },
            {
                "structure": "I2",
                "motif": "m1",
                "param_type": "Both",
                "param": "w0",
                "value": 0.35,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 0.35, "s": 0.1},
            },
            {
                "structure": "Pb2",
                "motif": "m2",
                "param_type": "Both",
                "param": "w1",
                "value": 0.4,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 0.4, "s": 0.1},
            },
        ],
        "active_parameter_indices": [0, 1, 2],
        "full_parameter_names": ["w2", "w0", "w1", "scale"],
        "fixed_parameter_values": [0.25, 0.35, 0.4, 1.0],
        "q_values": [0.1, 0.2],
        "experimental_intensities": [1.0, 0.8],
        "theoretical_intensities": [[1.0, 0.9]],
        "solvent_intensities": [0.0, 0.0],
    }
    (run_dir / "dream_runtime_metadata.json").write_text(
        json.dumps(metadata, indent=2) + "\n",
        encoding="utf-8",
    )
    np.save(
        run_dir / "dream_sampled_params.npy",
        np.asarray(
            [
                [
                    [0.25, 0.35, 0.40],
                    [0.26, 0.34, 0.41],
                ]
            ],
            dtype=float,
        ),
    )
    np.save(
        run_dir / "dream_log_ps.npy",
        np.asarray([[1.0, 2.0]], dtype=float),
    )
    return run_dir


def _write_stoichiometry_filter_dream_results(tmp_path):
    run_dir = tmp_path / "dream_stoichiometry_filter_test"
    run_dir.mkdir(parents=True)
    metadata = {
        "settings": {"burnin_percent": 0},
        "template_name": "template_pd_likelihood_monosq_decoupled",
        "parameter_map": [
            {
                "structure": "PbI2",
                "motif": "m1",
                "param_type": "Both",
                "param": "w0",
                "value": 1.0,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 1.0, "s": 0.1},
            },
            {
                "structure": "Pb2I5",
                "motif": "m2",
                "param_type": "Both",
                "param": "w1",
                "value": 0.0,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 1.0, "s": 0.1},
            },
            {
                "structure": "",
                "motif": "",
                "param_type": "SAXS",
                "param": "scale",
                "value": 1.0,
                "vary": False,
                "distribution": "norm",
                "dist_params": {"loc": 1.0, "scale": 0.1},
            },
        ],
        "active_parameter_entries": [
            {
                "structure": "PbI2",
                "motif": "m1",
                "param_type": "Both",
                "param": "w0",
                "value": 1.0,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 1.0, "s": 0.1},
            },
            {
                "structure": "Pb2I5",
                "motif": "m2",
                "param_type": "Both",
                "param": "w1",
                "value": 0.0,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 1.0, "s": 0.1},
            },
        ],
        "active_parameter_indices": [0, 1],
        "full_parameter_names": ["w0", "w1", "scale"],
        "fixed_parameter_values": [1.0, 0.0, 1.0],
        "q_values": [0.1, 0.2],
        "experimental_intensities": [1.0, 0.8],
        "theoretical_intensities": [[1.0, 0.9]],
        "solvent_intensities": [0.0, 0.0],
    }
    (run_dir / "dream_runtime_metadata.json").write_text(
        json.dumps(metadata, indent=2) + "\n",
        encoding="utf-8",
    )
    np.save(
        run_dir / "dream_sampled_params.npy",
        np.asarray(
            [
                [
                    [1.0, 0.0],
                    [0.9, 0.1],
                    [0.2, 0.8],
                    [0.0, 1.0],
                ]
            ],
            dtype=float,
        ),
    )
    np.save(
        run_dir / "dream_log_ps.npy",
        np.asarray([[4.0, 3.0, 2.0, 1.0]], dtype=float),
    )
    return run_dir


def _write_violin_mode_split_dream_results(tmp_path):
    run_dir = tmp_path / "dream_violin_mode_split_test"
    run_dir.mkdir(parents=True)
    metadata = {
        "settings": {"burnin_percent": 0},
        "template_name": "template_pd_likelihood_monosq_decoupled",
        "parameter_map": [
            {
                "structure": "A",
                "motif": "m1",
                "param_type": "Both",
                "param": "w0",
                "value": 0.35,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 0.35, "s": 0.1},
            },
            {
                "structure": "A",
                "motif": "m1",
                "param_type": "Both",
                "param": "r_eff_w0",
                "value": 9.0,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 9.0, "scale": 0.5},
            },
            {
                "structure": "B",
                "motif": "m2",
                "param_type": "Both",
                "param": "a_eff_w1",
                "value": 8.0,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 8.0, "scale": 0.5},
            },
            {
                "structure": "B",
                "motif": "m2",
                "param_type": "Both",
                "param": "b_eff_w1",
                "value": 10.0,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 10.0, "scale": 0.5},
            },
            {
                "structure": "B",
                "motif": "m2",
                "param_type": "Both",
                "param": "c_eff_w1",
                "value": 12.0,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 12.0, "scale": 0.5},
            },
            {
                "structure": "",
                "motif": "",
                "param_type": "SAXS",
                "param": "scale",
                "value": 1.0,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 1.0, "scale": 0.1},
            },
            {
                "structure": "",
                "motif": "",
                "param_type": "SAXS",
                "param": "offset",
                "value": 0.05,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 0.05, "scale": 0.01},
            },
            {
                "structure": "",
                "motif": "",
                "param_type": "SAXS",
                "param": "phi_int",
                "value": 0.12,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 0.12, "scale": 0.01},
            },
        ],
        "active_parameter_entries": [
            {
                "structure": "A",
                "motif": "m1",
                "param_type": "Both",
                "param": "w0",
                "value": 0.35,
                "vary": True,
                "distribution": "lognorm",
                "dist_params": {"loc": 0.0, "scale": 0.35, "s": 0.1},
            },
            {
                "structure": "A",
                "motif": "m1",
                "param_type": "Both",
                "param": "r_eff_w0",
                "value": 9.0,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 9.0, "scale": 0.5},
            },
            {
                "structure": "B",
                "motif": "m2",
                "param_type": "Both",
                "param": "a_eff_w1",
                "value": 8.0,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 8.0, "scale": 0.5},
            },
            {
                "structure": "B",
                "motif": "m2",
                "param_type": "Both",
                "param": "b_eff_w1",
                "value": 10.0,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 10.0, "scale": 0.5},
            },
            {
                "structure": "B",
                "motif": "m2",
                "param_type": "Both",
                "param": "c_eff_w1",
                "value": 12.0,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 12.0, "scale": 0.5},
            },
            {
                "structure": "",
                "motif": "",
                "param_type": "SAXS",
                "param": "scale",
                "value": 1.0,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 1.0, "scale": 0.1},
            },
            {
                "structure": "",
                "motif": "",
                "param_type": "SAXS",
                "param": "offset",
                "value": 0.05,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 0.05, "scale": 0.01},
            },
            {
                "structure": "",
                "motif": "",
                "param_type": "SAXS",
                "param": "phi_int",
                "value": 0.12,
                "vary": True,
                "distribution": "norm",
                "dist_params": {"loc": 0.12, "scale": 0.01},
            },
        ],
        "active_parameter_indices": list(range(8)),
        "full_parameter_names": [
            "w0",
            "r_eff_w0",
            "a_eff_w1",
            "b_eff_w1",
            "c_eff_w1",
            "scale",
            "offset",
            "phi_int",
        ],
        "fixed_parameter_values": [
            0.35,
            9.0,
            8.0,
            10.0,
            12.0,
            1.0,
            0.05,
            0.12,
        ],
        "q_values": [0.1, 0.2],
        "experimental_intensities": [1.0, 0.8],
        "theoretical_intensities": [[1.0, 0.9]],
        "solvent_intensities": [0.0, 0.0],
    }
    (run_dir / "dream_runtime_metadata.json").write_text(
        json.dumps(metadata, indent=2) + "\n",
        encoding="utf-8",
    )
    np.save(
        run_dir / "dream_sampled_params.npy",
        np.asarray(
            [
                [
                    [0.35, 9.0, 8.0, 10.0, 12.0, 1.0, 0.05, 0.12],
                    [0.36, 9.2, 8.1, 10.1, 12.1, 1.1, 0.07, 0.14],
                ]
            ],
            dtype=float,
        ),
    )
    np.save(
        run_dir / "dream_log_ps.npy",
        np.asarray([[1.0, 2.0]], dtype=float),
    )
    return run_dir


@pytest.fixture(scope="module")
def qapp():
    os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")
    app = QApplication.instance()
    if app is None:
        app = QApplication([])
    yield app


def _wait_for_dream_refresh(qapp, delay_ms: int = 120) -> None:
    qapp.processEvents()
    QTest.qWait(delay_ms)
    qapp.processEvents()


def _apply_dream_filter_changes(
    window: SAXSMainWindow,
    qapp,
    delay_ms: int = 120,
) -> None:
    window.dream_tab.apply_filter_button.click()
    _wait_for_dream_refresh(qapp, delay_ms=delay_ms)


def test_saxs_main_window_loads_project_prefit_and_dream_tabs(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    brand_widget = window.tabs.cornerWidget(Qt.Corner.TopLeftCorner)

    assert window.windowTitle() == "SAXSShell"
    assert not window.windowIcon().isNull()
    assert brand_widget is not None
    brand_labels = [
        label.text() for label in brand_widget.findChildren(QLabel)
    ]
    brand_icon_labels = [
        label
        for label in brand_widget.findChildren(QLabel)
        if label.pixmap() is not None and not label.pixmap().isNull()
    ]
    assert "SAXSShell" in brand_labels
    assert len(brand_icon_labels) == 1
    assert brand_icon_labels[0].height() >= 34
    assert brand_icon_labels[0].width() > brand_icon_labels[0].height()
    assert brand_widget.height() <= 40
    assert window.project_setup_tab.forward_model_group.isEnabled()
    assert window.project_setup_tab.model_group.isEnabled()
    assert window.project_setup_tab.template_combo.count() >= 1
    assert (
        window.project_setup_tab.selected_template_name()
        == "template_pd_likelihood_monosq_decoupled"
    )
    assert window.project_setup_tab.template_combo.currentText().startswith(
        "MonoSQ Decoupled"
    )
    assert (
        "Structure Factor:"
        in window.project_setup_tab.template_combo.toolTip()
    )
    assert (
        "keith.white@colorado.edu"
        in window.project_setup_tab.template_help_button.toolTip()
    )
    assert window.project_setup_tab.project_name_edit.text() == "saxs_project"
    assert window.project_setup_tab.component_log_x_checkbox.isChecked()
    assert window.project_setup_tab.component_log_y_checkbox.isChecked()
    assert window.project_setup_tab.project_dir_edit.text() == str(
        project_dir.parent
    )
    assert window.project_setup_tab.open_project_dir_edit.text().endswith(
        "saxs_project"
    )
    assert window.prefit_tab.parameter_table.rowCount() == 6
    assert window.prefit_tab.template_combo.count() >= 1
    assert window.prefit_tab.selected_template_name() == (
        "template_pd_likelihood_monosq_decoupled"
    )
    assert window.prefit_tab.template_combo.currentText().startswith(
        "MonoSQ Decoupled"
    )
    assert "Form Factor:" in window.prefit_tab.template_combo.toolTip()
    assert (
        "keith.white@colorado.edu"
        in window.prefit_tab.template_help_button.toolTip()
    )
    assert "Refine scale before the other model parameters" in (
        window.prefit_tab.prefit_help_button.toolTip()
    )
    assert window.prefit_tab.show_experimental_trace_checkbox.isChecked()
    assert window.prefit_tab.show_model_trace_checkbox.isChecked()
    assert not window.prefit_tab.show_solvent_trace_checkbox.isChecked()
    assert (
        not window.prefit_tab.show_structure_factor_trace_checkbox.isChecked()
    )
    assert window.prefit_tab.log_x_checkbox.isChecked()
    assert window.prefit_tab.log_y_checkbox.isChecked()
    assert window.prefit_tab.plot_toolbar is not None
    assert not window.prefit_tab.autosave_checkbox.isChecked()
    assert window.prefit_tab.recommended_scale_button is not None
    assert isinstance(window.prefit_tab._left_scroll_area, QScrollArea)
    assert (
        window.prefit_tab._left_scroll_area.widget()
        is window.prefit_tab._left_panel
    )
    assert (
        window.prefit_tab._pane_splitter.widget(0)
        is window.prefit_tab._left_scroll_area
    )
    assert window.prefit_tab.set_best_button is not None
    assert window.prefit_tab.reset_best_button is not None
    assert window.prefit_tab.restore_state_button is not None
    assert not window.prefit_tab.restore_state_button.isEnabled()
    assert "Best Prefit preset" in window.prefit_tab.set_best_button.toolTip()
    assert "template-default prefit preset" in (
        window.prefit_tab.reset_button.toolTip()
    )
    assert "timestamped prefit snapshot folders" in (
        window.prefit_tab.saved_state_combo.toolTip()
    )
    assert "Prefit summary:" in window.prefit_tab.summary_box.toPlainText()
    assert (
        "Template: template_pd_likelihood_monosq_decoupled"
        in window.prefit_tab.summary_box.toPlainText()
    )
    assert "Prefit Console" in window.prefit_tab.summary_box.toPlainText()
    assert window.dream_tab.export_model_report_button is not None
    assert (
        window.dream_tab.export_model_report_button.text()
        == "Export Model Report (PPTX)"
    )
    assert window.dream_tab.recycle_button is not None
    assert window.dream_tab.recycle_button.text() == "Recycle"


def test_launch_saxs_ui_shows_and_finishes_startup_splash(qapp, monkeypatch):
    calls: list[object] = []

    class _Splash:
        def show(self):
            calls.append("splash_show")

        def finish(self, window):
            calls.append(("splash_finish", window))

        def close(self):
            calls.append("splash_close")

    class _Window:
        def __init__(self, initial_project_dir=None):
            self.initial_project_dir = initial_project_dir
            calls.append(("window_init", initial_project_dir))

        def show(self):
            calls.append("window_show")

    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.configure_saxshell_application",
        lambda app: calls.append(("configure", app)),
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.create_saxshell_startup_splash",
        lambda: _Splash(),
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.SAXSMainWindow",
        _Window,
    )

    project_path = Path("/tmp/brand_test_project")
    result = launch_saxs_ui(project_path)

    assert result == 0
    assert calls[0][0] == "configure"
    assert calls[1] == "splash_show"
    assert calls[2] == ("window_init", project_path)
    assert calls[3] == "window_show"
    assert calls[4][0] == "splash_finish"
    assert isinstance(calls[4][1], _Window)


def test_prefit_ionic_radius_help_button_shows_citation(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    messages: list[tuple[str, str]] = []
    monkeypatch.setattr(
        "saxshell.saxs.ui.prefit_tab.QMessageBox.information",
        lambda _parent, title, message: messages.append((title, message)),
    )

    assert (
        window.prefit_tab.cluster_geometry_ionic_radius_help_button is not None
    )
    window.prefit_tab.cluster_geometry_ionic_radius_help_button.click()

    assert messages
    assert messages[-1][0] == "Ionic Radius Estimate Help"
    assert "Shannon, R. D. (1976)" in messages[-1][1]
    assert "https://doi.org/10.1107/S0567739476001551" in messages[-1][1]
    assert "+0.14 A offset" in messages[-1][1]
    window.close()


def test_prefit_solute_volume_fraction_help_button_shows_citation(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    messages: list[tuple[str, str]] = []
    monkeypatch.setattr(
        "saxshell.saxs.ui.prefit_tab.QMessageBox.information",
        lambda _parent, title, message: messages.append((title, message)),
    )

    window.prefit_tab.solute_volume_fraction_help_button.click()

    assert messages
    assert messages[-1][0] == "Solution Scattering Estimate Help"
    assert "Physical solute-associated volume fraction" in messages[-1][1]
    assert "SAXS-effective interaction contrast ratio" in messages[-1][1]
    assert "Fluorescence background proxy" in messages[-1][1]
    assert SOLUTE_VOLUME_FRACTION_CITATION_URL in messages[-1][1]
    window.close()


def test_prefit_solute_volume_fraction_estimator_is_template_aware_and_applies(
    qapp,
    tmp_path,
):
    del qapp
    base_project_dir, _base_paths = _build_minimal_saxs_project(
        tmp_path / "basic"
    )
    base_window = SAXSMainWindow(initial_project_dir=base_project_dir)

    assert not base_window.prefit_tab._solute_volume_fraction_group.isHidden()
    assert "solv_w" in (
        base_window.prefit_tab.solute_volume_fraction_status_label.text()
    )

    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(
        tmp_path / "poly"
    )
    poly_window = SAXSMainWindow(initial_project_dir=poly_project_dir)

    assert not poly_window.prefit_tab._solute_volume_fraction_group.isHidden()
    assert "phi_solute" in (
        poly_window.prefit_tab.solute_volume_fraction_status_label.text()
    )
    assert "solvent_scale" in (
        poly_window.prefit_tab.solute_volume_fraction_status_label.text()
    )
    assert poly_window.prefit_tab.solute_volume_fraction_is_collapsed()

    poly_window.prefit_tab.solute_volume_fraction_collapse_button.click()
    widget = poly_window.prefit_tab.solute_volume_fraction_widget
    assert not poly_window.prefit_tab.solute_volume_fraction_is_collapsed()
    assert widget.output_is_collapsed()
    widget.solution_density_spin.setValue(1.0)
    widget.solute_stoich_edit.setText("Cs1Pb1I3")
    widget.solvent_stoich_edit.setText("H2O")
    widget.molar_mass_solute_spin.setValue(620.0)
    widget.molar_mass_solvent_spin.setValue(18.015)
    widget.mass_solute_spin.setValue(1.0)
    widget.mass_solvent_spin.setValue(9.0)
    widget.solute_density_spin.setValue(2.0)
    widget.solvent_density_spin.setValue(1.0)

    expected_settings = widget.current_estimator_settings()
    expected_estimate = calculate_solution_scattering_estimate(
        SolutionScatteringEstimatorSettings(
            solution=expected_settings.solution,
            solute_density_g_per_ml=expected_settings.solute_density_g_per_ml,
            solvent_density_g_per_ml=expected_settings.solvent_density_g_per_ml,
            calculate_number_density=expected_settings.calculate_number_density,
            calculate_solute_volume_fraction=(
                expected_settings.calculate_solute_volume_fraction
            ),
            calculate_solvent_scattering_contribution=(
                expected_settings.calculate_solvent_scattering_contribution
            ),
            calculate_sample_fluorescence_yield=(
                expected_settings.calculate_sample_fluorescence_yield
            ),
            beam=expected_settings.beam,
        )
    )

    widget.calculate_button.click()

    phi_row = poly_window.prefit_tab.find_parameter_row("phi_solute")
    assert phi_row >= 0
    assert float(
        poly_window.prefit_tab.parameter_table.item(phi_row, 3).text()
    ) == pytest.approx(
        expected_estimate.interaction_contrast_estimate.saxs_effective_solute_interaction_ratio,
        rel=1e-3,
    )
    assert (
        poly_window.prefit_tab.parameter_table.item(phi_row, 4).checkState()
        == Qt.CheckState.Unchecked
    )
    solvent_row = poly_window.prefit_tab.find_parameter_row("solvent_scale")
    assert solvent_row >= 0
    assert float(
        poly_window.prefit_tab.parameter_table.item(solvent_row, 3).text()
    ) == pytest.approx(
        expected_estimate.attenuation_estimate.solvent_scattering_scale_factor,
        rel=1e-3,
    )
    assert (
        poly_window.prefit_tab.parameter_table.item(
            solvent_row, 4
        ).checkState()
        == Qt.CheckState.Unchecked
    )
    assert "Applied phi_solute =" in widget.output_box.toPlainText()
    assert "Applied solvent_scale =" in widget.output_box.toPlainText()
    assert (
        "Estimated solute volume: 0.500 cm^3"
        in widget.output_box.toPlainText()
    )
    assert (
        "Estimated solvent volume: 9.000 cm^3"
        in widget.output_box.toPlainText()
    )
    assert (
        "Recommended solvent scattering scale factor:"
        in widget.output_box.toPlainText()
    )
    assert (
        "SAXS-effective solute interaction ratio:"
        in widget.output_box.toPlainText()
    )
    assert not widget.output_is_collapsed()
    widget.output_toggle_button.click()
    assert widget.output_is_collapsed()
    widget.output_toggle_button.click()
    assert not widget.output_is_collapsed()

    base_window.close()
    poly_window.close()


def test_prefit_single_solvent_weight_uses_combined_saxs_effective_multiplier(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.prefit_tab.solute_volume_fraction_collapse_button.click()
    widget = window.prefit_tab.solute_volume_fraction_widget
    widget.solution_density_spin.setValue(1.0)
    widget.solute_stoich_edit.setText("Cs1Pb1I3")
    widget.solvent_stoich_edit.setText("H2O")
    widget.molar_mass_solute_spin.setValue(620.0)
    widget.molar_mass_solvent_spin.setValue(18.015)
    widget.mass_solute_spin.setValue(1.0)
    widget.mass_solvent_spin.setValue(9.0)
    widget.solute_density_spin.setValue(2.0)
    widget.solvent_density_spin.setValue(1.0)

    expected_settings = widget.current_estimator_settings()
    expected_estimate = calculate_solution_scattering_estimate(
        SolutionScatteringEstimatorSettings(
            solution=expected_settings.solution,
            solute_density_g_per_ml=expected_settings.solute_density_g_per_ml,
            solvent_density_g_per_ml=expected_settings.solvent_density_g_per_ml,
            calculate_number_density=expected_settings.calculate_number_density,
            calculate_solute_volume_fraction=(
                expected_settings.calculate_solute_volume_fraction
            ),
            calculate_solvent_scattering_contribution=(
                expected_settings.calculate_solvent_scattering_contribution
            ),
            calculate_sample_fluorescence_yield=(
                expected_settings.calculate_sample_fluorescence_yield
            ),
            beam=expected_settings.beam,
        )
    )

    widget.calculate_button.click()

    solvent_row = window.prefit_tab.find_parameter_row("solv_w")
    assert solvent_row >= 0
    assert float(
        window.prefit_tab.parameter_table.item(solvent_row, 3).text()
    ) == pytest.approx(
        expected_estimate.attenuation_estimate.solvent_scattering_scale_factor
        * expected_estimate.interaction_contrast_estimate.saxs_effective_solvent_background_ratio,
        rel=1e-3,
    )
    assert "Applied solv_w =" in widget.output_box.toPlainText()
    window.close()


def test_prefit_solute_volume_fraction_widget_hides_solute_density_in_molarity_mode(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.prefit_tab.solute_volume_fraction_collapse_button.click()
    widget = window.prefit_tab.solute_volume_fraction_widget
    molarity_index = widget.solution_mode_combo.findData("molarity_per_liter")
    assert molarity_index >= 0
    widget.solution_mode_combo.setCurrentIndex(molarity_index)

    assert not widget.solution_density_spin.isHidden()
    assert widget.solute_density_label.isHidden()
    assert widget.solute_density_spin.isHidden()
    assert not widget.solvent_density_label.isHidden()
    assert not widget.solvent_density_spin.isHidden()
    assert widget.current_estimator_settings().solute_density_g_per_ml is None
    assert (
        widget.current_estimator_settings().solvent_density_g_per_ml
        == pytest.approx(widget.solvent_density_spin.value())
    )
    assert "solvent density remains active" in (
        widget.solution_mode_hint_label.text().lower()
    )
    window.close()


def test_solute_volume_fraction_widget_loads_builtin_solvent_density_presets(
    qapp,
):
    del qapp
    widget = SoluteVolumeFractionWidget()

    dmf_index = widget.solution_preset_combo.findData("PbI2 - DMF - 0.49 M")
    assert dmf_index >= 0
    widget.solution_preset_combo.setCurrentIndex(dmf_index)
    widget._load_selected_solution_preset()
    assert widget.solvent_density_spin.value() == pytest.approx(0.944)

    dmso_index = widget.solution_preset_combo.findData("PbI2 - DMSO - 0.405 M")
    assert dmso_index >= 0
    widget.solution_preset_combo.setCurrentIndex(dmso_index)
    widget._load_selected_solution_preset()
    assert widget.solvent_density_spin.value() == pytest.approx(1.10)


def test_solution_scattering_widget_loads_builtin_beam_presets(qapp):
    del qapp
    widget = SoluteVolumeFractionWidget()

    assert (
        widget.beam_preset_combo.currentData() == "NSLS-II 28-ID-1 (default)"
    )
    assert widget.incident_energy_spin.value() == pytest.approx(74.0)
    assert widget.capillary_size_spin.value() == pytest.approx(1.0)
    assert widget.beam_footprint_width_spin.value() == pytest.approx(0.4)
    assert widget.beam_footprint_height_spin.value() == pytest.approx(0.4)

    focused_index = widget.beam_preset_combo.findData(
        "APS 5-IDD (default - focused)"
    )
    assert focused_index >= 0
    widget.beam_preset_combo.setCurrentIndex(focused_index)
    widget._load_selected_beam_preset()

    assert widget.incident_energy_spin.value() == pytest.approx(17.0)
    assert widget.beam_footprint_width_spin.value() == pytest.approx(0.05)
    assert widget.beam_footprint_height_spin.value() == pytest.approx(1.0)

    unfocused_index = widget.beam_preset_combo.findData(
        "APS 5-IDD (default - unfocused)"
    )
    assert unfocused_index >= 0
    widget.beam_preset_combo.setCurrentIndex(unfocused_index)
    widget._load_selected_beam_preset()

    assert widget.incident_energy_spin.value() == pytest.approx(17.5)
    assert widget.beam_footprint_width_spin.value() == pytest.approx(1.0)
    assert widget.beam_footprint_height_spin.value() == pytest.approx(1.0)


def test_solution_scattering_widget_saves_and_deletes_custom_beam_presets(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    preset_path = tmp_path / "beam_presets.json"
    monkeypatch.setenv(
        "SAXSHELL_BEAM_GEOMETRY_PRESETS_PATH",
        str(preset_path),
    )
    widget = SoluteVolumeFractionWidget()

    widget.incident_energy_spin.setValue(12.4)
    widget.capillary_size_spin.setValue(2.0)
    widget.beam_footprint_width_spin.setValue(0.25)
    widget.beam_footprint_height_spin.setValue(0.75)

    monkeypatch.setattr(
        "saxshell.saxs.ui.solution_scattering_widget.QInputDialog.getText",
        lambda *args, **kwargs: ("Custom Beam", True),
    )

    widget._save_current_beam_preset()

    saved_payload = json.loads(preset_path.read_text(encoding="utf-8"))
    assert "Custom Beam" in saved_payload["presets"]
    assert saved_payload["presets"]["Custom Beam"][
        "incident_energy_kev"
    ] == pytest.approx(12.4)
    assert widget.beam_preset_combo.findData("Custom Beam") >= 0

    widget.beam_preset_combo.setCurrentIndex(
        widget.beam_preset_combo.findData("Custom Beam")
    )
    widget._delete_selected_beam_preset()

    reloaded_payload = json.loads(preset_path.read_text(encoding="utf-8"))
    assert "Custom Beam" not in reloaded_payload["presets"]
    assert widget.beam_preset_combo.findData("Custom Beam") < 0


def test_solution_scattering_widget_wavelength_dialog_tracks_energy(qapp):
    del qapp
    widget = SoluteVolumeFractionWidget()

    widget.incident_energy_spin.setValue(17.0)
    widget._show_wavelength_dialog()

    dialog = widget._wavelength_dialog
    assert dialog is not None
    assert dialog.windowTitle() == "Beam Energy and Wavelength"
    assert dialog.energy_value_label.text() == "17"
    assert float(dialog.wavelength_value_label.text()) == pytest.approx(
        12.398419843320026 / 17.0,
        rel=1e-6,
    )

    widget.incident_energy_spin.setValue(74.0)
    QApplication.processEvents()

    assert dialog.energy_value_label.text() == "74"
    assert float(dialog.wavelength_value_label.text()) == pytest.approx(
        12.398419843320026 / 74.0,
        rel=1e-6,
    )
    dialog.close()


def test_solution_scattering_widget_places_output_below_run_button(qapp):
    del qapp
    widget = SoluteVolumeFractionWidget()
    widget.resize(960, 720)
    widget.set_output_collapsed(False)
    widget.show()
    QApplication.processEvents()

    assert (
        widget.output_panel.parentWidget()
        is widget.calculate_button.parentWidget()
    )
    assert widget.output_panel.y() > widget.calculate_button.y()


def test_prefit_fallback_preserves_selected_template_when_workflow_not_ready(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, paths = _build_poly_lma_geometry_project(tmp_path)
    (paths.project_dir / "md_prior_weights.json").unlink()

    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert window.prefit_workflow is None
    assert (
        window.prefit_tab.selected_template_name() == POLY_LMA_HS_MIX_TEMPLATE
    )
    assert not window.prefit_tab._solute_volume_fraction_group.isHidden()
    assert not window.prefit_tab._cluster_geometry_group.isHidden()
    assert window.prefit_tab.parameter_table.rowCount() == 0
    assert "Prefit workflow is not ready yet." in (
        window.prefit_tab.output_box.toPlainText()
    )
    window.close()


def test_prefit_cluster_geometry_section_is_template_aware(qapp, tmp_path):
    del qapp
    base_project_dir, _base_paths = _build_minimal_saxs_project(
        tmp_path / "basic"
    )
    base_window = SAXSMainWindow(initial_project_dir=base_project_dir)

    assert not base_window.prefit_workflow.supports_cluster_geometry_metadata()
    assert base_window.prefit_tab._cluster_geometry_group.isHidden()
    assert "waiting on template metadata" not in (
        base_window.prefit_tab._summary_text.lower()
    )

    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(
        tmp_path / "poly"
    )
    poly_window = SAXSMainWindow(initial_project_dir=poly_project_dir)

    assert poly_window.prefit_workflow.supports_cluster_geometry_metadata()
    assert not poly_window.prefit_tab._cluster_geometry_group.isHidden()
    assert "Compute cluster geometry metadata" in (
        poly_window.prefit_tab.cluster_geometry_status_label.text()
    )
    assert (
        poly_window.prefit_tab.cluster_geometry_radii_type_combo.currentData()
        == "ionic"
    )
    assert (
        poly_window.prefit_tab.cluster_geometry_ionic_radius_type_combo.currentData()
        == "effective"
    )
    assert poly_window.prefit_tab.cluster_geometry_table.rowCount() == 0
    assert "waiting on template metadata" in (
        poly_window.prefit_tab._summary_text.lower()
    )
    assert "cluster geometry metadata" in (
        poly_window.prefit_tab._summary_text.lower()
    )

    poly_window.compute_prefit_cluster_geometry()

    assert poly_window.prefit_tab.cluster_geometry_table.rowCount() == 1
    assert poly_window.prefit_tab._current_evaluation is not None
    assert (
        poly_window.prefit_tab.cluster_geometry_progress_label.text()
        == "Cluster geometry metadata ready."
    )
    assert (
        poly_window.prefit_tab.cluster_geometry_progress_bar.value()
        == poly_window.prefit_tab.cluster_geometry_progress_bar.maximum()
    )
    assert poly_window._progress_dialog is not None
    assert not poly_window._progress_dialog.isVisible()
    table = poly_window.prefit_tab.cluster_geometry_table
    sf_column = _table_column_index(table, "S.F. Approx.")
    map_to_column = _table_column_index(table, "Map To")
    mapping_combo = poly_window.prefit_tab.cluster_geometry_table.cellWidget(
        0,
        map_to_column,
    )
    assert mapping_combo is not None
    assert str(mapping_combo.currentData()) == "w0"
    sf_combo = poly_window.prefit_tab.cluster_geometry_table.cellWidget(
        0,
        sf_column,
    )
    assert sf_combo is not None
    assert isinstance(sf_combo, TableCellComboBox)
    assert str(sf_combo.currentData()) in {"sphere", "ellipsoid"}
    assert "effective_radii" in (
        poly_window.prefit_tab.cluster_geometry_status_label.text()
    )
    assert "waiting on template metadata" not in (
        poly_window.prefit_tab._summary_text.lower()
    )
    assert "Prefit summary:" in poly_window.prefit_tab._summary_text
    assert poly_window.dream_tab.chains_spin.value() == 4
    assert (
        poly_window.dream_tab.settings_preset_combo.currentText()
        == poly_window.dream_tab.ACTIVE_SETTINGS_LABEL
    )
    assert poly_window.dream_tab.model_name_edit is not None
    assert poly_window.dream_tab.nseedchains_spin.value() == 40
    assert poly_window.dream_tab.output_box is poly_window.dream_tab.log_box
    assert (
        poly_window.dream_tab.output_box is poly_window.dream_tab.summary_box
    )
    assert poly_window.dream_tab.model_toolbar is not None
    assert poly_window.dream_tab.violin_toolbar is not None
    assert poly_window.dream_tab.settings_preset_combo.toolTip()
    assert poly_window.dream_tab.model_name_edit.toolTip()
    assert poly_window.dream_tab.chains_spin.toolTip()
    assert poly_window.dream_tab.run_button.toolTip()
    assert "DREAM Summary" in poly_window.dream_tab.output_box.toPlainText()
    assert "DREAM Console" in poly_window.dream_tab.output_box.toPlainText()
    base_window.close()
    poly_window.close()


def test_strict_poly_lma_hs_ui_only_offers_sphere_cluster_geometry_shape(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_poly_lma_geometry_project(
        tmp_path,
        template_name=POLY_LMA_HS_TEMPLATE,
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    sf_column = _table_column_index(table, "S.F. Approx.")
    sf_combo = table.cellWidget(0, sf_column)

    assert sf_combo is not None
    assert sf_combo.count() == 1
    assert str(sf_combo.currentData()) == "sphere"
    assert sf_combo.findData("ellipsoid") < 0
    window.close()


def test_prefit_cluster_geometry_invalid_radii_alert_and_revert(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    errors: list[tuple[str, str]] = []
    monkeypatch.setattr(
        window,
        "_show_error",
        lambda title, message: errors.append((title, message)),
    )

    updated_rows = window.prefit_workflow.cluster_geometry_rows()
    updated_rows[0].bond_length_sphere_effective_radius = 0.0
    updated_rows[0].bond_length_ellipsoid_semiaxis_a = 0.0
    updated_rows[0].bond_length_ellipsoid_semiaxis_b = 0.0
    updated_rows[0].bond_length_ellipsoid_semiaxis_c = 0.0
    window.prefit_workflow.set_cluster_geometry_rows(updated_rows)
    window._refresh_prefit_cluster_geometry_section()

    window.prefit_tab.toggle_cluster_geometry_radii_button.click()

    assert errors
    assert errors[-1][0] == "Invalid cluster geometry radii"
    assert "positive active radii" in errors[-1][1]
    assert (
        window.prefit_tab.cluster_geometry_radii_type_combo.currentData()
        == ("ionic")
    )
    assert window.prefit_workflow.cluster_geometry_active_radii_type() == (
        "ionic"
    )
    assert window.prefit_tab._current_evaluation is not None
    window.close()


def test_prefit_cluster_geometry_manual_radius_updates_refresh_model(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    radius_column = _table_column_index(table, "Effective Radius")
    semiaxis_x_column = _table_column_index(table, "Semiaxis X")
    semiaxis_y_column = _table_column_index(table, "Semiaxis Y")
    semiaxis_z_column = _table_column_index(table, "Semiaxis Z")
    sf_column = _table_column_index(table, "S.F. Approx.")
    sf_combo = table.cellWidget(0, sf_column)
    assert sf_combo is not None

    sphere_index = sf_combo.findData("sphere")
    assert sphere_index >= 0
    sf_combo.setCurrentIndex(sphere_index)

    original_model = np.asarray(
        window.prefit_tab.current_evaluation().model_intensities,
        dtype=float,
    ).copy()
    table.item(0, radius_column).setText("3.5")
    window.prefit_tab.update_cluster_geometry_button.click()

    updated_row = window.prefit_workflow.cluster_geometry_rows()[0]
    assert updated_row.radii_type_used == "ionic"
    assert updated_row.sf_approximation == "sphere"
    assert updated_row.ionic_sphere_effective_radius == pytest.approx(3.5)
    assert updated_row.effective_radius == pytest.approx(3.5)
    radius_param_row = window.prefit_tab.find_parameter_row("r_eff_w0")
    assert radius_param_row >= 0
    assert float(
        window.prefit_tab.parameter_table.item(radius_param_row, 3).text()
    ) == pytest.approx(3.5)
    updated_model = np.asarray(
        window.prefit_tab.current_evaluation().model_intensities,
        dtype=float,
    )
    assert not np.allclose(updated_model, original_model)

    window.prefit_tab.toggle_cluster_geometry_radii_button.click()
    bond_length_radius = float(table.item(0, radius_column).text())
    window.prefit_tab.toggle_cluster_geometry_radii_button.click()
    assert bond_length_radius != pytest.approx(3.5)
    assert float(table.item(0, radius_column).text()) == pytest.approx(
        3.5,
        rel=1e-3,
    )

    ellipsoid_index = sf_combo.findData("ellipsoid")
    assert ellipsoid_index >= 0
    sf_combo.setCurrentIndex(ellipsoid_index)
    ellipsoid_model = np.asarray(
        window.prefit_tab.current_evaluation().model_intensities,
        dtype=float,
    ).copy()
    table.item(0, semiaxis_x_column).setText("4.0")
    table.item(0, semiaxis_y_column).setText("3.0")
    table.item(0, semiaxis_z_column).setText("2.0")
    window.prefit_tab.update_cluster_geometry_button.click()

    updated_row = window.prefit_workflow.cluster_geometry_rows()[0]
    expected_effective_radius = float(np.cbrt(4.0 * 3.0 * 2.0))
    assert updated_row.sf_approximation == "ellipsoid"
    assert updated_row.active_semiaxis_a == pytest.approx(4.0)
    assert updated_row.active_semiaxis_b == pytest.approx(3.0)
    assert updated_row.active_semiaxis_c == pytest.approx(2.0)
    assert updated_row.effective_radius == pytest.approx(
        expected_effective_radius
    )
    for parameter_name, expected_value in (
        ("a_eff_w0", 4.0),
        ("b_eff_w0", 3.0),
        ("c_eff_w0", 2.0),
    ):
        parameter_row = window.prefit_tab.find_parameter_row(parameter_name)
        assert parameter_row >= 0
        assert float(
            window.prefit_tab.parameter_table.item(parameter_row, 3).text()
        ) == pytest.approx(expected_value)
    refreshed_model = np.asarray(
        window.prefit_tab.current_evaluation().model_intensities,
        dtype=float,
    )
    assert not np.allclose(refreshed_model, ellipsoid_model)
    window.close()


def test_update_model_uses_regenerated_mixed_geometry_parameters(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    semiaxis_x_column = _table_column_index(table, "Semiaxis X")
    semiaxis_y_column = _table_column_index(table, "Semiaxis Y")
    semiaxis_z_column = _table_column_index(table, "Semiaxis Z")
    sf_column = _table_column_index(table, "S.F. Approx.")
    sf_combo = table.cellWidget(0, sf_column)
    assert sf_combo is not None
    assert str(sf_combo.currentData()) == "ellipsoid"

    original_model = np.asarray(
        window.prefit_tab.current_evaluation().model_intensities,
        dtype=float,
    ).copy()

    table.item(0, semiaxis_x_column).setText("8.0")
    table.item(0, semiaxis_y_column).setText("2.0")
    table.item(0, semiaxis_z_column).setText("1.0")
    window.update_prefit_model()

    updated_model = np.asarray(
        window.prefit_tab.current_evaluation().model_intensities,
        dtype=float,
    )
    assert not np.allclose(updated_model, original_model)
    for parameter_name, expected_value in (
        ("a_eff_w0", 8.0),
        ("b_eff_w0", 2.0),
        ("c_eff_w0", 1.0),
    ):
        parameter_row = window.prefit_tab.find_parameter_row(parameter_name)
        assert parameter_row >= 0
        assert float(
            window.prefit_tab.parameter_table.item(parameter_row, 3).text()
        ) == pytest.approx(expected_value)

    sphere_index = sf_combo.findData("sphere")
    assert sphere_index >= 0
    sf_combo.setCurrentIndex(sphere_index)
    radius_column = _table_column_index(table, "Effective Radius")
    table.item(0, radius_column).setText("4.5")
    ellipsoid_model = updated_model.copy()
    window.update_prefit_model()

    sphere_model = np.asarray(
        window.prefit_tab.current_evaluation().model_intensities,
        dtype=float,
    )
    assert not np.allclose(sphere_model, ellipsoid_model)
    radius_param_row = window.prefit_tab.find_parameter_row("r_eff_w0")
    assert radius_param_row >= 0
    assert float(
        window.prefit_tab.parameter_table.item(radius_param_row, 3).text()
    ) == pytest.approx(4.5)
    assert window.prefit_tab.find_parameter_row("a_eff_w0") < 0
    window.close()


def test_update_model_preserves_manual_geometry_parameter_edits_when_cluster_geometry_is_unchanged(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    parameter_name = "a_eff_w0"
    parameter_row = window.prefit_tab.find_parameter_row(parameter_name)
    assert parameter_row >= 0

    original_value = float(
        window.prefit_tab.parameter_table.item(parameter_row, 3).text()
    )
    custom_value = original_value + 0.75
    custom_minimum = max(original_value * 0.5, 1e-6)
    custom_maximum = original_value + 2.0

    window.prefit_tab.set_parameter_row(
        parameter_name,
        value=custom_value,
        minimum=custom_minimum,
        maximum=custom_maximum,
        vary=True,
    )

    window.update_prefit_model()

    assert float(
        window.prefit_tab.parameter_table.item(parameter_row, 3).text()
    ) == pytest.approx(custom_value)
    assert float(
        window.prefit_tab.parameter_table.item(parameter_row, 5).text()
    ) == pytest.approx(custom_minimum)
    assert float(
        window.prefit_tab.parameter_table.item(parameter_row, 6).text()
    ) == pytest.approx(custom_maximum)
    assert (
        window.prefit_tab.parameter_table.item(parameter_row, 4).checkState()
        == Qt.CheckState.Checked
    )
    window.close()


def test_update_prefit_model_batches_cluster_geometry_project_save(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    save_calls: list[bool] = []
    original_save_project = window.prefit_workflow.project_manager.save_project

    def record_save_project(settings, *, refresh_registered_paths=True):
        save_calls.append(bool(refresh_registered_paths))
        return original_save_project(
            settings,
            refresh_registered_paths=refresh_registered_paths,
        )

    monkeypatch.setattr(
        window.prefit_workflow.project_manager,
        "save_project",
        record_save_project,
    )

    window.update_prefit_model()

    assert save_calls == [False]
    window.close()


def test_loading_experimental_data_autosave_skips_registered_path_refresh(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    frames_dir = tmp_path / "frames"
    frames_dir.mkdir()
    (frames_dir / "frame_0000.xyz").write_text(
        "1\nframe\nPb 0.0 0.0 0.0\n",
        encoding="utf-8",
    )
    project_dir, _paths = _build_minimal_saxs_project(
        tmp_path,
        frames_dir=frames_dir,
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)
    experimental_path = tmp_path / "replacement_exp.txt"
    q_values = np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float)
    np.savetxt(
        experimental_path,
        np.column_stack(
            [q_values, np.asarray([100.0, 80.0, 55.0, 30.0], dtype=float)]
        ),
    )
    summary = load_experimental_data_file(experimental_path)

    save_calls: list[bool] = []
    original_save_project = window.project_manager.save_project

    def record_save_project(settings, *, refresh_registered_paths=True):
        save_calls.append(bool(refresh_registered_paths))
        return original_save_project(
            settings,
            refresh_registered_paths=refresh_registered_paths,
        )

    monkeypatch.setattr(
        window.project_manager,
        "save_project",
        record_save_project,
    )

    window.project_setup_tab._apply_experimental_file(
        experimental_path,
        summary,
    )

    assert save_calls == [False]
    window.close()


def test_loading_solvent_data_autosave_skips_registered_path_refresh(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    frames_dir = tmp_path / "frames"
    frames_dir.mkdir()
    (frames_dir / "frame_0000.xyz").write_text(
        "1\nframe\nPb 0.0 0.0 0.0\n",
        encoding="utf-8",
    )
    project_dir, _paths = _build_minimal_saxs_project(
        tmp_path,
        frames_dir=frames_dir,
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)
    solvent_path = tmp_path / "replacement_solvent.txt"
    q_values = np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float)
    np.savetxt(
        solvent_path,
        np.column_stack(
            [q_values, np.asarray([12.0, 11.0, 10.0, 9.0], dtype=float)]
        ),
    )
    summary = load_experimental_data_file(solvent_path)

    save_calls: list[bool] = []
    original_save_project = window.project_manager.save_project

    def record_save_project(settings, *, refresh_registered_paths=True):
        save_calls.append(bool(refresh_registered_paths))
        return original_save_project(
            settings,
            refresh_registered_paths=refresh_registered_paths,
        )

    monkeypatch.setattr(
        window.project_manager,
        "save_project",
        record_save_project,
    )

    window.project_setup_tab._apply_solvent_file(solvent_path, summary)

    assert save_calls == [False]
    window.close()


def test_model_only_mode_autosave_skips_registered_path_refresh(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    frames_dir = tmp_path / "frames"
    frames_dir.mkdir()
    (frames_dir / "frame_0000.xyz").write_text(
        "1\nframe\nPb 0.0 0.0 0.0\n",
        encoding="utf-8",
    )
    project_dir, _paths = _build_minimal_saxs_project(
        tmp_path,
        frames_dir=frames_dir,
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)

    save_calls: list[bool] = []
    original_save_project = window.project_manager.save_project

    def record_save_project(settings, *, refresh_registered_paths=True):
        save_calls.append(bool(refresh_registered_paths))
        return original_save_project(
            settings,
            refresh_registered_paths=refresh_registered_paths,
        )

    monkeypatch.setattr(
        window.project_manager,
        "save_project",
        record_save_project,
    )

    window.project_setup_tab.model_only_mode_checkbox.setChecked(True)

    assert save_calls == [False]
    window.close()


def test_updating_frames_folder_reference_skips_registered_path_refresh(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    frames_dir = tmp_path / "updated_frames"
    frames_dir.mkdir()
    (frames_dir / "frame_0000.xyz").write_text(
        "1\nframe\nPb 0.0 0.0 0.0\n",
        encoding="utf-8",
    )

    save_calls: list[bool] = []
    original_save_project = window.project_manager.save_project

    def record_save_project(settings, *, refresh_registered_paths=True):
        save_calls.append(bool(refresh_registered_paths))
        return original_save_project(
            settings,
            refresh_registered_paths=refresh_registered_paths,
        )

    monkeypatch.setattr(
        window.project_manager,
        "save_project",
        record_save_project,
    )

    window.project_setup_tab.frames_dir_edit.setText(str(frames_dir))
    window.project_setup_tab._on_frames_dir_edited()

    assert save_calls == [False]
    saved_settings = SAXSProjectManager().load_project(project_dir)
    assert saved_settings.frames_dir == str(frames_dir)
    assert saved_settings.frames_dir_snapshot is None
    window.close()


def test_prefit_mixed_shape_switch_logs_when_equivalent_sphere_curve_is_unchanged(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    sf_column = _table_column_index(table, "S.F. Approx.")
    sf_combo = table.cellWidget(0, sf_column)
    assert sf_combo is not None
    assert str(sf_combo.currentData()) == "ellipsoid"

    previous_model = np.asarray(
        window.prefit_tab.current_evaluation().model_intensities,
        dtype=float,
    ).copy()

    sphere_index = sf_combo.findData("sphere")
    assert sphere_index >= 0
    sf_combo.setCurrentIndex(sphere_index)
    window.update_prefit_model()

    current_model = np.asarray(
        window.prefit_tab.current_evaluation().model_intensities,
        dtype=float,
    )
    assert np.allclose(previous_model, current_model, rtol=1e-5, atol=1e-8)
    assert "equivalent-sphere approximation" in (
        window.prefit_tab.output_box.toPlainText().lower()
    )
    assert "changed approximation: ellipsoid -> sphere" in (
        window.prefit_tab.output_box.toPlainText().lower()
    )
    window.close()


def test_prefit_cluster_geometry_notes_are_collapsed_with_expand_and_tooltip(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    notes_column = _table_column_index(table, "Notes")
    notes_item = table.item(0, notes_column)
    assert notes_item is not None
    expected_collapsed_height = max(
        table.verticalHeader().defaultSectionSize(),
        30,
    )
    collapsed_height = table.rowHeight(0)

    assert "\n" not in notes_item.text()
    assert collapsed_height == expected_collapsed_height
    assert (
        "Active sphere radius" in notes_item.toolTip()
        or "Active ellipsoid semiaxes" in notes_item.toolTip()
    )

    window.prefit_tab._on_cluster_geometry_cell_double_clicked(0, notes_column)
    expanded_item = table.item(0, notes_column)
    assert expanded_item is not None
    expanded_height = table.rowHeight(0)
    assert "\n" in expanded_item.text()
    assert expanded_height >= collapsed_height

    window.prefit_tab._on_cluster_geometry_cell_double_clicked(0, notes_column)
    recollapsed_item = table.item(0, notes_column)
    assert recollapsed_item is not None
    assert "\n" not in recollapsed_item.text()
    assert table.rowHeight(0) == expected_collapsed_height
    assert table.rowHeight(0) <= expanded_height
    window.close()


def test_prefit_cluster_geometry_notes_reload_collapsed(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    notes_column = _table_column_index(table, "Notes")
    window.prefit_tab._on_cluster_geometry_cell_double_clicked(0, notes_column)
    assert "\n" in table.item(0, notes_column).text()
    window.close()

    reloaded_window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    reloaded_table = reloaded_window.prefit_tab.cluster_geometry_table
    reloaded_notes_item = reloaded_table.item(0, notes_column)
    assert reloaded_notes_item is not None
    assert "\n" not in reloaded_notes_item.text()
    assert reloaded_table.rowHeight(0) == max(
        reloaded_table.verticalHeader().defaultSectionSize(),
        30,
    )
    reloaded_window.close()


def test_prefit_cluster_geometry_paths_are_collapsed_with_expand_and_tooltip(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    path_column = _table_column_index(table, "Path")
    path_item = table.item(0, path_column)
    assert path_item is not None
    expected_collapsed_height = max(
        table.verticalHeader().defaultSectionSize(),
        30,
    )
    collapsed_height = table.rowHeight(0)

    assert "..." in path_item.text()
    assert path_item.toolTip().endswith("/clusters/A")
    assert collapsed_height == expected_collapsed_height

    window.prefit_tab._on_cluster_geometry_cell_double_clicked(0, path_column)
    expanded_item = table.item(0, path_column)
    assert expanded_item is not None
    assert "..." not in expanded_item.text()
    assert expanded_item.text() == expanded_item.toolTip()
    assert table.rowHeight(0) >= collapsed_height

    window.prefit_tab._on_cluster_geometry_cell_double_clicked(0, path_column)
    recollapsed_item = table.item(0, path_column)
    assert recollapsed_item is not None
    assert "..." in recollapsed_item.text()
    assert table.rowHeight(0) == expected_collapsed_height
    window.close()


def test_prefit_cluster_geometry_paths_reload_collapsed(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    path_column = _table_column_index(table, "Path")
    window.prefit_tab._on_cluster_geometry_cell_double_clicked(0, path_column)
    assert "..." not in table.item(0, path_column).text()
    window.close()

    reloaded_window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    reloaded_table = reloaded_window.prefit_tab.cluster_geometry_table
    reloaded_path_item = reloaded_table.item(0, path_column)
    assert reloaded_path_item is not None
    assert "..." in reloaded_path_item.text()
    assert reloaded_table.rowHeight(0) == max(
        reloaded_table.verticalHeader().defaultSectionSize(),
        30,
    )
    reloaded_window.close()


def test_prefit_cluster_geometry_ionic_radius_type_switch_updates_display(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    radius_column = _table_column_index(table, "Effective Radius")
    radii_type_column = _table_column_index(table, "Radii Type")
    ionic_combo = window.prefit_tab.cluster_geometry_ionic_radius_type_combo

    assert ionic_combo.currentData() == "effective"
    effective_radius = float(table.item(0, radius_column).text())
    assert table.item(0, radii_type_column).text() == "Ionic (effective)"

    crystal_index = ionic_combo.findData("crystal")
    assert crystal_index >= 0
    ionic_combo.setCurrentIndex(crystal_index)

    updated_row = window.prefit_workflow.cluster_geometry_rows()[0]
    crystal_radius = float(table.item(0, radius_column).text())

    assert ionic_combo.currentData() == "crystal"
    assert (
        window.prefit_workflow.cluster_geometry_active_ionic_radius_type()
        == ("crystal")
    )
    assert updated_row.ionic_radius_type_used == "crystal"
    assert table.item(0, radii_type_column).text() == "Ionic (crystal)"
    assert crystal_radius == pytest.approx(
        updated_row.crystal_ionic_sphere_effective_radius,
        rel=1e-3,
    )
    assert crystal_radius > effective_radius
    if updated_row.sf_approximation == "sphere":
        radius_param_row = window.prefit_tab.find_parameter_row("r_eff_w0")
        assert radius_param_row >= 0
        assert float(
            window.prefit_tab.parameter_table.item(radius_param_row, 3).text()
        ) == pytest.approx(updated_row.crystal_ionic_sphere_effective_radius)
    else:
        for parameter_name, expected_value in (
            ("a_eff_w0", updated_row.crystal_ionic_ellipsoid_semiaxis_a),
            ("b_eff_w0", updated_row.crystal_ionic_ellipsoid_semiaxis_b),
            ("c_eff_w0", updated_row.crystal_ionic_ellipsoid_semiaxis_c),
        ):
            parameter_row = window.prefit_tab.find_parameter_row(
                parameter_name
            )
            assert parameter_row >= 0
            assert float(
                window.prefit_tab.parameter_table.item(
                    parameter_row,
                    3,
                ).text()
            ) == pytest.approx(expected_value)

    window.prefit_tab.toggle_cluster_geometry_radii_button.click()
    assert not ionic_combo.isEnabled()
    window.prefit_tab.toggle_cluster_geometry_radii_button.click()
    assert ionic_combo.isEnabled()
    assert ionic_combo.currentData() == "crystal"
    window.close()


def test_update_model_reports_all_nonpositive_cluster_geometry_radii(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    semiaxis_x_column = _table_column_index(table, "Semiaxis X")
    semiaxis_y_column = _table_column_index(table, "Semiaxis Y")
    sf_column = _table_column_index(table, "S.F. Approx.")
    sf_combo = table.cellWidget(0, sf_column)
    assert sf_combo is not None
    ellipsoid_index = sf_combo.findData("ellipsoid")
    assert ellipsoid_index >= 0
    sf_combo.setCurrentIndex(ellipsoid_index)

    errors: list[tuple[str, str]] = []
    monkeypatch.setattr(
        window,
        "_show_error",
        lambda title, message: errors.append((title, message)),
    )

    table.item(0, semiaxis_x_column).setText("-1.0")
    table.item(0, semiaxis_y_column).setText("0.0")
    window.prefit_tab.update_button.click()

    assert errors
    assert errors[-1][0] == "Update model failed"
    assert "Cluster geometry radii must be positive" in errors[-1][1]
    assert "A Semiaxis X=-1" in errors[-1][1]
    assert "A Semiaxis Y=0" in errors[-1][1]
    assert window.prefit_tab.current_evaluation() is not None
    window.close()


def test_prefit_cluster_geometry_controls_switch_radii_mode_and_shape(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    radius_column = _table_column_index(table, "Effective Radius")
    radii_type_column = _table_column_index(table, "Radii Type")
    semiaxis_x_column = _table_column_index(table, "Semiaxis X")
    semiaxis_y_column = _table_column_index(table, "Semiaxis Y")
    semiaxis_z_column = _table_column_index(table, "Semiaxis Z")
    sf_column = _table_column_index(table, "S.F. Approx.")
    initial_sf_combo = table.cellWidget(0, sf_column)
    parameter_names = [
        window.prefit_tab.parameter_table.item(row, 2).text()
        for row in range(window.prefit_tab.parameter_table.rowCount())
    ]

    initial_radius_text = table.item(0, radius_column).text()
    assert (
        window.prefit_tab.cluster_geometry_radii_type_combo.currentData()
        == ("ionic")
    )
    assert (
        window.prefit_tab.cluster_geometry_ionic_radius_type_combo.currentData()
        == "effective"
    )
    assert table.item(0, radii_type_column).text() == "Ionic (effective)"
    assert initial_sf_combo is not None
    if str(initial_sf_combo.currentData()) == "sphere":
        assert "r_eff_w0" in parameter_names
        assert "a_eff_w0" not in parameter_names
    else:
        assert "r_eff_w0" not in parameter_names
        assert "a_eff_w0" in parameter_names
        assert "b_eff_w0" in parameter_names
        assert "c_eff_w0" in parameter_names

    window.prefit_tab.toggle_cluster_geometry_radii_button.click()

    assert (
        window.prefit_tab.cluster_geometry_radii_type_combo.currentData()
        == ("bond_length")
    )
    assert table.item(0, radii_type_column).text() == "Bond length"
    sf_combo = table.cellWidget(0, sf_column)
    assert sf_combo is not None
    sphere_index = sf_combo.findData("sphere")
    assert sphere_index >= 0
    sf_combo.setCurrentIndex(sphere_index)
    sphere_parameter_names = [
        window.prefit_tab.parameter_table.item(row, 2).text()
        for row in range(window.prefit_tab.parameter_table.rowCount())
    ]

    assert table.item(0, radius_column).text() != initial_radius_text
    assert (
        table.item(0, radius_column).foreground().color().name().lower()
        == window.prefit_tab.ACTIVE_CLUSTER_GEOMETRY_COLOR.name().lower()
    )
    assert "r_eff_w0" in sphere_parameter_names
    assert "a_eff_w0" not in sphere_parameter_names
    assert "b_eff_w0" not in sphere_parameter_names
    assert "c_eff_w0" not in sphere_parameter_names

    ellipsoid_index = sf_combo.findData("ellipsoid")
    assert ellipsoid_index >= 0
    sf_combo.setCurrentIndex(ellipsoid_index)
    ellipsoid_parameter_names = [
        window.prefit_tab.parameter_table.item(row, 2).text()
        for row in range(window.prefit_tab.parameter_table.rowCount())
    ]

    updated_row = window.prefit_workflow.cluster_geometry_rows()[0]
    assert updated_row.sf_approximation == "ellipsoid"
    assert updated_row.radii_type_used == "bond_length"
    assert float(table.item(0, radius_column).text()) == pytest.approx(
        updated_row.effective_radius,
        rel=1e-3,
    )
    assert float(table.item(0, semiaxis_x_column).text()) == pytest.approx(
        updated_row.active_semiaxis_a,
        rel=1e-3,
    )
    assert float(table.item(0, semiaxis_y_column).text()) == pytest.approx(
        updated_row.active_semiaxis_b,
        rel=1e-3,
    )
    assert float(table.item(0, semiaxis_z_column).text()) == pytest.approx(
        updated_row.active_semiaxis_c,
        rel=1e-3,
    )
    assert (
        table.item(0, radius_column).foreground().color().name().lower()
        == window.prefit_tab.INACTIVE_CLUSTER_GEOMETRY_COLOR.name().lower()
    )
    for column in (
        semiaxis_x_column,
        semiaxis_y_column,
        semiaxis_z_column,
    ):
        assert (
            table.item(0, column).foreground().color().name().lower()
            == window.prefit_tab.ACTIVE_CLUSTER_GEOMETRY_COLOR.name().lower()
        )
    assert "r_eff_w0" not in ellipsoid_parameter_names
    assert "a_eff_w0" in ellipsoid_parameter_names
    assert "b_eff_w0" in ellipsoid_parameter_names
    assert "c_eff_w0" in ellipsoid_parameter_names
    window.close()


def test_prefit_cluster_geometry_single_atom_can_switch_to_bond_length_mode(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(
        tmp_path,
        template_name=POLY_LMA_HS_TEMPLATE,
        single_atom=True,
    )
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    radius_column = _table_column_index(table, "Effective Radius")
    errors: list[tuple[str, str]] = []
    window._show_error = lambda title, message: errors.append((title, message))

    ionic_radius = float(table.item(0, radius_column).text())
    window.prefit_tab.toggle_cluster_geometry_radii_button.click()

    assert not errors
    assert (
        window.prefit_tab.cluster_geometry_radii_type_combo.currentData()
        == ("bond_length")
    )
    assert window.prefit_workflow.cluster_geometry_active_radii_type() == (
        "bond_length"
    )
    bond_length_radius = float(table.item(0, radius_column).text())
    assert bond_length_radius == pytest.approx(1.39, rel=1e-3)
    assert 0.0 < bond_length_radius < ionic_radius
    assert window.prefit_tab.current_evaluation() is not None
    window.close()


def test_poly_lma_dream_parameter_map_tracks_cluster_geometry_shape(
    qapp,
    tmp_path,
):
    del qapp
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    sf_column = _table_column_index(
        window.prefit_tab.cluster_geometry_table,
        "S.F. Approx.",
    )
    sf_combo = window.prefit_tab.cluster_geometry_table.cellWidget(
        0, sf_column
    )
    assert sf_combo is not None

    sphere_index = sf_combo.findData("sphere")
    assert sphere_index >= 0
    sf_combo.setCurrentIndex(sphere_index)
    workflow = window._load_dream_workflow()
    sphere_entries = workflow.create_default_parameter_map(persist=False)
    sphere_names = [entry.param for entry in sphere_entries]
    assert "r_eff_w0" in sphere_names
    assert "a_eff_w0" not in sphere_names
    sphere_radius_entry = next(
        entry for entry in sphere_entries if entry.param == "r_eff_w0"
    )
    assert sphere_radius_entry.distribution == "lognorm"
    assert set(sphere_radius_entry.dist_params) == {"loc", "scale", "s"}

    ellipsoid_index = sf_combo.findData("ellipsoid")
    assert ellipsoid_index >= 0
    sf_combo.setCurrentIndex(ellipsoid_index)
    workflow = window._load_dream_workflow()
    ellipsoid_entries = workflow.create_default_parameter_map(persist=False)
    ellipsoid_names = [entry.param for entry in ellipsoid_entries]
    assert "r_eff_w0" not in ellipsoid_names
    assert "a_eff_w0" in ellipsoid_names
    assert "b_eff_w0" in ellipsoid_names
    assert "c_eff_w0" in ellipsoid_names
    ellipsoid_a_entry = next(
        entry for entry in ellipsoid_entries if entry.param == "a_eff_w0"
    )
    assert ellipsoid_a_entry.distribution == "norm"

    window.close()


def test_main_window_menus_expose_project_tools_and_help(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert window.file_menu.title() == "File"
    assert window.create_project_action.text() == "Create Project"
    assert window.open_project_action.text() == "Open Existing Project..."
    assert window.open_recent_menu.title() == "Open Recent Project"
    assert window.save_project_action.text() == "Save Project"
    assert window.save_project_action.shortcuts()
    assert window.save_project_as_action.text() == "Save Project As..."
    assert all(
        button.text() != "Save Project State"
        for button in window.project_setup_tab.project_group.findChildren(
            QPushButton
        )
    )

    assert window.tools_menu.title() == "Tools"
    assert window.md_extraction_menu.title() == "MD Extraction"
    assert window.structure_analysis_menu.title() == "Structure Analysis"
    assert window.visualization_menu.title() == "Visualization"
    assert window.mdtrajectory_action.text() == "Open MD Trajectory Extraction"
    assert window.xyz2pdb_action.text() == "Open XYZ -> PDB Conversion"
    assert window.cluster_action.text() == "Open Cluster Extraction"
    assert window.bondanalysis_action.text() == "Open Bond Analysis"
    assert (
        window.debye_waller_analysis_action.text()
        == "Open Debye-Waller Analysis"
    )
    assert [action.text() for action in window.tools_menu.actions()] == [
        "MD Extraction",
        "Structure Analysis",
        "Cluster Dynamics",
        "PDF",
        "Visualization",
        "SAXS Calculation Preview",
        "X-ray Toolkit",
    ]
    assert [
        action.text() for action in window.md_extraction_menu.actions()
    ] == [
        "Open MD Trajectory Extraction",
        "Open XYZ -> PDB Conversion",
        "Open Cluster Extraction",
    ]
    assert [
        action.text() for action in window.structure_analysis_menu.actions()
    ] == [
        "Open Bond Analysis",
        "Open Debye-Waller Analysis",
    ]
    assert (
        window.clusterdynamics_action.text() == "Open Cluster Dynamics (only)"
    )
    assert (
        window.clusterdynamicsml_action.text() == "Open Cluster Dynamics (ML)"
    )
    assert window.fullrmc_action.text() == "Open fullrmc Setup"
    assert window.structure_viewer_action.text() == "Structure Viewer"
    assert window.blenderxyz_action.text() == "Open Blender XYZ Renderer"
    assert window.component_calculation_preview_menu.title() == (
        "SAXS Calculation Preview"
    )
    assert window.contrast_mode_action.text() == "Open SAXS Contrast Mode"
    assert (
        window.electron_density_mapping_action.text()
        == "Open Electron Density Mapping"
    )
    assert window.xray_toolkit_menu.title() == "X-ray Toolkit"
    assert (
        window.volume_fraction_action.text() == "Open Volume Fraction Estimate"
    )
    assert (
        window.number_density_action.text() == "Open Number Density Estimate"
    )

    assert window.settings_menu.title() == "Settings"
    assert (
        window.console_autoscroll_action.text() == "Autoscroll Console Output"
    )
    assert window.dream_output_settings_action.text() == "Main UI Settings..."
    assert window.window_presets_menu.title() == "Window Presets"
    assert window.auto_fit_window_action.text() == "Auto Fit Current Screen"
    preset_labels = [
        action.text()
        for action in window.window_presets_menu.actions()
        if action.text()
    ]
    assert "13-inch Laptop (Compact)" in preset_labels
    assert "14-inch Laptop / MacBook Pro" in preset_labels
    assert "15-inch / 16-inch Laptop" in preset_labels
    assert "External Display (1080p)" in preset_labels
    assert "External Display (1440p / QHD)" in preset_labels


def test_project_setup_shows_prep_help_tooltips(qapp):
    del qapp
    window = SAXSMainWindow()

    tooltip = window.project_setup_tab.open_mdtrajectory_help_button.toolTip()

    assert (
        window.project_setup_tab.open_mdtrajectory_button.text()
        == "Open MD Trajectory Extraction"
    )
    assert (
        window.project_setup_tab.open_xyz2pdb_button.text()
        == "Open XYZ -> PDB Conversion"
    )
    assert (
        window.project_setup_tab.open_cluster_button.text()
        == "Open Cluster Extraction"
    )
    assert "first extract frames" in tooltip.lower()
    assert "optionally convert" in tooltip.lower()
    assert "cluster extraction" in tooltip.lower()
    assert (
        window.project_setup_tab.open_cluster_help_button.toolTip() == tooltip
    )
    assert (
        window.project_setup_tab.open_xyz2pdb_help_button.toolTip() == tooltip
    )
    window.close()

    assert window.help_menu.title() == "Help"
    assert window.version_info_action.text() == "Version Information"
    assert window.github_action.text() == "Open GitHub Repository"
    assert window.contact_action.text() == "Contact Developer"

    version_text = window._version_information_text()
    assert __version__
    assert "Package version:" in version_text
    assert "GitHub repository:" in version_text
    assert "Developer contact:" in version_text
    assert "keith.white@colorado.edu" in version_text


def test_project_setup_prep_buttons_span_form_row(qapp):
    del qapp
    window = SAXSMainWindow()
    window.show()
    QTest.qWait(0)

    layout = window.project_setup_tab.forward_model_group.layout()
    prep_row = layout.itemAt(0, QFormLayout.ItemRole.SpanningRole).widget()
    frames_row = layout.itemAt(1, QFormLayout.ItemRole.FieldRole).widget()

    assert isinstance(layout, QFormLayout)
    assert (
        prep_row
        is window.project_setup_tab.open_mdtrajectory_button.parentWidget()
    )
    assert frames_row is not None
    assert prep_row.x() < frames_row.x()

    window.close()


def test_console_autoscroll_setting_controls_tab_output_scroll(
    qapp, tmp_path, monkeypatch
):
    class _FakeSettings:
        def __init__(self):
            self.values: dict[str, object] = {}

        def value(self, key, default=None):
            return self.values.get(key, default)

        def setValue(self, key, value):
            self.values[key] = value

    settings_store = _FakeSettings()
    monkeypatch.setattr(
        SAXSMainWindow,
        "_recent_projects_settings",
        lambda self: settings_store,
    )

    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.resize(920, 720)
    window.project_setup_tab.summary_box.setFixedHeight(90)
    window.prefit_tab.output_box.setFixedHeight(90)
    window.dream_tab.output_box.setFixedHeight(90)
    window.show()
    qapp.processEvents()

    assert window.console_autoscroll_action.isChecked()

    def _activate_tab(tab):
        window.tabs.setCurrentWidget(tab)
        qapp.processEvents()

    def _append_project_lines(start: int, stop: int):
        _activate_tab(window.project_setup_tab)
        for index in range(start, stop):
            window.project_setup_tab.append_summary(
                f"project summary line {index}"
            )
        qapp.processEvents()
        return window.project_setup_tab.summary_box.verticalScrollBar()

    def _append_prefit_lines(start: int, stop: int):
        _activate_tab(window.prefit_tab)
        for index in range(start, stop):
            window.prefit_tab.append_log(f"prefit log line {index}")
        qapp.processEvents()
        return window.prefit_tab.output_box.verticalScrollBar()

    def _append_dream_lines(start: int, stop: int):
        _activate_tab(window.dream_tab)
        for index in range(start, stop):
            window.dream_tab.append_log(f"dream log line {index}")
        qapp.processEvents()
        return window.dream_tab.output_box.verticalScrollBar()

    def _is_near_bottom(scrollbar) -> bool:
        return scrollbar.maximum() - scrollbar.value() <= max(
            int(scrollbar.pageStep()), 4
        )

    project_scrollbar = _append_project_lines(0, 80)
    prefit_scrollbar = _append_prefit_lines(0, 80)
    dream_scrollbar = _append_dream_lines(0, 80)

    assert project_scrollbar.maximum() > 0
    assert prefit_scrollbar.maximum() > 0
    assert dream_scrollbar.maximum() > 0
    assert _is_near_bottom(project_scrollbar)
    assert _is_near_bottom(prefit_scrollbar)
    assert _is_near_bottom(dream_scrollbar)

    window.console_autoscroll_action.trigger()
    qapp.processEvents()

    assert settings_store.values["console_autoscroll_enabled"] is False
    assert not window.console_autoscroll_action.isChecked()

    project_scrollbar.setValue(0)
    prefit_scrollbar.setValue(0)
    dream_scrollbar.setValue(0)
    qapp.processEvents()

    project_scrollbar = _append_project_lines(80, 100)
    prefit_scrollbar = _append_prefit_lines(80, 100)
    dream_scrollbar = _append_dream_lines(80, 100)

    assert not _is_near_bottom(project_scrollbar)
    assert not _is_near_bottom(prefit_scrollbar)
    assert not _is_near_bottom(dream_scrollbar)

    window.console_autoscroll_action.trigger()
    qapp.processEvents()

    assert settings_store.values["console_autoscroll_enabled"] is True
    project_scrollbar = _append_project_lines(100, 101)
    prefit_scrollbar = _append_prefit_lines(100, 101)
    dream_scrollbar = _append_dream_lines(100, 101)
    assert _is_near_bottom(project_scrollbar)
    assert _is_near_bottom(prefit_scrollbar)
    assert _is_near_bottom(dream_scrollbar)
    window.close()


def test_volume_fraction_tool_window_opens_with_citation_and_target(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window._open_solute_volume_fraction_tool()

    tool_window = window._solute_volume_fraction_tool_window
    assert tool_window is not None
    assert tool_window.windowTitle() == "Volume Fraction Estimate"
    labels = [label.text() for label in tool_window.findChildren(QLabel)]
    assert any(SOLUTE_VOLUME_FRACTION_CITATION_URL in text for text in labels)
    assert "phi_solute" in tool_window.estimator_widget.target_label.text()
    assert "solvent_scale" in tool_window.estimator_widget.target_label.text()
    assert (
        tool_window.estimator_widget.calculate_number_density_checkbox.isChecked()
    )
    molarity_index = tool_window.estimator_widget.solution_mode_combo.findData(
        "molarity_per_liter"
    )
    assert molarity_index >= 0
    tool_window.estimator_widget.solution_mode_combo.setCurrentIndex(
        molarity_index
    )
    assert tool_window.estimator_widget.solute_density_label.isHidden()
    assert tool_window.estimator_widget.solute_density_spin.isHidden()
    assert not tool_window.estimator_widget.solvent_density_label.isHidden()
    assert not tool_window.estimator_widget.solvent_density_spin.isHidden()
    assert "volume-closure calculations" in (
        tool_window.estimator_widget.solution_mode_hint_label.text().lower()
    )

    tool_window.estimator_widget.solution_density_spin.setValue(1.0)
    tool_window.estimator_widget.solute_stoich_edit.setText("Cs1Pb1I3")
    tool_window.estimator_widget.solvent_stoich_edit.setText("H2O")
    tool_window.estimator_widget.molar_mass_solute_spin.setValue(620.0)
    tool_window.estimator_widget.molar_mass_solvent_spin.setValue(18.015)
    tool_window.estimator_widget.molarity_spin.setValue(0.5)
    tool_window.estimator_widget.molarity_element_edit.setText("Pb")
    tool_window.estimator_widget.solvent_density_spin.setValue(1.0)
    tool_window.estimator_widget.calculate_button.click()

    assert "Number density estimate" in (
        tool_window.estimator_widget.output_box.toPlainText()
    )
    assert "atoms/A^3" in tool_window.estimator_widget.output_box.toPlainText()

    tool_window.close()
    window.close()


def test_number_density_attenuation_and_fluorescence_tool_windows_open_with_defaults(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window._open_number_density_tool()
    number_density_window = window._number_density_tool_window
    assert number_density_window is not None
    assert number_density_window.windowTitle() == "Number Density Estimate"
    assert (
        number_density_window.estimator_widget.calculate_number_density_checkbox.isChecked()
    )
    assert (
        not number_density_window.estimator_widget.calculate_volume_fraction_checkbox.isChecked()
    )
    assert (
        not number_density_window.estimator_widget.calculate_attenuation_checkbox.isChecked()
    )
    assert (
        not number_density_window.estimator_widget.calculate_fluorescence_checkbox.isChecked()
    )

    window._open_attenuation_tool()
    attenuation_window = window._attenuation_tool_window
    assert attenuation_window is not None
    assert attenuation_window.windowTitle() == "Attenuation Estimate"
    assert (
        not attenuation_window.estimator_widget.calculate_number_density_checkbox.isChecked()
    )
    assert (
        attenuation_window.estimator_widget.calculate_volume_fraction_checkbox.isChecked()
    )
    assert (
        attenuation_window.estimator_widget.calculate_attenuation_checkbox.isChecked()
    )
    assert (
        not attenuation_window.estimator_widget.calculate_fluorescence_checkbox.isChecked()
    )
    assert (
        "solvent_scale"
        in attenuation_window.estimator_widget.target_label.text()
    )

    window._open_fluorescence_tool()
    fluorescence_window = window._fluorescence_tool_window
    assert fluorescence_window is not None
    assert fluorescence_window.windowTitle() == "Fluorescence Estimate"
    assert (
        not fluorescence_window.estimator_widget.calculate_number_density_checkbox.isChecked()
    )
    assert (
        not fluorescence_window.estimator_widget.calculate_volume_fraction_checkbox.isChecked()
    )
    assert (
        not fluorescence_window.estimator_widget.calculate_attenuation_checkbox.isChecked()
    )
    assert (
        fluorescence_window.estimator_widget.calculate_fluorescence_checkbox.isChecked()
    )
    assert (
        "solvent_scale"
        in fluorescence_window.estimator_widget.target_label.text()
    )

    number_density_window.close()
    attenuation_window.close()
    fluorescence_window.close()
    window.close()


def test_contact_action_opens_developer_contact_window(qapp, monkeypatch):
    del qapp
    window = SAXSMainWindow()
    captured: dict[str, str] = {}

    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.QMessageBox.information",
        lambda parent, title, message: captured.update(
            {
                "title": title,
                "message": message,
            }
        ),
    )

    window._show_contact_information()

    assert captured["title"] == "Developer Contact"
    assert "contact the developer" in captured["message"].lower()
    assert "keith.white@colorado.edu" in captured["message"]


def test_bondanalysis_tool_uses_active_project_clusters_dir(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    clusters_dir = tmp_path / "clusters"
    clusters_dir.mkdir()
    window.current_settings.clusters_dir = str(clusters_dir)
    window.project_manager.save_project(window.current_settings)
    launched: dict[str, object] = {}

    class FakeBondAnalysisWindow:
        def __init__(
            self, initial_clusters_dir=None, initial_project_dir=None
        ):
            launched["clusters_dir"] = initial_clusters_dir
            launched["project_dir"] = initial_project_dir
            launched["instance"] = self

        def show(self):
            launched["shown"] = True

        def raise_(self):
            launched["raised"] = True

    monkeypatch.setattr(
        "saxshell.bondanalysis.ui.main_window.BondAnalysisMainWindow",
        FakeBondAnalysisWindow,
    )

    window._open_bondanalysis_tool()

    assert launched["clusters_dir"] == str(clusters_dir.resolve())
    assert launched["project_dir"] == Path(project_dir).resolve()
    assert launched["shown"] is True
    assert launched["raised"] is True
    assert launched["instance"] in window._child_tool_windows
    window.close()


def test_mdtrajectory_tool_uses_active_project_references(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    trajectory_file = tmp_path / "traj.xyz"
    topology_file = tmp_path / "topology.pdb"
    energy_file = tmp_path / "traj.ener"
    trajectory_file.write_text("1\nframe\nPb 0.0 0.0 0.0\n", encoding="utf-8")
    topology_file.write_text("MODEL        1\nENDMDL\n", encoding="utf-8")
    energy_file.write_text(
        "# step time kinetic temperature potential\n"
        "1 0.0 1.0 300.0 -10.0\n",
        encoding="utf-8",
    )
    window.current_settings.trajectory_file = str(trajectory_file)
    window.current_settings.topology_file = str(topology_file)
    window.current_settings.energy_file = str(energy_file)
    window.project_manager.save_project(window.current_settings)
    launched: dict[str, object] = {}

    class FakeMDTrajectoryWindow:
        pass

    def fake_launch_mdtrajectory_app(**kwargs):
        launched.update(kwargs)
        launched["instance"] = FakeMDTrajectoryWindow()
        return launched["instance"]

    monkeypatch.setattr(
        "saxshell.mdtrajectory.ui.main_window.launch_mdtrajectory_app",
        fake_launch_mdtrajectory_app,
    )

    window._open_mdtrajectory_tool()

    assert launched["project_dir"] == Path(project_dir).resolve()
    assert launched["trajectory_file"] == trajectory_file.resolve()
    assert launched["topology_file"] == topology_file.resolve()
    assert launched["energy_file"] == energy_file.resolve()
    assert launched["instance"] in window._child_tool_windows
    window.close()


def test_debye_waller_tool_uses_active_project_clusters_dir(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    clusters_dir = tmp_path / "clusters"
    clusters_dir.mkdir()
    window.current_settings.clusters_dir = str(clusters_dir.resolve())
    window.project_manager.save_project(window.current_settings)
    launched: dict[str, object] = {}

    class FakeDebyeWallerWindow(QWidget):
        def __init__(self):
            super().__init__()
            launched["instance"] = self

    def fake_launch_debye_waller_analysis_ui(**kwargs):
        launched.update(kwargs)
        return FakeDebyeWallerWindow()

    monkeypatch.setattr(
        "saxshell.saxs.debye_waller.ui.main_window.launch_debye_waller_analysis_ui",
        fake_launch_debye_waller_analysis_ui,
    )

    window._open_debye_waller_analysis_tool()

    assert launched["initial_project_dir"] == Path(project_dir).resolve()
    assert launched["initial_clusters_dir"] == clusters_dir.resolve()
    assert launched["instance"] in window._child_tool_windows
    window.close()


def _configure_debye_waller_project_clusters(
    project_dir: Path,
    tmp_path: Path,
    *,
    dirname: str = "project_dw_clusters",
) -> Path:
    def pdb_atom_line(
        atom_id: int,
        atom_name: str,
        residue_name: str,
        residue_number: int,
        x_coord: float,
        y_coord: float,
        z_coord: float,
        element: str,
    ) -> str:
        return (
            f"ATOM  {atom_id:5d} {atom_name:<4} {residue_name:>3}  "
            f"{residue_number:4d}    "
            f"{x_coord:8.3f}{y_coord:8.3f}{z_coord:8.3f}"
            f"  1.00  0.00          {element:>2}\n"
        )

    def write_water_frame(
        path: Path,
        *,
        mol1_oxygen_x: float,
        mol1_hydrogen_x: float,
        mol2_oxygen_x: float,
        mol2_hydrogen_x: float,
    ) -> None:
        lines = [
            pdb_atom_line(1, "O", "HOH", 1, mol1_oxygen_x, 0.0, 0.0, "O"),
            pdb_atom_line(2, "H1", "HOH", 1, mol1_hydrogen_x, 0.0, 0.0, "H"),
            pdb_atom_line(3, "O", "HOH", 2, mol2_oxygen_x, 0.0, 0.0, "O"),
            pdb_atom_line(4, "H1", "HOH", 2, mol2_hydrogen_x, 0.0, 0.0, "H"),
            "END\n",
        ]
        path.write_text("".join(lines), encoding="utf-8")

    clusters_dir = tmp_path / dirname
    structure_dir = clusters_dir / "H2O2"
    structure_dir.mkdir(parents=True, exist_ok=True)
    write_water_frame(
        structure_dir / "water_frame_0001.pdb",
        mol1_oxygen_x=0.0,
        mol1_hydrogen_x=1.00,
        mol2_oxygen_x=5.0,
        mol2_hydrogen_x=6.00,
    )
    write_water_frame(
        structure_dir / "water_frame_0002.pdb",
        mol1_oxygen_x=0.0,
        mol1_hydrogen_x=1.10,
        mol2_oxygen_x=5.2,
        mol2_hydrogen_x=6.35,
    )

    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.clusters_dir = str(clusters_dir.resolve())
    manager.save_project(settings)
    return clusters_dir


def _write_debye_waller_project_result(project_dir: Path, tmp_path: Path):
    from saxshell.saxs.debye_waller.workflow import (
        DebyeWallerWorkflow,
        save_debye_waller_analysis_to_project,
    )

    clusters_dir = _configure_debye_waller_project_clusters(
        project_dir,
        tmp_path,
    )

    result = DebyeWallerWorkflow(
        clusters_dir,
        project_dir=project_dir,
        output_dir=tmp_path / "dw_external_out",
        output_basename="dw_ui",
    ).run()
    return save_debye_waller_analysis_to_project(result, project_dir)


def test_project_setup_debye_waller_button_requires_active_pdb_clusters_and_tracks_saved_result(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert not window.project_setup_tab.debye_waller_button.isEnabled()
    assert "active PDB clusters folder" in (
        window.project_setup_tab.debye_waller_button.toolTip()
    )
    assert "#6b7280" in (
        window.project_setup_tab.debye_waller_ready_indicator.styleSheet()
    )

    _write_debye_waller_project_result(project_dir, tmp_path)
    window.load_project(project_dir)

    assert window.project_setup_tab.debye_waller_button.isEnabled()
    assert "#16a34a" in (
        window.project_setup_tab.debye_waller_ready_indicator.styleSheet()
    )
    assert (
        "computed and saved"
        in (
            window.project_setup_tab.debye_waller_ready_indicator.toolTip()
        ).lower()
    )
    help_tooltip = (
        window.project_setup_tab.debye_waller_help_button.toolTip().lower()
    )
    assert "optional" in help_tooltip
    assert "before building saxs components" in help_tooltip
    assert "pdb cluster" in help_tooltip

    _configure_debye_waller_project_clusters(
        project_dir,
        tmp_path,
        dirname="project_dw_clusters_alt",
    )
    window.load_project(project_dir)

    assert window.project_setup_tab.debye_waller_button.isEnabled()
    assert "#6b7280" in (
        window.project_setup_tab.debye_waller_ready_indicator.styleSheet()
    )
    assert (
        "different clusters folder"
        in (
            window.project_setup_tab.debye_waller_ready_indicator.toolTip()
        ).lower()
    )
    window.close()


def test_project_setup_debye_waller_button_launches_linked_tool_and_refreshes_status(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    from saxshell.saxs.debye_waller.workflow import (
        DebyeWallerWorkflow,
        save_debye_waller_analysis_to_project,
    )

    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    clusters_dir = _configure_debye_waller_project_clusters(
        project_dir,
        tmp_path,
        dirname="project_dw_clusters_launch",
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)
    launched: dict[str, object] = {}

    class FakeDebyeWallerWindow(QWidget):
        project_paths_registered = Signal(object)
        project_analysis_saved = Signal(object)

        def __init__(self):
            super().__init__()
            launched["instance"] = self

    def fake_launch_debye_waller_analysis_ui(**kwargs):
        launched.update(kwargs)
        return FakeDebyeWallerWindow()

    monkeypatch.setattr(
        "saxshell.saxs.debye_waller.ui.main_window.launch_debye_waller_analysis_ui",
        fake_launch_debye_waller_analysis_ui,
    )

    assert window.project_setup_tab.debye_waller_button.isEnabled()
    assert "#6b7280" in (
        window.project_setup_tab.debye_waller_ready_indicator.styleSheet()
    )

    window.project_setup_tab.debye_waller_button.click()

    assert launched["initial_project_dir"] == Path(project_dir).resolve()
    assert launched["initial_clusters_dir"] == clusters_dir.resolve()
    assert launched["instance"] in window._child_tool_windows

    result = DebyeWallerWorkflow(
        clusters_dir,
        project_dir=project_dir,
        output_dir=tmp_path / "dw_linked_out",
        output_basename="linked_refresh",
    ).run()
    save_debye_waller_analysis_to_project(result, project_dir)
    launched["instance"].project_analysis_saved.emit(
        {"project_dir": str(Path(project_dir).resolve())}
    )
    QApplication.processEvents()

    assert "#16a34a" in (
        window.project_setup_tab.debye_waller_ready_indicator.styleSheet()
    )
    window.close()


def test_project_setup_debye_waller_button_reports_startup_progress_in_popup(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    from saxshell.saxs.debye_waller.ui.main_window import (
        DEBYE_WALLER_WINDOW_LOAD_TOTAL_STEPS,
    )

    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    clusters_dir = _configure_debye_waller_project_clusters(
        project_dir,
        tmp_path,
        dirname="project_dw_clusters_popup",
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)
    launched: dict[str, object] = {}

    class FakeDebyeWallerWindow(QWidget):
        project_paths_registered = Signal(object)
        project_analysis_saved = Signal(object)

    def fake_launch_debye_waller_analysis_ui(**kwargs):
        launched.update(kwargs)
        progress_callback = kwargs.get("startup_progress_callback")
        log_callback = kwargs.get("startup_log_callback")
        assert progress_callback is not None
        assert log_callback is not None
        progress_callback(
            1,
            DEBYE_WALLER_WINDOW_LOAD_TOTAL_STEPS,
            "Preparing Debye-Waller analysis window...",
        )
        log_callback("Preparing Debye-Waller analysis window.")
        progress_callback(
            DEBYE_WALLER_WINDOW_LOAD_TOTAL_STEPS,
            DEBYE_WALLER_WINDOW_LOAD_TOTAL_STEPS,
            "Finalizing Debye-Waller window...",
        )
        log_callback("Debye-Waller analysis window is ready.")
        return FakeDebyeWallerWindow()

    monkeypatch.setattr(
        "saxshell.saxs.debye_waller.ui.main_window.launch_debye_waller_analysis_ui",
        fake_launch_debye_waller_analysis_ui,
    )

    window.project_setup_tab.debye_waller_button.click()

    assert launched["initial_project_dir"] == Path(project_dir).resolve()
    assert launched["initial_clusters_dir"] == clusters_dir.resolve()
    assert window._progress_dialog is not None
    assert not window._progress_dialog.isVisible()
    dialog_output = window._progress_dialog.output_box.toPlainText()
    assert "Loading Debye-Waller analysis from" in dialog_output
    assert "Preparing Debye-Waller analysis window." in dialog_output
    assert "Debye-Waller analysis window is ready." in dialog_output
    window.close()


def test_debye_waller_window_inherits_clusters_dir_from_project_reference(
    qapp,
    tmp_path,
):
    del qapp
    from saxshell.saxs.debye_waller.ui.main_window import (
        DebyeWallerAnalysisMainWindow,
    )

    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    clusters_dir = tmp_path / "project_clusters"
    clusters_dir.mkdir()

    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.clusters_dir = str(clusters_dir.resolve())
    manager.save_project(settings)

    window = DebyeWallerAnalysisMainWindow(initial_project_dir=project_dir)

    assert window.clusters_dir_edit.text() == str(clusters_dir.resolve())
    assert (
        "inherited from the active project"
        in window.summary_box.toPlainText().lower()
    )
    window.close()


def test_debye_waller_window_browsing_clusters_reference_skips_registered_path_refresh(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    from saxshell.saxs.debye_waller.ui.main_window import (
        DebyeWallerAnalysisMainWindow,
    )

    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    initial_clusters_dir = tmp_path / "initial_clusters"
    initial_clusters_dir.mkdir()
    updated_clusters_dir = tmp_path / "updated_clusters"
    updated_clusters_dir.mkdir()

    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.clusters_dir = str(initial_clusters_dir.resolve())
    manager.save_project(settings)
    assert manager.load_project(project_dir).clusters_dir_snapshot is not None

    window = DebyeWallerAnalysisMainWindow(initial_project_dir=project_dir)

    save_calls: list[bool] = []
    original_save_project = window._project_manager.save_project

    def record_save_project(settings, *, refresh_registered_paths=True):
        save_calls.append(bool(refresh_registered_paths))
        return original_save_project(
            settings,
            refresh_registered_paths=refresh_registered_paths,
        )

    monkeypatch.setattr(
        window._project_manager,
        "save_project",
        record_save_project,
    )
    monkeypatch.setattr(
        "saxshell.saxs.debye_waller.ui.main_window.QFileDialog.getExistingDirectory",
        lambda *args, **kwargs: str(updated_clusters_dir.resolve()),
    )

    window._browse_clusters_dir()

    assert save_calls == [False]
    saved_settings = SAXSProjectManager().load_project(project_dir)
    assert saved_settings.clusters_dir == str(updated_clusters_dir.resolve())
    assert saved_settings.clusters_dir_snapshot is None
    window.close()


def test_debye_waller_window_exposes_stoichiometry_info_tab(
    qapp,
    tmp_path,
):
    del qapp
    from saxshell.saxs.debye_waller.ui.main_window import (
        DebyeWallerAnalysisMainWindow,
    )

    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = DebyeWallerAnalysisMainWindow(initial_project_dir=project_dir)

    tab_labels = [
        window.results_tabs.tabText(index)
        for index in range(window.results_tabs.count())
    ]

    assert "Aggregated Pairs" in tab_labels
    assert "Stoichiometries" in tab_labels
    assert window.stoichiometry_info_table.columnCount() >= 1
    assert window.aggregated_pair_table.columnCount() >= 1
    window.close()


def test_debye_waller_window_loads_saved_project_analysis(
    qapp,
    tmp_path,
):
    del qapp
    from saxshell.saxs.debye_waller.ui.main_window import (
        DebyeWallerAnalysisMainWindow,
    )

    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    saved_result = _write_debye_waller_project_result(project_dir, tmp_path)

    window = DebyeWallerAnalysisMainWindow(initial_project_dir=project_dir)

    assert window._last_result is not None
    assert window.aggregated_pair_table.rowCount() > 0
    assert window.pair_summary_table.rowCount() > 0
    assert window.stoichiometry_info_table.rowCount() > 0
    assert "loaded saved debye-waller analysis" in (
        window.log_box.toPlainText().lower()
    )
    assert window.output_dir_edit.text() == str(saved_result.output_dir)
    window.close()


def test_debye_waller_window_does_not_resave_project_during_startup(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    from saxshell.saxs.debye_waller.ui.main_window import (
        DebyeWallerAnalysisMainWindow,
    )

    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    clusters_dir = _configure_debye_waller_project_clusters(
        project_dir,
        tmp_path,
        dirname="project_dw_clusters_no_resave",
    )
    save_calls: list[str | None] = []
    original_save_project = SAXSProjectManager.save_project

    def wrapped_save_project(self, settings, *args, **kwargs):
        save_calls.append(settings.clusters_dir)
        return original_save_project(self, settings, *args, **kwargs)

    monkeypatch.setattr(
        SAXSProjectManager,
        "save_project",
        wrapped_save_project,
    )

    window = DebyeWallerAnalysisMainWindow(
        initial_project_dir=project_dir,
        initial_clusters_dir=clusters_dir,
    )

    assert save_calls == []
    window.close()


def test_debye_waller_window_autosaves_first_project_analysis(
    qapp,
    tmp_path,
):
    del qapp
    from saxshell.saxs.debye_waller.ui.main_window import (
        DebyeWallerAnalysisMainWindow,
    )
    from saxshell.saxs.debye_waller.workflow import (
        DebyeWallerWorkflow,
        find_saved_project_debye_waller_analysis,
    )

    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    saved_result = _write_debye_waller_project_result(project_dir, tmp_path)
    saved_summary_path = (
        saved_result.artifacts.summary_json_path
        if saved_result.artifacts is not None
        else None
    )
    assert saved_summary_path is not None
    saved_summary_path.unlink()

    pair_csv = saved_result.artifacts.pair_summary_csv_path
    scope_csv = saved_result.artifacts.scope_summary_csv_path
    segment_csv = saved_result.artifacts.segment_csv_path
    aggregated_pair_csv = (
        saved_result.artifacts.aggregated_pair_summary_csv_path
    )
    for artifact_path in (
        aggregated_pair_csv,
        pair_csv,
        scope_csv,
        segment_csv,
    ):
        artifact_path.unlink()
    saved_result.output_dir.rmdir()

    clusters_dir = Path(
        SAXSProjectManager().load_project(project_dir).clusters_dir
    ).resolve()
    result = DebyeWallerWorkflow(
        clusters_dir,
        project_dir=project_dir,
        output_dir=tmp_path / "dw_new_run",
        output_basename="new_run",
    ).run()

    window = DebyeWallerAnalysisMainWindow(initial_project_dir=project_dir)
    window._project_had_saved_analysis_before_run = False
    window._finish_run(result)

    restored_path = find_saved_project_debye_waller_analysis(project_dir)
    assert restored_path is not None
    assert restored_path.exists()
    assert "auto-saved the first debye-waller analysis" in (
        window.log_box.toPlainText().lower()
    )
    assert window.save_project_button.isEnabled()
    window.close()


def test_cluster_tool_uses_active_project_frames_dir_and_project_dir(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    frames_dir = tmp_path / "splitxyz_f0fs"
    frames_dir.mkdir()
    window.current_settings.frames_dir = str(frames_dir)
    window.project_manager.save_project(window.current_settings)
    launched: dict[str, object] = {}

    class FakeClusterWindow:
        def __init__(self, initial_frames_dir=None, initial_project_dir=None):
            launched["frames_dir"] = initial_frames_dir
            launched["project_dir"] = initial_project_dir
            launched["instance"] = self

        def show(self):
            launched["shown"] = True

        def raise_(self):
            launched["raised"] = True

    monkeypatch.setattr(
        "saxshell.cluster.ui.main_window.ClusterMainWindow",
        FakeClusterWindow,
    )

    window._open_cluster_tool()

    assert launched["frames_dir"] == frames_dir.resolve()
    assert launched["project_dir"] == Path(project_dir).resolve()
    assert launched["shown"] is True
    assert launched["raised"] is True
    assert launched["instance"] in window._child_tool_windows
    window.close()


def test_xyz2pdb_tool_uses_active_project_frames_dir_and_project_dir(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    frames_dir = tmp_path / "splitxyz_f0fs"
    frames_dir.mkdir()
    window.current_settings.frames_dir = str(frames_dir)
    window.project_manager.save_project(window.current_settings)
    launched: dict[str, object] = {}

    class FakeXYZToPDBWindow:
        pass

    def fake_launch_xyz2pdb_ui(**kwargs):
        launched.update(kwargs)
        launched["instance"] = FakeXYZToPDBWindow()
        return launched["instance"]

    monkeypatch.setattr(
        "saxshell.xyz2pdb.ui.main_window.launch_xyz2pdb_ui",
        fake_launch_xyz2pdb_ui,
    )

    window._open_xyz2pdb_tool()

    assert launched["input_path"] == frames_dir.resolve()
    assert launched["project_dir"] == Path(project_dir).resolve()
    assert launched["instance"] in window._child_tool_windows
    window.close()


def test_mdtrajectory_tool_updates_main_project_frames_dir_from_child(
    qapp,
    tmp_path,
    monkeypatch,
):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    exported_frames_dir = tmp_path / "splitxyz_f5fs"
    exported_frames_dir.mkdir()

    saved_settings = window.project_manager.load_project(project_dir)
    saved_settings.frames_dir = str(exported_frames_dir.resolve())
    window.project_manager.save_project(saved_settings)

    class FakeMDTrajectoryWindow(QObject):
        project_paths_registered = Signal(object)

    fake_window = FakeMDTrajectoryWindow()

    def fake_launch_mdtrajectory_app(**kwargs):
        del kwargs
        return fake_window

    monkeypatch.setattr(
        "saxshell.mdtrajectory.ui.main_window.launch_mdtrajectory_app",
        fake_launch_mdtrajectory_app,
    )

    window._open_mdtrajectory_tool()
    fake_window.project_paths_registered.emit(
        {
            "project_dir": Path(project_dir).resolve(),
            "frames_dir": exported_frames_dir.resolve(),
        }
    )
    qapp.processEvents()

    assert (
        window.project_setup_tab.frames_dir() == exported_frames_dir.resolve()
    )
    assert window.current_settings is not None
    assert window.current_settings.resolved_frames_dir == (
        exported_frames_dir.resolve()
    )
    assert "linked tool" in window.project_setup_tab.summary_box.toPlainText()
    window.close()


def test_xyz2pdb_tool_updates_main_project_pdb_folder_from_child(
    qapp,
    tmp_path,
    monkeypatch,
):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    pdb_frames_dir = tmp_path / "xyz2pdb_splitxyz_f5fs"
    pdb_frames_dir.mkdir()
    saved_settings = window.project_manager.load_project(project_dir)
    saved_settings.pdb_frames_dir = str(pdb_frames_dir.resolve())
    window.project_manager.save_project(saved_settings)

    class FakeXYZToPDBWindow(QObject):
        project_paths_registered = Signal(object)

    fake_window = FakeXYZToPDBWindow()

    def fake_launch_xyz2pdb_ui(**kwargs):
        del kwargs
        return fake_window

    monkeypatch.setattr(
        "saxshell.xyz2pdb.ui.main_window.launch_xyz2pdb_ui",
        fake_launch_xyz2pdb_ui,
    )

    window._open_xyz2pdb_tool()
    fake_window.project_paths_registered.emit(
        {
            "project_dir": Path(project_dir).resolve(),
            "pdb_frames_dir": pdb_frames_dir.resolve(),
        }
    )
    qapp.processEvents()

    assert (
        window.project_setup_tab.pdb_frames_dir() == pdb_frames_dir.resolve()
    )
    assert window.current_settings is not None
    assert window.current_settings.resolved_pdb_frames_dir == (
        pdb_frames_dir.resolve()
    )
    window.close()


def test_cluster_tool_updates_main_project_folder_refs_from_child(
    qapp,
    tmp_path,
    monkeypatch,
):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    frames_dir = tmp_path / "splitxyz_f10fs"
    frames_dir.mkdir()
    clusters_dir = tmp_path / "clusters_splitxyz_f10fs"
    clusters_dir.mkdir()

    saved_settings = window.project_manager.load_project(project_dir)
    saved_settings.frames_dir = str(frames_dir.resolve())
    saved_settings.clusters_dir = str(clusters_dir.resolve())
    window.project_manager.save_project(saved_settings)

    scan_calls: list[bool] = []
    window.project_setup_tab.request_cluster_scan = lambda: scan_calls.append(
        True
    )

    class FakeClusterWindow(QWidget):
        project_paths_registered = Signal(object)

        def __init__(self, initial_frames_dir=None, initial_project_dir=None):
            super().__init__()
            self.initial_frames_dir = initial_frames_dir
            self.initial_project_dir = initial_project_dir

        def show(self):
            return None

        def raise_(self):
            return None

    fake_window: FakeClusterWindow | None = None

    def fake_cluster_window(*args, **kwargs):
        nonlocal fake_window
        fake_window = FakeClusterWindow(*args, **kwargs)
        return fake_window

    monkeypatch.setattr(
        "saxshell.cluster.ui.main_window.ClusterMainWindow",
        fake_cluster_window,
    )

    window._open_cluster_tool()
    assert fake_window is not None
    fake_window.project_paths_registered.emit(
        {
            "project_dir": Path(project_dir).resolve(),
            "frames_dir": frames_dir.resolve(),
            "clusters_dir": clusters_dir.resolve(),
        }
    )
    qapp.processEvents()

    assert window.project_setup_tab.frames_dir() == frames_dir.resolve()
    assert window.project_setup_tab.clusters_dir() == clusters_dir.resolve()
    assert window.current_settings is not None
    assert window.current_settings.resolved_frames_dir == frames_dir.resolve()
    assert window.current_settings.resolved_clusters_dir == (
        clusters_dir.resolve()
    )
    assert scan_calls == [True]
    window.close()


def test_main_window_refuses_close_when_child_tool_refuses_close(
    qapp,
    tmp_path,
):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.show()

    class BusyChild(QWidget):
        def close(self):
            return False

    child = BusyChild()
    window._track_child_tool_window(child)
    qapp.processEvents()

    assert not window.close()
    assert window.isVisible()
    assert (
        "linked tool is still busy"
        in window.statusBar().currentMessage().lower()
    )
    window.hide()


def test_fullrmc_tool_uses_active_project_dir(qapp, tmp_path, monkeypatch):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    launched: dict[str, object] = {}

    class FakeRMCSetupWindow:
        def __init__(self, initial_project_dir=None):
            launched["project_dir"] = initial_project_dir
            launched["instance"] = self

        def show(self):
            launched["shown"] = True

        def raise_(self):
            launched["raised"] = True

    monkeypatch.setattr(
        "saxshell.fullrmc.ui.main_window.RMCSetupMainWindow",
        FakeRMCSetupWindow,
    )

    window._open_fullrmc_tool()

    assert launched["project_dir"] == window.current_settings.project_dir
    assert launched["shown"] is True
    assert launched["raised"] is True
    assert launched["instance"] in window._child_tool_windows


def test_cluster_dynamics_tool_uses_active_project_dir(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    frames_dir = tmp_path / "splitxyz_f0fs"
    frames_dir.mkdir()
    energy_file = tmp_path / "traj.ener"
    energy_file.write_text(
        "# step time kinetic temperature potential\n"
        "1 0.0 1.0 300.0 -10.0\n",
        encoding="utf-8",
    )
    window.current_settings.frames_dir = str(frames_dir)
    window.current_settings.energy_file = str(energy_file)
    window.project_manager.save_project(window.current_settings)
    launched: dict[str, object] = {}

    class FakeClusterDynamicsWindow:
        def __init__(
            self,
            initial_frames_dir=None,
            initial_energy_file=None,
            initial_project_dir=None,
        ):
            launched["frames_dir"] = initial_frames_dir
            launched["energy_file"] = initial_energy_file
            launched["project_dir"] = initial_project_dir
            launched["instance"] = self

        def show(self):
            launched["shown"] = True

        def raise_(self):
            launched["raised"] = True

    monkeypatch.setattr(
        "saxshell.clusterdynamics.ui.main_window.ClusterDynamicsMainWindow",
        FakeClusterDynamicsWindow,
    )

    window._open_clusterdynamics_tool()

    assert launched["frames_dir"] == frames_dir.resolve()
    assert launched["energy_file"] == energy_file.resolve()
    assert (
        launched["project_dir"]
        == Path(window.current_settings.project_dir).resolve()
    )
    assert launched["shown"] is True
    assert launched["raised"] is True
    assert launched["instance"] in window._child_tool_windows
    window.close()


def test_cluster_dynamics_ml_tool_uses_active_project_dir(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    frames_dir = tmp_path / "splitxyz_f0fs"
    frames_dir.mkdir()
    clusters_dir = tmp_path / "clusters"
    clusters_dir.mkdir()
    energy_file = tmp_path / "traj.ener"
    energy_file.write_text(
        "# step time kinetic temperature potential\n"
        "1 0.0 1.0 300.0 -10.0\n",
        encoding="utf-8",
    )
    window.current_settings.frames_dir = str(frames_dir)
    window.current_settings.clusters_dir = str(clusters_dir)
    window.current_settings.energy_file = str(energy_file)
    window.project_manager.save_project(window.current_settings)
    launched: dict[str, object] = {}

    class FakeClusterDynamicsMLWindow:
        def __init__(
            self,
            initial_frames_dir=None,
            initial_energy_file=None,
            initial_project_dir=None,
            initial_clusters_dir=None,
            initial_experimental_data_file=None,
        ):
            launched["frames_dir"] = initial_frames_dir
            launched["energy_file"] = initial_energy_file
            launched["project_dir"] = initial_project_dir
            launched["clusters_dir"] = initial_clusters_dir
            launched["experimental_data_file"] = initial_experimental_data_file
            launched["instance"] = self

        def show(self):
            launched["shown"] = True

        def raise_(self):
            launched["raised"] = True

    monkeypatch.setattr(
        "saxshell.clusterdynamicsml.ui.main_window.ClusterDynamicsMLMainWindow",
        FakeClusterDynamicsMLWindow,
    )

    window._open_clusterdynamicsml_tool()

    assert launched["frames_dir"] == frames_dir.resolve()
    assert launched["energy_file"] == energy_file.resolve()
    assert launched["clusters_dir"] == clusters_dir.resolve()
    assert (
        launched["experimental_data_file"]
        == window.current_settings.resolved_experimental_data_path
    )
    assert (
        launched["project_dir"]
        == Path(window.current_settings.project_dir).resolve()
    )
    assert launched["shown"] is True
    assert launched["raised"] is True
    assert launched["instance"] in window._child_tool_windows
    window.close()


def test_project_setup_predict_structures_button_opens_cluster_dynamics_ml_tool(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    launched: dict[str, object] = {}

    class FakeClusterDynamicsMLWindow:
        def __init__(
            self,
            initial_frames_dir=None,
            initial_energy_file=None,
            initial_project_dir=None,
            initial_clusters_dir=None,
            initial_experimental_data_file=None,
        ):
            launched["project_dir"] = initial_project_dir
            launched["instance"] = self

        def show(self):
            launched["shown"] = True

        def raise_(self):
            launched["raised"] = True

    monkeypatch.setattr(
        "saxshell.clusterdynamicsml.ui.main_window.ClusterDynamicsMLMainWindow",
        FakeClusterDynamicsMLWindow,
    )

    window.project_setup_tab.predict_structures_button.click()

    assert (
        launched["project_dir"]
        == Path(window.current_settings.project_dir).resolve()
    )
    assert launched["shown"] is True
    assert launched["raised"] is True
    assert launched["instance"] in window._child_tool_windows
    window.close()


def test_contrast_mode_scaffold_window_populates_launch_context(
    qapp, tmp_path
):
    del qapp
    project_dir = tmp_path / "contrast_project"
    project_dir.mkdir()
    clusters_dir = tmp_path / "clusters"
    clusters_dir.mkdir()
    experimental_data_file = tmp_path / "exp_data.dat"
    experimental_data_file.write_text(
        "0.05 1.0\n0.08 0.8\n0.12 0.5\n",
        encoding="utf-8",
    )
    cluster_bin_dir = clusters_dir / "A2"
    cluster_bin_dir.mkdir()
    (cluster_bin_dir / "frame_0001.xyz").write_text(
        "2\nframe\nA 0.0 0.0 0.0\nA 1.0 0.0 0.0\n",
        encoding="utf-8",
    )

    window = ContrastModeMainWindow(
        initial_project_dir=project_dir.resolve(),
        initial_clusters_dir=clusters_dir.resolve(),
        initial_experimental_data_file=experimental_data_file.resolve(),
        initial_q_min=0.05,
        initial_q_max=0.8,
        initial_template_name="template_demo",
    )
    _flush_contrast_window_launch_preview(window)

    assert (
        window.windowTitle() == "SAXSShell (Contrast Debye Workflow) (Preview)"
    )
    assert window.mode_edit.text() == "Contrast (Debye)"
    assert window.project_dir_edit.text() == str(project_dir.resolve())
    assert window.clusters_dir_edit.text() == str(clusters_dir.resolve())
    assert window.experimental_data_edit.text() == str(
        experimental_data_file.resolve()
    )
    assert window.q_range_edit.text() == "0.05 to 0.8"
    assert window.q_min_spin.value() == pytest.approx(0.05)
    assert window.q_max_spin.value() == pytest.approx(0.8)
    assert window.template_edit.text() == "template_demo"
    assert (
        "Contrast-mode SAXS workspace"
        in window.workflow_summary_box.toPlainText()
    )
    assert window.output_path_edit.text() == str(
        (project_dir / "contrast_workflow").resolve()
    )
    assert window.representatives_output_edit.text() == str(
        (project_dir / "contrast_workflow" / "representatives").resolve()
    )
    assert window.screening_output_edit.text() == str(
        (
            project_dir / "contrast_workflow" / "representatives" / "screening"
        ).resolve()
    )
    assert window.summary_output_edit.text() == str(
        (
            project_dir
            / "contrast_workflow"
            / "representatives"
            / "selection_summary.json"
        ).resolve()
    )
    assert window.recognized_clusters_table.columnCount() == 7
    assert window.recognized_clusters_table.rowCount() == 1
    assert window.representative_table.columnCount() == 9
    assert window.representative_table.rowCount() == 0
    assert (
        window.add_representative_button.text()
        == "Add Representative (Custom)…"
    )
    assert window.solvent_preset_combo.findData("DMF") >= 0
    assert window.solvent_preset_combo.findData("DMSO") >= 0
    assert window.main_splitter.count() == 2
    assert window.left_scroll_area.widget() is not None
    assert window.right_scroll_area.widget() is not None
    assert window.right_splitter.count() == 2
    assert window.trace_figure.axes
    assert (
        window.trace_figure.axes[0].lines[0].get_label() == "Experimental data"
    )
    window.close()


def _flush_contrast_window_launch_preview(
    window: ContrastModeMainWindow,
) -> None:
    qapp = QApplication.instance()
    assert qapp is not None
    deadline = time.monotonic() + 2.0
    while time.monotonic() < deadline:
        qapp.processEvents()
        if not window._preview_refresh_timer.isActive():
            qapp.processEvents()
            return
        QTest.qWait(10)
    raise AssertionError(
        "Contrast window launch preview did not finish in time."
    )


def test_contrast_mode_window_defers_preview_rebuild_until_event_loop(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    clusters_dir = tmp_path / "clusters"
    clusters_dir.mkdir()
    cluster_bin_dir = clusters_dir / "A2"
    cluster_bin_dir.mkdir()
    (cluster_bin_dir / "frame_0001.xyz").write_text(
        "2\nframe\nA 0.0 0.0 0.0\nA 1.0 0.0 0.0\n",
        encoding="utf-8",
    )

    rebuild_calls: list[dict[str, object]] = []

    def fake_rebuild_preview(self, *, log_to_console, status_message):
        rebuild_calls.append(
            {
                "log_to_console": log_to_console,
                "status_message": status_message,
            }
        )

    monkeypatch.setattr(
        ContrastModeMainWindow,
        "_rebuild_preview",
        fake_rebuild_preview,
    )

    window = ContrastModeMainWindow(
        initial_clusters_dir=clusters_dir.resolve(),
    )

    assert rebuild_calls == []
    assert window._preview_refresh_timer.isActive() is True
    window.close()


def test_contrast_mode_workspace_marks_selected_q_range_on_experimental_preview(
    qapp,
    tmp_path,
):
    del qapp
    project_dir = tmp_path / "contrast_project"
    project_dir.mkdir()
    clusters_dir = tmp_path / "clusters"
    clusters_dir.mkdir()
    experimental_data_file = tmp_path / "exp_data.dat"
    experimental_data_file.write_text(
        "0.05 1.0\n0.08 0.8\n0.12 0.5\n",
        encoding="utf-8",
    )

    window = ContrastModeMainWindow(
        initial_project_dir=project_dir.resolve(),
        initial_clusters_dir=clusters_dir.resolve(),
        initial_experimental_data_file=experimental_data_file.resolve(),
        initial_q_min=0.06,
        initial_q_max=0.10,
    )
    _flush_contrast_window_launch_preview(window)

    axis = window.trace_figure.axes[0]
    labels = [line.get_label() for line in axis.lines]
    assert labels == ["Experimental data", "Selected q-range"]
    assert axis.lines[0].get_alpha() == pytest.approx(0.35)
    np.testing.assert_allclose(
        np.asarray(axis.lines[1].get_xdata(), dtype=float),
        np.asarray([0.08], dtype=float),
    )
    window.close()


def test_contrast_mode_tool_uses_active_project_context(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    clusters_dir = tmp_path / "contrast_clusters"
    clusters_dir.mkdir()
    experimental_data_file = tmp_path / "contrast_exp.dat"
    experimental_data_file.write_text("0.06 2.0\n", encoding="utf-8")
    window.current_settings.clusters_dir = str(clusters_dir.resolve())
    window.current_settings.experimental_data_path = str(
        experimental_data_file.resolve()
    )
    window.current_settings.q_min = 0.06
    window.current_settings.q_max = 0.75
    window.current_settings.selected_model_template = POLY_LMA_HS_TEMPLATE
    window.project_manager.save_project(window.current_settings)
    launched: dict[str, object] = {}

    class FakeContrastModeWindow(QWidget):
        def __init__(self):
            super().__init__()
            launched["instance"] = self

        def raise_(self):
            launched["raised"] = True

        def activateWindow(self):
            launched["activated"] = True

    def fake_launch_contrast_mode_ui(**kwargs):
        launched.update(kwargs)
        return FakeContrastModeWindow()

    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.launch_contrast_mode_ui",
        fake_launch_contrast_mode_ui,
    )

    window._open_contrast_mode_tool()
    saved_settings = window.project_manager.load_project(project_dir)

    assert launched["initial_project_dir"] == Path(project_dir).resolve()
    assert launched["initial_clusters_dir"] == clusters_dir.resolve()
    assert (
        launched["initial_experimental_data_file"]
        == saved_settings.resolved_experimental_data_path
    )
    assert launched["initial_q_min"] == pytest.approx(0.06)
    assert launched["initial_q_max"] == pytest.approx(0.75)
    assert launched["initial_template_name"] == POLY_LMA_HS_TEMPLATE
    assert launched["preview_mode"] is True
    assert launched["instance"] in window._child_tool_windows
    assert window._contrast_mode_tool_window is launched["instance"]
    window.close()


def test_contrast_mode_workspace_loads_predicted_structure_bins_from_project_settings(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    manager = SAXSProjectManager()
    project_dir = tmp_path / "contrast_predicted_project"
    settings = manager.create_project(project_dir)
    clusters_dir = project_dir / "clusters"
    structure_dir = clusters_dir / "PbI2"
    structure_dir.mkdir(parents=True, exist_ok=True)
    (structure_dir / "frame_0001.xyz").write_text(
        "\n".join(
            [
                "3",
                "frame_0001",
                "Pb 0.0 0.0 0.0",
                "I 2.8 0.0 0.0",
                "I 0.0 2.8 0.0",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    predicted_dir = tmp_path / "predicted_structures"
    predicted_dir.mkdir(parents=True, exist_ok=True)
    predicted_structure = predicted_dir / "04_rank01_PbI4.xyz"
    predicted_structure.write_text(
        "\n".join(
            [
                "5",
                "predicted_rank01",
                "Pb 0.0 0.0 0.0",
                "I 2.6 0.0 0.0",
                "I -2.6 0.0 0.0",
                "I 0.0 2.6 0.0",
                "I 0.0 0.0 2.6",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    settings.model_only_mode = True
    settings.clusters_dir = str(clusters_dir.resolve())
    settings.use_predicted_structure_weights = True
    settings.use_experimental_grid = False
    settings.q_min = 0.05
    settings.q_max = 0.30
    settings.q_points = 8
    settings.selected_model_template = POLY_LMA_HS_TEMPLATE
    manager.save_project(settings)

    def fake_predicted_payload(self, settings, *, cluster_inventory):
        del self, settings
        observed_weight = 0.75 / max(len(cluster_inventory.cluster_bins), 1)
        return (
            SimpleNamespace(dataset_file=project_dir / "mock_predicted.json"),
            {
                (cluster_bin.structure, cluster_bin.motif): observed_weight
                for cluster_bin in cluster_inventory.cluster_bins
            },
            [
                {
                    "prediction": SimpleNamespace(label="PbI4"),
                    "motif": "predicted_rank01",
                    "weight": 0.25,
                    "source_path": predicted_structure,
                }
            ],
            ["Pb", "I"],
        )

    monkeypatch.setattr(
        SAXSProjectManager,
        "_predicted_structure_weight_payload",
        fake_predicted_payload,
    )

    window = ContrastModeMainWindow(
        initial_project_dir=project_dir.resolve(),
        initial_clusters_dir=clusters_dir.resolve(),
        initial_q_min=settings.q_min,
        initial_q_max=settings.q_max,
        initial_template_name=settings.selected_model_template,
    )
    _flush_contrast_window_launch_preview(window)

    loaded_components = {
        window.recognized_clusters_table.item(row, 0).text()
        for row in range(window.recognized_clusters_table.rowCount())
    }
    predicted_rows = [
        row
        for row in range(window.recognized_clusters_table.rowCount())
        if window.recognized_clusters_table.item(row, 1).text() == "PbI4"
    ]

    assert window.recognized_clusters_table.rowCount() == 2
    assert loaded_components == {"PbI2", "PbI4 / predicted_rank01"}
    assert len(predicted_rows) == 1
    assert (
        window.recognized_clusters_table.item(predicted_rows[0], 2).text()
        == "Predicted structure"
    )
    assert (
        window.recognized_clusters_table.item(predicted_rows[0], 6).text()
        == predicted_structure.name
    )
    assert any(
        cluster_bin.structure == "PbI4"
        and cluster_bin.motif == "predicted_rank01"
        and cluster_bin.representative == predicted_structure.name
        for cluster_bin in window._recognized_cluster_bins
    )
    assert "predicted structure bin" in (
        window.cluster_table_status_label.text().lower()
    )
    window.close()


def test_project_setup_view_contrast_button_opens_active_distribution_context(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _settings, artifact_paths, _build_result = (
        _build_saved_contrast_distribution_project(tmp_path)
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)
    launched: dict[str, object] = {}

    class FakeContrastModeWindow(QWidget):
        def __init__(self):
            super().__init__()
            launched["instance"] = self

        def raise_(self):
            launched["raised"] = True

        def activateWindow(self):
            launched["activated"] = True

    def fake_launch_contrast_mode_ui(**kwargs):
        launched.update(kwargs)
        return FakeContrastModeWindow()

    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.launch_contrast_mode_ui",
        fake_launch_contrast_mode_ui,
    )

    assert (
        window.project_setup_tab.view_contrast_distribution_button.isEnabled()
    )
    window.project_setup_tab.view_contrast_distribution_button.click()

    assert (
        launched["initial_distribution_id"] == artifact_paths.distribution_id
    )
    assert launched["initial_distribution_root_dir"] == artifact_paths.root_dir
    assert (
        launched["initial_contrast_artifact_dir"]
        == artifact_paths.contrast_dir
    )
    assert launched["instance"] in window._child_tool_windows
    window.close()


def test_contrast_mode_tool_open_starts_clean_even_with_saved_distribution(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _settings, _artifact_paths, _build_result = (
        _build_saved_contrast_distribution_project(tmp_path)
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)
    launched: dict[str, object] = {}

    class FakeContrastModeWindow(QWidget):
        def __init__(self):
            super().__init__()
            launched["instance"] = self

        def raise_(self):
            launched["raised"] = True

        def activateWindow(self):
            launched["activated"] = True

    def fake_launch_contrast_mode_ui(**kwargs):
        launched.update(kwargs)
        return FakeContrastModeWindow()

    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.launch_contrast_mode_ui",
        fake_launch_contrast_mode_ui,
    )

    window._open_contrast_mode_tool()

    assert launched["initial_distribution_id"] is None
    assert launched["initial_distribution_root_dir"] is None
    assert launched["initial_contrast_artifact_dir"] is None
    assert launched["preview_mode"] is True
    window.close()


def test_electron_density_mapping_tool_uses_active_structure_folder(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    pdb_frames_dir = tmp_path / "pdb_frames"
    pdb_frames_dir.mkdir()
    xyz_frames_dir = tmp_path / "xyz_frames"
    xyz_frames_dir.mkdir()
    window.current_settings.pdb_frames_dir = str(pdb_frames_dir.resolve())
    window.current_settings.frames_dir = str(xyz_frames_dir.resolve())
    window.project_manager.save_project(window.current_settings)
    launched: dict[str, object] = {}

    class FakeElectronDensityWindow(QWidget):
        def __init__(self):
            super().__init__()
            launched["instance"] = self

    def fake_launch_electron_density_mapping_ui(**kwargs):
        launched.update(kwargs)
        return FakeElectronDensityWindow()

    monkeypatch.setattr(
        "saxshell.saxs.electron_density_mapping.ui.main_window.launch_electron_density_mapping_ui",
        fake_launch_electron_density_mapping_ui,
    )

    window._open_electron_density_mapping_tool()

    assert launched["initial_project_dir"] == Path(project_dir).resolve()
    assert launched["initial_input_path"] == pdb_frames_dir.resolve()
    assert launched["preview_mode"] is True
    assert launched["instance"] in window._child_tool_windows
    window.close()


def test_open_contrast_mode_tool_does_not_auto_start_build(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    launched: dict[str, object] = {"start_calls": 0}

    class FakeContrastModeWindow(QWidget):
        def __init__(self):
            super().__init__()
            launched["instance"] = self

        def start_contrast_component_build(self):
            launched["start_calls"] += 1

        def raise_(self):
            launched["raised"] = True

        def activateWindow(self):
            launched["activated"] = True

    def fake_launch_contrast_mode_ui(**kwargs):
        launched.update(kwargs)
        return FakeContrastModeWindow()

    monkeypatch.setattr(
        window,
        "_confirm_default_q_range_for_component_build",
        lambda: True,
    )
    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.launch_contrast_mode_ui",
        fake_launch_contrast_mode_ui,
    )

    window.project_setup_tab.set_component_build_mode(
        COMPONENT_BUILD_MODE_CONTRAST
    )
    window.build_project_components()

    assert launched["instance"] in window._child_tool_windows
    assert launched["start_calls"] == 0
    window.close()


def test_contrast_mode_window_can_close_while_workflow_thread_is_running(
    qapp,
    monkeypatch,
):
    del qapp
    window = ContrastModeMainWindow()

    class FakeRunningThread:
        def isRunning(self) -> bool:
            return True

    warnings: list[tuple[object, ...]] = []
    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.QMessageBox.warning",
        lambda *args, **kwargs: warnings.append(args),
    )
    window._workflow_thread = FakeRunningThread()

    assert window.close() is True
    assert warnings == []


def test_main_window_loads_contrast_distribution_after_tool_build_signal(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _settings, artifact_paths, _build_result = (
        _build_saved_contrast_distribution_project(tmp_path)
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)
    loaded_distribution_ids: list[str] = []

    monkeypatch.setattr(
        window,
        "_load_saved_distribution",
        lambda distribution_id: loaded_distribution_ids.append(
            distribution_id
        ),
    )

    window._on_contrast_components_built(
        {
            "project_dir": str(project_dir.resolve()),
            "distribution_id": artifact_paths.distribution_id,
            "distribution_dir": str(artifact_paths.root_dir),
            "component_dir": str(artifact_paths.component_dir),
            "component_map_path": str(artifact_paths.component_map_file),
        }
    )

    assert loaded_distribution_ids == [artifact_paths.distribution_id]
    assert "built and loaded" in window.statusBar().currentMessage().lower()
    window.close()


def test_contrast_mode_workspace_tracks_manual_representative_selection(
    qapp, tmp_path, monkeypatch
):
    del qapp
    representative_path = tmp_path / "representative_A2.pdb"
    representative_path.write_text(
        "ATOM      1  PB  RES A   1       0.000   0.000   0.000  1.00  0.00          Pb\n"
        "END\n",
        encoding="utf-8",
    )
    window = ContrastModeMainWindow()

    monkeypatch.setattr(
        QFileDialog,
        "getOpenFileNames",
        lambda *args, **kwargs: (
            [str(representative_path)],
            "Structure files (*.pdb *.xyz)",
        ),
    )
    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.QMessageBox.question",
        lambda *args, **kwargs: QMessageBox.StandardButton.Yes,
    )

    window.add_representative_button.click()

    assert window.representative_table.rowCount() == 1
    assert (
        window.representative_table.item(0, 2).text()
        == representative_path.name
    )
    assert window.representative_table.item(0, 3).text().startswith("#")
    assert window.structure_viewer.current_preview is not None
    assert (
        window.structure_viewer.current_preview.file_path
        == representative_path.resolve()
    )
    assert window.structure_viewer.current_preview.atom_count == 1
    assert "Previewing" in window.visualizer_status_label.text()
    assert (
        str(representative_path.resolve())
        in window.visualizer_details_box.toPlainText()
    )

    window.remove_representative_button.click()

    assert window.representative_table.rowCount() == 0
    assert window.structure_viewer.current_preview is None
    assert (
        "Selected representative: none"
        in window.visualizer_details_box.toPlainText()
    )
    window.close()


def test_contrast_mode_workspace_can_choose_representative_trace_color(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    representative_path = tmp_path / "representative_A2.pdb"
    representative_path.write_text(
        "ATOM      1  PB  RES A   1       0.000   0.000   0.000  1.00  0.00          Pb\n"
        "END\n",
        encoding="utf-8",
    )
    window = ContrastModeMainWindow()

    monkeypatch.setattr(
        QFileDialog,
        "getOpenFileNames",
        lambda *args, **kwargs: (
            [str(representative_path)],
            "Structure files (*.pdb *.xyz)",
        ),
    )
    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.QMessageBox.question",
        lambda *args, **kwargs: QMessageBox.StandardButton.Yes,
    )
    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.QColorDialog.getColor",
        lambda *args, **kwargs: QColor("#ff8800"),
    )

    window.add_representative_button.click()
    display_label = window._representative_row_metadata(0)["display_label"]
    window._generated_trace_profiles = [
        (
            str(display_label),
            np.asarray([0.08, 0.12], dtype=float),
            np.asarray([4.0, 2.0], dtype=float),
        )
    ]
    window._generated_traces_visible = True
    window._redraw_trace_plot()

    window._on_representative_table_cell_clicked(0, 3)

    metadata = window._representative_row_metadata(0)
    assert metadata["trace_color"] == "#ff8800"
    assert metadata["custom_trace_color"] is True
    assert window.representative_table.item(0, 3).text() == "#FF8800"
    assert (
        window.representative_table.item(0, 3).background().color().name()
        == "#ff8800"
    )
    assert (
        to_hex(window.trace_figure.axes[0].lines[0].get_color()).lower()
        == "#ff8800"
    )
    window.close()


def test_contrast_mode_workspace_warns_before_custom_representative_edits(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    representative_path = tmp_path / "representative_A2.pdb"
    representative_path.write_text(
        "ATOM      1  PB  RES A   1       0.000   0.000   0.000  1.00  0.00          Pb\n"
        "END\n",
        encoding="utf-8",
    )
    window = ContrastModeMainWindow()

    monkeypatch.setattr(
        QFileDialog,
        "getOpenFileNames",
        lambda *args, **kwargs: (
            [str(representative_path)],
            "Structure files (*.pdb *.xyz)",
        ),
    )
    prompts: list[str] = []
    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.QMessageBox.question",
        lambda *args, **kwargs: (
            prompts.append(
                str(args[2] if len(args) > 2 else kwargs.get("text", ""))
            )
            or QMessageBox.StandardButton.No
        ),
    )

    window.add_representative_button.click()

    assert window.representative_table.rowCount() == 0
    assert prompts
    assert "prior weights" in prompts[0].lower()

    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.QMessageBox.question",
        lambda *args, **kwargs: QMessageBox.StandardButton.Yes,
    )

    window.add_representative_button.click()
    assert window.representative_table.rowCount() == 1

    window.remove_representative_button.click()
    assert window.representative_table.rowCount() == 0
    window.close()


def test_contrast_mode_workspace_can_analyze_representative_structures(
    qapp, tmp_path
):
    del qapp
    project_dir = tmp_path / "contrast_project"
    project_dir.mkdir()
    clusters_dir = tmp_path / "clusters"
    clusters_dir.mkdir()
    cluster_bin_dir = clusters_dir / "PbI2"
    cluster_bin_dir.mkdir()
    (cluster_bin_dir / "frame_0001.xyz").write_text(
        "\n".join(
            [
                "5",
                "frame_0001",
                "Pb 0.0 0.0 0.0",
                "I 3.0 0.0 0.0",
                "I 0.7 1.8 0.0",
                "O 0.0 0.0 2.6",
                "O 0.0 0.0 5.1",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    (cluster_bin_dir / "frame_0002.xyz").write_text(
        "\n".join(
            [
                "5",
                "frame_0002",
                "Pb 0.0 0.0 0.0",
                "I 2.1 0.0 0.0",
                "I 0.0 2.1 0.0",
                "O 0.0 0.0 2.8",
                "O 0.0 0.0 5.2",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    (cluster_bin_dir / "frame_0003.xyz").write_text(
        "\n".join(
            [
                "4",
                "frame_0003",
                "Pb 0.0 0.0 0.0",
                "I 1.5 0.0 0.0",
                "I -1.4 1.6 0.0",
                "O 0.0 0.0 3.1",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    window = ContrastModeMainWindow(
        initial_project_dir=project_dir.resolve(),
        initial_clusters_dir=clusters_dir.resolve(),
    )
    _flush_contrast_window_launch_preview(window)

    assert window.representative_table.rowCount() == 0
    window.analyze_representatives_button.click()
    deadline = time.monotonic() + 10.0
    while time.monotonic() < deadline:
        QApplication.processEvents()
        if (
            window.representative_table.rowCount() == 1
            and window.analyze_representatives_button.isEnabled()
        ):
            break
        QTest.qWait(25)
    else:
        raise AssertionError(
            "Representative analysis did not finish within 10 seconds."
        )

    output_dir = (
        project_dir / "contrast_workflow" / "representatives"
    ).resolve()
    assert window.representative_table.rowCount() == 1
    assert window.representative_table.item(0, 1).text() == "Auto"
    assert (
        window.representative_table.item(0, 2).text() == "PbI2__frame_0002.xyz"
    )
    assert window.representative_table.item(0, 3).text().startswith("#")
    assert window.representative_table.item(0, 4).text() == "Direct 1; Outer 1"
    assert "Score" in window.representative_table.item(0, 8).text()
    assert "PbI2.json" in window.representative_table.item(0, 8).text()
    assert window.output_path_edit.text() == str(
        (project_dir / "contrast_workflow").resolve()
    )
    assert window.representatives_output_edit.text() == str(output_dir)
    assert window.screening_output_edit.text() == str(output_dir / "screening")
    assert window.summary_output_edit.text() == str(
        output_dir / "selection_summary.json"
    )
    assert output_dir.joinpath("selection_summary.json").is_file()
    assert output_dir.joinpath("selection_summary.tsv").is_file()
    assert output_dir.joinpath("screening", "PbI2.json").is_file()
    assert (
        "Contrast representative selection complete"
        in window.console_box.toPlainText()
    )
    assert (
        "Selected frame_0002.xyz for PbI2" in window.console_box.toPlainText()
    )
    assert (
        "representative analysis complete"
        in window.workflow_progress_label.text().lower()
    )
    assert "2/5 workflow stages" in window.workflow_progress_label.text()
    assert "1 bin(s) selected" in window.workflow_progress_label.text()
    assert window.workflow_progress_bar.value() == 2
    assert window.workflow_progress_bar.maximum() == 5
    window.close()


def test_contrast_mode_workspace_sampler_settings_are_collapsible_and_propagated(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir = tmp_path / "contrast_project"
    project_dir.mkdir()
    clusters_dir = tmp_path / "clusters"
    clusters_dir.mkdir()
    cluster_bin_dir = clusters_dir / "PbI2"
    cluster_bin_dir.mkdir()
    (cluster_bin_dir / "frame_0001.xyz").write_text(
        "\n".join(
            [
                "4",
                "frame_0001",
                "Pb 0.0 0.0 0.0",
                "I 2.1 0.0 0.0",
                "I 0.0 2.1 0.0",
                "O 0.0 0.0 2.8",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    window = ContrastModeMainWindow(
        initial_project_dir=project_dir.resolve(),
        initial_clusters_dir=clusters_dir.resolve(),
    )

    assert window.sampler_settings_widget.isHidden()
    window.sampler_settings_toggle_button.click()
    assert not window.sampler_settings_widget.isHidden()

    window.sampler_enabled_checkbox.setChecked(True)
    window.sampler_full_scan_threshold_spin.setValue(9)
    window.sampler_distribution_samples_spin.setValue(21)
    window.sampler_minimum_samples_spin.setValue(3)
    window.sampler_max_samples_spin.setValue(11)
    window.sampler_batch_size_spin.setValue(2)
    window.sampler_stratify_checkbox.setChecked(False)
    window.sampler_seed_spin.setValue(77)
    window.sampler_patience_spin.setValue(5)
    window.sampler_tolerance_spin.setValue(0.0125)

    sampler_settings = window._representative_sampler_settings()
    assert (
        sampler_settings
        == ContrastRepresentativeSamplerSettings.from_values(
            enabled=True,
            full_scan_threshold=9,
            target_distribution_samples=21,
            minimum_candidate_samples=3,
            max_candidate_samples=11,
            candidate_batch_size=2,
            random_seed=77,
            convergence_patience=5,
            improvement_tolerance=0.0125,
            stratify_sampling=False,
        )
    )

    captured: dict[str, object] = {}

    def fake_analyze(*args, sampler_settings=None, **kwargs):
        del args, kwargs
        captured["sampler_settings"] = sampler_settings
        return SimpleNamespace()

    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.analyze_contrast_representatives",
        fake_analyze,
    )

    settings = window._contrast_project_settings()
    window._run_representative_analysis_task(
        settings,
        sampler_settings,
        progress_callback=lambda *args: None,
        log_callback=lambda *args: None,
    )

    assert captured["sampler_settings"] == sampler_settings
    window.sampler_settings_toggle_button.click()
    assert window.sampler_settings_widget.isHidden()
    window.close()


def test_contrast_mode_workspace_loads_builtin_and_custom_solvent_presets(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    preset_path = tmp_path / "contrast_solvents.json"
    monkeypatch.setenv(
        "SAXSHELL_CONTRAST_SOLVENTS_PATH",
        str(preset_path),
    )

    window = ContrastModeMainWindow()

    dmf_index = window.solvent_preset_combo.findData("DMF")
    assert dmf_index >= 0
    window.solvent_preset_combo.setCurrentIndex(dmf_index)
    assert window.solvent_formula_edit.text() == "C3H7NO"
    assert window.solvent_density_spin.value() == pytest.approx(0.944)
    vacuum_index = window.solvent_preset_combo.findData("Vacuum")
    assert vacuum_index >= 0
    window.solvent_preset_combo.setCurrentIndex(vacuum_index)
    assert window.solvent_formula_edit.text() == "Vacuum"
    assert window.solvent_density_spin.value() == pytest.approx(0.0)

    reference_index = window.solvent_method_combo.findData(
        "reference_structure"
    )
    assert reference_index >= 0
    window.solvent_method_combo.setCurrentIndex(reference_index)
    assert window.reference_solvent_file_edit.isEnabled()
    assert not window.solvent_formula_edit.isEnabled()
    assert not window.direct_density_spin.isEnabled()

    neat_index = window.solvent_method_combo.findData("neat_solvent_estimate")
    assert neat_index >= 0
    window.solvent_method_combo.setCurrentIndex(neat_index)
    assert window.solvent_formula_edit.isEnabled()
    assert not window.reference_solvent_file_edit.isEnabled()
    assert not window.direct_density_spin.isEnabled()
    window.solvent_formula_edit.setText("C4H10O")
    window.solvent_density_spin.setValue(0.88)
    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.QInputDialog.getText",
        lambda *args, **kwargs: ("My Solvent", True),
    )
    monkeypatch.setattr(
        "saxshell.saxs.contrast.ui.main_window.QMessageBox.question",
        lambda *args, **kwargs: QMessageBox.StandardButton.Yes,
    )

    window.save_custom_solvent_button.click()

    assert preset_path.is_file()
    assert window.solvent_preset_combo.findData("My Solvent") >= 0
    direct_index = window.solvent_method_combo.findData(
        CONTRAST_SOLVENT_METHOD_DIRECT
    )
    assert direct_index >= 0
    window.solvent_method_combo.setCurrentIndex(direct_index)
    window.direct_density_spin.setValue(0.287)
    direct_settings = window._current_density_settings()
    assert direct_settings.solvent.method == CONTRAST_SOLVENT_METHOD_DIRECT
    assert (
        direct_settings.solvent.direct_electron_density_e_per_a3
        == pytest.approx(0.287)
    )
    assert window.direct_density_spin.isEnabled()
    assert not window.solvent_formula_edit.isEnabled()
    assert not window.reference_solvent_file_edit.isEnabled()
    window.close()

    reloaded = ContrastModeMainWindow()
    my_index = reloaded.solvent_preset_combo.findData("My Solvent")
    assert my_index >= 0
    reloaded.solvent_preset_combo.setCurrentIndex(my_index)
    assert reloaded.solvent_formula_edit.text() == "C4H10O"
    assert reloaded.solvent_density_spin.value() == pytest.approx(0.88)

    reloaded.delete_custom_solvent_button.click()
    assert reloaded.solvent_preset_combo.findData("My Solvent") == -1
    reloaded.close()


def test_contrast_mode_workspace_hydrogen_exclusion_feeds_density_settings(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = ContrastModeMainWindow(
        initial_project_dir=Path(project_dir).resolve(),
    )

    assert window.exclude_hydrogen_checkbox.isChecked() is False
    window.exclude_hydrogen_checkbox.setChecked(True)
    density_settings = window._current_density_settings()
    project_settings = window._contrast_project_settings()

    assert density_settings.exclude_elements == ("H",)
    assert "H" in project_settings.exclude_elements
    window.close()


def test_contrast_mode_workspace_viewer_loads_saved_mesh_and_density_metadata(
    qapp, tmp_path
):
    del qapp
    project_dir = tmp_path / "contrast_project"
    project_dir.mkdir()
    clusters_dir = tmp_path / "clusters"
    clusters_dir.mkdir()
    cluster_bin_dir = clusters_dir / "PbI2"
    cluster_bin_dir.mkdir()
    (cluster_bin_dir / "frame_0001.xyz").write_text(
        "\n".join(
            [
                "5",
                "frame_0001",
                "Pb 0.0 0.0 0.0",
                "I 3.0 0.0 0.0",
                "I 0.7 1.8 0.0",
                "O 0.0 0.0 2.6",
                "O 0.0 0.0 5.1",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    (cluster_bin_dir / "frame_0002.xyz").write_text(
        "\n".join(
            [
                "5",
                "frame_0002",
                "Pb 0.0 0.0 0.0",
                "I 2.1 0.0 0.0",
                "I 0.0 2.1 0.0",
                "O 0.0 0.0 2.8",
                "O 0.0 0.0 5.2",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    representative_result = analyze_contrast_representatives(
        project_dir,
        clusters_dir,
    )
    compute_contrast_geometry_and_electron_density(
        representative_result,
        ContrastGeometryDensitySettings(
            solvent=ContrastSolventDensitySettings.from_values(
                method=CONTRAST_SOLVENT_METHOD_NEAT,
                solvent_formula="H2O",
                solvent_density_g_per_ml=1.0,
            )
        ),
    )

    window = ContrastModeMainWindow(
        initial_project_dir=project_dir.resolve(),
        initial_clusters_dir=clusters_dir.resolve(),
    )
    window._representative_analysis_result = representative_result
    window._populate_representative_analysis_result(representative_result)

    assert window.representative_table.rowCount() == 1
    assert window.representative_table.item(0, 3).text().startswith("#")
    assert window.representative_table.item(0, 5).text().endswith("A^3")
    assert window.representative_table.item(0, 6).text().startswith("Δρ ")
    assert window.representative_table.item(0, 7).text() == "Density ready"
    assert window.structure_viewer.current_preview is not None
    assert window.structure_viewer.current_preview.has_mesh
    assert window.structure_viewer.current_preview.mesh_json_path is not None
    assert (
        window.structure_viewer.current_preview.density_json_path is not None
    )
    assert window.structure_viewer.minimumHeight() >= 430
    assert window.structure_viewer.canvas.minimumHeight() >= 420
    assert window.visualizer_details_box.maximumHeight() <= 112
    assert (
        "Contrast density term:" in window.visualizer_details_box.toPlainText()
    )

    axis = window.structure_viewer.figure.axes[0]
    initial_azimuth = float(axis.azim)
    initial_x_span = float(axis.get_xlim3d()[1] - axis.get_xlim3d()[0])
    window.rotate_left_button.click()
    assert float(axis.azim) != pytest.approx(initial_azimuth)
    window.zoom_in_button.click()
    assert float(axis.get_xlim3d()[1] - axis.get_xlim3d()[0]) < initial_x_span
    assert axis.get_legend() is not None
    assert axis.get_legend()._loc == 1
    assert any(
        collection.__class__.__name__ == "Poly3DCollection"
        for collection in axis.collections
    )
    window.close()


def test_contrast_mode_workspace_recomputes_displayed_trace_contrast_after_density_update(
    qapp,
    tmp_path,
):
    del qapp
    project_dir = tmp_path / "contrast_project"
    project_dir.mkdir()
    clusters_dir = tmp_path / "clusters"
    clusters_dir.mkdir()
    cluster_bin_dir = clusters_dir / "PbI2"
    cluster_bin_dir.mkdir()
    (cluster_bin_dir / "frame_0001.xyz").write_text(
        "\n".join(
            [
                "5",
                "frame_0001",
                "Pb 0.0 0.0 0.0",
                "I 3.0 0.0 0.0",
                "I 0.7 1.8 0.0",
                "O 0.0 0.0 2.6",
                "O 0.0 0.0 5.1",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    (cluster_bin_dir / "frame_0002.xyz").write_text(
        "\n".join(
            [
                "5",
                "frame_0002",
                "Pb 0.0 0.0 0.0",
                "I 2.1 0.0 0.0",
                "I 0.0 2.1 0.0",
                "O 0.0 0.0 2.8",
                "O 0.0 0.0 5.2",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    representative_result = analyze_contrast_representatives(
        project_dir,
        clusters_dir,
    )
    density_result_a = compute_contrast_geometry_and_electron_density(
        representative_result,
        ContrastGeometryDensitySettings(
            solvent=ContrastSolventDensitySettings.from_values(
                method=CONTRAST_SOLVENT_METHOD_NEAT,
                solvent_formula="H2O",
                solvent_density_g_per_ml=1.0,
            )
        ),
    )
    q_values = np.linspace(0.05, 0.30, 8)
    build_result_a = build_contrast_component_profiles(
        representative_result,
        density_result_a,
        q_values=q_values,
        output_dir=project_dir / "contrast_workflow" / "scattering_components",
        metadata_dir=project_dir / "contrast_workflow" / "debye",
        component_map_path=project_dir / "md_saxs_map.json",
    )

    window = ContrastModeMainWindow(
        initial_project_dir=project_dir.resolve(),
        initial_clusters_dir=clusters_dir.resolve(),
    )
    window._representative_analysis_result = representative_result
    window._populate_representative_analysis_result(representative_result)
    window._density_result = density_result_a
    window._load_generated_trace_profiles_from_trace_payloads(
        [
            trace_result.to_dict()
            for trace_result in build_result_a.trace_results
        ],
        default_component_dir=build_result_a.output_dir,
    )
    window._apply_saved_trace_metadata(
        [
            trace_result.to_dict()
            for trace_result in build_result_a.trace_results
        ]
    )

    first_status = window.representative_table.item(0, 7).text()
    first_scale = float(
        window._representative_row_metadata(0)["contrast_scale_factor"]
    )
    assert first_status.startswith("Trace ready (scale ")
    assert window._generated_trace_profiles

    density_result_b = compute_contrast_geometry_and_electron_density(
        representative_result,
        ContrastGeometryDensitySettings(
            solvent=ContrastSolventDensitySettings.from_values(
                method=CONTRAST_SOLVENT_METHOD_DIRECT,
                direct_electron_density_e_per_a3=0.05,
            )
        ),
    )
    window._on_density_finished(density_result_b)

    assert window.representative_table.item(0, 7).text() == "Density ready"
    assert window._generated_trace_profiles == []
    assert "contrast_scale_factor" not in window._representative_row_metadata(
        0
    )

    build_result_b = build_contrast_component_profiles(
        representative_result,
        density_result_b,
        q_values=q_values,
        output_dir=project_dir / "contrast_workflow" / "scattering_components",
        metadata_dir=project_dir / "contrast_workflow" / "debye",
        component_map_path=project_dir / "md_saxs_map.json",
    )
    window._load_generated_trace_profiles_from_trace_payloads(
        [
            trace_result.to_dict()
            for trace_result in build_result_b.trace_results
        ],
        default_component_dir=build_result_b.output_dir,
    )
    window._apply_saved_trace_metadata(
        [
            trace_result.to_dict()
            for trace_result in build_result_b.trace_results
        ]
    )

    second_status = window.representative_table.item(0, 7).text()
    second_scale = float(
        window._representative_row_metadata(0)["contrast_scale_factor"]
    )
    assert second_status.startswith("Trace ready (scale ")
    assert second_scale != pytest.approx(first_scale)
    assert second_status != first_status
    window.close()


def test_contrast_mode_workspace_reloads_saved_distribution_artifacts(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, settings, artifact_paths, _build_result = (
        _build_saved_contrast_distribution_project(tmp_path)
    )

    window = ContrastModeMainWindow(
        initial_project_dir=project_dir.resolve(),
        initial_clusters_dir=settings.resolved_clusters_dir,
        initial_q_min=settings.q_min,
        initial_q_max=settings.q_max,
        initial_template_name=settings.selected_model_template,
        initial_distribution_id=artifact_paths.distribution_id,
        initial_distribution_root_dir=artifact_paths.root_dir,
        initial_contrast_artifact_dir=artifact_paths.contrast_dir,
    )
    _flush_contrast_window_launch_preview(window)

    assert window.output_path_edit.text() == str(artifact_paths.root_dir)
    assert window.representatives_output_edit.text() == str(
        artifact_paths.contrast_dir
    )
    assert window.summary_output_edit.text() == str(
        artifact_paths.contrast_dir / "selection_summary.json"
    )
    assert window.representative_table.rowCount() == 1
    assert window.representative_table.item(0, 1).text() == "Auto"
    assert window.representative_table.item(0, 3).text().startswith("#")
    assert window.representative_table.item(0, 5).text().endswith("A^3")
    assert "Δρ" in window.representative_table.item(0, 6).text()
    assert (
        window.representative_table.item(0, 7)
        .text()
        .startswith("Trace ready (scale ")
    )
    assert window.structure_viewer.current_preview is not None
    assert window.structure_viewer.current_preview.has_mesh is True
    assert len(window._generated_trace_profiles) == 1
    assert window.trace_q_range_button.text() == "Autoscale to Model Range"
    assert (
        window.trace_generated_toggle_button.text() == "Hide Computed Traces"
    )
    assert window.trace_generated_toggle_button.isEnabled() is True
    assert "saved contrast distribution loaded" in (
        window.workflow_progress_label.text().lower()
    )
    assert "5/5 workflow stages" in window.workflow_progress_label.text()
    assert window.workflow_progress_bar.value() == 5
    assert window.workflow_progress_bar.maximum() == 5
    experimental_path = tmp_path / "contrast_saved_preview.txt"
    np.savetxt(
        experimental_path,
        np.column_stack(
            [
                np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float),
                np.asarray([100.0, 82.0, 58.0, 31.0], dtype=float),
            ]
        ),
    )
    window.experimental_data_edit.setText(str(experimental_path))
    assert window._load_experimental_preview(log_to_console=False) is True
    window._redraw_trace_plot()

    assert "generated contrast trace" in (
        window.trace_plot_status_label.text().lower()
    )
    assert len(window.trace_figure.axes) == 2
    experimental_axis, component_axis = window.trace_figure.axes
    assert (
        experimental_axis.get_title()
        == "Experimental Data and Contrast Traces"
    )
    assert experimental_axis.get_xlabel() == "q (Å⁻¹)"
    assert (
        experimental_axis.get_ylabel() == "Experimental Intensity (arb. units)"
    )
    assert component_axis.get_ylabel() == "Model Intensity (arb. units)"

    window.trace_generated_toggle_button.click()

    assert (
        window.trace_generated_toggle_button.text() == "Show Computed Traces"
    )
    assert all(
        not line.get_visible()
        for line in window.trace_figure.axes[1].get_lines()
    )

    window.trace_generated_toggle_button.click()

    assert (
        window.trace_generated_toggle_button.text() == "Hide Computed Traces"
    )
    assert any(
        line.get_visible() for line in window.trace_figure.axes[1].get_lines()
    )
    window.close()


def test_contrast_mode_workspace_reloaded_distribution_can_recompute_density(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, settings, artifact_paths, _build_result = (
        _build_saved_contrast_distribution_project(tmp_path)
    )

    window = ContrastModeMainWindow(
        initial_project_dir=project_dir.resolve(),
        initial_clusters_dir=settings.resolved_clusters_dir,
        initial_q_min=settings.q_min,
        initial_q_max=settings.q_max,
        initial_template_name=settings.selected_model_template,
        initial_distribution_id=artifact_paths.distribution_id,
        initial_distribution_root_dir=artifact_paths.root_dir,
        initial_contrast_artifact_dir=artifact_paths.contrast_dir,
    )
    _flush_contrast_window_launch_preview(window)

    started: dict[str, object] = {}

    def fake_start_workflow_task(**kwargs):
        started.update(kwargs)
        return True

    monkeypatch.setattr(
        window, "_start_workflow_task", fake_start_workflow_task
    )

    assert window._representative_analysis_result is not None

    window._compute_electron_density()

    assert started.get("task_name") == "compute_density"
    assert callable(started.get("task"))
    window.close()


def test_save_project_as_copies_project_and_rewrites_internal_paths(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    destination_parent = tmp_path / "copied_projects"
    destination_parent.mkdir()

    monkeypatch.setattr(
        QFileDialog,
        "getExistingDirectory",
        lambda *args, **kwargs: str(destination_parent),
    )
    monkeypatch.setattr(
        QInputDialog,
        "getText",
        lambda *args, **kwargs: ("renamed_project", True),
    )

    window.save_project_as()

    copied_dir = destination_parent / "renamed_project"
    copied_settings = SAXSProjectManager().load_project(copied_dir)

    assert copied_dir.is_dir()
    assert copied_settings.project_name == "renamed_project"
    assert copied_settings.project_dir == str(copied_dir.resolve())
    assert copied_settings.copied_experimental_data_file == str(
        (copied_dir / "experimental_data" / "exp_demo.txt").resolve()
    )
    assert copied_settings.experimental_data_path == str(
        (copied_dir / "experimental_data" / "exp_demo.txt").resolve()
    )
    assert window.current_settings is not None
    assert window.current_settings.project_dir == str(copied_dir.resolve())


def test_prefit_layout_uses_left_parameter_panel_and_combined_output(qapp):
    del qapp
    window = SAXSMainWindow()

    assert window.prefit_tab.parameter_table.parentWidget() is not None
    assert window.prefit_tab._scroll_area.widget() is not None
    assert window.prefit_tab._scroll_area.widgetResizable()
    assert window.prefit_tab._main_splitter.count() == 2
    assert (
        window.prefit_tab._main_splitter.widget(0)
        is window.prefit_tab._pane_splitter
    )
    assert (
        window.prefit_tab._main_splitter.widget(1)
        is window.prefit_tab._output_group
    )
    assert window.prefit_tab._pane_splitter.count() == 2
    assert (
        window.prefit_tab._pane_splitter.widget(0)
        is window.prefit_tab._left_scroll_area
    )
    assert (
        window.prefit_tab._pane_splitter.widget(1)
        is window.prefit_tab._plot_group
    )
    assert window.prefit_tab.output_box is window.prefit_tab.log_box
    assert window.prefit_tab.output_box is window.prefit_tab.summary_box
    assert window.prefit_tab.update_button.parentWidget() is not None
    assert window.prefit_tab.run_button.parentWidget() is not None
    assert window.prefit_tab.save_button.parentWidget() is not None
    assert window.prefit_tab.reset_button.parentWidget() is not None
    assert window.prefit_tab.set_best_button.parentWidget() is not None
    assert window.prefit_tab.reset_best_button.parentWidget() is not None
    assert window.prefit_tab.restore_state_button.parentWidget() is not None


def test_dream_layout_uses_combined_output_and_two_plot_panels(qapp):
    del qapp
    window = SAXSMainWindow()

    assert window.dream_tab._outer_scroll_area.widget() is not None
    assert window.dream_tab._outer_scroll_area.widgetResizable()
    assert window.dream_tab.output_box is window.dream_tab.log_box
    assert window.dream_tab.output_box is window.dream_tab.summary_box
    assert (
        window.dream_tab.output_box.lineWrapMode()
        == window.dream_tab.output_box.LineWrapMode.WidgetWidth
    )
    assert (
        window.dream_tab.output_box.wordWrapMode()
        == QTextOption.WrapMode.WrapAnywhere
    )
    assert (
        window.dream_tab.output_box.sizeAdjustPolicy()
        == QAbstractScrollArea.SizeAdjustPolicy.AdjustIgnored
    )
    assert (
        window.dream_tab.output_box.horizontalScrollBarPolicy()
        == Qt.ScrollBarPolicy.ScrollBarAlwaysOff
    )
    assert window.dream_tab.progress_label.wordWrap()
    assert (
        window.dream_tab.progress_label.sizePolicy().horizontalPolicy()
        == QSizePolicy.Policy.Ignored
    )
    assert window.dream_tab._main_splitter is window.dream_tab._top_splitter
    assert window.dream_tab._main_splitter.count() == 2
    assert window.dream_tab.verbose_checkbox.isChecked()
    assert window.dream_tab.verbose_interval_spin.value() == pytest.approx(5.0)


def test_dream_progress_label_wrap_does_not_resize_left_pane(qapp):
    del qapp
    window = SAXSMainWindow()
    window.show()
    QApplication.processEvents()

    before_sizes = window.dream_tab._main_splitter.sizes()
    message = (
        "DREAM runtime output path is writing a very long progress message "
        "that should wrap inside the left pane instead of expanding the "
        "splitter width to accommodate the text length. " * 4
    ).strip()

    window.dream_tab.start_progress(message)
    QApplication.processEvents()

    after_sizes = window.dream_tab._main_splitter.sizes()

    assert window.dream_tab.progress_label.wordWrap()
    assert (
        window.dream_tab.progress_label.sizePolicy().horizontalPolicy()
        == QSizePolicy.Policy.Ignored
    )
    assert abs(after_sizes[0] - before_sizes[0]) <= 2
    window.close()
    assert window.dream_tab.verbose_interval_spin.isEnabled()
    assert window.dream_tab.selected_search_filter_preset() == "medium"
    assert window.dream_tab.chains_spin.value() == 4
    assert window.dream_tab.iterations_spin.value() == 10000
    assert window.dream_tab.burnin_spin.value() == 20
    assert window.dream_tab.history_thin_spin.value() == 10
    assert window.dream_tab.nseedchains_spin.value() == 40
    assert window.dream_tab.lambda_spin.value() == pytest.approx(0.05)
    assert window.dream_tab.zeta_spin.value() == pytest.approx(1e-12)
    assert window.dream_tab.snooker_spin.value() == pytest.approx(0.1)
    assert window.dream_tab.p_gamma_unity_spin.value() == pytest.approx(0.2)
    assert (
        window.dream_tab._main_splitter.widget(0)
        is window.dream_tab._settings_scroll_area
    )
    assert (
        window.dream_tab._main_splitter.widget(1)
        is window.dream_tab._plot_scroll_area
    )
    assert (
        window.dream_tab._plot_scroll_area.widget()
        is window.dream_tab._plot_panel
    )
    assert window.dream_tab._plot_splitter.count() == 2
    assert window.dream_tab._output_group.parentWidget() is not None
    assert window.dream_tab.model_canvas is not None
    assert window.dream_tab.violin_canvas is not None
    assert window.dream_tab.show_experimental_trace_checkbox.isChecked()
    assert window.dream_tab.show_model_trace_checkbox.isChecked()
    assert not window.dream_tab.show_solvent_trace_checkbox.isChecked()
    assert (
        not window.dream_tab.show_structure_factor_trace_checkbox.isChecked()
    )
    assert window.dream_tab.model_log_x_checkbox.isChecked()
    assert window.dream_tab.model_log_y_checkbox.isChecked()
    assert (
        window.dream_tab._plot_splitter.widget(0)
        is window.dream_tab.model_canvas.parentWidget()
    )
    assert (
        window.dream_tab._plot_splitter.widget(1)
        is window.dream_tab.violin_canvas.parentWidget()
    )
    assert (
        window.dream_tab.posterior_filter_combo.currentData()
        == "all_post_burnin"
    )
    assert window.dream_tab.posterior_top_percent_spin.isEnabled()
    assert window.dream_tab.posterior_top_n_spin.isEnabled()
    assert window.dream_tab.auto_filter_assessment_checkbox.isChecked()
    assert (
        window.dream_tab.violin_sample_source_combo.currentData()
        == "filtered_posterior"
    )
    assert window.dream_tab.weight_order_combo.currentData() == "weight_index"
    assert (
        window.dream_tab.violin_value_scale_combo.currentData()
        == "parameter_value"
    )
    assert (
        window.dream_tab.violin_value_scale_combo.findData(
            "effective_radii_only"
        )
        >= 0
    )
    assert (
        window.dream_tab.violin_value_scale_combo.findData(
            "additional_parameters_only"
        )
        >= 0
    )
    assert (
        window.dream_tab.violin_mode_combo.findData("effective_radii_only")
        >= 0
    )
    assert (
        window.dream_tab.violin_mode_combo.findData(
            "additional_parameters_only"
        )
        >= 0
    )
    assert window.dream_tab.violin_palette_combo.currentData() == "Blues"
    assert not window.dream_tab.color_options_panel.isVisible()
    assert not window.dream_tab.color_options_toggle_button.isChecked()
    assert window.dream_tab.selected_violin_point_color() == to_hex(
        "tab:red",
        keep_alpha=False,
    )
    assert window.dream_tab.selected_violin_outline_color() == "#000000"
    assert window.dream_tab.violin_outline_width_spin.value() == pytest.approx(
        0.8
    )
    assert window.dream_tab.violin_custom_color_button.isEnabled()
    assert window.dream_tab.violin_point_color_button.isEnabled()
    assert (
        "automatic post-run filter assessment"
        in window.dream_tab.posterior_top_percent_spin.toolTip()
    )
    assert window.dream_tab.parameter_map_table.rowCount() == 0
    assert (
        window.dream_tab.parameter_map_table.editTriggers()
        == window.dream_tab.parameter_map_table.EditTrigger.NoEditTriggers
    )
    assert window.dream_tab.edit_button.parentWidget() is not None
    assert window.dream_tab.save_settings_button.parentWidget() is not None
    assert window.dream_tab.write_button.parentWidget() is not None
    assert window.dream_tab.preview_button.parentWidget() is not None
    assert window.dream_tab.run_button.parentWidget() is not None
    assert window.dream_tab.load_button.parentWidget() is not None
    assert window.dream_tab.report_button.parentWidget() is not None
    assert window.dream_tab.save_model_button.parentWidget() is not None
    assert window.dream_tab.save_violin_button.parentWidget() is not None
    assert window.dream_tab.setup_actions_group.title() == "DREAM Setup"
    assert window.dream_tab.filter_status_group.title() == "Filter Status"
    assert window.dream_tab.analysis_actions_group.title() == (
        "DREAM Analysis"
    )
    assert not window.dream_tab.apply_filter_button.isEnabled()
    assert "No DREAM dataset is loaded yet" in (
        window.dream_tab.filter_status_box.toPlainText()
    )


@pytest.mark.parametrize(
    ("tab_attr", "splitter_attr", "target_index"),
    [
        ("project_setup_tab", "_pane_splitter", 0),
        ("project_setup_tab", "_pane_splitter", 1),
        ("prefit_tab", "_pane_splitter", 0),
        ("prefit_tab", "_pane_splitter", 1),
        ("dream_tab", "_top_splitter", 0),
        ("dream_tab", "_top_splitter", 1),
    ],
)
def test_auto_snap_panes_resizes_clicked_supported_pane(
    qapp,
    tab_attr,
    splitter_attr,
    target_index,
):
    del qapp
    window = SAXSMainWindow()
    window.resize(1800, 980)
    window.show()
    tab = getattr(window, tab_attr)
    window.tabs.setCurrentWidget(tab)
    QApplication.processEvents()

    splitter = getattr(tab, splitter_attr)
    pane_snap_filter = tab._auto_snap_filter
    before_sizes, after_sizes = _pane_snap_resize_result(
        splitter,
        pane_snap_filter,
        target_index=target_index,
    )

    assert after_sizes[target_index] > before_sizes[target_index] + 20
    window.close()


def test_auto_snap_panes_disabled_preserves_manual_splitter_sizes(qapp):
    del qapp
    window = SAXSMainWindow()
    window.resize(1800, 980)
    window.show()
    window._set_auto_snap_panes_enabled(False, persist=False)
    window.tabs.setCurrentWidget(window.prefit_tab)
    QApplication.processEvents()

    splitter = window.prefit_tab._pane_splitter
    splitter.setSizes([900, 300])
    QApplication.processEvents()
    before_sizes = splitter.sizes()
    QTest.mouseClick(
        window.prefit_tab._plot_group,
        Qt.MouseButton.LeftButton,
        Qt.KeyboardModifier.NoModifier,
        window.prefit_tab._plot_group.rect().center(),
    )
    QApplication.processEvents()
    after_sizes = splitter.sizes()

    assert abs(after_sizes[0] - before_sizes[0]) <= 2
    assert abs(after_sizes[1] - before_sizes[1]) <= 2
    window.close()


@pytest.mark.parametrize("target_index", [0, 1])
def test_auto_snap_panes_focus_clicked_pane_when_its_preferred_width_is_small(
    qapp,
    target_index,
):
    del qapp
    window = SAXSMainWindow()
    window.resize(1800, 980)
    window.show()
    window.tabs.setCurrentWidget(window.project_setup_tab)
    QApplication.processEvents()

    before_sizes, after_sizes = _pane_snap_resize_result(
        window.project_setup_tab._pane_splitter,
        window.project_setup_tab._auto_snap_filter,
        target_index=target_index,
        desired_width=220,
        other_width=720,
    )

    assert after_sizes[target_index] > before_sizes[target_index] + 20
    assert after_sizes[1 - target_index] < before_sizes[1 - target_index] - 20
    window.close()


def test_dream_zeta_spin_uses_scientific_notation(qapp):
    del qapp
    window = SAXSMainWindow()

    window.dream_tab.set_settings(DreamRunSettings(zeta=1e-12))

    assert "e" in window.dream_tab.zeta_spin.text().lower()
    assert window.dream_tab.zeta_spin.text().lower().endswith("e-12")
    assert window.dream_tab.settings_payload().zeta == pytest.approx(1e-12)


def test_apply_dream_output_settings_updates_verbose_controls(qapp):
    del qapp
    window = SAXSMainWindow()

    window._apply_dream_output_settings(verbose=False, interval_seconds=2.5)

    assert not window.dream_tab.verbose_checkbox.isChecked()
    assert window.dream_tab.verbose_interval_spin.value() == pytest.approx(2.5)
    assert not window.dream_tab.verbose_interval_spin.isEnabled()
    assert "Updated DREAM output settings." in (
        window.dream_tab.output_box.toPlainText()
    )


def test_apply_powerpoint_export_settings_updates_project_state(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window._apply_powerpoint_export_settings(
        PowerPointExportSettings(
            font_family="Courier New",
            component_color_map="plasma",
            prior_histogram_color_map="cividis",
            solvent_sort_histogram_color_map="magma",
            include_output_summary=False,
            generate_manifest=False,
            export_figure_assets=False,
        )
    )

    settings = window.current_settings.powerpoint_export_settings
    assert settings.font_family == "Courier New"
    assert settings.component_color_map == "plasma"
    assert settings.prior_histogram_color_map == "cividis"
    assert settings.solvent_sort_histogram_color_map == "magma"
    assert not settings.include_output_summary
    assert not settings.generate_manifest
    assert not settings.export_figure_assets
    assert "Updated PowerPoint export settings." in (
        window.dream_tab.output_box.toPlainText()
    )


def test_dream_tab_batches_runtime_output_console_updates(qapp, monkeypatch):
    del qapp
    tab = DreamTab()
    render_calls: list[bool] = []
    original_render = tab._render_output

    def tracked_render(*, scroll_to_end: bool = False):
        render_calls.append(bool(scroll_to_end))
        return original_render(scroll_to_end=scroll_to_end)

    monkeypatch.setattr(tab, "_render_output", tracked_render)

    tab.append_runtime_output("line 1")
    tab.append_runtime_output("line 2")
    QApplication.processEvents()

    assert "line 1" not in tab.output_box.toPlainText()
    assert not render_calls

    QTest.qWait(tab.RUNTIME_OUTPUT_FLUSH_INTERVAL_MS + 75)
    QApplication.processEvents()

    text = tab.output_box.toPlainText()
    assert "DREAM Runtime Output" in text
    assert "line 1" in text
    assert "line 2" in text
    assert len(render_calls) == 1


def test_dream_tab_flushes_pending_runtime_output_before_regular_logs(qapp):
    del qapp
    tab = DreamTab()

    tab.set_log_text("Base log")
    tab.append_runtime_output("runtime line")
    tab.append_log("final log line")
    QApplication.processEvents()

    text = tab.output_box.toPlainText()
    assert "runtime line" in text
    assert "final log line" in text
    assert text.index("runtime line") < text.index("final log line")


def test_dream_search_filter_presets_update_and_fall_back_to_custom(qapp):
    del qapp
    window = SAXSMainWindow()

    aggressive_index = window.dream_tab.search_filter_preset_combo.findData(
        "more_aggressive"
    )
    window.dream_tab.search_filter_preset_combo.setCurrentIndex(
        aggressive_index
    )

    assert (
        window.dream_tab.selected_search_filter_preset() == "more_aggressive"
    )
    assert window.dream_tab.chains_spin.value() == 8
    assert window.dream_tab.iterations_spin.value() == 20000
    assert window.dream_tab.burnin_spin.value() == 25
    assert window.dream_tab.posterior_filter_combo.currentData() == (
        "top_percent_logp"
    )
    assert (
        window.dream_tab.posterior_top_percent_spin.value()
        == pytest.approx(5.0)
    )

    window.dream_tab.iterations_spin.setValue(21000)

    assert window.dream_tab.selected_search_filter_preset() == "custom"


def test_main_window_ui_scale_updates_fonts_and_minimum_sizes(qapp):
    del qapp
    window = SAXSMainWindow()

    base_font_size = window.font().pointSizeF()
    base_output_height = window.prefit_tab.output_box.minimumHeight()
    base_handle_width = window.dream_tab._top_splitter.handleWidth()

    window._set_ui_scale(1.2)

    assert window._ui_scale == 1.2
    assert window.font().pointSizeF() > base_font_size
    assert window.prefit_tab.output_box.minimumHeight() > base_output_height
    assert window.dream_tab._top_splitter.handleWidth() > base_handle_width

    window._reset_ui_scale()

    assert window._ui_scale == 1.0
    assert window.font().pointSizeF() == pytest.approx(base_font_size)
    assert window.prefit_tab.output_box.minimumHeight() == base_output_height
    assert window.dream_tab._top_splitter.handleWidth() == base_handle_width


def test_main_window_window_preset_resizes_and_scales_ui(qapp, monkeypatch):
    del qapp
    window = SAXSMainWindow()
    monkeypatch.setattr(
        window,
        "_current_available_geometry",
        lambda: QRect(0, 0, 1728, 1117),
    )

    window._apply_window_layout_preset("laptop_13")

    assert window.width() == 1180
    assert window.height() == 760
    assert window._ui_scale == pytest.approx(0.95)

    window._apply_window_layout_preset("display_1440p")

    assert window.width() == 1680
    assert window.height() == 980
    assert window._ui_scale == pytest.approx(1.1)


def test_dream_blink_highlight_does_not_override_button_stylesheet(qapp):
    del qapp
    window = SAXSMainWindow()

    window.dream_tab.blink_edit_priors_button()

    assert window.dream_tab.edit_button.styleSheet() == ""
    assert window.dream_tab.edit_button.graphicsEffect() is not None


def test_dream_blink_uses_per_button_effects(qapp):
    del qapp
    window = SAXSMainWindow()

    window.dream_tab.blink_edit_priors_button()
    first_effect = window.dream_tab.edit_button.graphicsEffect()
    window.dream_tab.blink_write_bundle_button()

    assert window.dream_tab.edit_button.graphicsEffect() is None
    assert window.dream_tab.write_button.graphicsEffect() is not None
    assert window.dream_tab.write_button.graphicsEffect() is not first_effect


def test_dream_prior_map_table_updates_after_saving_entries(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    entries = [
        DreamParameterEntry(
            structure="PbI2",
            motif="motif_A",
            param_type="SAXS",
            param="w0",
            value=0.6,
            vary=True,
            distribution="lognorm",
            dist_params={"loc": 0.0, "scale": 0.6, "s": 0.1},
        )
    ]

    window._save_distribution_entries(entries)

    table = window.dream_tab.parameter_map_table
    assert table.rowCount() == 1
    assert table.item(0, 0).text() == "PbI2"
    assert table.item(0, 3).text() == "w0"
    assert table.item(0, 5).text() == "Yes"


def test_distribution_window_prompts_before_quitting_without_first_save(
    qapp, monkeypatch
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="w0",
                value=0.6,
                vary=True,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 0.6, "s": 0.1},
            )
        ]
    )
    window.load_entries(
        window.current_entries(), has_existing_parameter_map=False
    )
    window.show()

    monkeypatch.setattr(
        "saxshell.saxs.ui.distribution_window.QMessageBox.question",
        lambda *args, **kwargs: QMessageBox.StandardButton.No,
    )

    window.close()
    QApplication.processEvents()

    assert window.isVisible()

    monkeypatch.setattr(
        "saxshell.saxs.ui.distribution_window.QMessageBox.question",
        lambda *args, **kwargs: QMessageBox.StandardButton.Yes,
    )

    window.close()
    QApplication.processEvents()

    assert not window.isVisible()


def test_distribution_window_layout_uses_splitters_for_plot_table_and_console(
    qapp,
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="w0",
                value=0.6,
                vary=True,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 0.6, "s": 0.1},
            )
        ]
    )
    window.show()
    QApplication.processEvents()

    assert isinstance(window._main_splitter, QSplitter)
    assert window._main_splitter.orientation() == Qt.Orientation.Horizontal
    assert window._main_splitter.count() == 2
    assert window._main_splitter.widget(0) is window._left_panel
    assert window._main_splitter.widget(1) is window._plot_panel
    assert all(size > 0 for size in window._main_splitter.sizes())

    assert isinstance(window._left_splitter, QSplitter)
    assert window._left_splitter.orientation() == Qt.Orientation.Vertical
    assert window._left_splitter.count() == 2
    assert window._left_splitter.widget(0) is window._editor_panel
    assert window._left_splitter.widget(1) is window.console
    assert all(size > 0 for size in window._left_splitter.sizes())


def test_distribution_window_interactive_center_lock_toggles_center_handle(
    qapp,
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
            )
        ]
    )
    window.show()
    QApplication.processEvents()

    assert window.lock_center_checkbox.isChecked()
    assert "currently locked" in window.interactive_hint_label.text()
    assert window._interactive_handles is not None
    assert window._interactive_handles.center is None

    window.lock_center_checkbox.setChecked(False)
    QApplication.processEvents()

    assert "red center handle" in window.interactive_hint_label.text()
    assert window._interactive_handles is not None
    assert window._interactive_handles.center is not None


def test_distribution_window_plot_is_square_and_shows_reset_baseline(qapp):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 1.0},
            )
        ]
    )
    window.show()
    QApplication.processEvents()

    axis = window.figure.axes[0]
    baseline_lines = _plot_lines_by_gid(axis, "reset-baseline")
    current_lines = _plot_lines_by_gid(axis, "current-distribution")

    assert float(axis.get_box_aspect()) == pytest.approx(1.0)
    assert len(baseline_lines) == 1
    assert len(current_lines) == 1
    assert np.allclose(
        baseline_lines[0].get_xdata(),
        current_lines[0].get_xdata(),
    )
    assert np.allclose(
        baseline_lines[0].get_ydata(),
        current_lines[0].get_ydata(),
    )


def test_distribution_window_width_drag_updates_norm_and_uniform_params():
    norm_entry = DreamParameterEntry(
        structure="PbI2",
        motif="motif_A",
        param_type="SAXS",
        param="scale",
        value=1.0,
        vary=True,
        distribution="norm",
        dist_params={"loc": 1.0, "scale": 0.2},
    )
    updated_norm = DistributionSetupWindow._width_drag_adjusted_entry(
        norm_entry,
        handle_kind="right_width",
        target_x=1.9,
    )
    assert updated_norm.value == pytest.approx(1.0)
    assert updated_norm.dist_params["loc"] == pytest.approx(1.0)
    assert updated_norm.dist_params["scale"] == pytest.approx(0.3)

    uniform_entry = DreamParameterEntry(
        structure="PbI2",
        motif="motif_A",
        param_type="SAXS",
        param="phi_solute",
        value=2.0,
        vary=True,
        distribution="uniform",
        dist_params={"loc": 1.6, "scale": 0.8},
    )
    updated_uniform = DistributionSetupWindow._width_drag_adjusted_entry(
        uniform_entry,
        handle_kind="left_width",
        target_x=1.2,
    )
    assert updated_uniform.value == pytest.approx(2.0)
    assert updated_uniform.dist_params["scale"] == pytest.approx(1.6)
    assert updated_uniform.dist_params["loc"] == pytest.approx(1.2)


def test_distribution_window_center_and_peak_drag_adjust_lognorm_cleanly():
    entry = DreamParameterEntry(
        structure="PbI2",
        motif="motif_A",
        param_type="Both",
        param="w0",
        value=0.6,
        vary=True,
        distribution="lognorm",
        dist_params={"loc": 0.0, "scale": 0.6, "s": 0.2},
    )

    moved = DistributionSetupWindow._center_drag_adjusted_entry(
        entry,
        target_center=0.9,
    )
    assert moved.value == pytest.approx(0.9)
    assert moved.dist_params["scale"] == pytest.approx(0.6)
    assert moved.dist_params["s"] == pytest.approx(0.2)
    assert moved.dist_params["loc"] == pytest.approx(0.3)

    narrowed = DistributionSetupWindow._peak_drag_adjusted_entry(
        entry,
        start_y=0.8,
        target_y=1.4,
        y_limits=(0.0, 2.0),
    )
    widened = DistributionSetupWindow._peak_drag_adjusted_entry(
        entry,
        start_y=0.8,
        target_y=0.2,
        y_limits=(0.0, 2.0),
    )
    assert narrowed.value == pytest.approx(entry.value)
    assert narrowed.dist_params["loc"] == pytest.approx(0.0)
    assert narrowed.dist_params["s"] < entry.dist_params["s"]
    assert widened.dist_params["loc"] == pytest.approx(0.0)
    assert widened.dist_params["s"] > entry.dist_params["s"]


def test_distribution_window_switches_lognorm_params_for_norm_and_uniform(
    qapp,
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="w0",
                value=0.6,
                vary=True,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 0.6, "s": 0.1},
            )
        ]
    )
    combo = window.table.cellWidget(0, 6)

    combo.setCurrentText("norm")
    QApplication.processEvents()

    norm_params = json.loads(window.table.item(0, 7).text())
    low_column = _table_column_index(window.table, "Guide Low")
    high_column = _table_column_index(window.table, "Guide High")
    assert "s" not in norm_params
    assert set(norm_params) == {"loc", "scale"}
    assert float(window.table.item(0, low_column).text()) == pytest.approx(
        -1.8
    )
    assert float(window.table.item(0, high_column).text()) == pytest.approx(
        1.8
    )
    assert window.figure.axes[0].get_title() == "w0: norm"

    combo.setCurrentText("uniform")
    QApplication.processEvents()

    uniform_params = json.loads(window.table.item(0, 7).text())
    assert "s" not in uniform_params
    assert set(uniform_params) == {"loc", "scale"}
    assert float(window.table.item(0, low_column).text()) == pytest.approx(
        0.0,
        abs=1e-12,
    )
    assert float(window.table.item(0, high_column).text()) == pytest.approx(
        0.6
    )
    assert window.figure.axes[0].get_title() == "w0: uniform"


def test_distribution_window_displays_distribution_guide_bounds(qapp):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
            ),
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_B",
                param_type="SAXS",
                param="phi_solute",
                value=0.45,
                vary=True,
                distribution="uniform",
                dist_params={"loc": 0.2, "scale": 0.5},
            ),
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_C",
                param_type="Both",
                param="w0",
                value=0.6,
                vary=True,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 0.6, "s": 0.1},
            ),
        ]
    )

    low_column = _table_column_index(window.table, "Guide Low")
    high_column = _table_column_index(window.table, "Guide High")
    q_low = stats.norm.cdf(-3.0)
    q_high = stats.norm.cdf(3.0)

    assert float(window.table.item(0, low_column).text()) == pytest.approx(0.4)
    assert float(window.table.item(0, high_column).text()) == pytest.approx(
        1.6
    )
    assert float(window.table.item(1, low_column).text()) == pytest.approx(0.2)
    assert float(window.table.item(1, high_column).text()) == pytest.approx(
        0.7
    )
    assert float(window.table.item(2, low_column).text()) == pytest.approx(
        stats.lognorm.ppf(q_low, s=0.1, loc=0.0, scale=0.6),
        rel=1e-6,
    )
    assert float(window.table.item(2, high_column).text()) == pytest.approx(
        stats.lognorm.ppf(q_high, s=0.1, loc=0.0, scale=0.6),
        rel=1e-6,
    )


def test_distribution_window_updates_distribution_guides_after_param_edit(
    qapp,
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
            )
        ]
    )

    low_column = _table_column_index(window.table, "Guide Low")
    high_column = _table_column_index(window.table, "Guide High")

    window.table.item(0, 7).setText(
        json.dumps({"loc": 1.0, "scale": 0.4}, sort_keys=True)
    )
    QApplication.processEvents()

    assert float(window.table.item(0, low_column).text()) == pytest.approx(
        -0.2
    )
    assert float(window.table.item(0, high_column).text()) == pytest.approx(
        2.2
    )


def test_distribution_window_manual_guide_edit_updates_norm_params(qapp):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
            )
        ]
    )

    low_column = _table_column_index(window.table, "Guide Low")
    high_column = _table_column_index(window.table, "Guide High")

    window.table.item(0, low_column).setText("0.1")
    QApplication.processEvents()

    entry = window.current_entries()[0]
    assert entry.value == pytest.approx(1.0)
    assert entry.dist_params["loc"] == pytest.approx(1.0)
    assert entry.dist_params["scale"] == pytest.approx(0.3)
    assert float(window.table.item(0, low_column).text()) == pytest.approx(0.1)
    assert float(window.table.item(0, high_column).text()) == pytest.approx(
        1.9
    )


def test_distribution_window_manual_guide_edit_updates_lognorm_params(qapp):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="Both",
                param="w0",
                value=0.6,
                vary=True,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 0.6, "s": 0.1},
            )
        ]
    )

    high_column = _table_column_index(window.table, "Guide High")

    window.table.item(0, high_column).setText("0.9")
    QApplication.processEvents()

    entry = window.current_entries()[0]
    assert entry.value == pytest.approx(0.6)
    assert float(window.table.item(0, high_column).text()) == pytest.approx(
        0.9
    )
    assert entry.dist_params["s"] != pytest.approx(0.1)


def test_distribution_window_interactive_drag_preview_is_throttled(
    monkeypatch,
    qapp,
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
            )
        ]
    )

    preview_calls: list[tuple[DreamParameterEntry, int | None, bool]] = []

    def _record_plot(
        entry: DreamParameterEntry,
        *,
        row: int | None = None,
        force_rescale: bool = False,
        interactive_preview: bool = False,
    ) -> None:
        del force_rescale
        preview_calls.append((entry, row, interactive_preview))

    monkeypatch.setattr(window, "_plot_entry", _record_plot)

    first_preview = DreamParameterEntry.from_dict(
        window.current_entries()[0].to_dict()
    )
    first_preview.dist_params["scale"] = 0.25
    second_preview = DreamParameterEntry.from_dict(first_preview.to_dict())
    second_preview.dist_params["scale"] = 0.35

    window._schedule_interactive_drag_preview(first_preview, row=0)
    window._schedule_interactive_drag_preview(second_preview, row=0)

    assert preview_calls == []

    QTest.qWait(INTERACTIVE_PREVIEW_THROTTLE_MS + 30)

    assert len(preview_calls) == 1
    entry, row, interactive_preview = preview_calls[0]
    assert row == 0
    assert interactive_preview is True
    assert entry.dist_params["scale"] == pytest.approx(0.35)


def test_distribution_window_tracks_current_row_and_rescales_plot(qapp):
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
            ),
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_B",
                param_type="Both",
                param="w0",
                value=1.0,
                vary=True,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 1.0, "s": 2.0},
            ),
        ]
    )

    axis = window.figure.axes[0]
    assert axis.get_title() == "scale: norm"
    assert axis.get_xscale() == "linear"

    window.table.setCurrentCell(1, 0)
    qapp.processEvents()

    axis = window.figure.axes[0]
    assert axis.get_title() == "w0: lognorm"
    assert axis.get_xscale() == "log"
    assert axis.get_xlim()[0] > 0.0
    assert axis.get_xlim()[1] < 100.0


def test_distribution_window_rescales_plot_after_distribution_mode_change(
    qapp,
):
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="Both",
                param="w0",
                value=1.0,
                vary=True,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 1.0, "s": 2.0},
            )
        ]
    )

    axis = window.figure.axes[0]
    assert axis.get_xscale() == "log"
    combo = window.table.cellWidget(0, 6)

    combo.setCurrentText("uniform")
    qapp.processEvents()

    axis = window.figure.axes[0]
    assert axis.get_title() == "w0: uniform"
    assert axis.get_xscale() == "linear"
    assert axis.get_xlim()[1] - axis.get_xlim()[0] < 2.0


def test_distribution_window_rescales_plot_after_param_edit(qapp):
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 1.0},
            )
        ]
    )

    axis = window.figure.axes[0]
    initial_xlim = axis.get_xlim()
    initial_ylim = axis.get_ylim()
    initial_baseline_y = _plot_lines_by_gid(axis, "reset-baseline")[
        0
    ].get_ydata()

    window.table.item(0, 7).setText(
        json.dumps({"loc": 1.0, "scale": 1.1}, sort_keys=True)
    )
    qapp.processEvents()

    axis = window.figure.axes[0]
    updated_xlim = axis.get_xlim()
    updated_ylim = axis.get_ylim()
    baseline_line = _plot_lines_by_gid(axis, "reset-baseline")[0]
    current_line = _plot_lines_by_gid(axis, "current-distribution")[0]
    assert axis.get_title() == "scale: norm"
    assert axis.get_xscale() == "linear"
    assert updated_xlim == pytest.approx(initial_xlim)
    assert updated_ylim == pytest.approx(initial_ylim)
    assert np.allclose(baseline_line.get_ydata(), initial_baseline_y)
    assert not np.allclose(
        current_line.get_ydata(),
        baseline_line.get_ydata(),
    )

    window.rescale_axes_button.click()
    qapp.processEvents()

    axis = window.figure.axes[0]
    rescaled_xlim = axis.get_xlim()
    rescaled_ylim = axis.get_ylim()
    assert (rescaled_xlim[1] - rescaled_xlim[0]) > (
        initial_xlim[1] - initial_xlim[0]
    )
    assert rescaled_ylim[1] < initial_ylim[1]


def test_distribution_window_auto_rescales_when_distribution_exits_window(
    qapp,
):
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 1.0},
            )
        ]
    )

    axis = window.figure.axes[0]
    initial_xlim = axis.get_xlim()
    initial_ylim = axis.get_ylim()

    window.table.item(0, 7).setText(
        json.dumps({"loc": 1.0, "scale": 1.35}, sort_keys=True)
    )
    qapp.processEvents()

    axis = window.figure.axes[0]
    assert (axis.get_xlim()[1] - axis.get_xlim()[0]) > (
        initial_xlim[1] - initial_xlim[0]
    )
    assert axis.get_ylim()[1] < initial_ylim[1]


def test_distribution_window_row_reset_restores_loaded_prior_settings(qapp):
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
                smart_preset_status="strict",
            )
        ]
    )

    reset_col = _table_column_index(window.table, "Reset")
    low_column = _table_column_index(window.table, "Guide Low")
    high_column = _table_column_index(window.table, "Guide High")
    distribution_combo = window.table.cellWidget(0, 6)
    vary_box = window.table.cellWidget(0, 5)
    reset_button = window.table.cellWidget(0, reset_col)

    window.table.item(0, 4).setText("2.5")
    vary_box.setChecked(True)
    distribution_combo.setCurrentText("uniform")
    window.table.item(0, 7).setText(
        json.dumps({"loc": 0.0, "scale": 4.0}, sort_keys=True)
    )
    qapp.processEvents()

    assert reset_button is not None
    reset_button.click()
    qapp.processEvents()

    entry = window.current_entries()[0]
    assert entry.value == pytest.approx(1.0)
    assert entry.vary is False
    assert entry.distribution == "norm"
    assert entry.dist_params == pytest.approx({"loc": 1.0, "scale": 0.2})
    assert entry.smart_preset_status == "strict"
    assert float(window.table.item(0, low_column).text()) == pytest.approx(0.4)
    assert float(window.table.item(0, high_column).text()) == pytest.approx(
        1.6
    )
    assert window.figure.axes[0].get_title() == "scale: norm"


def test_distribution_window_row_reset_preserves_original_baseline_after_preset(
    qapp,
):
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
            )
        ]
    )

    reset_col = _table_column_index(window.table, "Reset")
    window.smart_prior_preset_combo.setCurrentIndex(
        window.smart_prior_preset_combo.findData("lenient")
    )
    window.apply_smart_prior_preset_button.click()
    qapp.processEvents()

    adjusted_entry = window.current_entries()[0]
    assert adjusted_entry.dist_params["scale"] == pytest.approx(0.3)

    reset_button = window.table.cellWidget(0, reset_col)
    assert reset_button is not None
    reset_button.click()
    qapp.processEvents()

    reset_entry = window.current_entries()[0]
    assert reset_entry.dist_params["scale"] == pytest.approx(0.2)
    assert reset_entry.smart_preset_status == "custom"


def test_distribution_window_can_toggle_all_vary_flags(qapp):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="w0",
                value=0.6,
                vary=True,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 0.6, "s": 0.1},
            ),
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_B",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
            ),
        ]
    )

    window.set_all_vary_off_button.click()
    assert all(not entry.vary for entry in window.current_entries())

    window.set_all_vary_on_button.click()
    assert all(entry.vary for entry in window.current_entries())


def test_distribution_window_recommended_vary_selection_keeps_radius_params_off(
    qapp,
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="Both",
                param="w0",
                value=0.6,
                vary=False,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 0.6, "s": 0.1},
            ),
            DreamParameterEntry(
                structure="",
                motif="",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
            ),
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="r_eff_w0",
                value=4.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 4.0, "scale": 0.2},
            ),
        ]
    )

    window.select_recommended_vary_button.click()
    entries = window.current_entries()

    assert entries[0].vary is True
    assert entries[1].vary is True
    assert entries[2].vary is False


def test_distribution_window_smart_prior_preset_tightens_and_relaxes_spreads(
    qapp,
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="Both",
                param="w0",
                value=0.6,
                vary=True,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 0.6, "s": 0.2},
            ),
            DreamParameterEntry(
                structure="",
                motif="",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.3},
            ),
        ]
    )

    window.smart_prior_preset_combo.setCurrentText("Strict")
    window.apply_smart_prior_preset_button.click()
    strict_entries = window.current_entries()
    assert strict_entries[0].dist_params["s"] == pytest.approx(0.13)
    assert strict_entries[1].dist_params["scale"] == pytest.approx(0.195)

    window.smart_prior_preset_combo.setCurrentText("Lenient")
    window.apply_smart_prior_preset_button.click()
    lenient_entries = window.current_entries()
    assert lenient_entries[0].dist_params["s"] == pytest.approx(0.195)
    assert lenient_entries[1].dist_params["scale"] == pytest.approx(0.2925)


def test_distribution_window_smart_prior_preset_can_target_selected_structures(
    qapp,
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="Small",
                motif="m1",
                param_type="Both",
                param="w0",
                value=0.3,
                vary=True,
                distribution="norm",
                dist_params={"loc": 0.3, "scale": 0.1},
            ),
            DreamParameterEntry(
                structure="Small",
                motif="m1",
                param_type="SAXS",
                param="r_eff_w0",
                value=3.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 3.0, "scale": 0.2},
            ),
            DreamParameterEntry(
                structure="Large",
                motif="m2",
                param_type="Both",
                param="w1",
                value=0.7,
                vary=True,
                distribution="norm",
                dist_params={"loc": 0.7, "scale": 0.1},
            ),
        ]
    )

    scope_index = window.smart_prior_apply_scope_combo.findData("selected")
    assert scope_index >= 0
    window.smart_prior_apply_scope_combo.setCurrentIndex(scope_index)
    window.table.selectRow(0)
    window.smart_prior_preset_combo.setCurrentText("Strict")
    window.apply_smart_prior_preset_button.click()
    entries = window.current_entries()

    assert entries[0].dist_params["scale"] == pytest.approx(0.065)
    assert entries[1].dist_params["scale"] == pytest.approx(0.13)
    assert entries[2].dist_params["scale"] == pytest.approx(0.1)
    assert window.table.cellWidget(0, 8).currentText() == "Strict"
    assert window.table.cellWidget(1, 8).currentText() == "Strict"
    assert window.table.cellWidget(2, 8).currentText() == "Custom / Manual"


def test_distribution_window_size_aware_prior_preset_uses_effective_radii(
    qapp,
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="Small",
                motif="m1",
                param_type="Both",
                param="w0",
                value=0.3,
                vary=True,
                distribution="norm",
                dist_params={"loc": 0.3, "scale": 0.1},
            ),
            DreamParameterEntry(
                structure="Large",
                motif="m2",
                param_type="Both",
                param="w1",
                value=0.7,
                vary=True,
                distribution="norm",
                dist_params={"loc": 0.7, "scale": 0.1},
            ),
            DreamParameterEntry(
                structure="Small",
                motif="m1",
                param_type="SAXS",
                param="r_eff_w0",
                value=3.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 3.0, "scale": 0.2},
            ),
            DreamParameterEntry(
                structure="Large",
                motif="m2",
                param_type="SAXS",
                param="r_eff_w1",
                value=8.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 8.0, "scale": 0.2},
            ),
        ]
    )

    window.smart_prior_preset_combo.setCurrentText(
        "Strict Small / Lenient Large"
    )
    window.apply_smart_prior_preset_button.click()
    entries = window.current_entries()

    assert entries[0].dist_params["scale"] == pytest.approx(0.065)
    assert entries[1].dist_params["scale"] == pytest.approx(0.15)
    assert entries[2].dist_params["scale"] == pytest.approx(0.13)
    assert entries[3].dist_params["scale"] == pytest.approx(0.3)


def test_distribution_window_size_aware_preset_sets_row_statuses_for_all_structures(
    qapp,
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="Small",
                motif="m1",
                param_type="Both",
                param="w0",
                value=0.3,
                vary=True,
                distribution="norm",
                dist_params={"loc": 0.3, "scale": 0.1},
            ),
            DreamParameterEntry(
                structure="Large",
                motif="m2",
                param_type="Both",
                param="w1",
                value=0.7,
                vary=True,
                distribution="norm",
                dist_params={"loc": 0.7, "scale": 0.1},
            ),
            DreamParameterEntry(
                structure="Small",
                motif="m1",
                param_type="SAXS",
                param="r_eff_w0",
                value=3.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 3.0, "scale": 0.2},
            ),
            DreamParameterEntry(
                structure="Large",
                motif="m2",
                param_type="SAXS",
                param="r_eff_w1",
                value=8.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 8.0, "scale": 0.2},
            ),
            DreamParameterEntry(
                structure="",
                motif="",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.3},
            ),
        ]
    )

    scope_index = window.smart_prior_apply_scope_combo.findData("selected")
    assert scope_index >= 0
    window.smart_prior_apply_scope_combo.setCurrentIndex(scope_index)
    window.table.selectRow(0)
    window.smart_prior_preset_combo.setCurrentText(
        "Strict Small / Lenient Large"
    )
    window.apply_smart_prior_preset_button.click()

    assert window.table.cellWidget(0, 8).currentText() == "Strict"
    assert window.table.cellWidget(1, 8).currentText() == "Lenient"
    assert window.table.cellWidget(2, 8).currentText() == "Strict"
    assert window.table.cellWidget(3, 8).currentText() == "Lenient"
    assert window.table.cellWidget(4, 8).currentText() == "Proportional"


def test_distribution_window_row_status_can_override_single_structure_preset(
    qapp,
):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="Small",
                motif="m1",
                param_type="Both",
                param="w0",
                value=0.3,
                vary=True,
                distribution="norm",
                dist_params={"loc": 0.3, "scale": 0.1},
            ),
            DreamParameterEntry(
                structure="Large",
                motif="m2",
                param_type="Both",
                param="w1",
                value=0.7,
                vary=True,
                distribution="norm",
                dist_params={"loc": 0.7, "scale": 0.1},
            ),
            DreamParameterEntry(
                structure="Small",
                motif="m1",
                param_type="SAXS",
                param="r_eff_w0",
                value=3.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 3.0, "scale": 0.2},
            ),
            DreamParameterEntry(
                structure="Large",
                motif="m2",
                param_type="SAXS",
                param="r_eff_w1",
                value=8.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 8.0, "scale": 0.2},
            ),
        ]
    )

    window.smart_prior_preset_combo.setCurrentText(
        "Strict Small / Lenient Large"
    )
    window.apply_smart_prior_preset_button.click()

    small_status_combo = window.table.cellWidget(0, 8)
    very_lenient_index = small_status_combo.findData("very_lenient")
    assert very_lenient_index >= 0
    small_status_combo.setCurrentIndex(very_lenient_index)
    entries = window.current_entries()

    assert entries[0].dist_params["scale"] == pytest.approx(0.14625)
    assert entries[2].dist_params["scale"] == pytest.approx(0.2925)
    assert entries[1].dist_params["scale"] == pytest.approx(0.15)
    assert entries[3].dist_params["scale"] == pytest.approx(0.3)
    assert window.table.cellWidget(0, 8).currentText() == "Very Lenient"
    assert window.table.cellWidget(2, 8).currentText() == "Very Lenient"
    assert window.table.cellWidget(1, 8).currentText() == "Lenient"


def test_distribution_window_warns_when_effective_radius_is_set_to_vary(
    qapp,
    monkeypatch,
):
    del qapp
    monkeypatch.setattr(
        DistributionSetupWindow,
        "_session_skip_effective_radius_vary_warning",
        False,
    )
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="r_eff_w0",
                value=4.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 4.0, "scale": 0.2},
            )
        ]
    )

    warnings: list[tuple[str, str, str]] = []
    monkeypatch.setattr(
        "saxshell.saxs.ui.distribution_window.QMessageBox.exec",
        lambda dialog: warnings.append(
            (
                dialog.text(),
                dialog.informativeText(),
                (
                    dialog.checkBox().text()
                    if dialog.checkBox() is not None
                    else ""
                ),
            )
        )
        or int(QMessageBox.StandardButton.Ok),
    )

    vary_box = window.table.cellWidget(0, 5)
    assert isinstance(vary_box, QCheckBox)
    vary_box.setChecked(True)

    assert warnings
    assert "not recommended to vary effective-radius parameters" in (
        warnings[-1][0]
    )
    assert "r_eff_w0" in warnings[-1][1]
    assert (
        warnings[-1][2]
        == "Don't show this type of warning again during this session"
    )


def test_distribution_window_can_suppress_effective_radius_warning_for_session(
    qapp,
    monkeypatch,
):
    del qapp
    monkeypatch.setattr(
        DistributionSetupWindow,
        "_session_skip_effective_radius_vary_warning",
        False,
    )

    dialogs: list[tuple[str, str]] = []

    def _fake_exec(dialog):
        checkbox = dialog.checkBox()
        assert checkbox is not None
        checkbox.setChecked(True)
        dialogs.append((dialog.text(), dialog.informativeText()))
        return int(QMessageBox.StandardButton.Ok)

    monkeypatch.setattr(
        "saxshell.saxs.ui.distribution_window.QMessageBox.exec",
        _fake_exec,
    )

    first_window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="r_eff_w0",
                value=4.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 4.0, "scale": 0.2},
            )
        ]
    )
    first_vary_box = first_window.table.cellWidget(0, 5)
    assert isinstance(first_vary_box, QCheckBox)
    first_vary_box.setChecked(True)

    second_window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="SAXS",
                param="r_eff_w0",
                value=4.0,
                vary=False,
                distribution="norm",
                dist_params={"loc": 4.0, "scale": 0.2},
            )
        ]
    )
    second_vary_box = second_window.table.cellWidget(0, 5)
    assert isinstance(second_vary_box, QCheckBox)
    second_vary_box.setChecked(True)

    assert len(dialogs) == 1
    assert (
        "not recommended to vary effective-radius parameters" in dialogs[0][0]
    )
    assert "r_eff_w0" in dialogs[0][1]


def test_distribution_window_preview_defaults_to_weight_parameters(qapp):
    del qapp
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="Both",
                param="w0",
                value=0.6,
                vary=True,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 0.6, "s": 0.1},
            ),
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_B",
                param_type="Both",
                param="w1",
                value=0.4,
                vary=True,
                distribution="norm",
                dist_params={"loc": 0.4, "scale": 0.05},
            ),
            DreamParameterEntry(
                structure="",
                motif="",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
            ),
        ]
    )

    assert isinstance(window.toolbar, NavigationToolbar2QT)
    window.preview_weight_priors_button.click()

    assert window._weight_preview_window is not None
    assert window._weight_preview_window.isVisible()
    assert isinstance(
        window._weight_preview_window.toolbar, NavigationToolbar2QT
    )
    checkbox_states = {
        entry.param: checkbox.isChecked()
        for entry, checkbox in window._weight_preview_window._parameter_checkboxes
    }
    assert checkbox_states == {"w0": True, "w1": True, "scale": False}
    axis = window._weight_preview_window.figure.axes[0]
    assert axis.get_title() == "Prior distributions"
    assert axis.get_xlabel() == "Value"
    assert axis.get_ylabel() == "Density"
    plotted_labels = [line.get_label() for line in axis.get_lines()]
    assert plotted_labels == ["w0 (PbI2)", "w1 (PbI2)"]


def test_distribution_window_preview_can_toggle_non_weight_parameters(qapp):
    window = DistributionSetupWindow(
        [
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_A",
                param_type="Both",
                param="w0",
                value=0.6,
                vary=True,
                distribution="lognorm",
                dist_params={"loc": 0.0, "scale": 0.6, "s": 0.1},
            ),
            DreamParameterEntry(
                structure="PbI2",
                motif="motif_B",
                param_type="Both",
                param="w1",
                value=0.4,
                vary=True,
                distribution="norm",
                dist_params={"loc": 0.4, "scale": 0.05},
            ),
            DreamParameterEntry(
                structure="",
                motif="",
                param_type="SAXS",
                param="scale",
                value=1.0,
                vary=True,
                distribution="norm",
                dist_params={"loc": 1.0, "scale": 0.2},
            ),
        ]
    )

    window.preview_weight_priors_button.click()

    assert window._weight_preview_window is not None
    scale_checkbox = next(
        checkbox
        for entry, checkbox in window._weight_preview_window._parameter_checkboxes
        if entry.param == "scale"
    )
    scale_checkbox.setChecked(True)
    qapp.processEvents()

    axis = window._weight_preview_window.figure.axes[0]
    plotted_labels = [line.get_label() for line in axis.get_lines()]
    assert plotted_labels == ["w0 (PbI2)", "w1 (PbI2)", "scale"]


def test_template_dropdowns_use_display_names_and_tooltips(qapp):
    del qapp
    tab = ProjectSetupTab()
    basic_spec = load_template_spec("template_likelihood_monosq")
    decoupled_spec = load_template_spec(
        "template_pd_likelihood_monosq_decoupled"
    )

    tab.set_available_templates([basic_spec, decoupled_spec], basic_spec.name)

    assert tab.template_combo.itemText(0).startswith("MonoSQ Basic")
    assert (
        tab.template_combo.itemData(0, Qt.ItemDataRole.ToolTipRole)
        == basic_spec.description
    )
    assert tab.selected_template_name() == basic_spec.name


def test_project_setup_component_build_mode_defaults_and_round_trips(qapp):
    del qapp
    tab = ProjectSetupTab()
    settings = ProjectSettings(
        project_name="demo",
        project_dir="/tmp/demo",
        component_build_mode=COMPONENT_BUILD_MODE_CONTRAST,
    )

    assert tab.component_build_mode() == COMPONENT_BUILD_MODE_NO_CONTRAST
    assert (
        tab.component_build_mode_combo.currentText() == "No Contrast (Debye)"
    )

    tab.set_component_build_mode(COMPONENT_BUILD_MODE_CONTRAST)
    assert tab.component_build_mode() == COMPONENT_BUILD_MODE_CONTRAST
    assert tab.component_build_mode_combo.currentText() == "Contrast (Debye)"

    tab.set_component_build_mode(COMPONENT_BUILD_MODE_BORN_APPROXIMATION)
    assert (
        tab.component_build_mode() == COMPONENT_BUILD_MODE_BORN_APPROXIMATION
    )
    assert (
        tab.component_build_mode_combo.currentText()
        == "Born Approximation (Average)"
    )

    tab.set_project_settings(settings, [])
    assert tab.component_build_mode() == COMPONENT_BUILD_MODE_CONTRAST
    assert "Contrast (Debye)" in tab.component_build_mode_label.text()


def test_distribution_identity_and_label_include_component_build_mode(
    tmp_path,
):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    manager = SAXSProjectManager()
    no_contrast_settings = manager.load_project(project_dir)
    contrast_settings = ProjectSettings.from_dict(
        no_contrast_settings.to_dict()
    )
    contrast_settings.component_build_mode = COMPONENT_BUILD_MODE_CONTRAST

    assert project_module.distribution_id_for_settings(
        no_contrast_settings
    ) != project_module.distribution_id_for_settings(contrast_settings)
    assert "Build: No Contrast (Debye)" in (
        project_module.distribution_label_for_settings(no_contrast_settings)
    )
    assert "Build: Contrast (Debye)" in (
        project_module.distribution_label_for_settings(contrast_settings)
    )


def test_template_dropdowns_hide_deprecated_by_default_but_load_selected_deprecated(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    project_template_names = {
        str(window.project_setup_tab.template_combo.itemData(index) or "")
        for index in range(window.project_setup_tab.template_combo.count())
    }
    prefit_template_names = {
        str(window.prefit_tab.template_combo.itemData(index) or "")
        for index in range(window.prefit_tab.template_combo.count())
    }

    assert not window.project_setup_tab.show_deprecated_templates()
    assert not window.prefit_tab.show_deprecated_templates()
    assert window.project_setup_tab.selected_template_name() == (
        "template_pd_likelihood_monosq_decoupled"
    )
    assert window.prefit_tab.selected_template_name() == (
        "template_pd_likelihood_monosq_decoupled"
    )
    assert "template_pd_likelihood_monosq_decoupled" in project_template_names
    assert "template_pd_likelihood_monosq_decoupled" in prefit_template_names
    assert "template_likelihood_monosq" not in project_template_names
    assert "template_likelihood_monosq" not in prefit_template_names
    assert "template_pydream_poly_lma_hs_legacy" not in project_template_names
    assert "template_pydream_poly_lma_hs_legacy" not in prefit_template_names
    window.close()


def test_template_dropdowns_can_show_deprecated_and_stay_synced(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.project_setup_tab.show_deprecated_templates_checkbox.setChecked(
        True
    )

    project_template_names = {
        str(window.project_setup_tab.template_combo.itemData(index) or "")
        for index in range(window.project_setup_tab.template_combo.count())
    }
    prefit_template_names = {
        str(window.prefit_tab.template_combo.itemData(index) or "")
        for index in range(window.prefit_tab.template_combo.count())
    }

    assert window.project_setup_tab.show_deprecated_templates()
    assert window.prefit_tab.show_deprecated_templates()
    assert "template_likelihood_monosq" in project_template_names
    assert "template_likelihood_monosq" in prefit_template_names
    assert "template_pydream_poly_lma_hs_legacy" in project_template_names
    assert "template_pydream_poly_lma_hs_legacy" in prefit_template_names

    window.prefit_tab.show_deprecated_templates_checkbox.setChecked(False)

    project_template_names = {
        str(window.project_setup_tab.template_combo.itemData(index) or "")
        for index in range(window.project_setup_tab.template_combo.count())
    }
    prefit_template_names = {
        str(window.prefit_tab.template_combo.itemData(index) or "")
        for index in range(window.prefit_tab.template_combo.count())
    }

    assert not window.project_setup_tab.show_deprecated_templates()
    assert not window.prefit_tab.show_deprecated_templates()
    assert window.project_setup_tab.selected_template_name() == (
        "template_pd_likelihood_monosq_decoupled"
    )
    assert window.prefit_tab.selected_template_name() == (
        "template_pd_likelihood_monosq_decoupled"
    )
    assert "template_likelihood_monosq" not in project_template_names
    assert "template_likelihood_monosq" not in prefit_template_names
    assert "template_pd_likelihood_monosq_decoupled" in project_template_names
    assert "template_pd_likelihood_monosq_decoupled" in prefit_template_names
    window.close()


def test_install_model_dialog_collects_model_inputs(qapp, tmp_path):
    del qapp
    candidate_template = tmp_path / "candidate_install_model.py"
    candidate_template.write_text(
        "import numpy as np\n"
        "# model_lmfit: lmfit_model_profile\n"
        "# model_pydream: log_likelihood_candidate\n"
        "# inputs_lmfit: q, solvent_data, model_data, params\n"
        "# inputs_pydream: q_values, experimental_intensities, "
        "solvent_intensities, theoretical_intensities, params\n"
        "# param_columns: Structure, Motif, Param, Value, Vary, Min, Max\n"
        "# param: scale,1.0,True,0.0,10.0\n"
        "def lmfit_model_profile(q, solvent_data, model_data, **params):\n"
        "    del q, solvent_data\n"
        "    return params['scale'] * np.asarray(model_data[0], dtype=float)\n"
        "def log_likelihood_candidate(params):\n"
        "    del params\n"
        "    return -1.0\n",
        encoding="utf-8",
    )
    dialog = InstallModelDialog()
    dialog.model_name_edit.setText("Dialog Candidate Model")
    dialog.template_path_edit.setText(str(candidate_template))
    dialog.description_edit.setPlainText("Dialog-installed candidate.")

    request = dialog.selected_request()

    assert request == TemplateInstallRequest(
        model_name="Dialog Candidate Model",
        template_path=candidate_template.resolve(),
        model_description="Dialog-installed candidate.",
    )


def test_install_model_template_installs_and_refreshes_template_lists(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    install_dir = tmp_path / "installed_templates"
    request = TemplateInstallRequest(
        model_name="Installed Candidate Model",
        template_path=(
            Path(
                "tests/template_candidates/valid_installable_model.py"
            ).resolve()
        ),
        model_description="Installed from the Project Setup dialog.",
    )
    info_messages: list[tuple[str, str]] = []

    monkeypatch.setattr(
        window,
        "_prompt_template_install_request",
        lambda: request,
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.install_template_candidate",
        lambda template_path, **kwargs: install_template_candidate(
            template_path,
            destination_dir=install_dir,
            **kwargs,
        ),
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.list_template_specs",
        lambda **kwargs: list_template_specs(install_dir, **kwargs),
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.QMessageBox.information",
        lambda _parent, title, message: info_messages.append((title, message)),
    )

    window.project_setup_tab.install_model_button.click()

    assert (install_dir / "template_installed_candidate_model.py").is_file()
    assert (install_dir / "template_installed_candidate_model.json").is_file()
    assert any(
        window.project_setup_tab.template_combo.itemText(index)
        == "Installed Candidate Model"
        for index in range(window.project_setup_tab.template_combo.count())
    )
    assert any(
        window.prefit_tab.template_combo.itemText(index)
        == "Installed Candidate Model"
        for index in range(window.prefit_tab.template_combo.count())
    )
    assert info_messages
    assert info_messages[0][0] == "Model installed"
    assert "Installed Candidate Model" in info_messages[0][1]


def test_install_model_template_surfaces_validation_failures(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    install_dir = tmp_path / "installed_templates"
    request = TemplateInstallRequest(
        model_name="Broken Dream Callable Model",
        template_path=(
            Path(
                "tests/template_candidates/fail_missing_dream_callable_model.py"
            ).resolve()
        ),
        model_description="Expected to fail in the validation step.",
    )
    errors: list[tuple[str, str]] = []

    monkeypatch.setattr(
        window,
        "_prompt_template_install_request",
        lambda: request,
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.install_template_candidate",
        lambda template_path, **kwargs: install_template_candidate(
            template_path,
            destination_dir=install_dir,
            **kwargs,
        ),
    )
    monkeypatch.setattr(
        window,
        "_show_error",
        lambda title, message: errors.append((title, message)),
    )

    window.project_setup_tab.install_model_button.click()

    assert errors
    assert errors[0][0] == "Install model failed"
    assert "Missing callable log_likelihood_candidate" in errors[0][1]


def test_project_setup_empty_preview_message_is_wrapped(qapp):
    del qapp
    tab = ProjectSetupTab()
    tab.draw_component_plot(None)

    preview_axis = tab.component_figure.axes[0]
    text_labels = [text.get_text() for text in preview_axis.texts]

    assert any(
        "Select experimental data and build SAXS" in label
        and "averaged cluster profiles." in label
        and "\n" in label
        for label in text_labels
    )


def test_run_dream_requires_parameter_map_saved_in_session(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)
    workflow.create_default_parameter_map()
    window = SAXSMainWindow(initial_project_dir=project_dir)

    blinked = {"value": False}
    monkeypatch.setattr(
        window.dream_tab,
        "blink_edit_priors_button",
        lambda: blinked.update({"value": True}),
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.SAXSDreamWorkflow.run_bundle",
        lambda self, bundle: (_ for _ in ()).throw(
            AssertionError("run_bundle should not be called")
        ),
    )

    window.run_dream_bundle()

    assert blinked["value"]
    assert "Review the priors in Edit Priors" in (
        window.dream_tab.output_box.toPlainText()
    )


def test_run_dream_requires_written_runtime_bundle(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)
    entries = workflow.create_default_parameter_map(persist=True)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window._save_distribution_entries(entries)

    blinked = {"value": False}
    errors: list[tuple[str, str]] = []
    monkeypatch.setattr(
        window.dream_tab,
        "blink_write_bundle_button",
        lambda: blinked.update({"value": True}),
    )
    monkeypatch.setattr(
        window,
        "_show_error",
        lambda title, message: errors.append((title, message)),
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.SAXSDreamWorkflow.run_bundle",
        lambda self, bundle: (_ for _ in ()).throw(
            AssertionError("run_bundle should not be called")
        ),
    )

    window.run_dream_bundle()

    assert blinked["value"]
    assert errors == [
        (
            "Runtime Bundle not generated",
            "Runtime Bundle not generated. Click Write Runtime Bundle before running DREAM.",
        )
    ]
    assert "Runtime Bundle not generated" in (
        window.dream_tab.output_box.toPlainText()
    )


def test_preview_runtime_bundle_opens_latest_written_script(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)
    entries = workflow.create_default_parameter_map(persist=True)
    bundle = workflow.create_runtime_bundle(entries=entries)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window._last_written_dream_bundle = bundle

    opener = RuntimeBundleOpener(
        label="Fake Editor",
        stored_value="/Applications/FakeEditor.app",
        launch_target="/Applications/FakeEditor.app",
        launch_mode="mac_app",
    )
    launched: dict[str, object] = {}

    monkeypatch.setattr(
        window,
        "_available_runtime_bundle_openers",
        lambda: [opener],
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.QInputDialog.getItem",
        lambda *args, **kwargs: ("Fake Editor", True),
    )
    monkeypatch.setattr(
        window,
        "_launch_runtime_bundle_with_opener",
        lambda script_path, selected_opener: launched.update(
            {
                "path": str(script_path),
                "label": selected_opener.label,
                "stored_value": selected_opener.stored_value,
            }
        ),
    )

    window.preview_dream_runtime_bundle()

    assert launched["path"] == str(bundle.runtime_script_path)
    assert launched["label"] == "Fake Editor"
    assert "Opened DREAM runtime bundle preview" in (
        window.dream_tab.output_box.toPlainText()
    )
    reloaded_settings = SAXSProjectManager().load_project(project_dir)
    assert reloaded_settings.runtime_bundle_opener == opener.stored_value


def test_preview_runtime_bundle_reuses_saved_project_opener(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)
    entries = workflow.create_default_parameter_map(persist=True)
    bundle = workflow.create_runtime_bundle(entries=entries)

    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.runtime_bundle_opener = "/Applications/SavedEditor.app"
    manager.save_project(settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)
    window._last_written_dream_bundle = bundle
    launched: dict[str, object] = {}

    monkeypatch.setattr(
        window,
        "_available_runtime_bundle_openers",
        lambda: [
            RuntimeBundleOpener(
                label="Saved Editor",
                stored_value="/Applications/SavedEditor.app",
                launch_target="/Applications/SavedEditor.app",
                launch_mode="mac_app",
            )
        ],
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.QInputDialog.getItem",
        lambda *args, **kwargs: pytest.fail(
            "The opener chooser should not appear for a saved project opener."
        ),
    )
    monkeypatch.setattr(
        window,
        "_launch_runtime_bundle_with_opener",
        lambda script_path, selected_opener: launched.update(
            {
                "path": str(script_path),
                "label": selected_opener.label,
                "stored_value": selected_opener.stored_value,
            }
        ),
    )

    window.preview_dream_runtime_bundle()

    assert launched["path"] == str(bundle.runtime_script_path)
    assert launched["stored_value"] == "/Applications/SavedEditor.app"


def test_run_dream_shows_progress_and_popup_can_be_closed(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)
    entries = workflow.create_default_parameter_map(persist=True)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window._save_distribution_entries(entries)
    window.write_dream_bundle()

    def _fake_run_bundle(
        self,
        bundle,
        *,
        output_callback=None,
        output_interval_seconds=None,
    ):
        del output_interval_seconds
        if output_callback is not None:
            output_callback("DREAM sampler: initialization complete")
        time.sleep(0.15)
        if output_callback is not None:
            output_callback("DREAM sampler: collecting posterior samples")
        metadata = json.loads(bundle.metadata_path.read_text(encoding="utf-8"))
        active_count = len(metadata["active_parameter_entries"])
        active_values = np.asarray(
            [
                float(entry["value"])
                for entry in metadata["active_parameter_entries"]
            ],
            dtype=float,
        )
        sampled_params = []
        log_ps = []
        for chain_index in range(2):
            chain_samples = []
            chain_logps = []
            for step_index in range(4):
                adjustment = (chain_index + 1) * (step_index + 1) * 0.01
                chain_samples.append(active_values[:active_count] + adjustment)
                chain_logps.append(-5.0 + chain_index + step_index * 0.5)
            sampled_params.append(chain_samples)
            log_ps.append(chain_logps)
        np.save(
            bundle.run_dir / "dream_sampled_params.npy",
            np.asarray(sampled_params, dtype=float),
        )
        np.save(
            bundle.run_dir / "dream_log_ps.npy",
            np.asarray(log_ps, dtype=float)[..., np.newaxis],
        )
        return {
            "sampled_params_path": str(
                bundle.run_dir / "dream_sampled_params.npy"
            ),
            "log_ps_path": str(bundle.run_dir / "dream_log_ps.npy"),
        }

    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.SAXSDreamWorkflow.run_bundle",
        _fake_run_bundle,
    )

    window.run_dream_bundle()
    QApplication.processEvents()

    assert window.dream_tab.progress_bar.minimum() == 0
    assert window.dream_tab.progress_bar.maximum() == 0
    assert not window.dream_tab.run_button.isEnabled()
    assert window._dream_progress_dialog is not None
    assert window._dream_progress_dialog.isVisible()

    deadline = time.time() + 2.0
    while (
        "DREAM sampler: initialization complete"
        not in window.dream_tab.output_box.toPlainText()
        and time.time() < deadline
    ):
        QApplication.processEvents()
        time.sleep(0.02)

    assert "DREAM Runtime Output" in window.dream_tab.output_box.toPlainText()
    assert (
        "DREAM sampler: initialization complete"
        in window.dream_tab.output_box.toPlainText()
    )

    window._dream_progress_dialog.close()
    QApplication.processEvents()
    assert not window._dream_progress_dialog.isVisible()

    deadline = time.time() + 5.0
    while window._dream_task_thread is not None and time.time() < deadline:
        QApplication.processEvents()
        time.sleep(0.02)

    assert window._dream_task_thread is None
    assert window.dream_tab.run_button.isEnabled()
    assert window.dream_tab.progress_bar.maximum() == 1
    assert window.dream_tab.progress_bar.value() == 1
    assert (
        "DREAM sampler: collecting posterior samples"
        in window.dream_tab.output_box.toPlainText()
    )
    assert "DREAM run complete" in window.dream_tab.output_box.toPlainText()


def test_dream_results_loader_normalizes_singleton_logp_axis(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    bundle = _write_minimal_dream_results(project_dir)
    log_ps_path = bundle.run_dir / "dream_log_ps.npy"
    log_ps = np.load(log_ps_path)
    np.save(log_ps_path, np.asarray(log_ps, dtype=float)[..., np.newaxis])

    loader = SAXSDreamResultsLoader(bundle.run_dir, burnin_percent=0)
    summary = loader.get_summary()

    assert loader.log_ps.ndim == 2
    assert loader.log_ps.shape == (2, 4)
    assert summary.posterior_sample_count == 8


def test_dream_workflow_normalizes_stale_distribution_params_from_saved_map(
    qapp, tmp_path
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)
    workflow.parameter_map_path.write_text(
        json.dumps(
            {
                "entries": [
                    {
                        "structure": "PbI2",
                        "motif": "motif_A",
                        "param_type": "SAXS",
                        "param": "w0",
                        "value": 0.6,
                        "vary": True,
                        "distribution": "norm",
                        "dist_params": {
                            "loc": 0.6,
                            "scale": 0.2,
                            "s": 0.1,
                        },
                    }
                ]
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )

    entries = workflow.load_parameter_map()

    assert entries[0].distribution == "norm"
    assert entries[0].dist_params == {"loc": 0.6, "scale": 0.2}
    saved_project_map = json.loads(
        workflow.parameter_map_path.read_text(encoding="utf-8")
    )
    assert saved_project_map["entries"][0]["dist_params"] == {
        "loc": 0.6,
        "scale": 0.2,
    }

    bundle = workflow.create_runtime_bundle(entries=entries)
    saved_map = json.loads(
        bundle.parameter_map_path.read_text(encoding="utf-8")
    )
    assert saved_map["entries"][0]["dist_params"] == {
        "loc": 0.6,
        "scale": 0.2,
    }
    metadata = json.loads(bundle.metadata_path.read_text(encoding="utf-8"))
    assert metadata["active_parameter_entries"][0]["dist_params"] == {
        "loc": 0.6,
        "scale": 0.2,
    }


def test_dream_runtime_module_is_pickleable_for_parallel_execution(
    qapp, tmp_path
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)
    entries = workflow.create_default_parameter_map(persist=True)
    bundle = workflow.create_runtime_bundle(entries=entries)

    module_name, module, added_sys_path = workflow._load_runtime_module(bundle)
    try:
        pickled = pickle.dumps(module.active_log_likelihood)
    finally:
        workflow._unload_runtime_module(
            module_name,
            added_sys_path=added_sys_path,
        )

    assert pickled


def test_dream_runtime_module_saves_squeezed_log_ps(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)
    entries = workflow.create_default_parameter_map(persist=True)
    bundle = workflow.create_runtime_bundle(entries=entries)

    module_name, module, added_sys_path = workflow._load_runtime_module(bundle)
    try:
        monkeypatch.setattr(
            module,
            "run_dream",
            lambda **kwargs: (
                [
                    np.asarray(
                        [[0.6, 0.0, 0.05, 9.0, 0.0, 5e-4]], dtype=float
                    ),
                    np.asarray(
                        [[0.61, 0.0, 0.05, 9.1, 0.0, 5e-4]], dtype=float
                    ),
                ],
                [
                    np.asarray([[-5.0]], dtype=float),
                    np.asarray([[-4.8]], dtype=float),
                ],
            ),
        )
        module.run_sampler()
    finally:
        workflow._unload_runtime_module(
            module_name,
            added_sys_path=added_sys_path,
        )

    saved_log_ps = np.load(bundle.run_dir / "dream_log_ps.npy")
    assert saved_log_ps.shape == (2, 1)


def test_dream_runtime_module_handles_array_shaped_crossover_values(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)
    entries = workflow.create_default_parameter_map(persist=True)
    bundle = workflow.create_runtime_bundle(entries=entries)

    runtime_source = bundle.runtime_script_path.read_text(encoding="utf-8")
    assert "def _build_runtime_mp_context()" in runtime_source
    assert "mp_context=mp_context" in runtime_source

    module_name, module, added_sys_path = workflow._load_runtime_module(bundle)
    try:
        shared_vars = SimpleNamespace(
            cross_probs=mp.Array("d", [0.5, 0.5]),
            ncr_updates=mp.Array("d", [0.0, 0.0]),
            current_positions=mp.Array("d", [1.0, 2.0, 1.2, 2.4]),
            delta_m=mp.Array("d", [0.0, 0.0]),
            gamma_level_probs=mp.Array("d", [1.0]),
            ngamma_updates=mp.Array("d", [0.0]),
            delta_m_gamma=mp.Array("d", [0.0]),
        )
        monkeypatch.setattr(module, "Dream_shared_vars", shared_vars)
        module._PyDreamDream.estimate_crossover_probabilities.__globals__[
            "Dream_shared_vars"
        ] = shared_vars
        module._PyDreamDream.estimate_gamma_level_probs.__globals__[
            "Dream_shared_vars"
        ] = shared_vars

        fake_dream = SimpleNamespace(
            nCR=2,
            nchains=2,
            CR_values=np.asarray([0.5, 1.0], dtype=float),
            CR_probabilities=np.asarray([0.5, 0.5], dtype=float),
            ngamma=1,
            gamma_level_values=np.asarray([1], dtype=int),
        )

        crossover_value = module._PyDreamDream.set_CR(
            fake_dream,
            np.asarray([1.0, 0.0], dtype=float),
            np.asarray([0.5, 1.0], dtype=float),
        )
        assert float(crossover_value) == pytest.approx(0.5)

        cross_probs = module._PyDreamDream.estimate_crossover_probabilities(
            fake_dream,
            2,
            np.asarray([1.0, 2.0], dtype=float),
            np.asarray([1.1, 2.3], dtype=float),
            np.asarray([0.5], dtype=float),
        )
        assert np.asarray(cross_probs, dtype=float).shape == (2,)
        assert list(shared_vars.ncr_updates[:]) == [1.0, 0.0]

        gamma_probs = module._PyDreamDream.estimate_gamma_level_probs(
            fake_dream,
            2,
            np.asarray([1.0, 2.0], dtype=float),
            np.asarray([1.1, 2.3], dtype=float),
            np.asarray([1], dtype=int),
        )
        assert np.asarray(gamma_probs, dtype=float).shape == (1,)
        assert list(shared_vars.ngamma_updates[:]) == [1.0]
    finally:
        workflow._unload_runtime_module(
            module_name,
            added_sys_path=added_sys_path,
        )


def test_save_dream_settings_creates_named_preset_and_restores_active_state(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.dream_tab.chains_spin.setValue(6)
    window.dream_tab.iterations_spin.setValue(2500)
    window.dream_tab.posterior_filter_combo.setCurrentIndex(1)
    window.dream_tab.posterior_top_percent_spin.setValue(12.5)
    window.dream_tab.credible_interval_low_spin.setValue(10.0)
    window.dream_tab.credible_interval_high_spin.setValue(90.0)
    window.dream_tab.violin_sample_source_combo.setCurrentIndex(1)

    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.QInputDialog.getText",
        lambda *args, **kwargs: ("Preset A", True),
    )

    window.save_dream_settings()

    preset_index = window.dream_tab.settings_preset_combo.findText("Preset A")
    assert preset_index >= 0
    assert (
        window.dream_tab.settings_preset_combo.currentText()
        == window.dream_tab.ACTIVE_SETTINGS_LABEL
    )

    window.dream_tab.chains_spin.setValue(11)
    window.dream_tab.iterations_spin.setValue(3333)
    window.dream_tab.posterior_filter_combo.setCurrentIndex(2)
    window.dream_tab.posterior_top_n_spin.setValue(7)
    window.dream_tab.violin_sample_source_combo.setCurrentIndex(0)

    window.dream_tab.settings_preset_combo.setCurrentIndex(preset_index)
    QApplication.processEvents()

    assert window.dream_tab.chains_spin.value() == 6
    assert window.dream_tab.iterations_spin.value() == 2500
    assert (
        window.dream_tab.posterior_filter_combo.currentData()
        == "top_percent_logp"
    )
    assert window.dream_tab.posterior_top_percent_spin.value() == 12.5
    assert window.dream_tab.credible_interval_low_spin.value() == 10.0
    assert window.dream_tab.credible_interval_high_spin.value() == 90.0
    assert (
        window.dream_tab.violin_sample_source_combo.currentData()
        == "map_chain_only"
    )

    active_index = window.dream_tab.settings_preset_combo.findText(
        window.dream_tab.ACTIVE_SETTINGS_LABEL
    )
    window.dream_tab.settings_preset_combo.setCurrentIndex(active_index)
    QApplication.processEvents()

    assert window.dream_tab.chains_spin.value() == 11
    assert window.dream_tab.iterations_spin.value() == 3333
    assert (
        window.dream_tab.posterior_filter_combo.currentData() == "top_n_logp"
    )
    assert window.dream_tab.posterior_top_n_spin.value() == 7
    assert (
        window.dream_tab.violin_sample_source_combo.currentData()
        == "filtered_posterior"
    )


def test_dream_posterior_filter_controls_keep_default_thresholds_editable(
    qapp, tmp_path
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.dream_tab.posterior_filter_combo.setCurrentIndex(1)
    QApplication.processEvents()

    assert window.dream_tab.posterior_top_percent_spin.isEnabled()
    assert window.dream_tab.posterior_top_n_spin.isEnabled()

    window.dream_tab.posterior_filter_combo.setCurrentIndex(2)
    QApplication.processEvents()

    assert window.dream_tab.posterior_top_percent_spin.isEnabled()
    assert window.dream_tab.posterior_top_n_spin.isEnabled()


def test_load_latest_dream_results_updates_both_plot_panels(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.dream_tab.bestfit_method_combo.setCurrentIndex(1)
    window.dream_tab.violin_mode_combo.setCurrentIndex(2)
    window.dream_tab.posterior_filter_combo.setCurrentIndex(2)
    window.dream_tab.posterior_top_n_spin.setValue(1)
    window.dream_tab.credible_interval_low_spin.setValue(5.0)
    window.dream_tab.credible_interval_high_spin.setValue(95.0)
    window.dream_tab.violin_sample_source_combo.setCurrentIndex(1)

    window.load_latest_results()

    assert "Best-fit method: chain_mean" in (
        window.dream_tab.output_box.toPlainText()
    )
    assert "Posterior filter: top_n_logp" in (
        window.dream_tab.output_box.toPlainText()
    )
    assert "Posterior samples kept: 1" in (
        window.dream_tab.output_box.toPlainText()
    )
    assert "Violin data mode: weights_only" in (
        window.dream_tab.output_box.toPlainText()
    )
    assert "Violin sample source: map_chain_only" in (
        window.dream_tab.output_box.toPlainText()
    )
    assert "p5=" in window.dream_tab.output_box.toPlainText()
    assert "p95=" in window.dream_tab.output_box.toPlainText()
    assert (
        window.dream_tab.model_figure.axes[0]
        .get_title()
        .startswith("DREAM refinement:")
    )
    metric_text = "\n".join(
        text.get_text() for text in window.dream_tab.model_figure.axes[0].texts
    )
    assert "RMSE:" in metric_text
    assert "Mean |res|:" in metric_text
    assert "R²:" in metric_text
    assert (
        window.dream_tab.violin_figure.axes[0].get_title()
        == "Posterior parameter distributions"
    )
    tick_labels = [
        label.get_text()
        for label in window.dream_tab.violin_figure.axes[0].get_xticklabels()
    ]
    assert "w0 (A)" in tick_labels


def test_dream_analysis_saved_run_dropdown_loads_selected_run_state(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)

    older_entries = workflow.create_default_parameter_map()
    older_entries[0] = DreamParameterEntry(
        structure=older_entries[0].structure,
        motif=older_entries[0].motif,
        param_type=older_entries[0].param_type,
        param=older_entries[0].param,
        value=0.11,
        vary=older_entries[0].vary,
        distribution=older_entries[0].distribution,
        dist_params=dict(older_entries[0].dist_params),
        smart_preset_status=older_entries[0].smart_preset_status,
    )
    older_settings = DreamRunSettings(
        nchains=3,
        niterations=1234,
        burnin_percent=7,
        model_name="older_model",
        run_label="older",
    )
    older_bundle = _write_minimal_dream_results(
        project_dir,
        settings=older_settings,
        entries=older_entries,
    )

    newer_entries = workflow.create_default_parameter_map()
    newer_entries[0] = DreamParameterEntry(
        structure=newer_entries[0].structure,
        motif=newer_entries[0].motif,
        param_type=newer_entries[0].param_type,
        param=newer_entries[0].param,
        value=0.77,
        vary=newer_entries[0].vary,
        distribution=newer_entries[0].distribution,
        dist_params=dict(newer_entries[0].dist_params),
        smart_preset_status=newer_entries[0].smart_preset_status,
    )
    newer_settings = DreamRunSettings(
        nchains=8,
        niterations=4321,
        burnin_percent=22,
        model_name="newer_model",
        run_label="newer",
    )
    newer_bundle = _write_minimal_dream_results(
        project_dir,
        settings=newer_settings,
        entries=newer_entries,
    )

    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert window.dream_tab.saved_runs_combo.count() == 2
    assert (
        Path(window.dream_tab.saved_runs_combo.currentData()).resolve()
        == newer_bundle.run_dir.resolve()
    )

    older_index = window.dream_tab.saved_runs_combo.findData(
        str(older_bundle.run_dir)
    )
    assert older_index >= 0
    window.dream_tab.saved_runs_combo.setCurrentIndex(older_index)
    QApplication.processEvents()

    window.load_selected_results()

    loaded_settings = load_dream_settings(older_bundle.settings_path)
    assert window._last_results_loader is not None
    assert (
        window._last_results_loader.run_dir == older_bundle.run_dir.resolve()
    )
    assert window.dream_tab.chains_spin.value() == loaded_settings.nchains
    assert (
        window.dream_tab.iterations_spin.value() == loaded_settings.niterations
    )
    assert (
        window.dream_tab.burnin_spin.value() == loaded_settings.burnin_percent
    )
    assert float(
        window.dream_tab.parameter_map_table.item(0, 4).text()
    ) == pytest.approx(older_entries[0].value)
    assert (
        str(older_bundle.run_dir) in window.dream_tab.output_box.toPlainText()
    )


def test_dream_analysis_saved_run_preview_is_non_destructive(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)
    workflow = SAXSDreamWorkflow(project_dir)

    preview_entries = workflow.create_default_parameter_map()
    preview_entries[0] = DreamParameterEntry(
        structure=preview_entries[0].structure,
        motif=preview_entries[0].motif,
        param_type=preview_entries[0].param_type,
        param=preview_entries[0].param,
        value=0.42,
        vary=preview_entries[0].vary,
        distribution=preview_entries[0].distribution,
        dist_params=dict(preview_entries[0].dist_params),
        smart_preset_status=preview_entries[0].smart_preset_status,
    )
    preview_settings = DreamRunSettings(
        nchains=7,
        niterations=2468,
        burnin_percent=13,
        parallel=False,
        adapt_crossover=False,
        restart=True,
        verbose=False,
        history_file="saved_history.npy",
        model_name="preview_model",
        run_label="preview",
        search_filter_preset="less_aggressive",
        bestfit_method="median",
        posterior_filter_mode="top_n_logp",
        posterior_top_percent=8.0,
        posterior_top_n=321,
        stoichiometry_target_elements_text="Pb, I",
        stoichiometry_target_ratio_text="1:2",
        stoichiometry_filter_enabled=True,
        stoichiometry_tolerance_percent=3.5,
    )
    preview_bundle = _write_minimal_dream_results(
        project_dir,
        settings=preview_settings,
        entries=preview_entries,
    )

    window = SAXSMainWindow(initial_project_dir=project_dir)
    original_chains = window.dream_tab.chains_spin.value()
    original_iterations = window.dream_tab.iterations_spin.value()

    preview_index = window.dream_tab.saved_runs_combo.findData(
        str(preview_bundle.run_dir)
    )
    assert preview_index >= 0
    window.dream_tab.saved_runs_combo.setCurrentIndex(preview_index)
    QApplication.processEvents()

    window.preview_selected_dream_run()
    QApplication.processEvents()

    dialog = window._dream_saved_run_preview_dialog
    assert dialog is not None
    assert dialog.isVisible()
    assert window._last_results_loader is None
    assert window.dream_tab.chains_spin.value() == original_chains
    assert window.dream_tab.iterations_spin.value() == original_iterations

    preview_text = dialog.summary_box.toPlainText()
    assert "preview_model" in preview_text
    assert "Chains: 7" in preview_text
    assert "Iterations: 2468" in preview_text
    assert "Best-fit method: Median" in preview_text
    assert "Saved prior parameter map:" in preview_text

    assert dialog.parameter_map_table.rowCount() == len(preview_entries)
    assert (
        dialog.parameter_map_table.item(0, 3).text()
        == preview_entries[0].param
    )
    assert float(
        dialog.parameter_map_table.item(0, 4).text()
    ) == pytest.approx(preview_entries[0].value)

    dialog.close()
    window.close()


def test_dream_model_metrics_box_updates_with_bestfit_method(qapp, tmp_path):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.load_latest_results()
    axis = window.dream_tab.model_figure.axes[0]
    first_metrics = "\n".join(text.get_text() for text in axis.texts)

    window.dream_tab.bestfit_method_combo.setCurrentIndex(2)
    _apply_dream_filter_changes(window, qapp)

    axis = window.dream_tab.model_figure.axes[0]
    second_metrics = "\n".join(text.get_text() for text in axis.texts)

    assert "RMSE:" in second_metrics
    assert "Mean |res|:" in second_metrics
    assert "R²:" in second_metrics
    assert first_metrics != second_metrics


def test_dream_model_plot_includes_residual_subplot(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.load_latest_results()

    assert len(window.dream_tab.model_figure.axes) == 2
    top_axis = window.dream_tab.model_figure.axes[0]
    residual_axis = window.dream_tab.model_figure.axes[1]
    assert top_axis.get_title().startswith("DREAM refinement:")
    assert residual_axis.get_ylabel() == "Residual"
    assert residual_axis.get_xlabel() == "q (Å⁻¹)"
    assert residual_axis.get_xscale() == top_axis.get_xscale()

    residual_line = residual_axis.get_lines()[-1]
    plot_data = window.dream_tab._current_model_plot_data
    assert plot_data is not None
    expected = np.asarray(
        plot_data.model_intensities - plot_data.experimental_intensities,
        dtype=float,
    )
    assert np.allclose(
        np.asarray(residual_line.get_ydata(), dtype=float),
        expected,
    )
    window.close()


def test_dream_model_plot_redraw_on_log_x_avoids_nonpositive_xlim_warning(
    qapp,
    tmp_path,
):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        window.load_latest_results()
        window.dream_tab.model_log_x_checkbox.setChecked(True)
        qapp.processEvents()
        window.dream_tab.bestfit_method_combo.setCurrentIndex(1)
        _apply_dream_filter_changes(window, qapp)

    warning_messages = [str(item.message) for item in caught]
    assert not any(
        "Attempt to set non-positive xlim on a log-scaled axis" in message
        for message in warning_messages
    )
    window.close()


def test_prefit_model_metrics_box_updates_with_model_changes(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    axis = window.prefit_tab.figure.axes[0]
    first_metrics = "\n".join(text.get_text() for text in axis.texts)

    assert "RMSE:" in first_metrics
    assert "Mean |res|:" in first_metrics
    assert "R²:" in first_metrics

    window.prefit_tab.set_parameter_row("scale", value=1e-3)
    window.update_prefit_model()

    axis = window.prefit_tab.figure.axes[0]
    second_metrics = "\n".join(text.get_text() for text in axis.texts)

    assert "RMSE:" in second_metrics
    assert "Mean |res|:" in second_metrics
    assert "R²:" in second_metrics
    assert first_metrics != second_metrics
    window.close()


def test_dream_violin_scale_modes_and_palette_controls(qapp, tmp_path):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    palette_index = window.dream_tab.violin_palette_combo.findData("plasma")
    window.load_latest_results()
    window.dream_tab.violin_value_scale_combo.setCurrentIndex(1)
    window.dream_tab.violin_palette_combo.setCurrentIndex(palette_index)
    window.dream_tab._configure_plot_color_button(
        window.dream_tab.violin_point_color_button,
        "tab:blue",
        label="Point",
    )
    window.dream_tab.visualization_settings_changed.emit()
    _apply_dream_filter_changes(window, qapp)

    axis = window.dream_tab.violin_figure.axes[0]
    tick_labels = [label.get_text() for label in axis.get_xticklabels()]
    assert tick_labels == ["w0 (A)"]
    assert axis.get_ylabel() == "Weight fraction"
    assert axis.get_title() == "Posterior weight distributions"
    assert axis.get_ylim() == pytest.approx((0.0, 1.0))
    body = next(
        collection
        for collection in axis.collections
        if isinstance(collection, PolyCollection)
    )
    assert to_hex(body.get_facecolor()[0], keep_alpha=False) == to_hex(
        colormaps["plasma"](0.72),
        keep_alpha=False,
    )
    assert to_hex(
        axis.collections[-1].get_facecolor()[0], keep_alpha=False
    ) == to_hex(
        "tab:blue",
        keep_alpha=False,
    )

    window.dream_tab.violin_value_scale_combo.setCurrentIndex(2)
    _apply_dream_filter_changes(window, qapp)

    axis = window.dream_tab.violin_figure.axes[0]
    assert axis.get_ylabel() == "Normalized parameter value"
    assert axis.get_title() == "Posterior parameter distributions (normalized)"
    assert axis.get_ylim() == pytest.approx((0.0, 1.0))
    normalized_labels = [
        label.get_text()
        for label in axis.get_xticklabels()
        if label.get_text()
    ]
    assert "w0 (A)" in normalized_labels
    assert "solv_w" in normalized_labels


def test_dream_violin_custom_color_controls_apply_to_plot(
    qapp, tmp_path, monkeypatch
):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    chosen_colors = iter(
        [
            QColor("#123456"),
            QColor("#fedcba"),
            QColor("#654321"),
            QColor("#abcdef"),
            QColor("#111111"),
        ]
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.dream_tab.QColorDialog.getColor",
        lambda *args, **kwargs: next(chosen_colors),
    )

    palette_index = window.dream_tab.violin_palette_combo.findData(
        "custom_solid"
    )
    window.load_latest_results()
    window.dream_tab.violin_palette_combo.setCurrentIndex(palette_index)
    window.dream_tab._choose_violin_custom_color()
    window.dream_tab._choose_violin_point_color()
    window.dream_tab._choose_interval_color()
    window.dream_tab._choose_median_color()
    window.dream_tab._choose_outline_color()
    window.dream_tab.violin_outline_width_spin.setValue(1.7)
    _apply_dream_filter_changes(window, qapp)

    axis = window.dream_tab.violin_figure.axes[0]
    body = next(
        collection
        for collection in axis.collections
        if isinstance(collection, PolyCollection)
    )
    assert to_hex(body.get_facecolor()[0], keep_alpha=False) == "#123456"
    assert to_hex(body.get_edgecolor()[0], keep_alpha=False) == "#111111"
    assert body.get_linewidths()[0] == pytest.approx(1.7)
    assert (
        to_hex(
            axis.collections[-1].get_facecolor()[0],
            keep_alpha=False,
        )
        == "#fedcba"
    )
    line_colors = [
        to_hex(color, keep_alpha=False)
        for collection in axis.collections
        if isinstance(collection, LineCollection)
        for color in collection.get_colors()
    ]
    assert "#654321" in line_colors
    assert "#abcdef" in line_colors


def test_dream_violin_custom_color_picker_switches_palette_and_updates_plot(
    qapp,
    tmp_path,
    monkeypatch,
):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    monkeypatch.setattr(
        "saxshell.saxs.ui.dream_tab.QColorDialog.getColor",
        lambda *args, **kwargs: QColor("#224466"),
    )

    window.load_latest_results()
    assert window.dream_tab.violin_palette_combo.currentData() == "Blues"

    window.dream_tab.violin_custom_color_button.click()
    _apply_dream_filter_changes(window, qapp)

    assert (
        window.dream_tab.violin_palette_combo.currentData() == "custom_solid"
    )
    assert window.dream_tab.selected_violin_custom_color() == "#224466"
    axis = window.dream_tab.violin_figure.axes[0]
    body = next(
        collection
        for collection in axis.collections
        if isinstance(collection, PolyCollection)
    )
    assert to_hex(body.get_facecolor()[0], keep_alpha=False) == "#224466"


def test_dream_default_violin_palette_starts_with_higher_contrast_color(
    qapp,
):
    del qapp
    window = SAXSMainWindow()

    colors = window.dream_tab._violin_body_colors(4)
    first_color = colors[0]
    expected = colormaps.get_cmap("Blues")(0.35)

    assert to_hex(first_color, keep_alpha=False) == to_hex(
        expected,
        keep_alpha=False,
    )


def test_dream_tab_limits_violin_display_sample_count(qapp):
    del qapp
    samples = np.arange(12_000, dtype=float).reshape(6_000, 2)

    limited = DreamTab._display_violin_samples(samples)

    assert limited.shape == (DreamTab.MAX_VIOLIN_PLOT_SAMPLES, 2)
    assert np.allclose(limited[0], samples[0])
    assert np.allclose(limited[-1], samples[-1])


def test_dream_results_loader_can_order_weight_violin_labels_by_structure(
    tmp_path,
):
    run_dir = _write_weight_order_dream_results(tmp_path)
    loader = SAXSDreamResultsLoader(run_dir, burnin_percent=0)

    weight_index_plot = loader.build_violin_data(
        mode="weights_only",
        weight_order="weight_index",
    )
    structure_order_plot = loader.build_violin_data(
        mode="weights_only",
        weight_order="structure_order",
    )

    assert weight_index_plot.parameter_names == ["w2", "w0", "w1"]
    assert weight_index_plot.display_names == [
        "w2 (PbI2O)",
        "w0 (I2)",
        "w1 (Pb2)",
    ]
    assert structure_order_plot.parameter_names == ["w0", "w1", "w2"]
    assert structure_order_plot.display_names == [
        "w0 (I2)",
        "w1 (Pb2)",
        "w2 (PbI2O)",
    ]


def test_dream_results_loader_splits_radius_and_additional_violin_modes(
    tmp_path,
):
    run_dir = _write_violin_mode_split_dream_results(tmp_path)
    loader = SAXSDreamResultsLoader(run_dir, burnin_percent=0)

    radius_plot = loader.build_violin_data(mode="effective_radii_only")
    additional_plot = loader.build_violin_data(
        mode="additional_parameters_only"
    )
    fit_plot = loader.build_violin_data(mode="fit_parameters")

    assert radius_plot.parameter_names == [
        "r_eff_w0",
        "a_eff_w1",
        "b_eff_w1",
        "c_eff_w1",
    ]
    assert additional_plot.parameter_names == [
        "scale",
        "offset",
        "phi_int",
    ]
    assert fit_plot.parameter_names == [
        "r_eff_w0",
        "a_eff_w1",
        "b_eff_w1",
        "c_eff_w1",
        "scale",
        "offset",
        "phi_int",
    ]


def test_effective_dream_violin_mode_honors_new_value_scale_overrides():
    radius_settings = DreamRunSettings(
        violin_parameter_mode="fit_parameters",
        violin_value_scale_mode="effective_radii_only",
    )
    additional_settings = DreamRunSettings(
        violin_parameter_mode="weights_only",
        violin_value_scale_mode="additional_parameters_only",
    )

    assert (
        SAXSMainWindow._effective_dream_violin_mode(radius_settings)
        == "effective_radii_only"
    )
    assert (
        SAXSMainWindow._effective_dream_violin_mode(additional_settings)
        == "additional_parameters_only"
    )


def test_dream_results_loader_filters_posterior_samples(tmp_path):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    bundle = _write_minimal_dream_results(project_dir)
    loader = SAXSDreamResultsLoader(bundle.run_dir, burnin_percent=0)

    summary = loader.get_summary(
        bestfit_method="median",
        posterior_filter_mode="top_n_logp",
        posterior_top_n=1,
        credible_interval_low=10.0,
        credible_interval_high=90.0,
    )
    violin_plot = loader.build_violin_data(
        mode="varying_parameters",
        posterior_filter_mode="top_percent_logp",
        posterior_top_percent=25.0,
        sample_source="map_chain_only",
    )

    assert summary.posterior_filter_mode == "top_n_logp"
    assert summary.posterior_sample_count == 1
    assert summary.credible_interval_low == 10.0
    assert summary.credible_interval_high == 90.0
    assert np.allclose(
        summary.interval_low_values,
        summary.interval_high_values,
    )
    assert np.allclose(
        summary.bestfit_params,
        summary.interval_low_values,
    )
    assert violin_plot.sample_source == "map_chain_only"
    assert violin_plot.sample_count == 2


def test_dream_results_loader_can_filter_posterior_by_stoichiometry(tmp_path):
    run_dir = _write_stoichiometry_filter_dream_results(tmp_path)
    loader = SAXSDreamResultsLoader(run_dir, burnin_percent=0)

    summary = loader.get_summary(
        bestfit_method="map",
        stoichiometry_target_elements_text="Pb, I",
        stoichiometry_target_ratio_text="1:2",
        stoichiometry_filter_enabled=True,
        stoichiometry_tolerance_percent=2.0,
    )

    assert summary.posterior_candidate_sample_count == 1
    assert summary.posterior_sample_count == 1
    assert summary.stoichiometry_target is not None
    assert summary.stoichiometry_evaluation is not None
    assert summary.stoichiometry_evaluation.is_valid is True
    assert (
        summary.stoichiometry_evaluation.max_deviation_percent
        == pytest.approx(0.0)
    )
    assert np.allclose(summary.bestfit_params[:2], [1.0, 0.0])


def test_dream_filter_changes_wait_for_apply_button(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    bundle = _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window._load_dream_results_from_run_dir(bundle.run_dir)
    QApplication.processEvents()

    assert "Posterior samples kept: 8" in window.dream_tab._summary_text
    assert not window.dream_tab.apply_filter_button.isEnabled()

    filter_index = window.dream_tab.posterior_filter_combo.findData(
        "top_n_logp"
    )
    assert filter_index >= 0
    window.dream_tab.posterior_filter_combo.setCurrentIndex(filter_index)
    window.dream_tab.posterior_top_n_spin.setValue(1)
    QApplication.processEvents()

    assert "Posterior samples kept: 8" in window.dream_tab._summary_text
    assert window.dream_tab.apply_filter_button.isEnabled()
    assert "Pending changes are not applied yet" in (
        window.dream_tab.filter_status_box.toPlainText()
    )

    window.dream_tab.apply_filter_button.click()
    QApplication.processEvents()

    assert "Posterior samples kept: 1" in window.dream_tab._summary_text
    assert not window.dream_tab.apply_filter_button.isEnabled()


def test_prefit_stoichiometry_status_updates_with_weight_edits(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    entries = [
        PrefitParameterEntry(
            structure="PbI2",
            motif="m1",
            name="w0",
            value=1.0,
            vary=True,
            minimum=0.0,
            maximum=2.0,
            category="weight",
        ),
        PrefitParameterEntry(
            structure="Pb2I5",
            motif="m2",
            name="w1",
            value=0.0,
            vary=True,
            minimum=0.0,
            maximum=2.0,
            category="weight",
        ),
        PrefitParameterEntry(
            structure="",
            motif="",
            name="scale",
            value=1.0,
            vary=False,
            minimum=0.0,
            maximum=10.0,
            category="fit",
        ),
    ]
    window.prefit_tab.populate_parameter_table(entries)

    window.dream_tab.stoichiometry_elements_edit.setText("Pb, I")
    window.dream_tab.stoichiometry_ratio_edit.setText("1:2")
    QApplication.processEvents()

    before = window.prefit_tab.stoichiometry_status_label.text()
    assert "Target: Pb:I = 1:2" in before
    assert "Deviation from target: 0%" in before

    weight_row = window.prefit_tab.find_parameter_row("w1")
    weight_item = window.prefit_tab.parameter_table.item(weight_row, 3)
    assert weight_item is not None
    weight_item.setText("1.0")
    QApplication.processEvents()

    after = window.prefit_tab.stoichiometry_status_label.text()
    assert after != before
    assert "Observed ratio:" in after
    assert "Deviation from target:" in after


def test_dream_results_loader_reuses_cached_plot_data_across_interval_changes(
    tmp_path,
):
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    bundle = _write_minimal_dream_results(project_dir)
    loader = SAXSDreamResultsLoader(bundle.run_dir, burnin_percent=0)

    first_model = loader.build_model_fit_data(
        bestfit_method="median",
        credible_interval_low=5.0,
        credible_interval_high=95.0,
    )
    second_model = loader.build_model_fit_data(
        bestfit_method="median",
        credible_interval_low=10.0,
        credible_interval_high=90.0,
    )
    first_violin = loader.build_violin_data(
        mode="weights_only",
        credible_interval_low=5.0,
        credible_interval_high=95.0,
    )
    second_violin = loader.build_violin_data(
        mode="weights_only",
        credible_interval_low=10.0,
        credible_interval_high=90.0,
    )

    assert first_model is second_model
    assert first_violin is second_violin


def test_dream_runtime_bundle_carries_selected_solvent_trace(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    prefit = SAXSPrefitWorkflow(project_dir)
    prefit.save_fit(prefit.parameter_entries)

    solvent_q = np.linspace(0.05, 0.3, 8)
    solvent_intensity = np.linspace(3.0, 4.4, 8)
    solvent_path = tmp_path / "solvent_reference_trace.dat"
    np.savetxt(solvent_path, np.column_stack([solvent_q, solvent_intensity]))

    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.solvent_data_path = str(solvent_path)
    settings.copied_solvent_data_file = None
    manager.save_project(settings)

    workflow = SAXSDreamWorkflow(project_dir)
    entries = workflow.create_default_parameter_map(persist=True)
    bundle = workflow.create_runtime_bundle(entries=entries)
    metadata = json.loads(bundle.metadata_path.read_text(encoding="utf-8"))

    assert np.allclose(metadata["solvent_intensities"], solvent_intensity)


def test_dream_plot_data_exports_save_into_exported_results_data(
    qapp, tmp_path
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    class _AcceptedDialog:
        def __init__(self, *args, **kwargs):
            del args, kwargs
            self.selected_options = SimpleNamespace(
                output_dir=paths.exported_data_dir,
                base_name="dream_violin_export_test",
                save_csv=True,
                save_pkl=True,
            )

        def exec(self):
            return QDialog.DialogCode.Accepted

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.DreamViolinExportDialog",
        _AcceptedDialog,
    )

    window.load_latest_results()
    window.dream_tab.violin_mode_combo.setCurrentIndex(2)
    window.save_dream_model_fit()
    window.save_dream_violin_data()
    monkeypatch.undo()

    model_exports = sorted(
        paths.exported_data_dir.glob("dream_model_fit_*.csv")
    )
    violin_csv_exports = sorted(
        paths.exported_data_dir.glob("dream_violin_*.csv")
    )
    violin_pkl_exports = sorted(
        paths.exported_data_dir.glob("dream_violin_*.pkl")
    )
    dialog_csv_export = (
        paths.exported_data_dir / "dream_violin_export_test.csv"
    )
    dialog_pkl_export = (
        paths.exported_data_dir / "dream_violin_export_test.pkl"
    )
    dialog_metadata_export = (
        paths.exported_data_dir / "dream_violin_export_test.metadata.json"
    )
    dialog_report_export = (
        paths.exported_data_dir / "dream_violin_export_test.report.txt"
    )

    assert model_exports
    assert violin_csv_exports
    assert violin_pkl_exports
    assert dialog_csv_export.is_file()
    assert dialog_pkl_export.is_file()
    assert dialog_metadata_export.is_file()
    assert dialog_report_export.is_file()
    assert (
        "q,experimental_intensity,model_intensity,"
        "solvent_contribution,structure_factor"
    ) in model_exports[-1].read_text(encoding="utf-8")
    model_metadata = json.loads(
        model_exports[-1]
        .with_name(f"{model_exports[-1].stem}.metadata.json")
        .read_text(encoding="utf-8")
    )
    model_report = model_exports[-1].with_name(
        f"{model_exports[-1].stem}.report.txt"
    )
    assert model_report.is_file()
    assert model_metadata["export_kind"] == "dream_model_fit"
    assert model_metadata["model_fit"]["includes_structure_factor"] is True
    assert model_metadata["model_fit"]["fit_metrics"]["rmse"] >= 0.0
    assert "Model fit metrics:" in model_report.read_text(encoding="utf-8")
    assert "w0 (A)" in dialog_csv_export.read_text(encoding="utf-8")
    violin_payload = pickle.loads(dialog_pkl_export.read_bytes())
    violin_metadata = json.loads(
        dialog_metadata_export.read_text(encoding="utf-8")
    )
    assert violin_payload["violin_plot"]["display_names"] == ["w0 (A)"]
    assert violin_payload["screening_metrics"]["posterior_filter_mode"]
    assert violin_metadata["export_kind"] == "dream_violin"
    assert violin_metadata["screening_metrics"]["posterior_filter_mode"]
    assert violin_payload["plot_payload"]["ylabel"] == "Parameter value"
    assert np.asarray(violin_payload["plot_payload"]["samples"]).shape[1] == 1
    assert "Posterior violin data:" in dialog_report_export.read_text(
        encoding="utf-8"
    )


def test_dream_model_report_export_builds_context_and_writes_pptx(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.load_latest_results()

    captured: dict[str, object] = {}
    errors: list[tuple[str, str]] = []
    monkeypatch.setattr(
        window,
        "_show_error",
        lambda title, message: errors.append((title, message)),
    )

    def _fake_export(context, *, progress_callback=None):
        captured["context"] = context
        if progress_callback is not None:
            progress_callback(1, 3, "Rendering report plots...")
            progress_callback(2, 3, "Building report slides...")
            progress_callback(3, 3, "Saving PowerPoint report...")
        context.asset_dir.mkdir(parents=True, exist_ok=True)
        context.output_path.write_bytes(b"pptx")
        manifest_path = context.asset_dir / "report_manifest.json"
        manifest_path.write_text("{}\n", encoding="utf-8")
        return SimpleNamespace(
            report_path=context.output_path,
            manifest_path=manifest_path,
            figure_paths=(),
        )

    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.export_dream_model_report_pptx",
        _fake_export,
    )

    window.export_dream_model_report()

    assert not errors
    report_exports = sorted(
        paths.reports_dir.glob("dream_model_report_*.pptx")
    )
    assert report_exports
    assert report_exports[-1].read_bytes() == b"pptx"
    assert "Generating DREAM model report PowerPoint. Please wait..." in (
        window.dream_tab.output_box.toPlainText()
    )
    assert (
        window.dream_tab.progress_label.text()
        == "DREAM model report exported."
    )
    assert window.dream_tab.progress_bar.maximum() == 3
    assert window.dream_tab.progress_bar.value() == 3
    context = captured["context"]
    assert context.project_name == "saxs_project"
    assert context.project_dir == project_dir.resolve()
    assert context.user_q_range_text == "0.05 to 0.3"
    assert context.component_plot_without_solvent is not None
    assert context.component_plot_with_solvent is not None
    assert len(context.prior_histograms) == 4
    assert context.prefit_parameter_entries
    assert context.dream_parameter_map_entries
    assert context.dream_summary.posterior_sample_count > 0
    assert context.dream_filter_views
    assert "Posterior filter:" in "\n".join(context.output_summary_lines)
    assert (context.asset_dir / "report_manifest.json").is_file()


def test_dream_model_report_export_writes_real_pptx(qapp, tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.load_latest_results()

    context = window._build_dream_model_report_context(
        settings=window.dream_tab.settings_payload(),
        output_path=paths.reports_dir / "dream_model_report_test.pptx",
        asset_dir=paths.reports_dir / "dream_model_report_test_assets",
    )

    progress_updates: list[tuple[int, int, str]] = []

    result = export_dream_model_report_pptx(
        context,
        progress_callback=lambda processed, total, message: progress_updates.append(
            (processed, total, message)
        ),
    )

    assert result.report_path.is_file()
    assert result.report_path.stat().st_size > 0
    assert result.manifest_path.is_file()
    assert result.figure_paths
    figure_names = {path.name for path in result.figure_paths}
    assert "dream_filter_violin_comparison.png" in figure_names
    assert "dream_filter_violin_comparison_weights.png" in figure_names
    assert "dream_filter_violin_comparison_effective_radii.png" in figure_names
    assert "prefit_model_without_solvent.png" in figure_names
    if (
        context.prefit_evaluation is not None
        and context.prefit_evaluation.solvent_contribution is not None
        and np.any(
            np.isfinite(
                np.asarray(
                    context.prefit_evaluation.solvent_contribution,
                    dtype=float,
                )
            )
            & (
                np.asarray(
                    context.prefit_evaluation.solvent_contribution,
                    dtype=float,
                )
                > 0.0
            )
        )
    ):
        assert "prefit_model_with_solvent.png" in figure_names
    assert progress_updates
    assert progress_updates[0][0] == 0
    assert "Please wait" in progress_updates[0][2]
    assert progress_updates[-1][0] == progress_updates[-1][1]

    presentation = Presentation(str(result.report_path))
    assert presentation.slide_width / 914400 == pytest.approx(13.333, rel=1e-3)
    assert presentation.slide_height / 914400 == pytest.approx(7.5, rel=1e-3)
    slide_texts = [
        "\n".join(
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        for slide in presentation.slides
    ]

    for slide in presentation.slides:
        for shape in slide.shapes:
            assert shape.left >= 0
            assert shape.top >= 0
            assert shape.left + shape.width <= presentation.slide_width
            assert shape.top + shape.height <= presentation.slide_height

    title_shape = next(
        shape
        for shape in presentation.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
        and "SAXS Model Report" in shape.text
    )
    assert title_shape.text_frame.paragraphs[0].runs[0].font.name == "Arial"
    assert (
        sum(
            "Posterior Violin Comparison" in slide_text
            for slide_text in slide_texts
        )
        == 3
    )
    assert any(
        "Posterior Violin Comparison - Weights" in slide_text
        for slide_text in slide_texts
    )
    assert any(
        "Posterior Violin Comparison - Effective Radii" in slide_text
        for slide_text in slide_texts
    )
    if "prefit_model_with_solvent.png" in figure_names:
        assert any(
            any(
                getattr(shape, "has_text_frame", False)
                and "Prefit Model With Solvent" in shape.text
                for shape in slide.shapes
            )
            for slide in presentation.slides
        )


def test_dream_model_report_exports_model_information_from_template_context(
    qapp, tmp_path
):
    pytest.importorskip("pptx")
    from pptx import Presentation

    del qapp
    project_dir, paths = _build_poly_lma_geometry_project(
        tmp_path,
        template_name=POLY_LMA_HS_TEMPLATE,
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.compute_prefit_cluster_geometry()
    _write_minimal_dream_results(project_dir)
    window.load_latest_results()

    context = window._build_dream_model_report_context(
        settings=window.dream_tab.settings_payload(),
        output_path=paths.reports_dir / "dream_model_report_model_info.pptx",
        asset_dir=paths.reports_dir / "dream_model_report_model_info_assets",
    )

    assert context.template_display_name == "pyDREAM Poly LMA Hard-Sphere"
    assert context.template_module_path is not None
    assert (
        context.template_module_path.name == "template_pydream_poly_lma_hs.py"
    )
    assert context.model_equation_text is not None
    assert "I_model(q)" in context.model_equation_text
    assert "phi_solute" in context.model_equation_text

    definition_text = "\n".join(context.model_definition_lines)
    reference_text = "\n".join(context.model_reference_lines)
    assert "Structure Factor:" in definition_text
    assert "Parameter Definitions:" in definition_text
    assert "phi_solute" in definition_text
    assert "Pedersen review" in reference_text
    assert "sasview.org/docs/user/models/hardsphere.html" in reference_text

    result = export_dream_model_report_pptx(context)

    presentation = Presentation(str(result.report_path))
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )

    assert "Model Information" in slide_text
    assert "pyDREAM Poly LMA Hard-Sphere" in slide_text
    assert "Model equation:" in slide_text
    assert "I_model(q)" in slide_text
    assert "Term definitions:" in slide_text
    assert "References:" in slide_text
    assert "Pedersen review" in slide_text
    assert "Pedersen97.pdf" in slide_text


def test_dream_model_report_uses_selected_secondary_atom_for_solvent_sort_prior_histograms(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.load_latest_results()

    (paths.project_dir / "md_prior_weights.json").write_text(
        json.dumps(
            {
                "origin": "clusters",
                "total_files": 5,
                "available_elements": ["Pb", "I", "Br", "O"],
                "structures": {
                    "PbI2": {
                        "motif_A": {
                            "count": 2,
                            "weight": 0.4,
                            "profile_file": "A_no_motif.txt",
                            "secondary_atom_distributions": {
                                "Br": {"0": 1, "1": 1},
                                "O": {"0": 1, "2": 1},
                            },
                        }
                    },
                    "Pb2I4": {
                        "motif_B": {
                            "count": 3,
                            "weight": 0.6,
                            "profile_file": "A_no_motif.txt",
                            "secondary_atom_distributions": {
                                "Br": {"0": 1, "3": 2},
                                "O": {"1": 1, "2": 2},
                            },
                        }
                    },
                },
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    window.project_setup_tab.apply_cluster_import_data(
        ["Pb", "I", "Br", "O"],
        [
            {
                "structure": "PbI2",
                "motif": "motif_A",
                "count": 2,
                "weight": 0.4,
                "atom_fraction_percent": 40.0,
                "structure_fraction_percent": 40.0,
            },
            {
                "structure": "Pb2I4",
                "motif": "motif_B",
                "count": 3,
                "weight": 0.6,
                "atom_fraction_percent": 60.0,
                "structure_fraction_percent": 60.0,
            },
        ],
    )
    window.project_setup_tab.prior_mode_combo.setCurrentText(
        "Solvent Sort - Structure Fraction"
    )
    secondary_index = window.project_setup_tab.secondary_filter_combo.findText(
        "O"
    )
    assert secondary_index >= 0
    window.project_setup_tab.secondary_filter_combo.setCurrentIndex(
        secondary_index
    )
    window.project_setup_tab.prior_mode_combo.setCurrentText(
        "Structure Fraction"
    )

    context = window._build_dream_model_report_context(
        settings=window.dream_tab.settings_payload(),
        output_path=paths.reports_dir
        / "dream_model_report_secondary_atom.pptx",
        asset_dir=(
            paths.reports_dir / "dream_model_report_secondary_atom_assets"
        ),
    )

    assert (
        context.powerpoint_settings.solvent_sort_histogram_color_map
        == "summer"
    )
    assert (
        context.prior_histograms[2].mode == "solvent_sort_structure_fraction"
    )
    assert context.prior_histograms[2].secondary_element == "O"
    assert context.prior_histograms[2].cmap == "summer"
    assert context.prior_histograms[3].mode == "solvent_sort_atom_fraction"
    assert context.prior_histograms[3].secondary_element == "O"
    assert context.prior_histograms[3].cmap == "summer"


def test_dream_model_report_context_and_export_honor_powerpoint_settings(
    qapp,
    tmp_path,
):
    pytest.importorskip("pptx")

    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.load_latest_results()
    window.current_settings.powerpoint_export_settings = (
        PowerPointExportSettings(
            font_family="Courier New",
            component_color_map="plasma",
            prior_histogram_color_map="cividis",
            solvent_sort_histogram_color_map="magma",
            include_prior_histograms=False,
            include_directory_summary=False,
            generate_manifest=False,
            export_figure_assets=False,
        )
    )

    context = window._build_dream_model_report_context(
        settings=window.dream_tab.settings_payload(),
        output_path=paths.reports_dir / "dream_model_report_custom.pptx",
        asset_dir=paths.reports_dir / "dream_model_report_custom_assets",
    )

    assert context.powerpoint_settings.font_family == "Courier New"
    assert not context.powerpoint_settings.include_prior_histograms
    assert not context.powerpoint_settings.generate_manifest
    assert not context.powerpoint_settings.export_figure_assets
    assert not any(
        "Report assets:" in line or "Report manifest:" in line
        for line in context.directory_lines
    )
    assert context.prior_histograms[0].cmap == "cividis"
    assert context.prior_histograms[2].cmap == "magma"
    expected_component_color = to_hex(
        colormaps["plasma"](0.68),
        keep_alpha=False,
    )
    assert (
        context.component_plot_without_solvent.component_series[0].color
        == expected_component_color
    )

    result = export_dream_model_report_pptx(context)

    assert result.report_path.is_file()
    assert result.manifest_path is None
    assert not result.figure_paths
    assert not context.asset_dir.exists()


def test_dream_recycle_pushes_selected_best_fit_into_prefit(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.load_latest_results()
    update_calls: list[str] = []
    original_update_prefit_model = window.update_prefit_model

    def _record_update_prefit_model():
        update_calls.append("called")
        return original_update_prefit_model()

    window.update_prefit_model = _record_update_prefit_model

    settings = window.dream_tab.settings_payload()
    summary = window._last_results_loader.get_summary(
        bestfit_method=settings.bestfit_method,
        posterior_filter_mode=settings.posterior_filter_mode,
        posterior_top_percent=settings.posterior_top_percent,
        posterior_top_n=settings.posterior_top_n,
        credible_interval_low=settings.credible_interval_low,
        credible_interval_high=settings.credible_interval_high,
    )
    expected_values = {
        str(name): float(summary.bestfit_params[index])
        for index, name in enumerate(summary.full_parameter_names)
    }

    scale_row = window.prefit_tab.find_parameter_row("scale")
    assert scale_row >= 0
    vary_item = window.prefit_tab.parameter_table.item(scale_row, 4)
    assert vary_item is not None
    vary_item.setCheckState(Qt.CheckState.Unchecked)
    window.prefit_tab.set_parameter_row("scale", value=2e-3)
    window.prefit_tab.set_parameter_row("offset", value=0.333)

    window.recycle_dream_output_to_prefit()

    entries_by_name = {
        entry.name: entry for entry in window.prefit_tab.parameter_entries()
    }
    assert entries_by_name["scale"].value == pytest.approx(
        expected_values["scale"]
    )
    assert entries_by_name["offset"].value == pytest.approx(
        expected_values["offset"]
    )
    assert entries_by_name["w0"].value == pytest.approx(expected_values["w0"])
    assert not entries_by_name["scale"].vary
    assert window.tabs.currentWidget() is window.prefit_tab
    assert update_calls == []
    assert "Recycled DREAM output into Prefit." in (
        window.prefit_tab.output_box.toPlainText()
    )
    assert "Prefit preview refresh: deferred" in (
        window.prefit_tab.output_box.toPlainText()
    )
    saved_settings = SAXSProjectManager().load_project(project_dir)
    recycled_best_entries = {
        str(entry["name"]): float(entry["value"])
        for entry in saved_settings.best_prefit_parameter_entries
    }
    assert recycled_best_entries["scale"] == pytest.approx(
        expected_values["scale"]
    )
    reset_col = _table_column_index(window.prefit_tab.parameter_table, "Reset")
    reset_button = window.prefit_tab.parameter_table.cellWidget(
        scale_row,
        reset_col,
    )
    assert reset_button is not None
    window.prefit_tab.set_parameter_row("scale", value=2e-3)
    reset_button.click()
    recycled_entries = {
        entry.name: entry for entry in window.prefit_tab.parameter_entries()
    }
    assert recycled_entries["scale"].value == pytest.approx(
        expected_values["scale"]
    )
    dream_entries = {
        entry.param: entry
        for entry in SAXSDreamWorkflow(project_dir).load_parameter_map()
    }
    assert dream_entries["scale"].value == pytest.approx(
        expected_values["scale"]
    )
    assert "Recycled the current DREAM best fit into the Prefit tab." in (
        window.dream_tab.output_box.toPlainText()
    )


def test_dream_run_finish_auto_exports_condensed_results(qapp, tmp_path):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    bundle = _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window._on_dream_run_finished(
        str(bundle.run_dir),
        {
            "sampled_params_path": str(
                bundle.run_dir / "dream_sampled_params.npy"
            ),
            "log_ps_path": str(bundle.run_dir / "dream_log_ps.npy"),
        },
    )

    auto_model_exports = sorted(
        paths.exported_data_dir.glob("dream_model_fit_auto_*.csv")
    )
    auto_violin_exports = sorted(
        paths.exported_data_dir.glob("dream_violin_auto_*.csv")
    )

    assert auto_model_exports
    assert auto_violin_exports
    auto_model_base = auto_model_exports[-1].stem
    auto_violin_base = auto_violin_exports[-1].stem
    assert (
        paths.exported_data_dir / f"{auto_model_base}.metadata.json"
    ).is_file()
    assert (
        paths.exported_data_dir / f"{auto_model_base}.report.txt"
    ).is_file()
    assert (
        paths.exported_data_dir / f"{auto_violin_base}.metadata.json"
    ).is_file()
    assert (
        paths.exported_data_dir / f"{auto_violin_base}.report.txt"
    ).is_file()


def test_dream_overconstraint_warning_can_set_all_vary_to_yes(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    bundle = _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    workflow = SAXSDreamWorkflow(project_dir)
    entries = workflow.load_parameter_map()
    for entry in entries:
        if entry.param == "scale":
            entry.vary = False
        if entry.param == "w0":
            entry.vary = False
    workflow.save_parameter_map(entries)
    window._last_results_loader = SAXSDreamResultsLoader(
        bundle.run_dir,
        burnin_percent=0,
    )

    prefit_evaluation = window.prefit_workflow.evaluate()
    worse_model = (
        np.asarray(
            prefit_evaluation.model_intensities,
            dtype=float,
        )
        + 0.5
    )
    worse_metrics = window._fit_quality_metrics_from_curves(
        prefit_evaluation.experimental_intensities,
        worse_model,
    )
    monkeypatch.setattr(
        window._last_results_loader,
        "build_model_fit_data",
        lambda **kwargs: SimpleNamespace(
            q_values=prefit_evaluation.q_values,
            experimental_intensities=(
                prefit_evaluation.experimental_intensities
            ),
            model_intensities=worse_model,
            bestfit_method="map",
            template_name="template_pd_likelihood_monosq_decoupled",
            rmse=worse_metrics.rmse,
            mean_abs_residual=worse_metrics.mean_abs_residual,
            r_squared=worse_metrics.r_squared,
        ),
    )

    captured: dict[str, object] = {}

    def _accept_prompt(comparison):
        captured["comparison"] = comparison
        return True

    monkeypatch.setattr(
        window, "_prompt_dream_constraint_update", _accept_prompt
    )

    window._maybe_warn_about_dream_overconstraints(
        window.dream_tab.settings_payload()
    )

    updated_entries = workflow.load_parameter_map()
    assert all(entry.vary for entry in updated_entries)
    comparison = captured["comparison"]
    assert "scale" in comparison.fixed_non_weight_parameters
    assert "w0" in comparison.fixed_weight_parameters
    assert window._dream_parameter_map_saved_in_session is True
    assert "All DREAM parameters now have vary=yes" in (
        window.dream_tab.output_box.toPlainText()
    )


def test_dream_overconstraint_warning_skips_prompt_when_dream_not_worse(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    bundle = _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    workflow = SAXSDreamWorkflow(project_dir)
    entries = workflow.load_parameter_map()
    for entry in entries:
        if entry.param == "scale":
            entry.vary = False
    workflow.save_parameter_map(entries)
    window._last_results_loader = SAXSDreamResultsLoader(
        bundle.run_dir,
        burnin_percent=0,
    )

    prefit_evaluation = window.prefit_workflow.evaluate()
    better_model = np.asarray(
        prefit_evaluation.experimental_intensities,
        dtype=float,
    )
    better_metrics = window._fit_quality_metrics_from_curves(
        prefit_evaluation.experimental_intensities,
        better_model,
    )
    monkeypatch.setattr(
        window._last_results_loader,
        "build_model_fit_data",
        lambda **kwargs: SimpleNamespace(
            q_values=prefit_evaluation.q_values,
            experimental_intensities=(
                prefit_evaluation.experimental_intensities
            ),
            model_intensities=better_model,
            bestfit_method="map",
            template_name="template_pd_likelihood_monosq_decoupled",
            rmse=better_metrics.rmse,
            mean_abs_residual=better_metrics.mean_abs_residual,
            r_squared=better_metrics.r_squared,
        ),
    )

    monkeypatch.setattr(
        window,
        "_prompt_dream_constraint_update",
        lambda comparison: pytest.fail(
            f"Unexpected overconstraint prompt: {comparison}"
        ),
    )

    window._maybe_warn_about_dream_overconstraints(
        window.dream_tab.settings_payload()
    )

    assert any(not entry.vary for entry in workflow.load_parameter_map())


def test_dream_run_finish_auto_applies_recommended_filter(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    bundle = _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    monkeypatch.setattr(
        window,
        "_evaluate_dream_posterior_filters",
        lambda settings: (
            [
                {
                    "mode": "all_post_burnin",
                    "description": "all_post_burnin",
                    "rmse": 1.2,
                    "mean_abs_residual": 1.0,
                    "r_squared": 0.1,
                    "posterior_sample_count": 8,
                },
                {
                    "mode": "top_n_logp",
                    "description": "top_n_logp (top 500 samples by log-posterior)",
                    "rmse": 0.2,
                    "mean_abs_residual": 0.15,
                    "r_squared": 0.95,
                    "posterior_sample_count": 4,
                },
            ],
            {
                "mode": "top_n_logp",
                "description": "top_n_logp (top 500 samples by log-posterior)",
                "rmse": 0.2,
                "mean_abs_residual": 0.15,
                "r_squared": 0.95,
                "posterior_sample_count": 4,
            },
        ),
    )

    window._on_dream_run_finished(
        str(bundle.run_dir),
        {
            "sampled_params_path": str(
                bundle.run_dir / "dream_sampled_params.npy"
            ),
            "log_ps_path": str(bundle.run_dir / "dream_log_ps.npy"),
        },
    )

    assert window.dream_tab.posterior_filter_combo.currentData() == (
        "top_n_logp"
    )
    assert window.dream_tab.selected_search_filter_preset() == "custom"
    assert "Recommended posterior filter by fit quality" in (
        window.dream_tab.output_box.toPlainText()
    )
    assert "Applied recommended posterior filter automatically." in (
        window.dream_tab.output_box.toPlainText()
    )


def test_prefit_template_selection_stays_pending_until_change_requested(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    if window.prefit_tab.template_combo.count() < 2:
        pytest.skip("Need at least two templates to test template switching.")

    original_template = window.prefit_workflow.template_spec.name
    alternative = next(
        str(window.prefit_tab.template_combo.itemData(index) or "")
        for index in range(window.prefit_tab.template_combo.count())
        if str(window.prefit_tab.template_combo.itemData(index) or "")
        != original_template
    )

    window.prefit_tab.set_selected_template(alternative, emit_signal=True)

    assert window.prefit_workflow.template_spec.name == original_template
    assert window.prefit_tab.active_template_name() == original_template
    assert window.project_setup_tab.active_template_name() == original_template
    assert window.prefit_tab.selected_template_name() == alternative
    assert window.project_setup_tab.selected_template_name() == alternative
    assert window.prefit_tab.change_template_button.isEnabled()
    assert window.project_setup_tab.change_template_button.isEnabled()
    window.close()


def test_settings_from_project_tab_uses_active_template_when_selection_differs(
    qapp, tmp_path
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    if window.prefit_tab.template_combo.count() < 2:
        pytest.skip("Need at least two templates to test template switching.")

    original_template = window.prefit_workflow.template_spec.name
    alternative = next(
        str(window.prefit_tab.template_combo.itemData(index) or "")
        for index in range(window.prefit_tab.template_combo.count())
        if str(window.prefit_tab.template_combo.itemData(index) or "")
        != original_template
    )

    window.prefit_tab.set_selected_template(alternative, emit_signal=True)

    captured_settings = window._settings_from_project_tab()

    assert captured_settings.selected_model_template == original_template
    assert window.project_setup_tab.selected_template_name() == alternative
    assert window.project_setup_tab.active_template_name() == original_template
    window.close()


def test_change_template_creates_template_scoped_distribution_copy(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    _settings, original_artifacts = _seed_saved_distribution_from_root(
        project_dir,
        include_cluster_geometry=True,
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)
    if window.project_setup_tab.template_combo.count() < 2:
        pytest.skip("Need at least two templates to test template switching.")

    original_template = window.current_settings.selected_model_template
    alternative = next(
        str(window.project_setup_tab.template_combo.itemData(index) or "")
        for index in range(window.project_setup_tab.template_combo.count())
        if str(window.project_setup_tab.template_combo.itemData(index) or "")
        != original_template
    )

    monkeypatch.setattr(
        window,
        "_confirm_prefit_template_change",
        lambda current, new: (True, False),
    )

    window.project_setup_tab.set_selected_template(
        alternative, emit_signal=True
    )
    window.project_setup_tab.change_template_button.click()
    QApplication.processEvents()

    current_artifacts = project_artifact_paths(window.current_settings)

    assert window.current_settings.selected_model_template == alternative
    assert window.project_setup_tab.active_template_name() == alternative
    assert window.prefit_tab.active_template_name() == alternative
    assert (
        current_artifacts.distribution_id != original_artifacts.distribution_id
    )
    assert current_artifacts.component_map_file.is_file()
    assert current_artifacts.prior_weights_file.is_file()
    assert not current_artifacts.cluster_geometry_metadata_file.exists()
    assert window.project_setup_tab.computed_distribution_combo.count() == 2
    assert (
        window.project_setup_tab.selected_distribution_id()
        == current_artifacts.distribution_id
    )
    assert original_artifacts.cluster_geometry_metadata_file.is_file()
    window.close()


def test_project_setup_inputs_stay_locked_until_project_selected(
    qapp, tmp_path
):
    del qapp
    window = SAXSMainWindow()

    assert not window.project_setup_tab.forward_model_group.isEnabled()
    assert not window.project_setup_tab.model_group.isEnabled()
    assert not window.project_setup_tab.prior_mode_combo.isEnabled()
    assert not window.project_setup_tab.generate_prior_plot_button.isEnabled()
    assert window.prefit_tab.template_combo.count() >= 1
    assert window.project_setup_tab.use_experimental_grid()
    assert not window.project_setup_tab.resample_points_spin.isEnabled()

    project_dir = tmp_path / "Created SAXS Project"
    window.project_setup_tab.project_name_edit.setText("Created SAXS Project")
    window.project_setup_tab.project_dir_edit.setText(str(tmp_path))
    window.create_project_from_tab()

    assert window.project_setup_tab.forward_model_group.isEnabled()
    assert window.project_setup_tab.model_group.isEnabled()
    assert window.project_setup_tab.prior_mode_combo.isEnabled()
    assert window.project_setup_tab.generate_prior_plot_button.isEnabled()
    assert window.project_setup_tab.project_dir_edit.text() == str(tmp_path)
    assert window.project_setup_tab.project_name_edit.text() == (
        "Created SAXS Project"
    )
    assert window.project_setup_tab.open_project_dir_edit.text().endswith(
        "Created SAXS Project"
    )
    assert window.current_settings.project_name == "Created SAXS Project"
    assert Path(window.current_settings.project_dir) == project_dir.resolve()

    window.project_setup_tab.use_experimental_grid_checkbox.setChecked(False)

    assert window.project_setup_tab.resample_points_spin.isEnabled()


def test_recognized_cluster_table_is_resizable_scrollable_and_saves_colors(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.project_setup_tab.apply_cluster_import_data(
        ["A"],
        [
            {
                "structure": "A",
                "motif": "no_motif",
                "count": 1,
                "weight": 1.0,
                "atom_fraction_percent": 100.0,
                "structure_fraction_percent": 100.0,
            }
        ],
    )

    table = window.project_setup_tab.recognized_clusters_table
    assert (
        table.horizontalHeader().sectionResizeMode(0)
        == QHeaderView.ResizeMode.Interactive
    )
    assert (
        table.horizontalScrollBarPolicy()
        == Qt.ScrollBarPolicy.ScrollBarAsNeeded
    )
    assert (
        table.verticalScrollBarPolicy() == Qt.ScrollBarPolicy.ScrollBarAsNeeded
    )

    monkeypatch.setattr(
        "saxshell.saxs.ui.project_setup_tab.QColorDialog.getColor",
        lambda *args, **kwargs: QColor("#123456"),
    )

    window.project_setup_tab._on_recognized_cluster_cell_clicked(0, 7)

    assert table.item(0, 7).text() == "#123456"
    assert window.project_setup_tab.component_trace_colors() == {
        "A_no_motif": "#123456"
    }

    window.save_project_state()

    settings = SAXSProjectManager().load_project(project_dir)
    assert settings.component_trace_colors == {"A_no_motif": "#123456"}


def test_model_and_build_layout_keeps_cluster_table_below_template_fields(
    qapp,
):
    del qapp
    tab = ProjectSetupTab()

    model_layout = tab.model_group.layout()
    assert model_layout is not None
    assert model_layout.itemAt(0).widget() is tab._model_build_header_widget
    assert model_layout.itemAt(1).layout() is tab._model_build_lower_layout
    assert (
        tab._model_build_lower_layout.itemAt(0).widget()
        is tab._model_build_button_widget
    )
    assert (
        tab._model_build_lower_layout.itemAt(1).widget()
        is tab._recognized_clusters_group
    )


def test_open_project_uses_existing_project_field(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow()

    window.project_setup_tab.open_project_dir_edit.setText(str(project_dir))
    window.open_project_from_dialog()

    assert window.current_settings is not None
    assert window.current_settings.project_dir == str(project_dir.resolve())
    assert window.project_setup_tab.project_name_edit.text() == "saxs_project"


def test_existing_project_browser_starts_in_recent_project_parent(
    qapp, tmp_path, monkeypatch
):
    del qapp

    class _FakeSettings:
        def __init__(self, values: dict[str, object]):
            self.values = dict(values)

        def value(self, key, default=None):
            return self.values.get(key, default)

        def setValue(self, key, value):
            self.values[key] = value

    recent_root = tmp_path / "recent_projects"
    recent_root.mkdir()
    project_dir, _paths = _build_minimal_saxs_project(recent_root)
    settings_store = _FakeSettings(
        {"recent_project_dirs": [str(project_dir.resolve())]}
    )
    captured: dict[str, str] = {}

    monkeypatch.setattr(
        "saxshell.saxs.ui.project_setup_tab.QSettings",
        lambda *args, **kwargs: settings_store,
    )

    def _capture_existing_directory(*args, **kwargs):
        captured["start_dir"] = args[2]
        return ""

    monkeypatch.setattr(
        "saxshell.saxs.ui.project_setup_tab.QFileDialog.getExistingDirectory",
        _capture_existing_directory,
    )

    tab = ProjectSetupTab()
    tab._browse_existing_project_directory()

    assert captured["start_dir"] == str(project_dir.resolve().parent)


def test_open_project_dialog_uses_recent_project_parent_when_field_is_empty(
    qapp, tmp_path, monkeypatch
):
    del qapp

    class _FakeSettings:
        def __init__(self, values: dict[str, object]):
            self.values = dict(values)

        def value(self, key, default=None):
            return self.values.get(key, default)

        def setValue(self, key, value):
            self.values[key] = value

    recent_root = tmp_path / "recent_projects"
    recent_root.mkdir()
    project_dir, _paths = _build_minimal_saxs_project(recent_root)
    settings_store = _FakeSettings(
        {"recent_project_dirs": [str(project_dir.resolve())]}
    )
    captured: dict[str, str] = {}

    monkeypatch.setattr(
        SAXSMainWindow,
        "_recent_projects_settings",
        lambda self: settings_store,
    )

    def _capture_existing_directory(*args, **kwargs):
        captured["start_dir"] = args[2]
        return ""

    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.QFileDialog.getExistingDirectory",
        _capture_existing_directory,
    )

    window = SAXSMainWindow()
    window.project_setup_tab.open_project_dir_edit.clear()

    window.open_project_from_dialog()

    assert captured["start_dir"] == str(project_dir.resolve().parent)


def test_auto_snap_panes_setting_defaults_enabled_and_persists(
    qapp,
    monkeypatch,
):
    del qapp

    class _FakeSettings:
        def __init__(self, values: dict[str, object] | None = None):
            self.values = {} if values is None else dict(values)

        def value(self, key, default=None):
            return self.values.get(key, default)

        def setValue(self, key, value):
            self.values[key] = value

    settings_store = _FakeSettings()

    monkeypatch.setattr(
        SAXSMainWindow,
        "_recent_projects_settings",
        lambda self: settings_store,
    )

    first_window = SAXSMainWindow()

    assert first_window.auto_snap_panes_action.isChecked()
    assert first_window.project_setup_tab._auto_snap_filter.is_enabled()
    assert first_window.prefit_tab._auto_snap_filter.is_enabled()
    assert first_window.dream_tab._auto_snap_filter.is_enabled()

    first_window.auto_snap_panes_action.trigger()

    assert settings_store.values[AUTO_SNAP_PANES_KEY] is False
    assert not first_window.auto_snap_panes_action.isChecked()
    assert not first_window.project_setup_tab._auto_snap_filter.is_enabled()
    assert not first_window.prefit_tab._auto_snap_filter.is_enabled()
    assert not first_window.dream_tab._auto_snap_filter.is_enabled()

    second_window = SAXSMainWindow()

    assert not second_window.auto_snap_panes_action.isChecked()
    assert not second_window.project_setup_tab._auto_snap_filter.is_enabled()
    assert not second_window.prefit_tab._auto_snap_filter.is_enabled()
    assert not second_window.dream_tab._auto_snap_filter.is_enabled()

    first_window.close()
    second_window.close()


def test_open_project_reports_loader_progress_and_output(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow()

    window.project_setup_tab.open_project_dir_edit.setText(str(project_dir))
    window.open_project_from_dialog()

    assert (
        window.project_setup_tab.activity_progress_bar.maximum()
        == PROJECT_LOAD_TOTAL_STEPS
    )
    assert (
        window.project_setup_tab.activity_progress_bar.value()
        == PROJECT_LOAD_TOTAL_STEPS
    )
    assert (
        window.project_setup_tab.activity_progress_label.text()
        == "Project load complete."
    )

    summary_text = window.project_setup_tab.summary_box.toPlainText()
    assert "Loading project settings from" in summary_text
    assert "Experimental data preview:" in summary_text
    assert "Loaded 1 SAXS component trace." in summary_text
    assert "Loaded project saxs_project" in summary_text

    assert window._progress_dialog is not None
    assert not window._progress_dialog.isVisible()
    dialog_output = window._progress_dialog.output_box.toPlainText()
    assert "Loading project settings from" in dialog_output
    assert "Loading Prefit workflow and preview." in dialog_output


def test_open_project_prepares_loader_payload_off_main_thread(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow()
    main_thread_id = threading.get_ident()
    build_thread_id: int | None = None
    payload_flags: tuple[bool, bool] | None = None
    original_build = window._build_project_load_payload
    original_apply = window._apply_project_settings

    def wrapped_build(*args, **kwargs):
        nonlocal build_thread_id
        build_thread_id = threading.get_ident()
        return original_build(*args, **kwargs)

    def wrapped_apply(settings, **kwargs):
        nonlocal payload_flags
        payload_flags = (
            kwargs.get("prefit_payload") is not None,
            kwargs.get("dream_payload") is not None,
        )
        return original_apply(settings, **kwargs)

    monkeypatch.setattr(window, "_build_project_load_payload", wrapped_build)
    monkeypatch.setattr(window, "_apply_project_settings", wrapped_apply)

    window.load_project(project_dir)

    assert build_thread_id is not None
    assert build_thread_id != main_thread_id
    assert payload_flags == (True, True)


def test_loaded_project_prefers_original_experimental_reference_in_ui(
    qapp, tmp_path
):
    del qapp
    manager = SAXSProjectManager()
    source_dir = tmp_path / "source_data"
    source_dir.mkdir()
    source_path = source_dir / "exp_source.txt"
    np.savetxt(
        source_path,
        np.column_stack(
            [
                np.asarray([0.05, 0.10], dtype=float),
                np.asarray([10.0, 9.0], dtype=float),
            ]
        ),
    )

    project_dir = tmp_path / "project_with_copy"
    settings = manager.create_project(project_dir)
    paths = build_project_paths(project_dir)
    copied_path = paths.experimental_data_dir / source_path.name
    copied_path.write_text(
        source_path.read_text(encoding="utf-8"), encoding="utf-8"
    )
    settings.experimental_data_path = str(source_path)
    settings.copied_experimental_data_file = str(copied_path)
    manager.save_project(settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert window.project_setup_tab.experimental_data_edit.text() == str(
        source_path
    )


def test_prefit_recommended_scale_button_updates_scale_bounds(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    scale_row = window.prefit_tab.find_parameter_row("scale")
    assert scale_row >= 0
    window.prefit_tab.parameter_table.item(scale_row, 3).setText("1e-6")
    window.prefit_tab.parameter_table.item(scale_row, 5).setText("1e-7")
    window.prefit_tab.parameter_table.item(scale_row, 6).setText("1e-5")

    window.apply_recommended_scale_settings()

    scale_entry = next(
        entry
        for entry in window.prefit_tab.parameter_entries()
        if entry.name == "scale"
    )
    offset_entry = next(
        entry
        for entry in window.prefit_tab.parameter_entries()
        if entry.name == "offset"
    )
    assert scale_entry.vary
    assert scale_entry.value == pytest.approx(5e-4)
    assert scale_entry.minimum == pytest.approx(5e-5)
    assert scale_entry.maximum == pytest.approx(5e-3)
    assert offset_entry.value == pytest.approx(0.05)
    assert "Applied autoscale settings." in (
        window.prefit_tab.output_box.toPlainText()
    )


def test_prefit_tab_reorders_controls_and_parameter_actions(qapp):
    del qapp
    tab = PrefitTab()

    assert tab._parameter_action_layout.itemAt(0).widget() is (
        tab.recommended_scale_button
    )
    assert tab._parameter_action_layout.itemAt(1).widget() is (
        tab.update_button
    )
    assert tab._parameter_action_layout.itemAt(2).widget() is (
        tab.auto_update_checkbox
    )
    assert tab._parameter_action_layout.itemAt(3).widget() is (
        tab.scrollable_parameter_checkbox
    )
    assert tab._prefit_control_button_grid.itemAtPosition(0, 0).widget() is (
        tab._run_button_cell
    )
    assert tab._prefit_control_button_grid.itemAtPosition(0, 1).widget() is (
        tab.autosave_checkbox
    )
    assert tab._prefit_control_button_grid.itemAtPosition(1, 0).widget() is (
        tab.save_button
    )
    assert tab._prefit_control_button_grid.itemAtPosition(1, 1).widget() is (
        tab.reset_button
    )
    assert tab._prefit_control_button_grid.itemAtPosition(2, 0).widget() is (
        tab.set_best_button
    )
    assert tab._prefit_control_button_grid.itemAtPosition(2, 1).widget() is (
        tab.reset_best_button
    )
    assert tab.sequence_history_checkbox.text() == "Sequence history logger"


def test_prefit_sequence_history_logger_records_fit_actions(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert window.prefit_workflow is not None
    monkeypatch.setattr(window, "_load_dream_workflow", lambda: None)
    monkeypatch.setattr(window, "update_prefit_model", lambda: None)
    monkeypatch.setattr(
        window,
        "_confirm_large_prefit_parameter_count",
        lambda names: (True, False),
    )
    monkeypatch.setattr(
        window.prefit_tab, "plot_evaluation", lambda value: None
    )
    monkeypatch.setattr(
        window.prefit_tab, "set_summary_text", lambda text: None
    )
    monkeypatch.setattr(
        window,
        "_refresh_saved_prefit_states",
        lambda selected_name=None: None,
    )
    window.prefit_tab.sequence_history_checkbox.setChecked(True)

    scale_row = window.prefit_tab.find_parameter_row("scale")
    scale_item = window.prefit_tab.parameter_table.item(scale_row, 3)
    assert scale_item is not None
    scale_item.setText("0.0025")

    window.apply_recommended_scale_settings()

    def fake_run_fit(entries, *, method="leastsq", max_nfev=10000):
        del max_nfev
        evaluation = window.prefit_workflow.evaluate(entries)
        residuals = np.asarray(evaluation.residuals, dtype=float)
        return PrefitFitResult(
            parameter_entries=entries,
            evaluation=evaluation,
            fit_report="fake lmfit report",
            method=method,
            nfev=12,
            chi_square=float(np.sum(residuals**2)),
            reduced_chi_square=float(np.mean(residuals**2)),
            r_squared=0.5,
            optimization_strategy=f"lmfit {method}",
            grid_evaluations=0,
        )

    monkeypatch.setattr(window.prefit_workflow, "run_fit", fake_run_fit)

    window.run_prefit()

    history_path = window.prefit_workflow.sequence_history_path()
    assert history_path.is_file()
    payload = json.loads(history_path.read_text(encoding="utf-8"))
    assert payload["format_version"] == 1
    assert payload["sequence_logger_enabled"] is True

    event_types = [event["event_type"] for event in payload["events"]]
    assert "sequence_history_logger_toggled" in event_types
    assert "manual_parameter_update" in event_types
    assert "autoscale_applied" in event_types
    assert "prefit_run_complete" in event_types

    manual_event = next(
        event
        for event in payload["events"]
        if event["event_type"] == "manual_parameter_update"
    )
    assert manual_event["details"]["trigger"] == "apply_autoscale"
    assert any(
        change["parameter_name"] == "scale"
        for change in manual_event["details"]["changed_parameters"]
    )

    autoscale_event = next(
        event
        for event in payload["events"]
        if event["event_type"] == "autoscale_applied"
    )
    assert autoscale_event["details"]["points_used"] > 0
    assert any(
        entry["name"] == "scale"
        for entry in autoscale_event["parameter_entries"]
    )

    run_event = next(
        event
        for event in payload["events"]
        if event["event_type"] == "prefit_run_complete"
    )
    assert run_event["details"]["method"] == "leastsq"
    assert run_event["details"]["statistics"]["nfev"] == 12
    assert any(
        entry["parameter_name"] == "scale"
        for entry in run_event["details"]["varying_parameters"]
    )
    window.close()


def test_solute_volume_fraction_help_text_and_labels_omit_fullrmc(qapp):
    del qapp
    widget = SoluteVolumeFractionWidget()

    assert "fullrmc" not in SOLUTE_VOLUME_FRACTION_HELP_TEXT.lower()
    label_text = "\n".join(
        label.text().lower() for label in widget.findChildren(QLabel)
    )
    assert "fullrmc" not in label_text


def test_prefit_tab_auto_updates_model_on_value_change_when_enabled(qapp):
    del qapp
    tab = PrefitTab()
    events: list[str] = []
    tab.update_model_requested.connect(lambda: events.append("update"))
    tab.populate_parameter_table(
        [
            PrefitParameterEntry(
                structure="A",
                motif="motif",
                name="scale",
                value=1.0,
                vary=True,
                minimum=0.1,
                maximum=10.0,
                category="fit",
            )
        ]
    )

    row = tab.find_parameter_row("scale")
    assert row >= 0

    tab.parameter_table.item(row, 3).setText("2.5")
    assert events == []

    tab.auto_update_checkbox.setChecked(True)
    tab.parameter_table.item(row, 3).setText("3.5")
    assert events == ["update"]

    tab.set_parameter_row("scale", value=4.5)
    assert events == ["update"]


def test_prefit_tab_supports_linked_parameter_expressions(qapp):
    del qapp
    tab = PrefitTab()
    events: list[str] = []
    tab.update_model_requested.connect(lambda: events.append("update"))
    tab.populate_parameter_table(
        [
            PrefitParameterEntry(
                structure="A",
                motif="motif",
                name="scale",
                value=2.5,
                vary=True,
                minimum=0.1,
                maximum=10.0,
                category="fit",
            ),
            PrefitParameterEntry(
                structure="A",
                motif="motif",
                name="offset",
                value=0.0,
                vary=True,
                minimum=-10.0,
                maximum=10.0,
                category="fit",
            ),
        ]
    )

    tab.auto_update_checkbox.setChecked(True)
    offset_row = tab.find_parameter_row("offset")
    assert offset_row >= 0

    tab.parameter_table.item(offset_row, 3).setText("*scale")

    entries = {entry.name: entry for entry in tab.parameter_entries()}
    offset_entry = entries["offset"]
    vary_item = tab.parameter_table.item(offset_row, 4)

    assert events == ["update"]
    assert offset_entry.initial_value_expression == "*scale"
    assert offset_entry.value_expression is None
    assert offset_entry.value == pytest.approx(2.5)
    assert offset_entry.vary is True
    assert vary_item.checkState() == Qt.CheckState.Checked
    assert vary_item.flags() & Qt.ItemFlag.ItemIsUserCheckable


def test_prefit_tab_supports_dependent_parameter_expressions_with_vary_off(
    qapp,
):
    del qapp
    tab = PrefitTab()
    tab.populate_parameter_table(
        [
            PrefitParameterEntry(
                structure="A",
                motif="motif",
                name="scale",
                value=2.5,
                vary=True,
                minimum=0.1,
                maximum=10.0,
                category="fit",
            ),
            PrefitParameterEntry(
                structure="A",
                motif="motif",
                name="offset",
                value=0.0,
                vary=False,
                minimum=-10.0,
                maximum=10.0,
                category="fit",
            ),
        ]
    )

    offset_row = tab.find_parameter_row("offset")
    assert offset_row >= 0

    tab.parameter_table.item(offset_row, 3).setText("*scale")

    entries = {entry.name: entry for entry in tab.parameter_entries()}
    offset_entry = entries["offset"]
    vary_item = tab.parameter_table.item(offset_row, 4)

    assert offset_entry.value_expression == "*scale"
    assert offset_entry.initial_value_expression is None
    assert offset_entry.value == pytest.approx(2.5)
    assert offset_entry.vary is False
    assert vary_item.checkState() == Qt.CheckState.Unchecked
    assert vary_item.flags() & Qt.ItemFlag.ItemIsUserCheckable


def test_prefit_tab_caches_resolved_parameter_entries_until_table_changes(
    qapp,
    monkeypatch,
):
    del qapp
    tab = PrefitTab()
    tab.populate_parameter_table(
        [
            PrefitParameterEntry(
                structure="A",
                motif="motif",
                name="scale",
                value=2.5,
                vary=True,
                minimum=0.1,
                maximum=10.0,
                category="fit",
            ),
            PrefitParameterEntry(
                structure="A",
                motif="motif",
                name="offset",
                value=0.0,
                vary=False,
                minimum=-10.0,
                maximum=10.0,
                category="fit",
            ),
        ]
    )

    resolve_calls = {"count": 0}
    original_resolve = prefit_tab_module.resolve_prefit_parameter_entries

    def record_resolve(entries):
        resolve_calls["count"] += 1
        return original_resolve(entries)

    monkeypatch.setattr(
        prefit_tab_module,
        "resolve_prefit_parameter_entries",
        record_resolve,
    )

    first_entries = {entry.name: entry for entry in tab.parameter_entries()}
    second_entries = {entry.name: entry for entry in tab.parameter_entries()}

    assert resolve_calls["count"] == 0
    assert first_entries["scale"].value == pytest.approx(2.5)
    assert second_entries["offset"].value == pytest.approx(0.0)

    tab.parameter_table.item(tab.find_parameter_row("offset"), 3).setText(
        "*scale"
    )
    updated_entries = {entry.name: entry for entry in tab.parameter_entries()}

    assert resolve_calls["count"] == 1
    assert updated_entries["offset"].value_expression == "*scale"
    assert updated_entries["offset"].value == pytest.approx(2.5)


def test_prefit_tab_scrollable_parameter_supports_expression_seed_ranges(qapp):
    del qapp
    tab = PrefitTab()
    tab.show()
    tab.populate_parameter_table(
        [
            PrefitParameterEntry(
                structure="A",
                motif="motif",
                name="scale",
                value=2.5,
                vary=True,
                minimum=0.1,
                maximum=10.0,
                category="fit",
            ),
            PrefitParameterEntry(
                structure="A",
                motif="motif",
                name="offset",
                value=0.0,
                vary=True,
                minimum=-10.0,
                maximum=10.0,
                category="fit",
            ),
        ]
    )

    tab.auto_update_checkbox.setChecked(True)
    tab.scrollable_parameter_checkbox.setChecked(True)
    offset_row = tab.find_parameter_row("offset")
    assert offset_row >= 0
    tab.parameter_table.setCurrentCell(offset_row, 3)
    tab.parameter_table.item(offset_row, 3).setText("*scale")

    assert tab.parameter_scroll_panel.isVisible()
    assert tab.parameter_scroll_bar.isEnabled()
    assert "Initial expression seed" in tab.parameter_scroll_mode_label.text()
    assert (
        "no numeric range"
        not in tab.parameter_scroll_name_label.text().lower()
    )


def test_prefit_metrics_note_non_positive_model_points(qapp):
    del qapp
    evaluation = PrefitEvaluation(
        q_values=np.asarray([0.1, 0.2, 0.3], dtype=float),
        experimental_intensities=np.asarray([1.0, 0.8, 0.6], dtype=float),
        model_intensities=np.asarray([1.1, -0.1, 0.5], dtype=float),
        residuals=np.asarray([0.1, -0.9, -0.1], dtype=float),
    )

    metric_lines = PrefitTab._prefit_metric_lines(evaluation)

    assert "Model <= 0 at 1 q-points" in metric_lines


def test_prefit_tab_scrollable_parameter_uses_bounds_and_updates_value(qapp):
    tab = PrefitTab()
    tab.show()
    qapp.processEvents()
    events: list[str] = []
    tab.update_model_requested.connect(lambda: events.append("update"))
    tab.populate_parameter_table(
        [
            PrefitParameterEntry(
                structure="A",
                motif="motif",
                name="scale",
                value=1e-4,
                vary=True,
                minimum=1e-5,
                maximum=5e-3,
                category="fit",
            ),
            PrefitParameterEntry(
                structure="A",
                motif="motif",
                name="offset",
                value=0.0,
                vary=True,
                minimum=-20.0,
                maximum=30.0,
                category="fit",
            ),
        ]
    )

    assert not tab.scrollable_parameter_checkbox.isEnabled()
    tab.auto_update_checkbox.setChecked(True)
    assert tab.scrollable_parameter_checkbox.isEnabled()
    tab.scrollable_parameter_checkbox.setChecked(True)

    scale_row = tab.find_parameter_row("scale")
    offset_row = tab.find_parameter_row("offset")
    assert scale_row >= 0
    assert offset_row >= 0

    tab.parameter_table.setCurrentCell(scale_row, 3)
    qapp.processEvents()
    assert tab.parameter_scroll_panel.isVisible()
    assert tab.parameter_scroll_mode_label.text() == "Log scroll"

    tab.parameter_scroll_bar.setValue(tab.PARAMETER_SCROLL_RESOLUTION)
    qapp.processEvents()
    assert events
    assert float(
        tab.parameter_table.item(scale_row, 3).text()
    ) == pytest.approx(5e-3)

    tab.parameter_table.setCurrentCell(offset_row, 3)
    qapp.processEvents()
    assert tab.parameter_scroll_mode_label.text() == "Linear scroll"


def test_scrollable_parameter_preserves_manual_geometry_bounds(
    qapp,
    tmp_path,
):
    poly_project_dir, _poly_paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=poly_project_dir)
    window.compute_prefit_cluster_geometry()

    parameter_name = "a_eff_w0"
    parameter_row = window.prefit_tab.find_parameter_row(parameter_name)
    assert parameter_row >= 0

    original_value = float(
        window.prefit_tab.parameter_table.item(parameter_row, 3).text()
    )
    custom_minimum = max(original_value * 0.5, 1e-6)
    custom_maximum = original_value + 2.0

    window.prefit_tab.set_parameter_row(
        parameter_name,
        minimum=custom_minimum,
        maximum=custom_maximum,
        vary=True,
    )
    window.prefit_tab.auto_update_checkbox.setChecked(True)
    window.prefit_tab.scrollable_parameter_checkbox.setChecked(True)
    window.prefit_tab.parameter_table.setCurrentCell(parameter_row, 3)

    qapp.processEvents()
    window.prefit_tab.parameter_scroll_bar.setValue(
        window.prefit_tab.PARAMETER_SCROLL_RESOLUTION
    )
    qapp.processEvents()

    assert float(
        window.prefit_tab.parameter_table.item(parameter_row, 5).text()
    ) == pytest.approx(custom_minimum)
    assert float(
        window.prefit_tab.parameter_table.item(parameter_row, 6).text()
    ) == pytest.approx(custom_maximum)
    assert float(
        window.prefit_tab.parameter_table.item(parameter_row, 3).text()
    ) == pytest.approx(custom_maximum)
    window.close()


def test_run_prefit_keeps_manual_weight_value_outside_previous_bounds(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    monkeypatch.setattr(
        window,
        "_confirm_large_prefit_parameter_count",
        lambda _names: (True, False),
    )

    weight_row = window.prefit_tab.find_parameter_row("w0")
    assert weight_row >= 0
    assert (
        float(window.prefit_tab.parameter_table.item(weight_row, 6).text())
        < 0.9
    )
    window.prefit_tab.parameter_table.item(weight_row, 3).setText("0.9")

    window.run_prefit()

    assert float(
        window.prefit_tab.parameter_table.item(weight_row, 3).text()
    ) == pytest.approx(0.9)
    assert (
        float(window.prefit_tab.parameter_table.item(weight_row, 6).text())
        >= 0.9
    )
    window.close()


def test_run_prefit_returns_to_editor_when_large_parameter_warning_is_rejected(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    for row in range(min(4, window.prefit_tab.parameter_table.rowCount())):
        vary_item = window.prefit_tab.parameter_table.item(row, 4)
        assert vary_item is not None
        vary_item.setCheckState(Qt.CheckState.Checked)

    warnings: list[list[str]] = []
    monkeypatch.setattr(
        window,
        "_confirm_large_prefit_parameter_count",
        lambda names: warnings.append(list(names)) or (False, False),
    )
    monkeypatch.setattr(
        window.prefit_workflow,
        "run_fit",
        lambda *args, **kwargs: (_ for _ in ()).throw(
            AssertionError("run_fit should not be called")
        ),
    )

    window.run_prefit()

    assert warnings
    assert len(warnings[0]) > 3
    assert "Prefit canceled." in window.prefit_tab.output_box.toPlainText()
    window.close()


def test_run_prefit_can_suppress_large_parameter_warning_for_session(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    for row in range(min(4, window.prefit_tab.parameter_table.rowCount())):
        vary_item = window.prefit_tab.parameter_table.item(row, 4)
        assert vary_item is not None
        vary_item.setCheckState(Qt.CheckState.Checked)

    helper_calls: list[list[str]] = []
    monkeypatch.setattr(
        window,
        "_confirm_large_prefit_parameter_count",
        lambda names: helper_calls.append(list(names)) or (True, True),
    )

    run_calls = {"count": 0}

    def fake_run_fit(entries, *, method="leastsq", max_nfev=10000):
        del max_nfev
        run_calls["count"] += 1
        evaluation = window.prefit_workflow.evaluate(entries)
        residuals = np.asarray(evaluation.residuals, dtype=float)
        return PrefitFitResult(
            parameter_entries=entries,
            evaluation=evaluation,
            fit_report="fake lmfit report",
            method=method,
            nfev=0,
            chi_square=float(np.sum(residuals**2)),
            reduced_chi_square=float(np.mean(residuals**2)),
            r_squared=0.0,
            optimization_strategy=f"lmfit {method}",
            grid_evaluations=0,
        )

    monkeypatch.setattr(window.prefit_workflow, "run_fit", fake_run_fit)

    window.run_prefit()
    window.run_prefit()

    assert run_calls["count"] == 2
    assert len(helper_calls) == 1
    window.close()


def test_best_prefit_preset_saves_resets_and_reloads(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.prefit_tab.set_parameter_row(
        "scale",
        value=7e-4,
        minimum=1e-5,
        maximum=8e-3,
        vary=True,
    )
    window.prefit_tab.set_parameter_row("offset", value=0.125)
    window.set_best_prefit_parameters()

    settings = SAXSProjectManager().load_project(project_dir)
    assert settings.best_prefit_template == (
        "template_pd_likelihood_monosq_decoupled"
    )
    assert settings.template_reset_template == (
        "template_pd_likelihood_monosq_decoupled"
    )
    assert settings.best_prefit_parameter_entries
    assert settings.template_reset_parameter_entries
    dream_entries = {
        entry.param: entry
        for entry in SAXSDreamWorkflow(project_dir).load_parameter_map()
    }
    assert dream_entries["scale"].value == pytest.approx(7e-4)
    assert dream_entries["offset"].value == pytest.approx(0.125)

    window.prefit_tab.set_parameter_row("scale", value=2e-3)
    window.prefit_tab.set_parameter_row("offset", value=0.333)
    window.reset_parameters_to_best_prefit()

    best_entries = {
        entry.name: entry for entry in window.prefit_tab.parameter_entries()
    }
    assert best_entries["scale"].value == pytest.approx(7e-4)
    assert best_entries["offset"].value == pytest.approx(0.125)

    window.prefit_tab.set_parameter_row("scale", value=2e-3)
    window.reset_prefit_entries()

    template_entries = {
        entry.name: entry for entry in window.prefit_tab.parameter_entries()
    }
    assert template_entries["scale"].value == pytest.approx(5e-4)
    assert template_entries["offset"].value == pytest.approx(0.0)

    reloaded_window = SAXSMainWindow(initial_project_dir=project_dir)
    reloaded_entries = {
        entry.name: entry
        for entry in reloaded_window.prefit_tab.parameter_entries()
    }
    assert reloaded_entries["scale"].value == pytest.approx(7e-4)
    assert reloaded_entries["offset"].value == pytest.approx(0.125)


def test_set_best_prefit_retargets_existing_dream_prior_centers(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    dream_workflow = SAXSDreamWorkflow(project_dir)
    parameter_map = dream_workflow.create_default_parameter_map(persist=False)
    scale_entry = next(
        entry for entry in parameter_map if entry.param == "scale"
    )
    scale_entry.vary = False
    scale_entry.distribution = "norm"
    scale_entry.dist_params = {"loc": 1.1e-3, "scale": 2.0e-4}
    dream_workflow.save_parameter_map(parameter_map)

    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.prefit_tab.set_parameter_row("scale", value=7e-4)
    window.prefit_tab.set_parameter_row("offset", value=0.125)

    window.set_best_prefit_parameters()

    updated_entries = {
        entry.param: entry
        for entry in SAXSDreamWorkflow(project_dir).load_parameter_map()
    }
    updated_scale = updated_entries["scale"]
    assert updated_scale.value == pytest.approx(7e-4)
    assert updated_scale.vary is False
    assert updated_scale.distribution == "norm"
    assert updated_scale.dist_params["loc"] == pytest.approx(7e-4)
    assert updated_scale.dist_params["scale"] == pytest.approx(2.0e-4)
    assert window._dream_parameter_map_saved_in_session is True
    window.close()


def test_individual_prefit_parameter_reset_button_restores_template_default(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.prefit_tab.set_parameter_row(
        "scale",
        value=2e-3,
        minimum=1e-6,
        maximum=9e-3,
        vary=True,
    )
    window.prefit_tab.set_parameter_row("offset", value=0.125)

    scale_row = window.prefit_tab.find_parameter_row("scale")
    assert scale_row >= 0
    reset_col = _table_column_index(window.prefit_tab.parameter_table, "Reset")
    reset_button = window.prefit_tab.parameter_table.cellWidget(
        scale_row,
        reset_col,
    )
    assert reset_button is not None

    reset_button.click()

    entries = {
        entry.name: entry for entry in window.prefit_tab.parameter_entries()
    }
    assert entries["scale"].value == pytest.approx(5e-4)
    assert entries["scale"].minimum == pytest.approx(1e-5)
    assert entries["scale"].maximum == pytest.approx(5e-3)
    assert entries["scale"].vary is False
    assert entries["offset"].value == pytest.approx(0.125)
    window.close()


def test_restore_prefit_state_recovers_saved_parameters_and_run_config(
    qapp, tmp_path
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.prefit_tab.set_parameter_row(
        "scale",
        value=7e-4,
        minimum=1e-5,
        maximum=8e-3,
        vary=True,
    )
    window.prefit_tab.set_parameter_row("offset", value=0.125)
    window.prefit_tab.set_run_config(method="powell", max_nfev=4321)
    window.prefit_tab.set_autosave(True)
    window.set_best_prefit_parameters()
    window.save_prefit()

    assert window.prefit_tab.restore_state_button.isEnabled()
    saved_state_name = window.prefit_tab.selected_saved_state_name()
    assert saved_state_name is not None
    assert saved_state_name.startswith("prefit_")
    assert "R^2=" in window.prefit_tab.saved_state_combo.currentText()

    window.prefit_tab.set_parameter_row("scale", value=2e-3)
    window.prefit_tab.set_parameter_row("offset", value=0.333)
    window.prefit_tab.set_run_config(method="nelder", max_nfev=9999)
    window.prefit_tab.set_autosave(False)
    window.restore_prefit_state()

    restored_entries = {
        entry.name: entry for entry in window.prefit_tab.parameter_entries()
    }
    assert restored_entries["scale"].value == pytest.approx(7e-4)
    assert restored_entries["offset"].value == pytest.approx(0.125)
    assert window.prefit_tab.run_config().method == "powell"
    assert window.prefit_tab.run_config().max_nfev == 4321
    assert window.prefit_tab.autosave_checkbox.isChecked()

    window.prefit_tab.set_parameter_row("scale", value=1.5e-3)
    window.reset_parameters_to_best_prefit()
    best_entries = {
        entry.name: entry for entry in window.prefit_tab.parameter_entries()
    }
    assert best_entries["scale"].value == pytest.approx(7e-4)


def test_restore_prefit_state_restores_cluster_geometry_mode_and_shape(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_poly_lma_geometry_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.compute_prefit_cluster_geometry()

    table = window.prefit_tab.cluster_geometry_table
    sf_column = _table_column_index(table, "S.F. Approx.")
    sf_combo = table.cellWidget(0, sf_column)
    assert sf_combo is not None

    window.prefit_tab.toggle_cluster_geometry_radii_button.click()
    ellipsoid_index = sf_combo.findData("ellipsoid")
    assert ellipsoid_index >= 0
    sf_combo.setCurrentIndex(ellipsoid_index)
    window.save_prefit()

    window.prefit_tab.toggle_cluster_geometry_radii_button.click()
    sphere_index = sf_combo.findData("sphere")
    assert sphere_index >= 0
    sf_combo.setCurrentIndex(sphere_index)

    window.restore_prefit_state()

    restored_table = window.prefit_tab.cluster_geometry_table
    restored_sf_combo = restored_table.cellWidget(0, sf_column)
    assert restored_sf_combo is not None
    assert (
        window.prefit_tab.cluster_geometry_radii_type_combo.currentData()
        == "bond_length"
    )
    assert str(restored_sf_combo.currentData()) == "ellipsoid"
    restored_row = window.prefit_workflow.cluster_geometry_rows()[0]
    assert restored_row.radii_type_used == "bond_length"
    assert restored_row.sf_approximation == "ellipsoid"

    window.close()


def test_generate_prior_weights_does_not_force_prefit_without_components(
    qapp, tmp_path, monkeypatch
):
    del qapp
    window = SAXSMainWindow()
    window.project_setup_tab.project_name_edit.setText("demo_project")
    window.project_setup_tab.project_dir_edit.setText(str(tmp_path))
    window.create_project_from_tab()

    scheduled: dict[str, object] = {}

    monkeypatch.setattr(
        window,
        "_start_project_task",
        lambda task_name, task_fn, start_message, settings=None: scheduled.update(
            {
                "task_name": task_name,
                "task_fn": task_fn,
                "start_message": start_message,
                "settings": settings,
            }
        ),
    )
    saved_paths: list[Path] = []
    monkeypatch.setattr(
        window.project_manager,
        "save_project",
        lambda settings: saved_paths.append(
            build_project_paths(settings.project_dir).project_file
        )
        or build_project_paths(settings.project_dir).project_file,
    )

    window.build_prior_weights()

    assert scheduled["task_name"] == "build_prior_weights"
    assert scheduled["start_message"] == "Generating prior weights..."
    assert scheduled["settings"] is window.current_settings
    assert saved_paths


def test_generate_plot_opens_standalone_prior_histogram_window(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.show_prior_histogram_window()

    prior_window = window._prior_histogram_windows[-1]
    assert isinstance(prior_window, PriorHistogramWindow)
    assert prior_window.parent() is None
    assert prior_window.isWindow()
    assert prior_window.toolbar is not None


def test_solvent_sort_mode_populates_secondary_atom_filter(qapp):
    del qapp
    tab = ProjectSetupTab()
    tab.set_project_selected(True)
    tab.apply_cluster_import_data(
        ["Pb", "I", "O"],
        [
            {
                "structure": "PbI2",
                "motif": "motif_A",
                "count": 2,
                "weight": 0.4,
                "atom_fraction_percent": 40.0,
                "structure_fraction_percent": 40.0,
            },
            {
                "structure": "Pb2I4",
                "motif": "motif_B",
                "count": 3,
                "weight": 0.6,
                "atom_fraction_percent": 60.0,
                "structure_fraction_percent": 60.0,
            },
        ],
    )

    tab.prior_mode_combo.setCurrentText("Solvent Sort - Structure Fraction")

    assert not tab.secondary_filter_combo.isHidden()
    assert tab.secondary_filter_combo.isEnabled()
    assert tab.secondary_filter_combo.count() == 1
    assert tab.secondary_filter_combo.currentText() == "O"


def test_save_prior_plot_png_writes_to_exported_results_plots(qapp, tmp_path):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.save_prior_plot_png()

    saved_images = list(paths.exported_plots_dir.glob("prior_histogram_*.png"))
    assert saved_images


def test_save_component_plot_data_exports_csv(qapp, tmp_path, monkeypatch):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    export_path = tmp_path / "component_plot_export.csv"

    monkeypatch.setattr(
        QFileDialog,
        "getSaveFileName",
        lambda *args, **kwargs: (
            str(export_path),
            "CSV files (*.csv)",
        ),
    )

    window.save_component_plot_data()

    assert export_path.is_file()
    contents = export_path.read_text(encoding="utf-8")
    assert "series,component_key,axis_index,axis_ylabel,color,visible,x,y" in (
        contents
    )
    assert "Experimental data" in contents
    assert "A_no_motif" in contents


def test_save_prior_plot_data_exports_npy(qapp, tmp_path, monkeypatch):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    export_path = tmp_path / "prior_plot_export.npy"

    monkeypatch.setattr(
        QFileDialog,
        "getSaveFileName",
        lambda *args, **kwargs: (
            str(export_path),
            "NumPy files (*.npy)",
        ),
    )

    window.save_prior_plot_data()

    assert export_path.is_file()
    payload = np.load(export_path, allow_pickle=True).item()
    assert payload["plot_mode"] == "structure_fraction"
    assert list(payload["labels"]) == ["A"]
    np.testing.assert_allclose(payload["matrix"], np.asarray([[100.0]]))


def test_prior_histogram_window_uses_secondary_filter_and_colormap(
    qapp, tmp_path
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    (paths.project_dir / "md_prior_weights.json").write_text(
        json.dumps(
            {
                "origin": "clusters",
                "total_files": 5,
                "available_elements": ["Pb", "I", "O"],
                "structures": {
                    "PbI2": {
                        "motif_A": {
                            "count": 2,
                            "weight": 0.4,
                            "profile_file": "A_no_motif.txt",
                            "secondary_atom_distributions": {
                                "O": {"0": 1, "1": 1}
                            },
                        }
                    },
                    "Pb2I4": {
                        "motif_B": {
                            "count": 3,
                            "weight": 0.6,
                            "profile_file": "A_no_motif.txt",
                            "secondary_atom_distributions": {
                                "O": {"0": 1, "2": 2}
                            },
                        }
                    },
                },
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.project_setup_tab.apply_cluster_import_data(
        ["Pb", "I", "O"],
        [
            {
                "structure": "PbI2",
                "motif": "motif_A",
                "count": 2,
                "weight": 0.4,
                "atom_fraction_percent": 40.0,
                "structure_fraction_percent": 40.0,
            },
            {
                "structure": "Pb2I4",
                "motif": "motif_B",
                "count": 3,
                "weight": 0.6,
                "atom_fraction_percent": 60.0,
                "structure_fraction_percent": 60.0,
            },
        ],
    )
    window.project_setup_tab.prior_mode_combo.setCurrentText(
        "Solvent Sort - Structure Fraction"
    )
    window.project_setup_tab.prior_color_combo.setCurrentText("viridis")

    window.show_prior_histogram_window()

    prior_window = window._prior_histogram_windows[-1]
    assert prior_window.mode == "solvent_sort_structure_fraction"
    assert prior_window.secondary_element == "O"
    assert prior_window.cmap == "viridis"


def test_prior_histogram_can_match_component_trace_colors(qapp, tmp_path):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    component_path = paths.scattering_components_dir / "A_no_motif.txt"

    plasma_index = (
        window.project_setup_tab.component_trace_color_scheme_combo.findData(
            "plasma"
        )
    )
    window.project_setup_tab.component_trace_color_scheme_combo.setCurrentIndex(
        plasma_index
    )
    window.project_setup_tab.draw_component_plot([component_path])
    window.project_setup_tab.draw_prior_plot(
        paths.project_dir / "md_prior_weights.json"
    )

    trace_color = to_hex(
        window.project_setup_tab._component_line_lookup[
            "A_no_motif"
        ].get_color()
    )
    bar_color = to_hex(
        window.project_setup_tab.prior_figure.axes[0]
        .patches[0]
        .get_facecolor()
    )

    assert (
        not window.project_setup_tab.prior_match_trace_colors_checkbox.isHidden()
    )
    assert bar_color == trace_color

    window.show_prior_histogram_window()

    prior_window = window._prior_histogram_windows[-1]
    assert prior_window.structure_motif_colors == {"A_no_motif": trace_color}


def test_prior_histogram_predicted_segments_match_component_trace_colors(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    _write_predicted_structure_artifacts(paths)
    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.use_predicted_structure_weights = True
    manager.save_project(settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)
    plasma_index = (
        window.project_setup_tab.component_trace_color_scheme_combo.findData(
            "plasma"
        )
    )
    window.project_setup_tab.component_trace_color_scheme_combo.setCurrentIndex(
        plasma_index
    )
    window.project_setup_tab.draw_prior_plot(
        window.project_setup_tab.current_prior_json_path()
    )

    observed_trace_color = to_hex(
        window.project_setup_tab._component_line_lookup[
            "A_no_motif"
        ].get_color()
    )
    predicted_trace_color = to_hex(
        window.project_setup_tab._component_line_lookup[
            "A2_predicted_rank01"
        ].get_color()
    )
    patches = [
        patch
        for patch in window.project_setup_tab.prior_figure.axes[0].patches
        if float(patch.get_height()) > 0.0
    ]

    assert len(patches) >= 2
    assert to_hex(patches[0].get_facecolor()) == observed_trace_color
    assert to_hex(patches[1].get_facecolor()) == predicted_trace_color


def test_solvent_sort_prior_histogram_does_not_auto_match_component_colors(
    qapp, tmp_path
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    component_path = paths.scattering_components_dir / "A_no_motif.txt"

    window.project_setup_tab.apply_cluster_import_data(
        ["A", "O"],
        [
            {
                "structure": "A",
                "motif": "no_motif",
                "count": 1,
                "weight": 1.0,
                "atom_fraction_percent": 100.0,
                "structure_fraction_percent": 100.0,
            }
        ],
    )
    window.project_setup_tab.prior_mode_combo.setCurrentText(
        "Solvent Sort - Structure Fraction"
    )
    plasma_index = (
        window.project_setup_tab.component_trace_color_scheme_combo.findData(
            "plasma"
        )
    )
    window.project_setup_tab.component_trace_color_scheme_combo.setCurrentIndex(
        plasma_index
    )
    window.project_setup_tab.prior_match_trace_colors_checkbox.setChecked(True)
    window.project_setup_tab.draw_component_plot([component_path])

    assert window.project_setup_tab.prior_structure_motif_colors() is None

    window.show_prior_histogram_window()

    prior_window = window._prior_histogram_windows[-1]
    assert prior_window.mode == "solvent_sort_structure_fraction"
    assert prior_window.structure_motif_colors is None


def test_create_project_warns_before_overwriting_existing_folder(
    qapp, tmp_path, monkeypatch
):
    del qapp
    existing_project_dir = tmp_path / "existing_project"
    existing_project_dir.mkdir()
    window = SAXSMainWindow()
    window.project_setup_tab.project_name_edit.setText("existing_project")
    window.project_setup_tab.project_dir_edit.setText(str(tmp_path))

    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.QMessageBox.warning",
        lambda *args, **kwargs: QMessageBox.StandardButton.No,
    )

    window.create_project_from_tab()

    assert window.current_settings is None
    assert not (existing_project_dir / "saxs_project.json").exists()


def test_selecting_clusters_directory_triggers_project_autosave(
    qapp, tmp_path, monkeypatch
):
    del qapp
    window = SAXSMainWindow()
    window.project_setup_tab.project_name_edit.setText("demo_project")
    window.project_setup_tab.project_dir_edit.setText(str(tmp_path))
    window.create_project_from_tab()

    clusters_dir = tmp_path / "clusters"
    (clusters_dir / "PbI2").mkdir(parents=True)
    (clusters_dir / "PbI2" / "frame_0001.xyz").write_text(
        "2\ncomment\nPb 0.0 0.0 0.0\nI 1.0 0.0 0.0\n",
        encoding="utf-8",
    )

    monkeypatch.setattr(
        "saxshell.saxs.ui.project_setup_tab.QFileDialog.getExistingDirectory",
        lambda *args, **kwargs: str(clusters_dir),
    )
    monkeypatch.setattr(
        window, "_start_project_task", lambda *args, **kwargs: None
    )
    saved_clusters: list[str | None] = []
    monkeypatch.setattr(
        window.project_manager,
        "save_project",
        lambda settings: saved_clusters.append(settings.clusters_dir)
        or build_project_paths(settings.project_dir).project_file,
    )

    window.project_setup_tab._choose_clusters_directory()

    assert saved_clusters[-1] == str(clusters_dir)


def test_selecting_frames_directory_triggers_project_autosave(
    qapp, tmp_path, monkeypatch
):
    del qapp
    window = SAXSMainWindow()
    window.project_setup_tab.project_name_edit.setText("demo_project")
    window.project_setup_tab.project_dir_edit.setText(str(tmp_path))
    window.create_project_from_tab()

    frames_dir = tmp_path / "splitxyz_f0fs"
    frames_dir.mkdir()
    (frames_dir / "frame_0000.xyz").write_text(
        "1\nframe\nPb 0.0 0.0 0.0\n",
        encoding="utf-8",
    )

    monkeypatch.setattr(
        "saxshell.saxs.ui.project_setup_tab.QFileDialog.getExistingDirectory",
        lambda *args, **kwargs: str(frames_dir),
    )
    saved_frames: list[str | None] = []
    monkeypatch.setattr(
        window.project_manager,
        "save_project",
        lambda settings: saved_frames.append(settings.frames_dir)
        or build_project_paths(settings.project_dir).project_file,
    )

    window.project_setup_tab._choose_frames_directory()

    assert saved_frames[-1] == str(frames_dir)


def test_selecting_pdb_frames_directory_triggers_project_autosave(
    qapp, tmp_path, monkeypatch
):
    del qapp
    window = SAXSMainWindow()
    window.project_setup_tab.project_name_edit.setText("demo_project")
    window.project_setup_tab.project_dir_edit.setText(str(tmp_path))
    window.create_project_from_tab()

    pdb_frames_dir = tmp_path / "xyz2pdb_splitxyz_f0fs"
    pdb_frames_dir.mkdir()

    monkeypatch.setattr(
        "saxshell.saxs.ui.project_setup_tab.QFileDialog.getExistingDirectory",
        lambda *args, **kwargs: str(pdb_frames_dir),
    )
    saved_pdb_frames: list[str | None] = []
    monkeypatch.setattr(
        window.project_manager,
        "save_project",
        lambda settings: saved_pdb_frames.append(settings.pdb_frames_dir)
        or build_project_paths(settings.project_dir).project_file,
    )

    window.project_setup_tab._choose_pdb_frames_directory()

    assert saved_pdb_frames[-1] == str(pdb_frames_dir)


def test_loading_project_reports_registered_folder_warnings(qapp, tmp_path):
    del qapp
    manager = SAXSProjectManager()
    project_dir = tmp_path / "saxs_project"
    settings = manager.create_project(project_dir)
    frames_dir = tmp_path / "splitxyz_f0fs"
    frames_dir.mkdir()
    (frames_dir / "frame_0000.xyz").write_text(
        "1\nframe\nPb 0.0 0.0 0.0\n",
        encoding="utf-8",
    )
    clusters_dir = tmp_path / "clusters"
    (clusters_dir / "PbI2").mkdir(parents=True)
    (clusters_dir / "PbI2" / "frame_0001.xyz").write_text(
        "2\ncomment\nPb 0.0 0.0 0.0\nI 1.0 0.0 0.0\n",
        encoding="utf-8",
    )
    settings.frames_dir = str(frames_dir)
    settings.clusters_dir = str(clusters_dir)
    trajectory_file = tmp_path / "traj.xyz"
    topology_file = tmp_path / "topology.pdb"
    energy_file = tmp_path / "traj.ener"
    trajectory_file.write_text(
        "1\nframe\nPb 0.0 0.0 0.0\n",
        encoding="utf-8",
    )
    topology_file.write_text("MODEL        1\nENDMDL\n", encoding="utf-8")
    energy_file.write_text(
        "# step time kinetic temperature potential\n"
        "1 0.0 1.0 300.0 -10.0\n",
        encoding="utf-8",
    )
    settings.trajectory_file = str(trajectory_file)
    settings.topology_file = str(topology_file)
    settings.energy_file = str(energy_file)
    manager.save_project(settings)

    (frames_dir / "frame_0001.xyz").write_text(
        "1\nframe\nI 1.0 0.0 0.0\n",
        encoding="utf-8",
    )
    (clusters_dir / "PbI2" / "frame_0001.xyz").unlink()
    (clusters_dir / "PbI2").rmdir()
    clusters_dir.rmdir()
    trajectory_file.write_text(
        "1\nframe\nI 1.0 0.0 0.0\n",
        encoding="utf-8",
    )
    topology_file.unlink()
    energy_file.write_text(
        "# step time kinetic temperature potential\n" "1 0.0 1.2 305.0 -9.5\n",
        encoding="utf-8",
    )

    window = SAXSMainWindow(initial_project_dir=project_dir)
    summary_text = window.project_setup_tab.summary_box.toPlainText()

    assert "Loaded project saxs_project" in summary_text
    assert "Registered frames folder contents changed" in summary_text
    assert "Registered clusters folder is missing" in summary_text
    assert "Registered trajectory file changed" in summary_text
    assert "Registered topology file is missing" in summary_text
    assert "Registered energy file changed" in summary_text
    window.close()


def test_selecting_experimental_file_triggers_project_autosave(
    qapp, tmp_path, monkeypatch
):
    del qapp
    window = SAXSMainWindow()
    window.project_setup_tab.project_name_edit.setText("demo_project")
    window.project_setup_tab.project_dir_edit.setText(str(tmp_path))
    window.create_project_from_tab()

    data_path = tmp_path / "exp_demo.txt"
    np.savetxt(
        data_path,
        np.column_stack(
            [
                np.asarray([0.05, 0.10], dtype=float),
                np.asarray([10.0, 9.5], dtype=float),
            ]
        ),
    )
    monkeypatch.setattr(
        "saxshell.saxs.ui.project_setup_tab.QFileDialog.getOpenFileName",
        lambda *args, **kwargs: (
            str(data_path),
            "Data files (*.txt *.dat *.iq)",
        ),
    )
    saved_data_paths: list[str | None] = []
    monkeypatch.setattr(
        window.project_manager,
        "save_project",
        lambda settings: saved_data_paths.append(
            settings.experimental_data_path
        )
        or build_project_paths(settings.project_dir).project_file,
    )

    window.project_setup_tab._choose_experimental_file()

    assert saved_data_paths[-1] == str(data_path)


def test_experimental_data_header_dialog_recovers_header_rows(qapp, tmp_path):
    del qapp
    data_path = tmp_path / "exp_with_header.txt"
    data_path.write_text(
        "Example SAXS export\n"
        "q intensity sigma\n"
        "0.05 10.0 0.1\n"
        "0.10 9.5 0.1\n",
        encoding="utf-8",
    )

    dialog = ExperimentalDataHeaderDialog(data_path)
    dialog.header_rows_spin.setValue(2)
    dialog._try_accept()

    assert dialog.accepted_summary is not None
    assert dialog.header_rows() == 2
    assert dialog.accepted_summary.header_rows == 2
    assert np.allclose(dialog.accepted_summary.q_values, [0.05, 0.10])
    assert np.allclose(
        dialog.accepted_summary.intensities,
        [10.0, 9.5],
    )


def test_load_experimental_data_file_detects_three_column_headers(
    tmp_path,
):
    data_path = tmp_path / "exp_three_columns.txt"
    data_path.write_text(
        "q_demo\tintensity_demo\terror_demo\n"
        "0.05\t10.0\t0.1\n"
        "0.10\t9.5\t0.2\n",
        encoding="utf-8",
    )

    summary = load_experimental_data_file(data_path)

    assert summary.header_rows == 1
    assert summary.column_names == [
        "q_demo",
        "intensity_demo",
        "error_demo",
    ]
    assert summary.q_column == 0
    assert summary.intensity_column == 1
    assert summary.error_column == 2
    assert summary.errors is not None
    assert np.allclose(summary.q_values, [0.05, 0.10])
    assert np.allclose(summary.intensities, [10.0, 9.5])
    assert np.allclose(summary.errors, [0.1, 0.2])


def test_experimental_data_header_dialog_allows_manual_column_selection(
    qapp, tmp_path
):
    del qapp
    data_path = tmp_path / "exp_swapped_columns.txt"
    data_path.write_text(
        "intensity_value\tq_value\tsigma_value\n"
        "10.0\t0.05\t0.1\n"
        "9.5\t0.10\t0.2\n",
        encoding="utf-8",
    )

    dialog = ExperimentalDataHeaderDialog(data_path)
    dialog.header_rows_spin.setValue(1)
    dialog.q_column_combo.setCurrentIndex(dialog.q_column_combo.findData(1))
    dialog.intensity_column_combo.setCurrentIndex(
        dialog.intensity_column_combo.findData(0)
    )
    dialog.error_column_combo.setCurrentIndex(
        dialog.error_column_combo.findData(2)
    )
    dialog._try_accept()

    assert dialog.accepted_summary is not None
    assert dialog.accepted_summary.column_names == [
        "intensity_value",
        "q_value",
        "sigma_value",
    ]
    assert dialog.q_column() == 1
    assert dialog.intensity_column() == 0
    assert dialog.error_column() == 2
    assert np.allclose(dialog.accepted_summary.q_values, [0.05, 0.10])
    assert np.allclose(
        dialog.accepted_summary.intensities,
        [10.0, 9.5],
    )
    assert np.allclose(dialog.accepted_summary.errors, [0.1, 0.2])


def test_project_setup_preview_updates_with_experimental_q_range(
    qapp, tmp_path
):
    del qapp
    data_path = tmp_path / "exp_preview.txt"
    np.savetxt(
        data_path,
        np.column_stack(
            [
                np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float),
                np.asarray([100.0, 80.0, 55.0, 30.0], dtype=float),
            ]
        ),
    )

    summary = load_experimental_data_file(data_path)
    tab = ProjectSetupTab()
    tab._apply_experimental_file(data_path, summary)
    tab.draw_component_plot(None)

    preview_axis = tab.component_figure.axes[0]
    assert tab.component_log_x_checkbox.isChecked()
    assert tab.component_log_y_checkbox.isChecked()
    assert preview_axis.get_title() == "Experimental Data Preview"
    assert preview_axis.get_xlabel() == "q (Å⁻¹)"
    assert preview_axis.get_ylabel() == "Intensity (arb. units)"
    assert preview_axis.get_xscale() == "log"
    assert preview_axis.get_yscale() == "log"

    tab.qmin_edit.setText("0.08")
    tab.qmax_edit.setText("0.12")
    tab.component_log_x_checkbox.setChecked(False)
    tab.draw_component_plot(None)

    preview_axis = tab.component_figure.axes[0]
    legend_labels = preview_axis.get_legend_handles_labels()[1]
    assert "Selected q-range" in legend_labels
    assert preview_axis.get_xscale() == "linear"
    assert preview_axis.get_yscale() == "log"


def test_save_project_state_reloads_prefit_and_dream_for_reduced_q_range(
    qapp, tmp_path
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.project_setup_tab.qmin_edit.setText("0.12")
    window.project_setup_tab.qmax_edit.setText("0.19")
    window.save_project_state()

    assert window.prefit_workflow is not None
    assert np.allclose(
        window.prefit_workflow.evaluate().q_values,
        [0.12142857142857144, 0.15714285714285714, 0.19285714285714284],
    )
    assert window.dream_workflow is not None
    assert np.allclose(
        window.dream_workflow.prefit_workflow.evaluate().q_values,
        [0.12142857142857144, 0.15714285714285714, 0.19285714285714284],
    )
    window.close()


def test_save_project_state_preserves_live_prefit_parameters(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    window.prefit_tab.set_parameter_row(
        "scale",
        value=7e-4,
        minimum=1e-5,
        maximum=8e-3,
        vary=True,
    )
    window.prefit_tab.set_parameter_row("offset", value=0.125)
    window.set_best_prefit_parameters()

    window.prefit_tab.set_parameter_row("scale", value=2e-3)
    window.prefit_tab.set_parameter_row("offset", value=0.333)
    window.project_setup_tab.qmin_edit.setText("0.12")
    window.project_setup_tab.qmax_edit.setText("0.19")

    window.save_project_state()

    entries = {
        entry.name: entry for entry in window.prefit_tab.parameter_entries()
    }
    assert entries["scale"].value == pytest.approx(2e-3)
    assert entries["offset"].value == pytest.approx(0.333)
    assert window.prefit_workflow is not None
    assert np.allclose(
        window.prefit_workflow.evaluate().q_values,
        [0.12142857142857144, 0.15714285714285714, 0.19285714285714284],
    )
    window.close()


def test_save_project_state_updates_workflows_in_place_without_resetting_tabs(
    qapp, tmp_path
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert window.prefit_workflow is not None
    assert window.dream_workflow is not None
    original_prefit_workflow = window.prefit_workflow
    original_dream_workflow = window.dream_workflow

    window.prefit_tab.set_parameter_row(
        "scale",
        value=1.7e-3,
        minimum=1e-5,
        maximum=8e-3,
        vary=True,
    )
    window.prefit_tab.set_parameter_row("offset", value=0.222)
    window.dream_tab.iterations_spin.setValue(4321)
    window.dream_tab.nseedchains_spin.setValue(33)
    window.project_setup_tab.qmin_edit.setText("0.12")
    window.project_setup_tab.qmax_edit.setText("0.19")

    window.save_project_state()

    assert window.prefit_workflow is original_prefit_workflow
    assert window.dream_workflow is original_dream_workflow
    entries = {
        entry.name: entry for entry in window.prefit_tab.parameter_entries()
    }
    assert entries["scale"].value == pytest.approx(1.7e-3)
    assert entries["offset"].value == pytest.approx(0.222)
    assert window.dream_tab.settings_payload().niterations == 4321
    assert window.dream_tab.settings_payload().nseedchains == 33
    assert np.allclose(
        window.prefit_workflow.evaluate().q_values,
        [0.12142857142857144, 0.15714285714285714, 0.19285714285714284],
    )
    assert np.allclose(
        window.dream_workflow.prefit_workflow.evaluate().q_values,
        [0.12142857142857144, 0.15714285714285714, 0.19285714285714284],
    )
    window.close()


def test_save_project_state_warns_when_q_range_expands_beyond_components(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    captured: dict[str, str] = {}

    monkeypatch.setattr(
        window,
        "_show_error",
        lambda title, message: captured.update(
            {
                "title": title,
                "message": message,
            }
        ),
    )

    window.project_setup_tab.qmin_edit.setText("0.04")
    window.project_setup_tab.qmax_edit.setText("0.31")
    window.save_project_state()

    assert (
        captured["title"]
        == "Expanded q-range requires rebuilding SAXS components"
    )
    assert "Recompute the SAXS model components" in captured["message"]
    window.close()


def test_build_components_warns_when_q_range_is_still_default(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    warnings: list[tuple[str, str]] = []
    start_calls: list[str] = []

    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.QMessageBox.warning",
        lambda _parent, title, message, *args, **kwargs: warnings.append(
            (title, message)
        )
        or QMessageBox.StandardButton.No,
    )
    monkeypatch.setattr(
        window,
        "_start_project_task",
        lambda task_name, *args, **kwargs: start_calls.append(task_name),
    )

    window.build_project_components()

    assert warnings
    assert warnings[0][0] == "Build SAXS components with default q-range?"
    assert "full experimental-data default" in warnings[0][1]
    assert not start_calls
    window.close()


def test_build_components_auto_refreshes_component_plot_for_built_distribution(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    external_data_dir = tmp_path / "external_data"
    external_data_dir.mkdir(parents=True, exist_ok=True)
    external_data_path = external_data_dir / "exp_external.txt"
    external_data_path.write_text(
        (paths.experimental_data_dir / "exp_demo.txt").read_text(
            encoding="utf-8"
        ),
        encoding="utf-8",
    )

    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    clusters_dir = paths.project_dir / "clusters"
    structure_dir = clusters_dir / "A"
    structure_dir.mkdir(parents=True, exist_ok=True)
    (structure_dir / "frame_0001.xyz").write_text(
        "1\nframe 1\nA 0.0 0.0 0.0\n",
        encoding="utf-8",
    )
    settings.clusters_dir = str(clusters_dir)
    settings.experimental_data_path = str(external_data_path)
    settings.copied_experimental_data_file = None
    manager.save_project(settings)

    component_trace = np.linspace(10.0, 17.0, 8)

    def fake_build_profiles(
        builder,
        *,
        cluster_bins,
        progress_callback=None,
        progress_total=None,
        **kwargs,
    ):
        del progress_callback, progress_total, kwargs
        output_path = builder.output_dir / "A_no_motif.txt"
        _write_component_file(
            output_path,
            builder.q_values,
            component_trace,
        )
        cluster_bin = cluster_bins[0]
        return [
            AveragedComponent(
                structure=cluster_bin.structure,
                motif=cluster_bin.motif,
                file_count=len(cluster_bin.files),
                representative=cluster_bin.representative,
                source_dir=cluster_bin.source_dir,
                q_values=np.asarray(builder.q_values, dtype=float),
                mean_intensity=np.asarray(component_trace, dtype=float),
                std_intensity=np.zeros_like(builder.q_values, dtype=float),
                se_intensity=np.zeros_like(builder.q_values, dtype=float),
                output_path=output_path,
            )
        ]

    monkeypatch.setattr(
        project_module.DebyeProfileBuilder,
        "build_profiles",
        fake_build_profiles,
    )

    window = SAXSMainWindow(initial_project_dir=project_dir)
    monkeypatch.setattr(
        window,
        "_confirm_default_q_range_for_component_build",
        lambda: True,
    )

    def run_task_sync(
        task_name,
        task_fn,
        *,
        start_message,
        settings=None,
    ):
        del start_message, settings
        result = task_fn(lambda *_args, **_kwargs: None)
        window._on_task_finished(task_name, result)

    monkeypatch.setattr(window, "_start_project_task", run_task_sync)

    expected_distribution_id = project_module.distribution_id_for_settings(
        window._settings_from_project_tab()
    )

    window.build_project_components()
    QApplication.processEvents()

    assert window.project_setup_tab.computed_distribution_combo.count() == 1
    assert (
        window.project_setup_tab.selected_distribution_id()
        == expected_distribution_id
    )
    assert "A_no_motif" in window.project_setup_tab._component_line_lookup
    np.testing.assert_allclose(
        window.project_setup_tab._component_line_lookup[
            "A_no_motif"
        ].get_ydata(),
        component_trace,
    )
    window.close()


def test_build_components_in_contrast_mode_launches_scaffold_instead_of_builder_task(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    start_calls: list[str] = []
    launched: list[str] = []

    monkeypatch.setattr(
        window,
        "_confirm_default_q_range_for_component_build",
        lambda: True,
    )
    monkeypatch.setattr(
        window,
        "_start_project_task",
        lambda task_name, *args, **kwargs: start_calls.append(task_name),
    )
    monkeypatch.setattr(
        window,
        "_open_contrast_mode_tool",
        lambda **kwargs: launched.append(kwargs),
    )

    window.project_setup_tab.set_component_build_mode(
        COMPONENT_BUILD_MODE_CONTRAST
    )
    window.build_project_components()

    saved_settings = SAXSProjectManager().load_project(project_dir)
    assert not start_calls
    assert launched == [{"preview_mode": False}]
    assert saved_settings.component_build_mode == COMPONENT_BUILD_MODE_CONTRAST
    assert window.current_settings is not None
    assert (
        window.current_settings.component_build_mode
        == COMPONENT_BUILD_MODE_CONTRAST
    )
    assert (
        "Contrast (Debye)"
        in window.project_setup_tab.summary_box.toPlainText()
    )
    window.close()


def test_build_components_in_born_approximation_launches_electron_density_workflow(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    launched: list[dict[str, object]] = []
    start_calls: list[str] = []
    clusters_dir = tmp_path / "clusters"
    (clusters_dir / "PbI2").mkdir(parents=True)
    (clusters_dir / "PbI2" / "frame_0001.xyz").write_text(
        "\n".join(
            [
                "3",
                "frame_0001",
                "Pb 0.0 0.0 0.0",
                "I 2.8 0.0 0.0",
                "I 0.0 2.8 0.0",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    window.project_setup_tab.clusters_dir_edit.setText(
        str(clusters_dir.resolve())
    )
    window.project_setup_tab.set_component_build_mode(
        COMPONENT_BUILD_MODE_BORN_APPROXIMATION
    )
    monkeypatch.setattr(
        window,
        "_confirm_default_q_range_for_component_build",
        lambda: True,
    )
    monkeypatch.setattr(
        window,
        "_start_project_task",
        lambda task_name, *args, **kwargs: start_calls.append(task_name),
    )
    monkeypatch.setattr(
        window,
        "_open_electron_density_mapping_tool",
        lambda **kwargs: launched.append(kwargs),
    )

    window.build_project_components()

    assert not start_calls
    assert launched == [{"preview_mode": False}]
    assert window.current_settings is not None
    assert (
        window.current_settings.component_build_mode
        == COMPONENT_BUILD_MODE_BORN_APPROXIMATION
    )
    assert (
        "Born Approximation (Average)"
        in window.project_setup_tab.summary_box.toPlainText()
    )
    window.close()


def test_saved_distributions_can_coexist_and_load_by_component_build_mode(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    manager = SAXSProjectManager()
    _no_contrast_settings, no_contrast_artifacts = (
        _seed_saved_distribution_from_root(
            project_dir,
            component_build_mode=COMPONENT_BUILD_MODE_NO_CONTRAST,
        )
    )
    contrast_settings, contrast_artifacts = _seed_saved_distribution_from_root(
        project_dir,
        component_build_mode=COMPONENT_BUILD_MODE_CONTRAST,
    )

    records = {
        record.distribution_id: record
        for record in manager.list_saved_distributions(project_dir)
    }
    contrast_metadata = json.loads(
        contrast_artifacts.distribution_metadata_file.read_text(
            encoding="utf-8"
        )
    )

    assert (
        no_contrast_artifacts.distribution_id
        != contrast_artifacts.distribution_id
    )
    assert set(records) == {
        no_contrast_artifacts.distribution_id,
        contrast_artifacts.distribution_id,
    }
    assert (
        records[no_contrast_artifacts.distribution_id].component_build_mode
        == COMPONENT_BUILD_MODE_NO_CONTRAST
    )
    assert (
        records[contrast_artifacts.distribution_id].component_build_mode
        == COMPONENT_BUILD_MODE_CONTRAST
    )
    assert "Build: No Contrast (Debye)" in (
        records[no_contrast_artifacts.distribution_id].label
    )
    assert "Build: Contrast (Debye)" in (
        records[contrast_artifacts.distribution_id].label
    )
    assert (
        contrast_metadata["component_build_mode"]
        == COMPONENT_BUILD_MODE_CONTRAST
    )
    assert contrast_metadata["label"] == (
        project_module.distribution_label_for_settings(contrast_settings)
    )

    loaded_settings = manager.settings_for_saved_distribution(
        project_dir,
        contrast_artifacts.distribution_id,
    )
    assert (
        loaded_settings.component_build_mode == COMPONENT_BUILD_MODE_CONTRAST
    )

    window = SAXSMainWindow(initial_project_dir=project_dir)
    target_index = (
        window.project_setup_tab.computed_distribution_combo.findData(
            contrast_artifacts.distribution_id
        )
    )
    assert target_index >= 0
    window.project_setup_tab.computed_distribution_combo.setCurrentIndex(
        target_index
    )
    window.project_setup_tab.load_distribution_button.click()
    QApplication.processEvents()

    assert window.project_setup_tab.selected_distribution_id() == (
        contrast_artifacts.distribution_id
    )
    assert (
        window.project_setup_tab.component_build_mode()
        == COMPONENT_BUILD_MODE_CONTRAST
    )
    assert window.current_settings is not None
    assert (
        window.current_settings.component_build_mode
        == COMPONENT_BUILD_MODE_CONTRAST
    )
    window.close()


def test_save_project_state_ignores_tiny_q_range_edge_mismatch(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    captured: dict[str, str] = {}

    monkeypatch.setattr(
        window,
        "_show_error",
        lambda title, message: captured.update(
            {
                "title": title,
                "message": message,
            }
        ),
    )

    window.project_setup_tab.qmin_edit.setText("0.04995")
    window.project_setup_tab.qmax_edit.setText("0.30005")
    window.save_project_state()

    assert captured == {}
    assert window.prefit_workflow is not None
    assert np.allclose(
        window.prefit_workflow.evaluate().q_values,
        np.linspace(0.05, 0.3, 8),
    )
    window.close()


def test_project_setup_preview_plots_solvent_data_in_green(qapp, tmp_path):
    del qapp
    experimental_path = tmp_path / "exp_preview.txt"
    solvent_path = tmp_path / "solvent_preview.txt"
    q_values = np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float)
    np.savetxt(
        experimental_path,
        np.column_stack(
            [q_values, np.asarray([100.0, 80.0, 55.0, 30.0], dtype=float)]
        ),
    )
    np.savetxt(
        solvent_path,
        np.column_stack(
            [q_values, np.asarray([12.0, 11.0, 10.0, 9.0], dtype=float)]
        ),
    )

    experimental_summary = load_experimental_data_file(experimental_path)
    solvent_summary = load_experimental_data_file(solvent_path)
    tab = ProjectSetupTab()
    tab._apply_experimental_file(experimental_path, experimental_summary)
    tab._apply_solvent_file(solvent_path, solvent_summary)

    preview_axis = tab.component_figure.axes[0]
    labels = [line.get_label() for line in preview_axis.get_lines()]

    assert "Experimental data" in labels
    assert "Solvent data" in labels
    solvent_line = next(
        line
        for line in preview_axis.get_lines()
        if line.get_label() == "Solvent data"
    )
    assert to_hex(solvent_line.get_color()) == "#008000"
    assert preview_axis.get_legend() is not None


def test_project_setup_data_trace_controls_toggle_and_persist(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    solvent_path = tmp_path / "solvent_control_trace.txt"
    q_values = np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float)
    np.savetxt(
        solvent_path,
        np.column_stack(
            [q_values, np.asarray([12.0, 11.0, 10.0, 9.0], dtype=float)]
        ),
    )

    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.solvent_data_path = str(solvent_path)
    settings.copied_solvent_data_file = None
    manager.save_project(settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)
    tab = window.project_setup_tab

    def _pick_color(*args, **kwargs):
        title = str(args[2] if len(args) > 2 else kwargs.get("title", ""))
        return QColor("#224466" if "Experimental" in title else "#228833")

    monkeypatch.setattr(
        "saxshell.saxs.ui.project_setup_tab.QColorDialog.getColor",
        _pick_color,
    )

    tab._choose_data_trace_color("experimental")
    tab._choose_data_trace_color("solvent")
    tab.solvent_trace_visible_checkbox.setChecked(False)

    preview_axis = tab.component_figure.axes[0]
    labels = [line.get_label() for line in preview_axis.get_lines()]
    experimental_line = next(
        line
        for line in preview_axis.get_lines()
        if line.get_label() == "Experimental data"
    )

    assert tab.experimental_trace_visible_checkbox.isChecked()
    assert not tab.solvent_trace_visible_checkbox.isChecked()
    assert to_hex(experimental_line.get_color()) == "#224466"
    assert "Solvent data" not in labels

    window.save_project_state()

    reloaded_settings = SAXSProjectManager().load_project(project_dir)
    assert reloaded_settings.experimental_trace_color == "#224466"
    assert reloaded_settings.solvent_trace_color == "#228833"
    assert reloaded_settings.experimental_trace_visible is True
    assert reloaded_settings.solvent_trace_visible is False


def test_project_setup_component_overlay_uses_secondary_y_axis(qapp, tmp_path):
    del qapp
    q_values = np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float)
    data_path = tmp_path / "exp_overlay.txt"
    np.savetxt(
        data_path,
        np.column_stack(
            [
                q_values,
                np.asarray([100.0, 80.0, 55.0, 30.0], dtype=float),
            ]
        ),
    )
    component_path = tmp_path / "A_no_motif.txt"
    _write_component_file(
        component_path,
        q_values,
        np.asarray([4.8, 4.6, 4.4, 4.2], dtype=float),
    )

    summary = load_experimental_data_file(data_path)
    tab = ProjectSetupTab()
    tab.component_log_y_checkbox.setChecked(False)
    tab._apply_experimental_file(data_path, summary)
    tab.draw_component_plot([component_path])

    assert len(tab.component_figure.axes) == 2
    experimental_axis, component_axis = tab.component_figure.axes
    assert (
        experimental_axis.get_title()
        == "Experimental Data and SAXS Components"
    )
    assert experimental_axis.get_xlabel() == "q (Å⁻¹)"
    assert (
        experimental_axis.get_ylabel() == "Experimental Intensity (arb. units)"
    )
    assert component_axis.get_ylabel() == "Model Intensity (arb. units)"
    assert experimental_axis.get_legend() is not None
    assert (
        component_axis.get_ylim()[1] - component_axis.get_ylim()[0] > 4.8 - 4.2
    )


def test_component_legend_toggle_hides_and_shows_legend(qapp, tmp_path):
    del qapp
    q_values = np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float)
    data_path = tmp_path / "exp_overlay_toggle.txt"
    np.savetxt(
        data_path,
        np.column_stack(
            [
                q_values,
                np.asarray([100.0, 80.0, 55.0, 30.0], dtype=float),
            ]
        ),
    )
    component_path = tmp_path / "PbI2_edge.txt"
    _write_component_file(
        component_path,
        q_values,
        np.asarray([4.8, 4.6, 4.4, 4.2], dtype=float),
    )

    summary = load_experimental_data_file(data_path)
    tab = ProjectSetupTab()
    tab._apply_experimental_file(data_path, summary)
    tab.draw_component_plot([component_path])

    assert tab.component_figure.axes[0].get_legend() is not None

    tab.component_legend_toggle_button.setChecked(False)

    assert tab.component_figure.axes[0].get_legend() is None


def test_component_autoscale_to_model_range_uses_model_q_limits(
    qapp, tmp_path
):
    del qapp
    experimental_q = np.asarray([0.03, 0.05, 0.08, 0.12, 0.18, 0.25])
    data_path = tmp_path / "exp_overlay_autoscale.txt"
    np.savetxt(
        data_path,
        np.column_stack(
            [
                experimental_q,
                np.asarray([150.0, 120.0, 95.0, 70.0, 40.0, 20.0]),
            ]
        ),
    )
    component_path = tmp_path / "PbI2_edge.txt"
    model_q = np.asarray([0.08, 0.12, 0.18], dtype=float)
    _write_component_file(
        component_path,
        model_q,
        np.asarray([4.8, 4.6, 4.2], dtype=float),
    )

    summary = load_experimental_data_file(data_path)
    tab = ProjectSetupTab()
    tab.component_log_y_checkbox.setChecked(False)
    tab._apply_experimental_file(data_path, summary)
    tab.draw_component_plot([component_path])

    tab.component_model_range_button.setChecked(True)

    x_limits = tab.component_figure.axes[0].get_xlim()
    assert x_limits[0] == pytest.approx(0.08)
    assert x_limits[1] == pytest.approx(0.18)


def test_component_show_hide_all_button_toggles_all_model_traces(
    qapp, tmp_path
):
    del qapp
    q_values = np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float)
    component_edge = tmp_path / "PbI2_edge.txt"
    component_face = tmp_path / "PbI2_face.txt"
    _write_component_file(
        component_edge,
        q_values,
        np.asarray([4.8, 4.6, 4.4, 4.2], dtype=float),
    )
    _write_component_file(
        component_face,
        q_values,
        np.asarray([3.8, 3.6, 3.4, 3.2], dtype=float),
    )

    tab = ProjectSetupTab()
    tab.apply_cluster_import_data(
        ["Pb", "I"],
        [
            {
                "structure": "PbI2",
                "motif": "edge",
                "count": 7,
                "weight": 0.55,
                "atom_fraction_percent": 55.0,
                "structure_fraction_percent": 55.0,
            },
            {
                "structure": "PbI2",
                "motif": "face",
                "count": 5,
                "weight": 0.45,
                "atom_fraction_percent": 45.0,
                "structure_fraction_percent": 45.0,
            },
        ],
    )
    tab.draw_component_plot([component_edge, component_face])

    assert tab.component_all_traces_button.text() == "Hide Computed Traces"

    tab.component_all_traces_button.click()

    assert not tab._component_line_lookup["PbI2_edge"].get_visible()
    assert not tab._component_line_lookup["PbI2_face"].get_visible()
    assert (
        tab.recognized_clusters_table.item(0, 6).checkState()
        == Qt.CheckState.Unchecked
    )
    assert (
        tab.recognized_clusters_table.item(1, 6).checkState()
        == Qt.CheckState.Unchecked
    )
    assert tab.component_all_traces_button.text() == "Show Computed Traces"

    tab.component_all_traces_button.click()

    assert tab._component_line_lookup["PbI2_edge"].get_visible()
    assert tab._component_line_lookup["PbI2_face"].get_visible()
    assert (
        tab.recognized_clusters_table.item(0, 6).checkState()
        == Qt.CheckState.Checked
    )
    assert (
        tab.recognized_clusters_table.item(1, 6).checkState()
        == Qt.CheckState.Checked
    )
    assert tab.component_all_traces_button.text() == "Hide Computed Traces"


def test_component_trace_color_scheme_applies_histogram_colormap(
    qapp, tmp_path
):
    del qapp
    q_values = np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float)
    component_edge = tmp_path / "PbI2_edge.txt"
    component_face = tmp_path / "PbI2_face.txt"
    _write_component_file(
        component_edge,
        q_values,
        np.asarray([4.8, 4.6, 4.4, 4.2], dtype=float),
    )
    _write_component_file(
        component_face,
        q_values,
        np.asarray([3.8, 3.6, 3.4, 3.2], dtype=float),
    )

    tab = ProjectSetupTab()
    tab.draw_component_plot([component_edge, component_face])
    plasma_index = tab.component_trace_color_scheme_combo.findData("plasma")
    tab.component_trace_color_scheme_combo.setCurrentIndex(plasma_index)

    expected_colors = tab._component_scheme_colors(
        [component_edge, component_face]
    )

    assert to_hex(tab._component_line_lookup["PbI2_edge"].get_color()) == (
        expected_colors["PbI2_edge"]
    )
    assert to_hex(tab._component_line_lookup["PbI2_face"].get_color()) == (
        expected_colors["PbI2_face"]
    )


def test_component_table_toggle_and_legend_pick_stay_synced(qapp, tmp_path):
    del qapp
    q_values = np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float)
    data_path = tmp_path / "exp_overlay_sync.txt"
    np.savetxt(
        data_path,
        np.column_stack(
            [
                q_values,
                np.asarray([100.0, 80.0, 55.0, 30.0], dtype=float),
            ]
        ),
    )
    component_path = tmp_path / "PbI2_edge.txt"
    _write_component_file(
        component_path,
        q_values,
        np.asarray([4.8, 4.6, 4.4, 4.2], dtype=float),
    )

    summary = load_experimental_data_file(data_path)
    tab = ProjectSetupTab()
    tab._apply_experimental_file(data_path, summary)
    tab.apply_cluster_import_data(
        ["Pb", "I"],
        [
            {
                "structure": "PbI2",
                "motif": "edge",
                "count": 7,
                "weight": 0.1239,
                "atom_fraction_percent": 44.44,
                "structure_fraction_percent": 12.39,
            }
        ],
    )
    tab.draw_component_plot([component_path])

    visibility_item = tab.recognized_clusters_table.item(0, 6)
    color_item = tab.recognized_clusters_table.item(0, 7)

    assert visibility_item.checkState() == Qt.CheckState.Checked
    assert color_item.text() != "--"

    visibility_item.setCheckState(Qt.CheckState.Unchecked)

    assert not tab._component_line_lookup["PbI2_edge"].get_visible()
    assert tab._component_legend_lookup[
        "PbI2_edge"
    ].get_alpha() == pytest.approx(0.25)

    event = type(
        "LegendEvent",
        (),
        {"artist": tab._component_legend_lookup["PbI2_edge"]},
    )()
    tab._handle_component_legend_pick(event)

    refreshed_item = tab.recognized_clusters_table.item(0, 6)
    assert tab._component_line_lookup["PbI2_edge"].get_visible()
    assert refreshed_item.checkState() == Qt.CheckState.Checked


def test_component_trace_color_scheme_persists_with_project_state(
    qapp, tmp_path
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    plasma_index = (
        window.project_setup_tab.component_trace_color_scheme_combo.findData(
            "plasma"
        )
    )
    window.project_setup_tab.component_trace_color_scheme_combo.setCurrentIndex(
        plasma_index
    )

    window.save_project_state()

    settings = SAXSProjectManager().load_project(project_dir)
    assert settings.component_trace_color_scheme == "plasma"

    reloaded_window = SAXSMainWindow(initial_project_dir=project_dir)
    assert (
        reloaded_window.project_setup_tab.component_trace_color_scheme()
        == "plasma"
    )


def test_prefit_plot_shows_solvent_contribution_and_legend_pick_toggles_model(
    qapp, tmp_path
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    solvent_q = np.linspace(0.05, 0.3, 8)
    solvent_intensity = np.linspace(1.5, 2.2, 8)
    solvent_path = tmp_path / "prefit_solvent_trace.dat"
    np.savetxt(solvent_path, np.column_stack([solvent_q, solvent_intensity]))

    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.solvent_data_path = str(solvent_path)
    settings.copied_solvent_data_file = None
    manager.save_project(settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)
    entries = window.prefit_workflow.load_parameter_entries()
    for entry in entries:
        if entry.name == "solv_w":
            entry.value = 0.5
        if entry.name == "scale":
            entry.value = 2e-3

    evaluation = window.prefit_workflow.evaluate(entries)
    window.prefit_tab.plot_evaluation(evaluation)
    window.prefit_tab.show_solvent_trace_checkbox.setChecked(True)

    top_axis = window.prefit_tab.figure.axes[0]
    labels = [line.get_label() for line in top_axis.get_lines()]
    model_line = next(
        line for line in top_axis.get_lines() if line.get_label() == "Model"
    )

    assert "Solvent contribution" in labels

    event = type(
        "LegendEvent",
        (),
        {"artist": window.prefit_tab._legend_handle_lookup["Model"]},
    )()
    window.prefit_tab._handle_legend_pick(event)

    assert not model_line.get_visible()
    assert window.prefit_tab._legend_handle_lookup[
        "Model"
    ].get_alpha() == pytest.approx(0.25)


def test_dream_plot_shows_solvent_contribution_and_legend_pick_toggles_model(
    qapp, tmp_path
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    solvent_q = np.linspace(0.05, 0.3, 8)
    solvent_intensity = np.linspace(1.5, 2.2, 8)
    solvent_path = tmp_path / "dream_solvent_trace.dat"
    np.savetxt(solvent_path, np.column_stack([solvent_q, solvent_intensity]))

    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.solvent_data_path = str(solvent_path)
    settings.copied_solvent_data_file = None
    manager.save_project(settings)

    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.load_latest_results()
    window.dream_tab.show_solvent_trace_checkbox.setChecked(True)

    top_axis = window.dream_tab.model_figure.axes[0]
    labels = [line.get_label() for line in top_axis.get_lines()]
    model_label = next(
        line.get_label()
        for line in top_axis.get_lines()
        if str(line.get_label()).startswith("Model (")
    )
    model_line = next(
        line
        for line in top_axis.get_lines()
        if line.get_label() == model_label
    )

    assert "Solvent contribution" in labels

    event = type(
        "LegendEvent",
        (),
        {"artist": window.dream_tab._model_legend_handle_lookup[model_label]},
    )()
    window.dream_tab._handle_model_legend_pick(event)

    assert not model_line.get_visible()
    assert window.dream_tab._model_legend_handle_lookup[
        model_label
    ].get_alpha() == pytest.approx(0.25)


def test_dream_plot_trace_toggles_control_visible_series(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    solvent_q = np.linspace(0.05, 0.3, 8)
    solvent_intensity = np.linspace(1.5, 2.2, 8)
    solvent_path = tmp_path / "dream_toggle_solvent_trace.dat"
    np.savetxt(solvent_path, np.column_stack([solvent_q, solvent_intensity]))

    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.solvent_data_path = str(solvent_path)
    settings.copied_solvent_data_file = None
    manager.save_project(settings)

    _write_minimal_dream_results(project_dir)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.load_latest_results()

    assert window.dream_tab.show_experimental_trace_checkbox.isEnabled()
    assert window.dream_tab.show_model_trace_checkbox.isEnabled()
    assert window.dream_tab.show_solvent_trace_checkbox.isEnabled()
    assert window.dream_tab.show_structure_factor_trace_checkbox.isEnabled()

    top_axis = window.dream_tab.model_figure.axes[0]
    labels = [
        line.get_label()
        for axis in window.dream_tab.model_figure.axes
        for line in axis.get_lines()
    ]
    assert "Solvent contribution" not in labels
    assert "Structure factor S(q)" not in labels

    window.dream_tab.show_solvent_trace_checkbox.setChecked(True)
    top_axis = window.dream_tab.model_figure.axes[0]
    labels = [
        line.get_label()
        for axis in window.dream_tab.model_figure.axes
        for line in axis.get_lines()
    ]
    assert "Solvent contribution" in labels

    window.dream_tab.show_structure_factor_trace_checkbox.setChecked(True)
    labels = [
        line.get_label()
        for axis in window.dream_tab.model_figure.axes
        for line in axis.get_lines()
    ]
    assert "Structure factor S(q)" in labels

    window.dream_tab.show_experimental_trace_checkbox.setChecked(False)
    top_axis = window.dream_tab.model_figure.axes[0]
    collection_labels = [
        collection.get_label() for collection in top_axis.collections
    ]
    assert "Experimental" not in collection_labels

    window.dream_tab.show_model_trace_checkbox.setChecked(False)
    top_axis = window.dream_tab.model_figure.axes[0]
    line_labels = [line.get_label() for line in top_axis.get_lines()]
    assert not any(str(label).startswith("Model (") for label in line_labels)


def test_prefit_plot_trace_toggles_control_visible_series(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    solvent_q = np.linspace(0.05, 0.3, 8)
    solvent_intensity = np.linspace(1.5, 2.2, 8)
    solvent_path = tmp_path / "prefit_toggle_solvent_trace.dat"
    np.savetxt(solvent_path, np.column_stack([solvent_q, solvent_intensity]))

    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.solvent_data_path = str(solvent_path)
    settings.copied_solvent_data_file = None
    manager.save_project(settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)
    entries = window.prefit_workflow.load_parameter_entries()
    for entry in entries:
        if entry.name == "solv_w":
            entry.value = 0.5
        if entry.name == "scale":
            entry.value = 2e-3

    evaluation = window.prefit_workflow.evaluate(entries)
    window.prefit_tab.plot_evaluation(evaluation)

    assert window.prefit_tab.show_experimental_trace_checkbox.isEnabled()
    assert window.prefit_tab.show_model_trace_checkbox.isEnabled()
    assert window.prefit_tab.show_solvent_trace_checkbox.isEnabled()
    assert window.prefit_tab.show_structure_factor_trace_checkbox.isEnabled()

    top_axis = window.prefit_tab.figure.axes[0]
    labels = [
        line.get_label()
        for axis in window.prefit_tab.figure.axes
        for line in axis.get_lines()
    ]
    assert "Solvent contribution" not in labels
    assert "Structure factor S(q)" not in labels

    window.prefit_tab.show_solvent_trace_checkbox.setChecked(True)
    top_axis = window.prefit_tab.figure.axes[0]
    labels = [
        line.get_label()
        for axis in window.prefit_tab.figure.axes
        for line in axis.get_lines()
    ]
    assert "Solvent contribution" in labels

    window.prefit_tab.show_structure_factor_trace_checkbox.setChecked(True)
    labels = [
        line.get_label()
        for axis in window.prefit_tab.figure.axes
        for line in axis.get_lines()
    ]
    assert "Structure factor S(q)" in labels

    window.prefit_tab.show_experimental_trace_checkbox.setChecked(False)
    top_axis = window.prefit_tab.figure.axes[0]
    line_labels = [line.get_label() for line in top_axis.get_lines()]
    assert "Experimental" not in line_labels


def test_prefit_field_interaction_warns_before_components_are_built(
    qapp, tmp_path, monkeypatch
):
    project_dir = tmp_path / "prefit_warning_project"
    manager = SAXSProjectManager()
    settings = manager.create_project(project_dir)
    settings.selected_model_template = (
        "template_pd_likelihood_monosq_decoupled"
    )
    manager.save_project(settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)
    window.show()
    qapp.processEvents()

    warnings: list[tuple[str, str]] = []
    monkeypatch.setattr(
        "saxshell.saxs.ui.main_window.QMessageBox.warning",
        lambda _parent, title, message, *args, **kwargs: warnings.append(
            (title, message)
        )
        or QMessageBox.StandardButton.Ok,
    )

    QTest.mouseClick(
        window.prefit_tab.parameter_table.viewport(),
        Qt.MouseButton.LeftButton,
    )
    qapp.processEvents()
    QTest.mouseClick(
        window.prefit_tab.parameter_table.viewport(),
        Qt.MouseButton.LeftButton,
    )
    qapp.processEvents()

    assert warnings
    assert warnings[0][0] == "Build SAXS components first"
    assert "Project Setup tab" in warnings[0][1]
    assert len(warnings) == 1
    window.close()


def test_model_only_mode_disables_fit_controls_and_dream(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    dream_index = window.tabs.indexOf(window.dream_tab)

    assert window.tabs.isTabEnabled(dream_index)
    assert window.project_setup_tab.experimental_file_button.isEnabled()
    assert window.project_setup_tab.experimental_data_edit.isEnabled()
    assert window.prefit_tab.run_button.isEnabled()

    window.project_setup_tab.model_only_mode_checkbox.setChecked(True)

    assert window.current_settings is not None
    assert window.current_settings.model_only_mode is True
    assert not window.project_setup_tab.experimental_file_button.isEnabled()
    assert not window.project_setup_tab.experimental_data_edit.isEnabled()
    assert not window.project_setup_tab.solvent_file_button.isEnabled()
    assert not window.project_setup_tab.solvent_data_edit.isEnabled()
    assert (
        not window.project_setup_tab.experimental_trace_visible_checkbox.isEnabled()
    )
    assert (
        not window.project_setup_tab.solvent_trace_visible_checkbox.isEnabled()
    )
    assert not window.prefit_tab.method_combo.isEnabled()
    assert not window.prefit_tab.nfev_spin.isEnabled()
    assert not window.prefit_tab.run_button.isEnabled()
    assert not window.prefit_tab.recommended_scale_button.isEnabled()
    assert window.prefit_tab.update_button.isEnabled()
    assert not window.tabs.isTabEnabled(dream_index)

    assert len(window.prefit_tab.figure.axes) == 1
    top_axis = window.prefit_tab.figure.axes[0]
    line_labels = [line.get_label() for line in top_axis.get_lines()]
    assert "Experimental" not in line_labels

    saved_settings = SAXSProjectManager().load_project(project_dir)
    assert saved_settings.model_only_mode is True

    window.project_setup_tab.model_only_mode_checkbox.setChecked(False)

    assert window.current_settings is not None
    assert window.current_settings.model_only_mode is False
    assert window.project_setup_tab.experimental_file_button.isEnabled()
    assert window.project_setup_tab.experimental_data_edit.isEnabled()
    assert window.prefit_tab.run_button.isEnabled()
    assert window.tabs.isTabEnabled(dream_index)

    window.prefit_tab.show_model_trace_checkbox.setChecked(False)
    top_axis = window.prefit_tab.figure.axes[0]
    line_labels = [line.get_label() for line in top_axis.get_lines()]
    assert "Model" not in line_labels


def test_save_prefit_plot_data_exports_csv_with_metadata(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    export_path = tmp_path / "prefit_plot_export.csv"

    monkeypatch.setattr(
        QFileDialog,
        "getSaveFileName",
        lambda *args, **kwargs: (
            str(export_path),
            "CSV files (*.csv)",
        ),
    )

    window.save_prefit_plot_data()

    assert export_path.is_file()
    contents = export_path.read_text(encoding="utf-8")
    assert "# fit_conditions:" in contents
    assert "# fit_metrics:" in contents
    assert (
        "q,experimental_intensity,model_intensity,residual,"
        "solvent_intensity,solvent_contribution,structure_factor" in contents
    )


def test_save_prefit_plot_data_exports_npy_with_metadata_sidecar(
    qapp, tmp_path, monkeypatch
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)
    export_path = tmp_path / "prefit_plot_export.npy"

    monkeypatch.setattr(
        QFileDialog,
        "getSaveFileName",
        lambda *args, **kwargs: (
            str(export_path),
            "NumPy files (*.npy)",
        ),
    )

    window.save_prefit_plot_data()

    assert export_path.is_file()
    metadata_path = export_path.with_suffix(".metadata.json")
    assert metadata_path.is_file()
    matrix = np.load(export_path)
    metadata = json.loads(metadata_path.read_text(encoding="utf-8"))
    assert matrix.shape[1] == 7
    assert metadata["columns"] == [
        "q",
        "experimental_intensity",
        "model_intensity",
        "residual",
        "solvent_intensity",
        "solvent_contribution",
        "structure_factor",
    ]
    assert "fit_conditions" in metadata
    assert "fit_metrics" in metadata


def test_experimental_status_note_is_wrapped_and_multiline(qapp):
    del qapp
    tab = ProjectSetupTab()

    assert tab.data_status_label.wordWrap()
    assert tab.data_status_label.minimumHeight() >= 72
    assert "\n" in tab.data_status_label.text()


def test_project_setup_template_controls_have_expanded_minimum_width(qapp):
    del qapp
    tab = ProjectSetupTab()

    assert tab.template_combo.minimumWidth() >= 420
    assert tab.active_template_edit.minimumWidth() >= 420


def test_project_setup_uses_scrollable_resizable_side_panes(qapp):
    del qapp
    tab = ProjectSetupTab()

    assert tab._pane_splitter.count() == 2
    assert tab._pane_splitter.widget(0) is tab._left_scroll_area
    assert tab._pane_splitter.widget(1) is tab._right_scroll_area
    assert tab._left_scroll_area.widget() is not None
    assert tab._right_scroll_area.widget() is not None
    assert (
        tab._right_scroll_area.verticalScrollBarPolicy()
        == Qt.ScrollBarPolicy.ScrollBarAsNeeded
    )
    assert (
        tab.component_group.parentWidget() is tab._right_scroll_area.widget()
    )
    assert tab.prior_group.parentWidget() is tab._right_scroll_area.widget()


def test_project_setup_plot_previews_expand_with_the_right_pane(qapp):
    del qapp
    tab = ProjectSetupTab()

    assert tab.component_canvas.parentWidget() is tab.component_group
    assert tab.prior_canvas.parentWidget() is tab.prior_group
    assert not bool(tab.component_canvas.property("_saxs_skip_scale"))
    assert not bool(tab.prior_canvas.property("_saxs_skip_scale"))
    assert (
        tab.component_group.sizePolicy().verticalPolicy()
        == QSizePolicy.Policy.Expanding
    )
    assert (
        tab.prior_group.sizePolicy().verticalPolicy()
        == QSizePolicy.Policy.Expanding
    )
    assert (
        tab.component_canvas.sizePolicy().horizontalPolicy()
        == QSizePolicy.Policy.Expanding
    )
    assert (
        tab.component_canvas.sizePolicy().verticalPolicy()
        == QSizePolicy.Policy.Expanding
    )
    assert (
        tab.prior_canvas.sizePolicy().horizontalPolicy()
        == QSizePolicy.Policy.Expanding
    )
    assert (
        tab.prior_canvas.sizePolicy().verticalPolicy()
        == QSizePolicy.Policy.Expanding
    )
    assert tab.component_canvas.minimumHeight() >= 320
    assert tab.prior_canvas.minimumHeight() >= 240


def test_project_setup_activity_progress_updates(qapp):
    del qapp
    tab = ProjectSetupTab()

    tab.start_activity_progress(8, "Importing cluster files...")
    assert tab.activity_progress_label.text() == "Importing cluster files..."
    assert tab.activity_progress_bar.maximum() == 8
    assert tab.activity_progress_bar.value() == 0

    tab.update_activity_progress(3, 8, "Building SAXS components...")
    assert tab.activity_progress_label.text() == "Building SAXS components..."
    assert tab.activity_progress_bar.value() == 3

    tab.finish_activity_progress("Prior-weight generation complete.")
    assert (
        tab.activity_progress_label.text()
        == "Prior-weight generation complete."
    )
    assert tab.activity_progress_bar.value() == 8

    tab.reset_activity_progress()
    assert tab.activity_progress_label.text() == "Progress: idle"
    assert tab.activity_progress_bar.maximum() == 1
    assert tab.activity_progress_bar.value() == 0


def test_recognized_clusters_weight_is_truncated_to_three_decimals(qapp):
    del qapp
    tab = ProjectSetupTab()
    tab._populate_recognized_clusters_table(
        [
            {
                "structure": "Pb2I4",
                "motif": "edge",
                "count": 7,
                "weight": 0.1239,
                "atom_fraction_percent": 44.44,
                "structure_fraction_percent": 12.39,
            }
        ]
    )

    assert tab.recognized_clusters_table.item(0, 3).text() == "0.123"


def test_project_setup_tracks_selected_experimental_columns(qapp, tmp_path):
    del qapp
    data_path = tmp_path / "exp_project_columns.txt"
    data_path.write_text(
        "intensity_trace\tq_trace\terror_trace\n"
        "10.0\t0.05\t0.1\n"
        "9.0\t0.10\t0.2\n",
        encoding="utf-8",
    )

    summary = load_experimental_data_file(
        data_path,
        q_column=1,
        intensity_column=0,
        error_column=2,
    )
    tab = ProjectSetupTab()
    tab._apply_experimental_file(data_path, summary)

    assert tab.experimental_q_column() == 1
    assert tab.experimental_intensity_column() == 0
    assert tab.experimental_error_column() == 2
    assert "q=q_trace" in tab.data_status_label.text()
    assert "intensity=intensity_trace" in tab.data_status_label.text()
    assert "error=error_trace" in tab.data_status_label.text()


def test_prefit_plot_log_axes_toggle_independently(qapp, tmp_path):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    top_axis, bottom_axis = window.prefit_tab.figure.axes
    assert top_axis.get_xscale() == "log"
    assert top_axis.get_yscale() == "log"
    assert bottom_axis.get_xscale() == "log"

    window.prefit_tab.log_x_checkbox.setChecked(False)
    window.prefit_tab.log_y_checkbox.setChecked(False)

    top_axis, bottom_axis = window.prefit_tab.figure.axes
    assert top_axis.get_xscale() == "linear"
    assert top_axis.get_yscale() == "linear"
    assert bottom_axis.get_xscale() == "linear"


def test_use_experimental_grid_crops_to_nearest_available_q_values(tmp_path):
    manager = SAXSProjectManager()
    settings = ProjectSettings(
        project_name="demo_project",
        project_dir=str(tmp_path / "demo_project"),
        use_experimental_grid=True,
        q_min=0.081,
        q_max=0.171,
    )
    summary = ExperimentalDataSummary(
        path=tmp_path / "exp_demo.txt",
        q_values=np.asarray([0.05, 0.08, 0.12, 0.18], dtype=float),
        intensities=np.asarray([10.0, 8.0, 6.0, 4.0], dtype=float),
    )

    q_grid = manager._build_q_grid(settings, summary)

    assert np.allclose(q_grid, [0.08, 0.12, 0.18])


def test_build_components_persists_effective_q_range_after_grid_crop(
    tmp_path,
    monkeypatch,
):
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    manager = SAXSProjectManager()
    clusters_dir = paths.project_dir / "clusters"
    structure_dir = clusters_dir / "A"
    structure_dir.mkdir(parents=True, exist_ok=True)
    (structure_dir / "frame_0001.xyz").write_text(
        "2\nframe 1\nA 0.0 0.0 0.0\nH 1.0 0.0 0.0\n",
        encoding="utf-8",
    )

    def fake_build_profiles(
        builder,
        *,
        cluster_bins,
        progress_callback=None,
        progress_total=None,
        **kwargs,
    ):
        del progress_callback, progress_total, kwargs
        values = np.linspace(10.0, 17.0, len(builder.q_values))
        output_path = builder.output_dir / "A_no_motif.txt"
        _write_component_file(output_path, builder.q_values, values)
        cluster_bin = cluster_bins[0]
        return [
            AveragedComponent(
                structure=cluster_bin.structure,
                motif=cluster_bin.motif,
                file_count=len(cluster_bin.files),
                representative=cluster_bin.representative,
                source_dir=cluster_bin.source_dir,
                q_values=np.asarray(builder.q_values, dtype=float),
                mean_intensity=np.asarray(values, dtype=float),
                std_intensity=np.zeros_like(builder.q_values, dtype=float),
                se_intensity=np.zeros_like(builder.q_values, dtype=float),
                output_path=output_path,
            )
        ]

    monkeypatch.setattr(
        project_module.DebyeProfileBuilder,
        "build_profiles",
        fake_build_profiles,
    )

    settings = manager.load_project(project_dir)
    settings.clusters_dir = str(clusters_dir)
    settings.q_min = 0.12
    settings.q_max = 0.19

    manager.build_scattering_components(settings)

    saved_settings = manager.load_project(project_dir)
    assert saved_settings.q_min == pytest.approx(0.12142857142857144)
    assert saved_settings.q_max == pytest.approx(0.19285714285714284)


def test_model_only_q_grid_uses_configured_range_without_experimental_data(
    tmp_path,
):
    manager = SAXSProjectManager()
    settings = ProjectSettings(
        project_name="demo_project",
        project_dir=str(tmp_path / "demo_project"),
        model_only_mode=True,
        use_experimental_grid=False,
        q_min=0.05,
        q_max=0.30,
        q_points=8,
    )

    q_grid = manager._build_q_grid(settings, None)

    assert np.allclose(q_grid, np.linspace(0.05, 0.30, 8))


def test_predicted_structure_mode_status_guides_user_when_bundle_is_missing(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, _paths = _build_minimal_saxs_project(tmp_path)
    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.use_predicted_structure_weights = True
    manager.save_project(settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert window.project_setup_tab.use_predicted_structure_weights()
    assert (
        window.project_setup_tab.use_predicted_structure_weights_checkbox.isEnabled()
    )
    assert "#6b7280" in (
        window.project_setup_tab.predicted_structure_ready_indicator.styleSheet()
    )
    assert "Predicted Structures mode is on" in (
        window.project_setup_tab.predicted_structure_status_label.text()
    )
    assert "Open Cluster Dynamics (ML)" in (
        window.project_setup_tab.predicted_structure_status_label.text()
    )
    window.close()


def test_predicted_structure_toggle_stays_disabled_until_predictions_exist(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert not window.project_setup_tab.use_predicted_structure_weights()
    assert (
        not window.project_setup_tab.use_predicted_structure_weights_checkbox.isEnabled()
    )
    assert "#6b7280" in (
        window.project_setup_tab.predicted_structure_ready_indicator.styleSheet()
    )

    _write_predicted_structure_artifacts(paths)
    window._refresh_predicted_structure_status()

    assert (
        window.project_setup_tab.use_predicted_structure_weights_checkbox.isEnabled()
    )
    assert "#16a34a" in (
        window.project_setup_tab.predicted_structure_ready_indicator.styleSheet()
    )
    assert "1 predicted structure is available" in (
        window.project_setup_tab.predicted_structure_ready_indicator.toolTip()
    )
    window.close()


def test_predicted_structure_mode_toggles_between_predicted_and_observed_artifacts(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    _write_predicted_structure_artifacts(paths)
    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.use_predicted_structure_weights = True
    manager.save_project(settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert window.project_setup_tab.use_predicted_structure_weights()
    assert (
        window.project_setup_tab.current_prior_json_path().name
        == "md_prior_weights_predicted_structures.json"
    )
    assert window.project_setup_tab.recognized_clusters_table.rowCount() == 2
    assert (
        window.project_setup_tab.recognized_clusters_table.item(1, 0).text()
        == "A2"
    )
    assert window.prefit_workflow is not None
    assert {
        (component.structure, component.motif)
        for component in window.prefit_workflow.components
    } == {("A", "no_motif"), ("A2", "predicted_rank01")}
    assert window.dream_workflow is not None
    assert {
        (component.structure, component.motif)
        for component in window.dream_workflow.prefit_workflow.components
    } == {("A", "no_motif"), ("A2", "predicted_rank01")}
    assert "observed + Predicted Structures" in (
        window.project_setup_tab.predicted_structure_status_label.text()
    )

    window.project_setup_tab.use_predicted_structure_weights_checkbox.setChecked(
        False
    )
    QApplication.processEvents()
    saved_settings = SAXSProjectManager().load_project(project_dir)

    assert not saved_settings.use_predicted_structure_weights
    assert (
        window.project_setup_tab.current_prior_json_path().name
        == "md_prior_weights.json"
    )
    assert window.project_setup_tab.recognized_clusters_table.rowCount() == 1
    assert (
        window.project_setup_tab.recognized_clusters_table.item(0, 0).text()
        == "A"
    )
    assert window.prefit_workflow is not None
    assert {
        (component.structure, component.motif)
        for component in window.prefit_workflow.components
    } == {("A", "no_motif")}
    dream_workflow = window._load_dream_workflow()
    assert {
        (component.structure, component.motif)
        for component in dream_workflow.prefit_workflow.components
    } == {("A", "no_motif")}
    assert "Predicted Structures mode is off" in (
        window.project_setup_tab.predicted_structure_status_label.text()
    )
    window.close()


def test_project_setup_predicted_mode_toggles_observed_and_predicted_traces(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    _write_predicted_structure_artifacts(paths)
    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.use_predicted_structure_weights = True
    manager.save_project(settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)
    tab = window.project_setup_tab

    assert not tab.component_observed_traces_button.isHidden()
    assert not tab.component_predicted_traces_button.isHidden()
    assert tab.component_observed_traces_button.text() == (
        "Hide Observed Traces"
    )
    assert tab.component_predicted_traces_button.text() == (
        "Hide Predicted Traces"
    )
    assert tab._component_line_lookup["A_no_motif"].get_visible()
    assert tab._component_line_lookup["A2_predicted_rank01"].get_visible()

    tab.component_predicted_traces_button.click()
    QApplication.processEvents()

    assert tab._component_line_lookup["A_no_motif"].get_visible()
    assert not tab._component_line_lookup["A2_predicted_rank01"].get_visible()
    assert tab.component_predicted_traces_button.text() == (
        "Show Predicted Traces"
    )
    assert (
        tab.recognized_clusters_table.item(0, 6).checkState()
        == Qt.CheckState.Checked
    )
    assert (
        tab.recognized_clusters_table.item(1, 6).checkState()
        == Qt.CheckState.Unchecked
    )

    tab.component_predicted_traces_button.click()
    QApplication.processEvents()

    assert tab._component_line_lookup["A_no_motif"].get_visible()
    assert tab._component_line_lookup["A2_predicted_rank01"].get_visible()

    tab.component_observed_traces_button.click()
    QApplication.processEvents()

    assert not tab._component_line_lookup["A_no_motif"].get_visible()
    assert tab._component_line_lookup["A2_predicted_rank01"].get_visible()
    assert tab.component_observed_traces_button.text() == (
        "Show Observed Traces"
    )
    assert (
        tab.recognized_clusters_table.item(0, 6).checkState()
        == Qt.CheckState.Unchecked
    )
    assert (
        tab.recognized_clusters_table.item(1, 6).checkState()
        == Qt.CheckState.Checked
    )

    tab.component_observed_traces_button.click()
    QApplication.processEvents()

    assert tab._component_line_lookup["A_no_motif"].get_visible()
    assert tab._component_line_lookup["A2_predicted_rank01"].get_visible()

    tab.use_predicted_structure_weights_checkbox.setChecked(False)
    QApplication.processEvents()

    assert tab.component_observed_traces_button.isHidden()
    assert tab.component_predicted_traces_button.isHidden()
    window.close()


def test_project_setup_remove_selected_elements_updates_exclude_list(qapp):
    del qapp
    tab = ProjectSetupTab()
    tab.set_project_selected(True)
    tab.set_available_elements(["H", "O", "Pb"])
    tab.exclude_elements_edit.setText("H O")

    tab.available_elements_list.clearSelection()
    tab.available_elements_list.item(0).setSelected(True)
    tab._remove_selected_elements_from_exclude()

    assert tab.exclude_elements() == ["O"]


def test_project_setup_exclude_buttons_use_clear_labels(qapp):
    del qapp
    tab = ProjectSetupTab()

    assert tab.exclude_selected_elements_button.text() == "Exclude Selected"
    assert tab.include_selected_elements_button.text() == "Include Selected"


def test_project_setup_predicted_mode_updates_prefit_geometry_table(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, paths = _build_poly_lma_geometry_project(tmp_path)
    observed_workflow = SAXSPrefitWorkflow(project_dir)
    observed_workflow.compute_cluster_geometry_table()

    _write_predicted_structure_artifacts(paths)
    predicted_dir = tmp_path / "ui_predicted_geometry"
    predicted_dir.mkdir(parents=True, exist_ok=True)
    predicted_structure = predicted_dir / "02_rank01_A2.xyz"
    predicted_structure.write_text(
        "2\npredicted\nA 0.0 0.0 0.0\nA 1.5 0.0 0.0\n",
        encoding="utf-8",
    )

    def fake_predicted_bins(self):
        if not self.settings.use_predicted_structure_weights:
            return []
        return [
            ClusterBin(
                structure="A2",
                motif="predicted_rank01",
                source_dir=predicted_structure.parent,
                files=(predicted_structure,),
                representative=predicted_structure.name,
            )
        ]

    monkeypatch.setattr(
        SAXSPrefitWorkflow,
        "_predicted_structure_cluster_bins_for_active_components",
        fake_predicted_bins,
    )

    manager = SAXSProjectManager()
    predicted_settings = manager.load_project(project_dir)
    predicted_settings.use_predicted_structure_weights = True
    manager.save_project(predicted_settings)
    predicted_workflow = SAXSPrefitWorkflow(project_dir)
    predicted_workflow.compute_cluster_geometry_table()

    observed_settings = manager.load_project(project_dir)
    observed_settings.use_predicted_structure_weights = False
    manager.save_project(observed_settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert {
        (row.structure, row.motif)
        for row in window.prefit_workflow.cluster_geometry_table.rows
    } == {("A", "no_motif")}
    assert all(
        entry.structure != "A2"
        for entry in window.prefit_tab.parameter_entries()
    )

    window.project_setup_tab.use_predicted_structure_weights_checkbox.setChecked(
        True
    )
    QApplication.processEvents()

    assert {
        (row.structure, row.motif)
        for row in window.prefit_workflow.cluster_geometry_table.rows
    } == {("A", "no_motif"), ("A2", "predicted_rank01")}
    assert any(
        entry.structure == "A2"
        for entry in window.prefit_tab.parameter_entries()
        if entry.category in {"weight", "geometry"}
    )

    window.project_setup_tab.use_predicted_structure_weights_checkbox.setChecked(
        False
    )
    QApplication.processEvents()

    assert {
        (row.structure, row.motif)
        for row in window.prefit_workflow.cluster_geometry_table.rows
    } == {("A", "no_motif")}
    assert all(
        entry.structure != "A2"
        for entry in window.prefit_tab.parameter_entries()
    )
    window.close()


def test_project_setup_predicted_histograms_use_observed_stoichiometry_bins(
    qapp,
    tmp_path,
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    _write_predicted_structure_artifacts(paths)
    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    settings.use_predicted_structure_weights = True
    manager.save_project(settings)

    window = SAXSMainWindow(initial_project_dir=project_dir)
    prior_path = window.project_setup_tab.current_prior_json_path()

    structure_payload = build_prior_histogram_export_payload(
        prior_path,
        mode="structure_fraction",
    )
    atom_payload = build_prior_histogram_export_payload(
        prior_path,
        mode="atom_fraction",
    )

    structure_matrix = np.asarray(structure_payload["matrix"], dtype=float)
    atom_matrix = np.asarray(atom_payload["matrix"], dtype=float)
    structure_by_label = {
        str(label): {
            str(segment): float(value)
            for segment, value in zip(
                structure_payload["segments"],
                structure_matrix[row_index],
                strict=False,
            )
        }
        for row_index, label in enumerate(structure_payload["labels"])
    }
    atom_by_label = {
        str(label): {
            str(segment): float(value)
            for segment, value in zip(
                atom_payload["segments"],
                atom_matrix[row_index],
                strict=False,
            )
        }
        for row_index, label in enumerate(atom_payload["labels"])
    }

    assert structure_payload["labels"] == ["A", "A2"]
    assert atom_payload["labels"] == ["A", "A2"]
    assert float(np.sum(structure_payload["totals"])) == pytest.approx(100.0)
    assert float(np.sum(atom_payload["totals"])) == pytest.approx(100.0)
    assert structure_by_label["A"]["no_motif"] == pytest.approx(75.0)
    assert structure_by_label["A2"]["predicted_rank01"] == pytest.approx(25.0)
    assert atom_by_label["A"]["no_motif"] == pytest.approx(60.0)
    assert atom_by_label["A2"]["predicted_rank01"] == pytest.approx(40.0)

    tick_labels = {
        tick.get_text().strip()
        for tick in window.project_setup_tab.prior_figure.axes[
            0
        ].get_xticklabels()
        if tick.get_text().strip()
    }
    assert tick_labels == {"A", "A$_{2}$"}
    window.close()


def test_project_setup_predicted_histogram_x_axis_updates_for_all_modes(
    qapp,
):
    del qapp
    payload = {
        "includes_predicted_structures": True,
        "available_elements": ["A", "O"],
        "structures": {
            "A": {
                "no_motif": {
                    "count": 3.0,
                    "normalized_weight": 0.75,
                    "source_kind": "cluster_dir",
                    "secondary_atom_distributions": {"O": {"0": 3.0}},
                }
            },
            "A2": {
                "predicted_rank01": {
                    "count": 1.0,
                    "normalized_weight": 0.25,
                    "source_kind": "predicted_structure",
                    "secondary_atom_distributions": {"O": {"1": 1.0}},
                }
            },
        },
    }
    figure = Figure()
    axis = figure.add_subplot(111)

    for mode, secondary in (
        ("structure_fraction", None),
        ("atom_fraction", None),
        ("solvent_sort_structure_fraction", "O"),
        ("solvent_sort_atom_fraction", "O"),
    ):
        plot_md_prior_histogram(
            payload,
            mode=mode,
            secondary_element=secondary,
            ax=axis,
        )
        tick_labels = {
            tick.get_text().strip()
            for tick in axis.get_xticklabels()
            if tick.get_text().strip()
        }
        assert tick_labels == {"A", "A$_{2}$"}


def test_predicted_structure_mode_switch_reuses_distribution_and_refreshes_tabs(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    manager = SAXSProjectManager()
    clusters_dir = paths.project_dir / "clusters"
    structure_dir = clusters_dir / "A"
    structure_dir.mkdir(parents=True, exist_ok=True)
    observed_structure = structure_dir / "frame_0001.xyz"
    observed_structure.write_text(
        "1\nframe 1\nA 0.0 0.0 0.0\n",
        encoding="utf-8",
    )

    observed_component = np.linspace(10.0, 17.0, 8)
    predicted_component = np.linspace(4.0, 11.0, 8)

    def fake_build_profiles(
        builder,
        *,
        cluster_bins,
        progress_callback=None,
        progress_total=None,
        **kwargs,
    ):
        del progress_callback, progress_total, kwargs
        output_path = builder.output_dir / "A_no_motif.txt"
        _write_component_file(
            output_path,
            builder.q_values,
            observed_component,
        )
        cluster_bin = cluster_bins[0]
        return [
            AveragedComponent(
                structure=cluster_bin.structure,
                motif=cluster_bin.motif,
                file_count=len(cluster_bin.files),
                representative=cluster_bin.representative,
                source_dir=cluster_bin.source_dir,
                q_values=np.asarray(builder.q_values, dtype=float),
                mean_intensity=np.asarray(observed_component, dtype=float),
                std_intensity=np.zeros_like(builder.q_values, dtype=float),
                se_intensity=np.zeros_like(builder.q_values, dtype=float),
                output_path=output_path,
            )
        ]

    monkeypatch.setattr(
        project_module.DebyeProfileBuilder,
        "build_profiles",
        fake_build_profiles,
    )
    monkeypatch.setattr(
        project_module,
        "compute_debye_intensity",
        lambda coordinates, elements, q_values: np.asarray(
            predicted_component,
            dtype=float,
        ),
    )

    dataset_dir = (
        paths.exported_data_dir
        / "clusterdynamicsml"
        / "saved_results"
        / "20260330_120000_demo"
    )
    dataset_dir.mkdir(parents=True, exist_ok=True)
    dataset_file = dataset_dir / "20260330_120000_demo_clusterdynamicsml.json"
    dataset_file.write_text("{}\n", encoding="utf-8")
    predicted_dir = dataset_dir / (f"{dataset_file.stem}_predicted_structures")
    predicted_dir.mkdir(parents=True, exist_ok=True)
    predicted_structure = predicted_dir / "02_rank01_A2.xyz"
    predicted_structure.write_text(
        "2\npredicted\nA 0.0 0.0 0.0\nA 1.0 0.0 0.0\n",
        encoding="utf-8",
    )
    loaded_dataset = SimpleNamespace(
        dataset_file=dataset_file,
        result=SimpleNamespace(
            training_observations=(
                ClusterDynamicsMLTrainingObservation(
                    label="A",
                    node_count=1,
                    cluster_size=1,
                    element_counts={"A": 1},
                    file_count=1,
                    representative_path=observed_structure,
                    structure_dir=structure_dir,
                    motifs=("no_motif",),
                    mean_atom_count=1.0,
                    mean_radius_of_gyration=1.0,
                    mean_max_radius=1.0,
                    mean_semiaxis_a=1.0,
                    mean_semiaxis_b=1.0,
                    mean_semiaxis_c=1.0,
                    total_observations=1,
                    occupied_frames=1,
                    mean_count_per_frame=1.0,
                    occupancy_fraction=1.0,
                    association_events=0,
                    dissociation_events=0,
                    association_rate_per_ps=0.0,
                    dissociation_rate_per_ps=0.0,
                    completed_lifetime_count=1,
                    window_truncated_lifetime_count=0,
                    mean_lifetime_fs=10.0,
                    std_lifetime_fs=0.0,
                ),
            ),
            predictions=(
                PredictedClusterCandidate(
                    target_node_count=2,
                    rank=1,
                    label="A2",
                    element_counts={"A": 2},
                    predicted_mean_count_per_frame=0.2,
                    predicted_occupancy_fraction=1.0,
                    predicted_mean_lifetime_fs=20.0,
                    predicted_association_rate_per_ps=0.0,
                    predicted_dissociation_rate_per_ps=0.0,
                    predicted_mean_radius_of_gyration=1.0,
                    predicted_mean_max_radius=1.0,
                    predicted_mean_semiaxis_a=1.0,
                    predicted_mean_semiaxis_b=1.0,
                    predicted_mean_semiaxis_c=1.0,
                    predicted_population_share=100.0,
                    predicted_stability_score=1.0,
                    source_label="A",
                    notes="",
                    generated_elements=("A", "A"),
                    generated_coordinates=np.asarray(
                        [[0.0, 0.0, 0.0], [1.0, 0.0, 0.0]],
                        dtype=float,
                    ),
                ),
            ),
            dynamics_result=SimpleNamespace(
                preview=SimpleNamespace(frame_timestep_fs=1.0)
            ),
            saxs_comparison=None,
        ),
    )
    monkeypatch.setattr(
        SAXSProjectManager,
        "_load_latest_predicted_structures_dataset",
        lambda self, project_dir: loaded_dataset,
    )

    observed_settings = manager.load_project(project_dir)
    observed_settings.clusters_dir = str(clusters_dir)
    observed_settings.q_min = None
    observed_settings.q_max = None
    observed_settings.use_predicted_structure_weights = False
    manager.save_project(observed_settings)
    manager.build_scattering_components(observed_settings)
    manager.generate_prior_weights(observed_settings)
    observed_artifacts = project_artifact_paths(
        observed_settings,
        storage_mode="distribution",
    )
    _write_distribution_history_artifacts(
        observed_artifacts,
        state_name="prefit_observed",
        run_name="dream_observed",
    )

    predicted_settings = manager.load_project(project_dir)
    predicted_settings.clusters_dir = str(clusters_dir)
    predicted_settings.q_min = None
    predicted_settings.q_max = None
    predicted_settings.use_predicted_structure_weights = True
    manager.save_project(predicted_settings)
    manager.build_scattering_components(predicted_settings)
    manager.generate_prior_weights(predicted_settings)
    predicted_artifacts = project_artifact_paths(
        predicted_settings,
        storage_mode="distribution",
    )
    _write_distribution_history_artifacts(
        predicted_artifacts,
        state_name="prefit_predicted",
        run_name="dream_predicted",
    )

    manager.save_project(observed_settings)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert window.project_setup_tab.selected_distribution_id() == (
        observed_artifacts.distribution_id
    )
    assert window.project_setup_tab.current_prior_json_path() == (
        observed_artifacts.prior_weights_file
    )
    assert window.project_setup_tab.recognized_clusters_table.rowCount() == 1
    assert window.prefit_workflow.prefit_dir == observed_artifacts.prefit_dir
    assert window.prefit_tab.saved_state_combo.itemText(0) == (
        "prefit_observed"
    )
    assert window.dream_tab.selected_saved_run_dir() == str(
        (observed_artifacts.dream_runtime_dir / "dream_observed").resolve()
    )
    observed_structure_payload = build_prior_histogram_export_payload(
        window.project_setup_tab.current_prior_json_path(),
        mode="structure_fraction",
    )
    observed_atom_payload = build_prior_histogram_export_payload(
        window.project_setup_tab.current_prior_json_path(),
        mode="atom_fraction",
    )

    window.project_setup_tab.use_predicted_structure_weights_checkbox.setChecked(
        True
    )
    QApplication.processEvents()

    assert window.project_setup_tab.selected_distribution_id() == (
        predicted_artifacts.distribution_id
    )
    assert window.project_setup_tab.current_prior_json_path() == (
        predicted_artifacts.prior_weights_file
    )
    assert window.project_setup_tab.recognized_clusters_table.rowCount() == 2
    assert window.prefit_workflow.prefit_dir == predicted_artifacts.prefit_dir
    assert window.prefit_tab.saved_state_combo.itemText(0) == (
        "prefit_predicted"
    )
    assert window.dream_tab.selected_saved_run_dir() == str(
        (predicted_artifacts.dream_runtime_dir / "dream_predicted").resolve()
    )
    assert {
        (component.structure, component.motif)
        for component in window.prefit_workflow.components
    } == {("A", "no_motif"), ("A2", "predicted_rank01")}
    assert window.dream_workflow is not None
    assert {
        (component.structure, component.motif)
        for component in window.dream_workflow.prefit_workflow.components
    } == {("A", "no_motif"), ("A2", "predicted_rank01")}

    structure_payload = build_prior_histogram_export_payload(
        window.project_setup_tab.current_prior_json_path(),
        mode="structure_fraction",
    )
    atom_payload = build_prior_histogram_export_payload(
        window.project_setup_tab.current_prior_json_path(),
        mode="atom_fraction",
    )
    structure_matrix = np.asarray(structure_payload["matrix"], dtype=float)
    atom_matrix = np.asarray(atom_payload["matrix"], dtype=float)
    structure_by_label = {
        str(label): {
            str(segment): float(value)
            for segment, value in zip(
                structure_payload["segments"],
                structure_matrix[row_index],
                strict=False,
            )
        }
        for row_index, label in enumerate(structure_payload["labels"])
    }
    atom_by_label = {
        str(label): {
            str(segment): float(value)
            for segment, value in zip(
                atom_payload["segments"],
                atom_matrix[row_index],
                strict=False,
            )
        }
        for row_index, label in enumerate(atom_payload["labels"])
    }

    assert structure_payload["labels"] == ["A", "A2"]
    assert atom_payload["labels"] == ["A", "A2"]
    assert observed_structure_payload["labels"] == ["A"]
    assert observed_atom_payload["labels"] == ["A"]
    assert float(np.sum(structure_payload["totals"])) == pytest.approx(100.0)
    assert float(np.sum(atom_payload["totals"])) == pytest.approx(100.0)
    assert structure_by_label["A2"]["predicted_rank01"] > 0.0
    assert atom_by_label["A2"]["predicted_rank01"] > 0.0
    assert structure_by_label["A"]["no_motif"] < 100.0
    assert atom_by_label["A"]["no_motif"] < 100.0
    np.testing.assert_allclose(
        np.asarray(observed_structure_payload["matrix"], dtype=float)[0],
        np.asarray([100.0], dtype=float),
    )
    np.testing.assert_allclose(
        np.asarray(observed_atom_payload["matrix"], dtype=float)[0],
        np.asarray([100.0], dtype=float),
    )

    window.project_setup_tab.use_predicted_structure_weights_checkbox.setChecked(
        False
    )
    QApplication.processEvents()

    assert window.project_setup_tab.selected_distribution_id() == (
        observed_artifacts.distribution_id
    )
    assert window.project_setup_tab.current_prior_json_path() == (
        observed_artifacts.prior_weights_file
    )
    assert window.project_setup_tab.recognized_clusters_table.rowCount() == 1
    assert window.prefit_workflow.prefit_dir == observed_artifacts.prefit_dir
    assert window.prefit_tab.saved_state_combo.itemText(0) == (
        "prefit_observed"
    )
    assert window.dream_tab.selected_saved_run_dir() == str(
        (observed_artifacts.dream_runtime_dir / "dream_observed").resolve()
    )
    window.close()


def test_project_setup_can_load_saved_distribution_and_scope_prefit_and_dream(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    manager = SAXSProjectManager()
    clusters_dir = paths.project_dir / "clusters"
    structure_dir = clusters_dir / "A"
    structure_dir.mkdir(parents=True, exist_ok=True)
    (structure_dir / "frame_0001.xyz").write_text(
        "2\nframe 1\nA 0.0 0.0 0.0\nH 1.0 0.0 0.0\n",
        encoding="utf-8",
    )

    observed_component = np.linspace(10.0, 17.0, 8)

    def fake_build_profiles(
        builder,
        *,
        cluster_bins,
        progress_callback=None,
        progress_total=None,
        **kwargs,
    ):
        del progress_callback, progress_total, kwargs
        scale = 0.5 if builder.exclude_elements else 1.0
        output_path = builder.output_dir / "A_no_motif.txt"
        _write_component_file(
            output_path,
            builder.q_values,
            observed_component * scale,
        )
        cluster_bin = cluster_bins[0]
        return [
            AveragedComponent(
                structure=cluster_bin.structure,
                motif=cluster_bin.motif,
                file_count=len(cluster_bin.files),
                representative=cluster_bin.representative,
                source_dir=cluster_bin.source_dir,
                q_values=np.asarray(builder.q_values, dtype=float),
                mean_intensity=np.asarray(
                    observed_component * scale,
                    dtype=float,
                ),
                std_intensity=np.zeros_like(builder.q_values, dtype=float),
                se_intensity=np.zeros_like(builder.q_values, dtype=float),
                output_path=output_path,
            )
        ]

    monkeypatch.setattr(
        project_module.DebyeProfileBuilder,
        "build_profiles",
        fake_build_profiles,
    )

    observed_settings = manager.load_project(project_dir)
    observed_settings.clusters_dir = str(clusters_dir)
    observed_settings.exclude_elements = []
    manager.save_project(observed_settings)
    manager.build_scattering_components(observed_settings)
    manager.generate_prior_weights(observed_settings)
    observed_artifacts = project_artifact_paths(observed_settings)
    _write_distribution_history_artifacts(
        observed_artifacts,
        state_name="prefit_observed",
        run_name="dream_observed",
    )

    excluded_settings = manager.load_project(project_dir)
    excluded_settings.clusters_dir = str(clusters_dir)
    excluded_settings.exclude_elements = ["H"]
    manager.save_project(excluded_settings)
    manager.build_scattering_components(excluded_settings)
    manager.generate_prior_weights(excluded_settings)
    excluded_artifacts = project_artifact_paths(excluded_settings)
    _write_distribution_history_artifacts(
        excluded_artifacts,
        state_name="prefit_excluded",
        run_name="dream_excluded",
    )

    manager.save_project(observed_settings)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    assert window.project_setup_tab.computed_distribution_combo.count() == 2
    assert window.project_setup_tab.current_prior_json_path() == (
        observed_artifacts.prior_weights_file
    )
    assert window.prefit_workflow.prefit_dir == observed_artifacts.prefit_dir
    assert window.prefit_tab.saved_state_combo.itemText(0) == (
        "prefit_observed"
    )
    assert window.dream_tab.selected_saved_run_dir() == str(
        (observed_artifacts.dream_runtime_dir / "dream_observed").resolve()
    )

    target_index = (
        window.project_setup_tab.computed_distribution_combo.findData(
            excluded_artifacts.distribution_id
        )
    )
    assert target_index >= 0
    window.project_setup_tab.computed_distribution_combo.setCurrentIndex(
        target_index
    )
    window.project_setup_tab.load_distribution_button.click()
    QApplication.processEvents()

    assert window.project_setup_tab.exclude_elements() == ["H"]
    assert window.project_setup_tab.current_prior_json_path() == (
        excluded_artifacts.prior_weights_file
    )
    assert window.prefit_workflow.prefit_dir == excluded_artifacts.prefit_dir
    assert window.prefit_tab.saved_state_combo.itemText(0) == (
        "prefit_excluded"
    )
    assert window.dream_tab.selected_saved_run_dir() == str(
        (excluded_artifacts.dream_runtime_dir / "dream_excluded").resolve()
    )
    window.close()


def test_no_contrast_distribution_paths_preserve_legacy_loading_without_build_mode(
    tmp_path,
):
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    manager = SAXSProjectManager()
    settings = manager.load_project(project_dir)
    template_legacy_distribution_id = (
        project_module._distribution_id_for_settings(
            settings,
            include_template=True,
            include_build_mode=False,
        )
    )
    oldest_legacy_distribution_id = (
        project_module._distribution_id_for_settings(
            settings,
            include_template=False,
            include_build_mode=False,
        )
    )

    def write_legacy_distribution(
        distribution_id: str,
        *,
        include_template_name: bool,
    ) -> Path:
        distribution_dir = (
            build_project_paths(project_dir).saved_distributions_dir
            / distribution_id
        )
        component_dir = distribution_dir / "scattering_components"
        component_dir.mkdir(parents=True, exist_ok=True)
        shutil.copytree(
            paths.scattering_components_dir,
            component_dir,
            dirs_exist_ok=True,
        )
        shutil.copy2(
            paths.project_dir / "md_saxs_map.json",
            distribution_dir / "md_saxs_map.json",
        )
        shutil.copy2(
            paths.project_dir / "md_prior_weights.json",
            distribution_dir / "md_prior_weights.json",
        )
        payload = {
            "schema_version": 1,
            "distribution_id": distribution_id,
            "label": (
                "Observed Only | Template: "
                f"{settings.selected_model_template} | Excluded: None | "
                "q-range: default | Grid: experimental grid"
                if include_template_name
                else "Observed Only | Excluded: None | q-range: default | "
                "Grid: experimental grid"
            ),
            "use_predicted_structure_weights": False,
            "exclude_elements": [],
            "clusters_dir": None,
            "q_min": None,
            "q_max": None,
            "use_experimental_grid": True,
            "q_points": None,
            "component_artifacts_ready": True,
            "prior_artifacts_ready": True,
        }
        if include_template_name:
            payload["template_name"] = settings.selected_model_template
        (distribution_dir / "distribution.json").write_text(
            json.dumps(payload, indent=2) + "\n",
            encoding="utf-8",
        )
        return distribution_dir

    template_legacy_dir = write_legacy_distribution(
        template_legacy_distribution_id,
        include_template_name=True,
    )
    oldest_legacy_dir = write_legacy_distribution(
        oldest_legacy_distribution_id,
        include_template_name=False,
    )

    resolved_template_legacy = project_artifact_paths(
        settings,
        storage_mode="distribution",
    )
    assert resolved_template_legacy.distribution_id == (
        template_legacy_distribution_id
    )
    assert resolved_template_legacy.root_dir == template_legacy_dir

    shutil.rmtree(template_legacy_dir)
    resolved_oldest_legacy = project_artifact_paths(
        settings,
        storage_mode="distribution",
    )
    assert resolved_oldest_legacy.distribution_id == (
        oldest_legacy_distribution_id
    )
    assert resolved_oldest_legacy.root_dir == oldest_legacy_dir

    contrast_settings = ProjectSettings.from_dict(settings.to_dict())
    contrast_settings.component_build_mode = COMPONENT_BUILD_MODE_CONTRAST
    contrast_paths = project_artifact_paths(
        contrast_settings,
        storage_mode="distribution",
    )
    assert contrast_paths.distribution_id not in {
        template_legacy_distribution_id,
        oldest_legacy_distribution_id,
    }

    loaded_settings = manager.settings_for_saved_distribution(
        project_dir,
        oldest_legacy_distribution_id,
        base_settings=contrast_settings,
    )
    assert (
        loaded_settings.component_build_mode
        == COMPONENT_BUILD_MODE_NO_CONTRAST
    )
    assert loaded_settings.selected_model_template == (
        contrast_settings.selected_model_template
    )


def test_project_setup_loads_saved_distribution_with_cropped_q_range(
    qapp,
    tmp_path,
    monkeypatch,
):
    del qapp
    project_dir, paths = _build_minimal_saxs_project(tmp_path)
    manager = SAXSProjectManager()
    clusters_dir = paths.project_dir / "clusters"
    structure_dir = clusters_dir / "A"
    structure_dir.mkdir(parents=True, exist_ok=True)
    (structure_dir / "frame_0001.xyz").write_text(
        "2\nframe 1\nA 0.0 0.0 0.0\nH 1.0 0.0 0.0\n",
        encoding="utf-8",
    )

    def fake_build_profiles(
        builder,
        *,
        cluster_bins,
        progress_callback=None,
        progress_total=None,
        **kwargs,
    ):
        del progress_callback, progress_total, kwargs
        values = np.linspace(10.0, 17.0, len(builder.q_values))
        output_path = builder.output_dir / "A_no_motif.txt"
        _write_component_file(output_path, builder.q_values, values)
        cluster_bin = cluster_bins[0]
        return [
            AveragedComponent(
                structure=cluster_bin.structure,
                motif=cluster_bin.motif,
                file_count=len(cluster_bin.files),
                representative=cluster_bin.representative,
                source_dir=cluster_bin.source_dir,
                q_values=np.asarray(builder.q_values, dtype=float),
                mean_intensity=np.asarray(values, dtype=float),
                std_intensity=np.zeros_like(builder.q_values, dtype=float),
                se_intensity=np.zeros_like(builder.q_values, dtype=float),
                output_path=output_path,
            )
        ]

    monkeypatch.setattr(
        project_module.DebyeProfileBuilder,
        "build_profiles",
        fake_build_profiles,
    )

    original_settings = manager.load_project(project_dir)
    original_settings.clusters_dir = str(clusters_dir)
    original_settings.q_min = None
    original_settings.q_max = None
    manager.save_project(original_settings)
    manager.build_scattering_components(original_settings)
    manager.generate_prior_weights(original_settings)
    reduced_settings = manager.load_project(project_dir)
    reduced_settings.clusters_dir = str(clusters_dir)
    reduced_settings.q_min = 0.12
    reduced_settings.q_max = 0.19
    manager.save_project(reduced_settings)
    manager.build_scattering_components(reduced_settings)
    manager.generate_prior_weights(reduced_settings)
    reduced_artifacts = project_artifact_paths(reduced_settings)

    manager.save_project(original_settings)
    window = SAXSMainWindow(initial_project_dir=project_dir)

    target_index = (
        window.project_setup_tab.computed_distribution_combo.findData(
            reduced_artifacts.distribution_id
        )
    )
    assert target_index >= 0
    window.project_setup_tab.computed_distribution_combo.setCurrentIndex(
        target_index
    )
    window.project_setup_tab.load_distribution_button.click()
    QApplication.processEvents()

    expected_q = np.asarray(
        [0.12142857142857144, 0.15714285714285714, 0.19285714285714284],
        dtype=float,
    )
    assert window.project_setup_tab.selected_distribution_id() == (
        reduced_artifacts.distribution_id
    )
    assert window.prefit_workflow is not None
    assert window.prefit_workflow.prefit_dir == reduced_artifacts.prefit_dir
    assert np.allclose(
        window.prefit_workflow.evaluate().q_values,
        expected_q,
    )
    assert window.dream_workflow is not None
    assert np.allclose(
        window.dream_workflow.prefit_workflow.evaluate().q_values,
        expected_q,
    )
    window.close()


def test_project_settings_roundtrip_preserves_powerpoint_export_settings(
    tmp_path,
):
    manager = SAXSProjectManager()
    settings = ProjectSettings(
        project_name="demo_project",
        project_dir=str(tmp_path / "demo_project"),
        powerpoint_export_settings=PowerPointExportSettings(
            font_family="Courier New",
            component_color_map="plasma",
            prior_histogram_color_map="cividis",
            solvent_sort_histogram_color_map="magma",
            generate_manifest=False,
            export_figure_assets=False,
        ),
    )

    manager.save_project(settings)
    loaded = manager.load_project(settings.project_dir)

    assert loaded.powerpoint_export_settings.font_family == "Courier New"
    assert loaded.powerpoint_export_settings.component_color_map == "plasma"
    assert (
        loaded.powerpoint_export_settings.prior_histogram_color_map
        == "cividis"
    )
    assert (
        loaded.powerpoint_export_settings.solvent_sort_histogram_color_map
        == "magma"
    )
    assert not loaded.powerpoint_export_settings.generate_manifest
    assert not loaded.powerpoint_export_settings.export_figure_assets


def test_scan_cluster_inventory_reports_progress_and_rows(tmp_path):
    cluster_dir = tmp_path / "clusters"
    structure_dir = cluster_dir / "Pb2I4"
    structure_dir.mkdir(parents=True)
    (structure_dir / "frame_0001.xyz").write_text(
        "2\ncomment\nPb 0.0 0.0 0.0\nI 1.0 0.0 0.0\n",
        encoding="utf-8",
    )
    manager = SAXSProjectManager()
    events: list[tuple[int, int, str]] = []

    result = manager.scan_cluster_inventory(
        cluster_dir,
        progress_callback=lambda processed, total, message: events.append(
            (processed, total, message)
        ),
    )

    assert isinstance(result, ClusterImportResult)
    assert result.available_elements == ["I", "Pb"]
    assert result.total_files == 1
    assert result.cluster_rows[0]["structure"] == "Pb2I4"
    assert events[0][2] == "Importing cluster files..."
    assert events[-1][2] == "Cluster import complete."


def test_cluster_inventory_cache_reuses_previous_scan(tmp_path, monkeypatch):
    cluster_dir = tmp_path / "clusters"
    structure_dir = cluster_dir / "Pb2I4"
    structure_dir.mkdir(parents=True)
    (structure_dir / "frame_0001.xyz").write_text(
        "3\ncomment\nPb 0.0 0.0 0.0\nI 1.0 0.0 0.0\nO 0.0 1.0 0.0\n",
        encoding="utf-8",
    )
    manager = SAXSProjectManager()

    initial = manager.scan_cluster_inventory(cluster_dir)

    def _fail_if_rescanned(_path):
        raise AssertionError("cluster inventory should have reused its cache")

    monkeypatch.setattr(
        project_module,
        "scan_structure_elements",
        _fail_if_rescanned,
    )

    cached = manager._collect_cluster_inventory(cluster_dir)

    assert cached.available_elements == initial.available_elements
    assert cached.total_files == initial.total_files
